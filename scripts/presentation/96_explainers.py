"""Backup slides that explain the findings behind each results slide.

These are appended to the author's own working deck
(`EWGT_26_Bienzeisler_new.pptx`, 41 slides) as a backup block: one to three
slides per results slide, each answering the question that slide provokes when
somebody in the audience looks closely.

Two style rules, both from the author:

* **Body in the v1 language** — the kicker, real bulleted lists, stat blocks,
  tables and native-shape schematics of `91_build_pptx.py`. Not the icon-badge
  house style of `94_build_house_deck.py`.
* **Headlines in the house form** — two lines in the title placeholder: the
  running section, then this slide's subject.

Every number is measured from the Stage-3 grid or the repository's constants;
the analysis behind the frequency-mix block is reproduced by `--audit`, which
recomputes each figure and fails loudly if it has moved.

The source deck is never written to. Output goes to a new file.

The 2026-08 revision withdrew the first chapter. It explained the bump of
consolidation at theta = 10 %, which the universal tour rule showed to be a
pricing artefact rather than a behaviour of the system; `block_mix_rev()`
replaces it with the withdrawal and the revision grid's own numbers. Pass
`--no-revision` to rebuild the submission-era backup unchanged.

The source deck is never written to, and neither is any deck that already
exists: the output goes through `_outguard.resolve()`.

Usage:
    python scripts/presentation/96_explainers.py [--out PATH] [--audit]
    python scripts/presentation/96_explainers.py --out-suffix _rev2026-08
"""
from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):      # theta and rho in the audit log
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

import numpy as np
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "revision"))

import _house as H                                                # noqa: E402
import _outguard as G                                             # noqa: E402
from _house import (AMBER, BODY_T, CRIM, DIM, GREEN, INK, INK2, L, LINE,
                    PANEL, RED, S1, S3, S4, S5, S6, SW, TEAL, W, WHITE,
                    hrule, hslide, label_box, pic, rect, txt)      # noqa: E402

import _conclusion as CONC                                        # noqa: E402
import _dummies as DUM                                            # noqa: E402

B = H.B
FIG = H.FIG / "tierB"          # the per-carrier and bulge panels
SRC = Path(r"C:/Users/bienzeisler/Documents/Präsentationen/EWGT/2026/"
           r"EWGT_26_Bienzeisler_new.pptx")
DEFAULT_OUT = SRC.parent / "EWGT_26_Bienzeisler_new_plus_explainers.pptx"

# The results slide each block unpacks, matched by a fragment of its title.
# Not by slide number: the working deck is edited between builds, and a
# hard-coded index silently starts pointing at the wrong slide. The number is
# resolved from the deck at build time and simply omitted if the title is gone.
EXPLAINS = {
    "mix": "delivery-frequency mix",
    "trade": "service improves faster",
    "range": "efficiency range",
    "maps": "where the delivery days land",
    "where": "pays where delivery is sparse",
    "valid": "validation",
}
_RESOLVED: dict = {}


def resolve_targets(prs) -> dict:
    """Find each block's results slide in the deck, by title fragment."""
    titles = []
    for i, sl in enumerate(prs.slides, 1):
        ph = [sh for sh in sl.shapes if sh.is_placeholder
              and sh.placeholder_format.idx == 0 and sh.has_text_frame]
        titles.append((i, ph[0].text_frame.text.replace("", " ") if ph else ""))
    out = {}
    for key, frag in EXPLAINS.items():
        hit = next(((i, t) for i, t in titles if frag in t.lower()), None)
        if hit is None:
            print(f"  ! no slide matches {frag!r} — tag will omit the number")
        else:
            print(f"  {key:6s} -> slide {hit[0]:2d}  {hit[1][:60]}")
        out[key] = hit
    return out


VAN = 189.15        # EUR per van-day, config/constants.py
CAP = 230           # parcels per van

# The revision grid, set by main() and read by block_mix_rev(). None when
# --no-revision is passed, in which case the submission-era block_mix() runs
# and the deck is exactly what it was.
_F = None
_RV = None
_TAG = True


def mark96(s, note, cite):
    """Speaker notes plus the provisional chip, when a grid is loaded."""
    if _RV is None:
        return s
    _RV.notes(s, note, cite=cite)
    _RV.provisional(s, enabled=_TAG)
    return s


def _MAXSAV(provider: str) -> float:
    """That LSP's best operator-lens saving anywhere on the theta = 1 grid."""
    import numpy as np
    c = _RV.D.load_costs_v2()
    base = float(c[np.isclose(c.penalty, 0) & np.isclose(c.share_willing, 0)
                   & (c.provider == provider)].operator_cost_eur.iloc[0])
    at1 = c[np.isclose(c.share_willing, 1.0) & (c.provider == provider)]
    return float((1 - at1.operator_cost_eur / base).max() * 100)


# ── slide scaffolding: house headline, v1 body ──────────────────────────────
def xslide(prs, key, section, subject, source=None):
    """A backup slide, tagged with the results slide it unpacks."""
    s = hslide(prs, section, subject, source)
    hit = _RESOLVED.get(key)
    tag = (f"explains slide {hit[0]}  ·  {hit[1]}" if hit
           else f"explains  ·  {EXPLAINS[key]}")
    txt(s, L, 1.00, W, 0.30, tag, 11.5, bold=True, color=RED, spc=1.4,
        caps=True)
    return s


def vbullets(slide, items, t, *, size=20.0, floor=17.0, l=L, w=W, label=""):
    """A v1 bulleted list that stops above this slide's source line.

    `B.bullets` takes a fixed height and happily writes past the bottom of the
    slide. Here the block is measured first and the type steps down until it
    fits the room the citation leaves; if even the floor is not enough the
    build says so rather than printing into the footer.
    """
    limit = getattr(slide, "body_bottom", H.BODY_B) - t
    gap = 12.0

    def need(sz):
        return sum(H.text_height(it, w, sz, 1.22) + gap / 72.0 for it in items)

    while size > floor and need(size) > limit:
        size -= 0.5
    over = need(size) - limit
    if over > 0.03:
        print(f"  ! {label or 'bullets'}: overruns by {over:.2f} in at "
              f"{size:g} pt", file=sys.stderr)
    return B.bullets(slide, items, t, l=l, w=w, size=size,
                     h=max(0.4, limit), gap=int(gap))


def divider(prs, num, kicker, headline, body):
    return B.divider(prs, num, kicker, headline, body)



# ═══════════════════════════════════════════════════════════════════════════
# the backup as a chain of questions
# ═══════════════════════════════════════════════════════════════════════════
# Each slide answers exactly one question, and the question is the headline.
# The order is the order somebody actually asks them in: first what looks odd,
# then whether it is real, then why, then whether the numbers can be trusted.
# `chapter()` records where each part starts so the contents slide can be
# filled in with real slide numbers once everything has been laid down.

_CHAPTERS: list = []
_CONTENTS = None


def chapter(prs, title, sub):
    _CHAPTERS.append((title, sub, len(prs.slides) + 1))


def block_contents(prs):
    """A map of the backup, filled in with real numbers once it is built."""
    global _CONTENTS
    s = hslide(prs, "Backup", "What is in here, and what each part answers",
               "Every slide in this section answers one question. Jump "
               "straight to the one you were asked.")
    _CONTENTS = s
    return s


def fill_contents():
    """Write the recorded chapters onto the contents slide."""
    if _CONTENTS is None:
        return
    y = BODY_T + 0.22
    for title, sub, num in _CHAPTERS:
        label_box(_CONTENTS, L, y, 0.95, 0.54, RED,
                  [(str(num), 21, True, WHITE)])
        txt(_CONTENTS, L + 1.18, y - 0.03, W - 1.18, 0.34, title, 20,
            bold=True, color=INK)
        txt(_CONTENTS, L + 1.18, y + 0.31, W - 1.18, 0.30, sub, 15, color=DIM)
        y += 0.70


# ═══════════════════════════════════════════════════════════════════════════
# part 1 · the odd thing in the frequency picture
# ═══════════════════════════════════════════════════════════════════════════
def block_mix(prs):
    chapter(prs, "The odd thing in the frequency picture",
            "Why bundling survives a huge fee when almost nobody takes part")

    # ── 1 · what looks odd ───────────────────────────────────────────────
    s = xslide(prs, "mix", "Part 1 · The odd thing", "What looks odd here?",
               "Share of the 312 delivery areas that give up daily delivery · "
               "Stage-3 grid.")
    txt(s, L + 2.30, BODY_T + 0.06, W - 2.30, 0.28,
        "share of customers willing to wait", 12, color=DIM, spc=0.8)
    y = B.table(s, ["Fee", "10 %", "20 %", "30 %", "40 %", "100 %"],
                [[("none", "key"), "87.5 %", "92.0 %", "94.2 %", "96.2 %",
                  "100 %"],
                 [("P = 5", "key"), ("49.7 %", "num"), ("42.9 %", "num"),
                  "14.7 %", "0 %", "0 %"],
                 [("P = 10", "key"), ("41.7 %", "num"), ("10.9 %", "num"),
                  "0 %", "0 %", "0 %"]],
                BODY_T + 0.40, widths=[2.3, 2.0, 2.0, 2.0, 2.0, 2.0],
                reserve=2.6)
    label_box(s, L, y + 0.30, W, 0.85, H.BLUSH,
              [("Follow the bottom row from left to right. The more people "
                "join in, the less gets bundled.", 22, True, RED)],
              line_col=RED)
    vbullets(s, ["That feels backwards. More volunteers should make bundling "
                 "easier, not harder.",
                 "And the top row, with no fee at all, does exactly what you "
                 "expect: it climbs.",
                 "The next slides work out why the two rows point in opposite "
                 "directions."],
             y + 1.30)

    # ── 2 + 3 · the argument, in pictures ────────────────────────────────
    DUM.slide_two_price_tags(prs, xslide)
    DUM.slide_the_proof(prs, xslide)

    # ── 3b · the fee is not a cost ───────────────────────────────────────
    # A reader stops here and asks: surely the shadow price is not added to
    # the bill? It is not. It steers the choice and nothing else, and the
    # slide has to say so before the next one shows what that steering costs.
    s = xslide(prs, "mix", "Part 1 · The odd thing",
               "Wait — is the fee added to the cost? No.",
               "Reported cost at the harshest fee, 10 % taking part: "
               "1 841 323 € with and without the penalty term. Identical.")
    for i, (head, body, col, fill, strike) in enumerate([
            ("What we count", "driving, vehicles, labour — real routing money",
             INK, PANEL, False),
            ("What we never count", "the service fee. Nobody is ever charged "
             "it", RED, H.BLUSH, True)]):
        x = L + i * (W / 2 + 0.20)
        cwd = W / 2 - 0.20
        rect(s, x, BODY_T + 0.35, cwd, 1.45, fill, line_col=LINE)
        txt(s, x + 0.30, BODY_T + 0.58, cwd - 0.60, 0.42, head, 24, bold=True,
            color=col)
        txt(s, x + 0.30, BODY_T + 1.06, cwd - 0.60, 0.60, body, 19, color=INK2,
            line=1.22)
        if strike:
            hrule(s, x + 0.30, BODY_T + 1.22, cwd - 0.60, RED, 2.5)
    label_box(s, L, BODY_T + 2.10, W, 0.95, WHITE,
              [("The fee lives only inside the optimiser. It decides which "
                "weekly plan wins — it is never money.", 22, True, INK)],
              line_col=LINE)
    vbullets(s, ["So why does a higher fee reduce the saving at all?",
                 "Because it makes the optimiser pick gentler plans. Gentler "
                 "plans deliver more often, and delivering more often costs "
                 "more.",
                 "The money that disappears is routing money never saved — not "
                 "a fee anybody hands over."],
             BODY_T + 3.25)

    # ── 3b2 · the same thing, with the numbers written out ───────────────
    DUM.slide_worked_example(prs, xslide)

    # ── 3c · why 10 % beats 20 % ─────────────────────────────────────────
    s = xslide(prs, "mix", "Part 1 · The odd thing",
               "Why is 10 % better than 20 %?",
               "Bar height is the most that could be saved; green is what is "
               "actually saved at the harshest fee.")
    pic(s, FIG / "figB5_prize_and_bill.png", L, BODY_T + 0.24, W, 3.35)
    vbullets(s, [[("Twice as many wait, so the fee steers twice as hard: the "
                   "areas that still bundle drop from ", False),
                  ("41.7 % to 10.9 %", True), (".", False)],
                 "A gentler plan saves less — the region gives up 76 616 € "
                 "at 10 %, but 144 037 € at 20 %.",
                 "But the prize grows only a tenth — so 68 425 € are left, "
                 "then 15 109 €, then nothing."],
             BODY_T + 3.66)

    # ── 4 · are they outliers? ───────────────────────────────────────────
    # An earlier version of this block claimed the saving came from dropping
    # whole van-days. It does not: 24 of 6 397 are saved and 107 of the 130
    # areas drop none. The slides below carry what is in the data.
    s = xslide(prs, "mix", "Part 1 · The odd thing",
               "Are these just a few tiny areas we should throw out?",
               "The areas that give up daily delivery at the harshest fee and "
               "10 % participation, by their weekly parcel volume.")
    pic(s, FIG / "figB1_who_consolidates.png", L, BODY_T + 0.24, W, 3.35)
    vbullets(s, ["No. There are 130 of them, and together they carry a "
                 "quarter of all parcels in the region.",
                 [("The smallest still handles ", False),
                  ("774 parcels a week", True),
                  (", the typical one 2 172. Not one is a curiosity.", False)],
                 "Throwing them out would delete a quarter of the study area."],
             BODY_T + 3.75)

    # ── 5 · do we save vans? ─────────────────────────────────────────────
    s = xslide(prs, "mix", "Part 1 · The odd thing",
               "So do we save delivery vans? About half of it, yes.",
               "Vehicle-days on the corrected fleet metric (2026-08-25: the "
               "express parcels of all non-delivering cells at a depot ride "
               "one pooled tour). Saving against the 1 909 748 € baseline.")
    pic(s, FIG / "figB2_where_the_money_is.png", L, BODY_T + 0.24, W, 3.35)
    vbullets(s, [[("180 of 6 397 van-days fall away — worth 34 047 €, which "
                   "is ", False), ("half", True),
                  (" of the 68 425 € saved there.", False)],
                 "The other half is shorter driving with the same vehicles.",
                 "At the operating point it is the other way round: 524 "
                 "van-days, but only 28 % of the money."],
             BODY_T + 3.75)

    # ── 6 · what decides it ──────────────────────────────────────────────
    s = xslide(prs, "mix", "Part 1 · The odd thing",
               "What decides it, then?",
               "Each dot is one setting of fee and participation from the "
               "Stage-3 grid.")
    pic(s, FIG / "figB3_ptheta_collapse.png", L, BODY_T + 0.24, W, 3.35)
    vbullets(s, ["Not the fee on its own, and not the participation on its "
                 "own — the two multiplied.",
                 "That product is simply the bill you end up paying.",
                 "A fee of 10 with one in ten joining behaves just like a fee "
                 "of 1 with everybody joining."],
             BODY_T + 3.75)

    # ── 7 · which areas ──────────────────────────────────────────────────
    s = xslide(prs, "mix", "Part 1 · The odd thing",
               "Why do small areas bundle and big ones not?",
               "Weekly parcel volume of the areas choosing each number of "
               "delivery days. The boxes cover the middle half of them.")
    pic(s, FIG / "figB4_size_vs_takt.png", L, BODY_T + 0.24, W, 3.35)
    vbullets(s, ["The fewer delivery days an area picks, the smaller it is. "
                 "Big areas stay daily.",
                 "Double the participation and the size limit halves, so the "
                 "two-day option empties out first.",
                 "Five days a week lasts longest: it saves a little and costs "
                 "hardly any waiting."],
             BODY_T + 3.75)


# ═══════════════════════════════════════════════════════════════════════════
# part 2 · what a fee buys
# ═══════════════════════════════════════════════════════════════════════════
def block_mix_rev(prs):
    """Replaces part 1, which explained an artefact.

    The old chapter took ten slides to work out why bundling survives a
    punitive fee when almost nobody takes part, and concluded that the real
    knob is the product P x theta. Both the phenomenon and the conclusion were
    consequences of the pre-revision express price: standard parcels of every
    non-delivering area of a hub rode ONE pooled tour, which no operator would
    dispatch and which only the scenario could ever have. With one universal
    tour rule the phenomenon is gone, so the chapter is not corrected -- it is
    withdrawn, and these three slides say so and show what replaced it.
    """
    f, RV = _F, _RV
    chapter(prs, "A finding this deck has withdrawn",
            "The bump at 10 % adoption was a pricing artefact, not a mechanism")

    def mark(s, note, cite):
        RV.notes(s, note, cite=cite)
        RV.provisional(s, enabled=_TAG)
        return s

    # ── 1 · what the old chapter claimed ─────────────────────────────────
    s = xslide(prs, "mix", "Part 1 · Withdrawn", "What this chapter used to say",
               "The claim, as it stood in the submission-era deck.")
    y = B.table(s, ["Fee", "10 % join", "20 % join", "30 % join"],
                [[("none", "key"), "87.5 %", "92.0 %", "94.2 %"],
                 [("P = 5", "key"), ("49.7 %", "num"), ("42.9 %", "num"),
                  "14.7 %"],
                 [("P = 10", "key"), ("41.7 %", "num"), ("10.9 %", "num"),
                  "0 %"]],
                BODY_T + 0.30, widths=[2.6, 3.2, 3.2, 3.2], reserve=2.60)
    label_box(s, L, y + 0.28, W, 0.85, H.BLUSH,
              [("The old reading: the effective knob is the product P x theta "
                "— bundling survives a punitive fee when few people join.",
                22, True, RED)], line_col=RED)
    vbullets(s, ["Those shares came from the submission grid, where the "
                 "standard parcels of every non-delivering area of a hub rode "
                 "one pooled express tour.",
                 "That tour was cheap because it was one tour — and no "
                 "operator would ever dispatch it."],
             y + 1.28, label="withdrawn/claim")
    mark(s, RV.bulge_notes(f), RV.cites("§40.7", "§40.8"))

    # ── 2 · what the revision grid says ──────────────────────────────────
    rows = RV.bulge_rows(f)
    s = xslide(prs, "mix", "Part 1 · Withdrawn",
               "What the same cells do under one tour rule",
               "Share of the 312 delivery areas that give up daily delivery, "
               "and the routing saving at the same cell, on the revision grid "
               "(routing-optimal plan).")
    y = B.table(s, ["Operating point", "Areas consolidating", "Routing saving"],
                [[(r[0], "key"), (r[1], "num"), r[2]] for r in rows],
                BODY_T + 0.30, widths=[4.0, 4.0, 3.2], reserve=2.60)
    label_box(s, L, y + 0.28, W, 0.85, H.BLUSH,
              [(f"41.7 % becomes {rows[2][1]}, and the saving that went with "
                f"it becomes {rows[2][2]}.", 22, True, RED)], line_col=RED)
    vbullets(s, ["Per-cell express with a 230-parcel minimum tour, applied "
                 "scenario-blind: there is no code branch left that could "
                 "price the baseline differently.",
                 "The theta < 1 column falls to an honest floor — and the "
                 "bump falls out of it."],
             y + 1.28, label="withdrawn/now")
    mark(s, RV.TOUR_RULE_NOTES, RV.cites("§40.7", "§40.8", "§40.9"))

    # ── 3 · what to say if asked ─────────────────────────────────────────
    s = xslide(prs, "mix", "Part 1 · Withdrawn", "What to say if you are asked",
               "The honest framing of a withdrawn result.")
    vbullets(s, ["The bump was real in the model and it is gone from the "
                 "model — because the model changed, not because the data did.",
                 "The change was a modelling correction we made ourselves: "
                 "one tour rule for baseline and scenario, and a minimum tour "
                 "size of one van-load.",
                 "It cost us the most interesting-looking finding in the "
                 "submission, and it is the reason the theta < 1 column is now "
                 "small rather than surprising.",
                 "Nothing in the full-adoption column depended on it: at "
                 "theta = 1 the express path is provably dead."],
             BODY_T + 0.35, label="withdrawn/framing")
    mark(s, "The withdrawal itself is the reviewer-facing point: the finding "
            "was retracted by our own correction, before review, and the "
            "theta = 1 anchor is untouched because the express path is "
            "provably inactive there (fast_share = 0).",
         ["§40.7", "§40.8"])


def block_trade(prs):
    chapter(prs, "What a fee actually buys",
            "Why the middle of the range beats both ends")

    _h = _F.headline if _F is not None else None
    s = xslide(prs, "trade", "Part 2 · What a fee buys",
               "What does each step of the fee cost us, and what does it buy?",
               (f"Everyone taking part, operator-polished plan (stage 2), "
                f"operator lens. Saving against the "
                f"{_RV.eur(_F.base_operator)} € weekly operator baseline; "
                f"waiting averaged over all parcels, including those that "
                f"never wait." if _h else
                "Everyone taking part. Saving measured against the "
                "1 909 748 € weekly baseline; waiting averaged over all "
                "parcels, including those that never wait."))
    if _h:
        _rows = []
        for _P in (0.0, 0.25, 0.5, 0.75, 1.0, 2.0):
            _x = _h[_P]
            _given = ("—" if _P == 0
                      else f"{_h[0.0]['op2'] - _x['op2']:.1f} points")
            _rm = ("—" if _P == 0 else
                   f"{(1 - _x['wait2'] / _h[0.0]['wait2']) * 100:.0f} %")
            _rows.append([("none" if _P == 0 else f"P = {_P:g}", "key"),
                          (f"{_x['op2']:.1f} %", "num"),
                          f"{_x['wait2']:.2f} d", _given,
                          (_rm, "good") if 0 < _P <= 0.5 else _rm])
    else:
        _rows = [[("none", "key"), ("22.8 %", "num"), "0.98 d", "—", "—"],
                 [("P = 0.25", "key"), ("18.5 %", "num"), "0.46 d",
                  "4.3 points", ("half of it", "good")],
                 [("P = 0.5", "key"), ("13.5 %", "num"), "0.23 d",
                  "9.3 points", ("three quarters", "good")],
                 [("P = 0.75", "key"), "10.2 %", "0.14 d", "12.6 points",
                  "86 %"],
                 [("P = 1", "key"), "7.5 %", "0.09 d", "15.3 points", "91 %"],
                 [("P = 2", "key"), "1.2 %", "0.01 d", "21.6 points", "99 %"]]
    y = B.table(s, ["Fee", "Operator saving", "People wait", "Saving given up",
                    "Waiting removed"], _rows,
                BODY_T + 0.20, widths=[2.2, 2.4, 2.2, 2.4, 2.4], reserve=1.6)
    txt(s, L, y + 0.24, W, 1.32,
        "The first step is still the bargain: give up under two points of "
        "operator saving and half the waiting disappears.\nEvery step after "
        "that buys less and costs more.", 22, bold=True, color=RED, line=1.28)
    mark96(s, "Restated in the operator lens on the operator-polished plan. "
              "The submission's column was routing euro on the "
              "routing-optimal plan, which this revision supersedes.",
           ["§40.15", "§40.18"])

    s = xslide(prs, "range", "Part 2 · What a fee buys",
               "Why not simply take the biggest saving?",
               ("Operator lens, operator-polished plan, revision grid. Fleet "
                "figures are the sum of hub peaks against the daily-delivery "
                "baseline of the same grid." if _h else
                "Fleet figures are for the balanced and smoothed schedules; "
                "18 of the 80 settings sit on the efficient front."))
    _cards = ([("No fee", f"{_h[0.0]['op2']:.1f} %",
                f"{_h[0.0]['wait2']:.2f} d",
                "the cheapest week for an operator — but the longest wait",
                False),
               ("P = 0.25", f"{_h[0.25]['op2']:.1f} %",
                f"{_h[0.25]['wait2']:.2f} d",
                "half the waiting, and the same peak-fleet cut", True),
               ("P = 0.5", f"{_h[0.5]['op2']:.1f} %",
                f"{_h[0.5]['wait2']:.2f} d",
                f"peak fleet down {abs(_h[0.5]['peak2_pct']):.0f} %, weekday "
                f"swings down 67 %", True)] if _h else [
              ("No fee", "22.8 %", "0.98 d",
               "the cheapest week — but a full day of waiting", False),
              ("P = 0.25", "18.5 %", "0.46 d",
               "half the waiting for a fifth of the saving", True),
              ("P = 0.5", "13.5 %", "0.23 d",
               "peak fleet down 12.9 %, weekday swings down 54 %", True)])
    for i, (nm, sav, wait, note, hot) in enumerate(_cards):
        x = L + i * (W / 3 + 0.02)
        cw = W / 3 - 0.30
        rect(s, x, BODY_T + 0.30, cw, 0.10, RED if hot else LINE)
        txt(s, x, BODY_T + 0.52, cw, 0.44, nm, 24, bold=True,
            color=RED if hot else INK)
        txt(s, x, BODY_T + 1.05, cw, 0.60, sav, 40, bold=True,
            color=RED if hot else INK)
        txt(s, x, BODY_T + 1.72, cw, 0.40, f"+{wait} waiting", 20, color=DIM)
        txt(s, x, BODY_T + 2.20, cw, 1.10, note, 20, color=INK2, line=1.22)
    vbullets(s, ["Push the fee to zero and you get the most money — and the "
                 "longest wait.",
                 "Push it high and the waiting vanishes, but so does the point "
                 "of doing it.",
                 "Which of the middle settings is right is a service decision, "
                 "not a modelling one."],
             BODY_T + 3.35)
    mark96(s, "Operator lens, operator-polished plan. In this lens P = 0 is "
              "still the largest number, but it costs 0.77 d of waiting "
              "against 0.39 d at P = 0.25 for the same peak-fleet cut -- "
              "which is why the knee stays at P = 0.25.",
           ["§40.15", "§40.18"])


# ═══════════════════════════════════════════════════════════════════════════
# part 3 · where it happens
# ═══════════════════════════════════════════════════════════════════════════
def block_maps(prs):
    chapter(prs, "Where it happens first",
            "Why the countryside changes and the city centre does not")

    s = xslide(prs, "maps", "Part 3 · Where it happens",
               "Why does the countryside change first?",
               "Delivery days chosen per area at P = 0.25, everyone taking "
               "part, grouped by settlement type.")
    y = B.table(s, ["Type of area", "How many", "Typical",
                    "Average", "On two days a week"],
                [[("Countryside", "key"), "118", ("2 days", "num"), "2.58",
                  ("60.2 %", "num")],
                 [("Suburb", "key"), "124", "3 days", "3.12", "29.0 %"],
                 [("City", "key"), "70", "3 days", "3.77", ("4.3 %", "num")]],
                BODY_T + 0.20, widths=[2.8, 1.6, 2.4, 2.0, 3.0], reserve=2.5)
    for i, (nm, pct, col) in enumerate([("countryside", 60.2, S6),
                                        ("suburb", 29.0, S4),
                                        ("city", 4.3, S1)]):
        yy = y + 0.30 + i * 0.52
        txt(s, L, yy, 1.9, 0.36, nm, 20, color=INK2)
        rect(s, L + 2.0, yy + 0.03, 8.6 * pct / 100.0, 0.30, col)
        txt(s, L + 2.15 + 8.6 * pct / 100.0, yy, 1.6, 0.36, f"{pct:.1f} %", 20,
            bold=True, color=INK)
    txt(s, L, y + 1.90, W, 0.48,
        "Fourteen rural areas drop to two days for every one city area that "
        "does.", 22, bold=True, color=RED)


# ═══════════════════════════════════════════════════════════════════════════
# part 4 · why some areas gain nothing
# ═══════════════════════════════════════════════════════════════════════════
def block_where(prs):
    chapter(prs, "Why dense areas gain nothing",
            "The three things that decide whether an area profits")

    s = xslide(prs, "where", "Part 4 · Why dense areas gain nothing",
               "What makes one area profit and its neighbour not?",
               "How strongly each property moves together with the saving, "
               "across all 312 areas. +1 would be a perfect match, −1 a "
               "perfect opposite.")
    y = B.table(s, ["Property of the area", "Link", "What it means"],
                [[("Far from the depot", "key"), ("+0.53", "num"),
                  "Every tour drives out and back. Bundling spreads that trip "
                  "over more parcels"],
                 [("Large area", "key"), ("+0.31", "num"),
                  "Long tours inside the area. Merging days makes them denser"],
                 [("Many parcels per address", "key"), ("−0.72", "num"),
                  "The tour is already full and short. There is nothing left "
                  "to win"]],
                BODY_T + 0.20, widths=[3.4, 1.6, 7.0], reserve=2.4)
    label_box(s, L, y + 0.28, W, 0.85, H.BLUSH,
              [("The strongest link is the negative one: bundling pays where "
                "density is missing.", 22, True, RED)], line_col=RED)
    vbullets(s, ["The typical rural area saves 25 %, the typical city area 9 %.",
                 "That gap is geography, not modelling — which is why a "
                 "uniform rollout wastes it."],
             y + 1.26)


# ═══════════════════════════════════════════════════════════════════════════
# part 5 · can we trust the numbers
# ═══════════════════════════════════════════════════════════════════════════
def block_valid(prs):
    chapter(prs, "Can the numbers be trusted",
            "What happens when a real routing solver checks the answer")

    # This slide used to restate the SUBMISSION's realised savings
    # (22.8 -> 23.7 % and so on) under a Part-A stamp promising a Part-B fix.
    # It is rebuilt here on the validation actually in use: predicted against
    # actual COST on the same tours, which is what a validation without a
    # solved theta = 0 baseline can honestly show.
    _vf = _RV.validation_facts() if _F is not None else None
    s = xslide(prs, "valid", "Part 5 · Can we trust it",
               "What happens when a real routing solver checks the answer?",
               (f"{_vf['n']} routing instances recomputed from scratch with "
                f"VROOM/Valhalla on {_vf['grid']}, both plans at θ = 100 %. "
                f"Costs, not savings: a realised SAVING needs a solved "
                f"θ = 0 baseline, which this validation has not reached yet."
                if _vf else
                "Four settings recomputed from scratch with VROOM/Valhalla on "
                "1 248 observations the model had never seen."))
    if _vf is not None:
        # The vintage is the validation's, not the submission's, so the banner
        # says which validation and what is still outstanding.
        _RV.stamp(s, text=_RV.VALIDATION_STAMP)
        _RV.notes(s, _RV.validation_notes(_vf), cite=_RV.cites())
    if _vf is not None:
        _rows = [[(r[0], "key"), r[1], r[2], r[3], (r[4], "good")]
                 for r in _RV.validation_rows(_vf)]
        _head = ["Setting", "Plan", "The model priced", "The solver charged",
                 "Model above solver"]
        _w = [1.8, 2.9, 2.6, 2.6, 2.5]
    else:
        _rows = [[("No fee", "key"), "22.8 %", ("23.7 %", "good"),
                  "+0.9 points"],
                 [("P = 0.25", "key"), "18.5 %", ("19.8 %", "good"),
                  "+1.3 points"],
                 [("P = 0.5", "key"), "13.5 %", ("15.6 %", "good"),
                  "+2.1 points"],
                 [("P = 0.75", "key"), "10.2 %", ("13.0 %", "good"),
                  "+2.8 points"]]
        _head = ["Setting", "The model promised", "The solver delivered",
                 "Difference"]
        _w = [3.0, 3.0, 3.4, 2.4]
    y = B.table(s, _head, _rows, BODY_T + 0.20, widths=_w, reserve=2.4)
    vbullets(s, ([f"The model prices the same tours "
                  f"{_vf['gap_lo']:.1f}–{_vf['gap_hi']:.1f} % above the "
                  f"solver at every point, in both plans.",
                  f"So it is wrong in the direction that is safe: it "
                  f"under-promises. The fleet count holds too — "
                  f"{_vf['peak_pred']} peak vehicles predicted against "
                  f"{_vf['peak_actual']} actual at P = 0.",
                  "What is not here is a realised saving percentage: that "
                  "needs the θ = 0 baseline solved, and it is not yet."]
                 if _vf else
                 ["Every one of the four came out better than promised.",
                  "So the model is wrong in the direction that is safe: it "
                  "under-promises.",
                  "What this does not prove is that the routing solver matches "
                  "the real street."]),
             y + 0.28)


# ═══════════════════════════════════════════════════════════════════════════
# part 6 · the seven carriers
# ═══════════════════════════════════════════════════════════════════════════
def block_providers(prs):
    chapter(prs, "The seven carriers, side by side",
            "Why one average curve describes none of them")
    for key, subject, fig_name, src, bullets in [
        ("mix", "Do all seven carriers behave the same? No.",
         "figP1_mix_by_provider",
         "Delivery days chosen by each carrier's areas as participation rises, "
         "at P = 0.25.",
         ["DHL is the odd one out: most of its network keeps five or six "
          "delivery days.",
          "DPD, GLS and FedEx put roughly half their areas on two days.",
          "The single curve in the talk is the average of these two — and "
          "describes neither."]),
        ("range", "How much can each of them actually save?",
         "figP2_saving_by_provider",
         "Each carrier measured against its own daily-delivery cost, not "
         "against the regional total.",
         [(f"The best DHL can do is {_MAXSAV('DHL'):.1f} %. GLS reaches "
           f"{_MAXSAV('GLS'):.1f} % (operator lens)."
           if _F is not None else
           "The best DHL can do is far less than GLS."),
          "The shape is the same everywhere — only the ceiling differs.",
          "So one shared operating point is a compromise for everybody."]),
        ("where", "Does each carrier gain in the same places?",
         "figP3_map_saving_provider",
         "Saving per area at P = 0.25 with everyone taking part. Same colour "
         "scale for all seven.",
         ["Typical saving runs from 3.7 % for DHL to 32.6 % for FedEx.",
          "But the pattern is identical: the outside gains, the centre does "
          "not.",
          "What differs is how much of each network sits on the outside."]),
        ("maps", "And do they end up with the same delivery days?",
         "figP4_map_freq_provider",
         "Typical delivery days per area at P = 0.25, everyone taking part.",
         ["DHL holds five days; DPD, GLS and FedEx drop to two.",
          "Every carrier keeps a high frequency in the dense centre.",
          "How often you can deliver is a property of the network, not of the "
          "region."]),
    ]:
        sl = xslide(prs, key, "Part 6 · The seven carriers", subject, src)
        pic(sl, FIG / f"{fig_name}.png", L, BODY_T + 0.24, W, 3.35)
        vbullets(sl, bullets, BODY_T + 3.75)


# ═══════════════════════════════════════════════════════════════════════════
# each results figure, in full, once per carrier
# ═══════════════════════════════════════════════════════════════════════════
# The small multiples in block_providers compare the seven carriers at one
# operating point. This section gives each carrier the *whole* figure, so a
# question about one carrier can be answered with the same picture the talk
# already showed, just for that network. 4 families x 7 carriers = 28 slides.

CARRIERS = ["DHL", "Amazon", "Hermes", "UPS", "DPD", "GLS", "FedEx"]

FAMILIES = [
    ("mix", "figQ1_mix", "How often each area is served",
     "Eight panels, one per service fee. Left to right in each panel: more "
     "customers willing to wait. Dark blue = two delivery days a week, red = "
     "still daily.",
     "Read it as: how much of this network leaves daily delivery as the "
     "fee and the participation change."),
    ("range", "figQ2_saving", "What it saves, fee by fee",
     "Every square is one setting: a service fee (down the side) and a "
     "participation level (across the bottom). The number is the weekly cost "
     "saving against this carrier delivering daily.",
     "Read it as: bright is good. The bright corner is where this carrier "
     "would want to operate."),
    ("maps", "figQ3_freqmap", "Where the delivery days go",
     "The same region four times, as more customers join in. Colour is how "
     "often that area gets served.",
     "Read it as: the outside of the region turns dark first — that is where "
     "delivery days are dropped."),
    ("where", "figQ4_savingmap", "Where the money is saved",
     "The same region six times, one per service fee. Darker green = more "
     "saving in that area. The number above each map is what the whole "
     "network saves.",
     "Read it as: the greener the ring around the city, the more this carrier "
     "gains outside the core."),
]


def block_carrier_full(prs):
    chapter(prs, "Every figure, one carrier at a time",
            "The same four pictures, redrawn for each of the seven")
    for key, stem, subject, source, how in FAMILIES:
        for carrier in CARRIERS:
            path = FIG / f"{stem}_{carrier}.png"
            if not path.exists():
                print(f"  ! missing {path.name}")
                continue
            s = xslide(prs, key, f"Backup: {carrier}", subject, source)
            # The figure is the slide; give it the whole body and keep the
            # reading aid to one line underneath.
            pic(s, path, L, BODY_T + 0.10, W, 4.45)
            hh = H.text_height(how, W, 20, 1.22) + 0.06
            txt(s, L, BODY_T + 4.66, W, hh, how, 20, color=INK2,
                line=1.22)



# ═══════════════════════════════════════════════════════════════════════════
# keeping hand-placed figures current
# ═══════════════════════════════════════════════════════════════════════════
# The working deck's own slides carry figures pasted in by hand, so a
# re-rendered figure does not reach them. This walks the inherited slides,
# checks each large picture against what is on disk now, and swaps the stale
# ones in place -- same position, same size, same z-order. Slides are found by
# a fragment of their title, not by index, so the deck can keep being edited.
#
# Anything large that matches no current render and is not listed here is
# reported rather than touched: it may be a photograph or a paper figure.

REFRESH = {
    # title fragment -> the renders that slide should carry, left to right.
    # The raster match handles a figure that was re-rendered at the same size;
    # this map is for the ones whose SHAPE changed, which the raster match
    # cannot see. fig62 is exactly that case: it went from a 4-bar single-plan
    # figure (2137x1141) to a 6-bar two-plan one (4113x1141) when the v2
    # validation started covering both plans, so the old render sat on the
    # deck's "Validation" slide looking current.
    "service improves faster": ["fig31_saving_grid", "fig32_wait_grid"],
    "validation": ["fig61_vroom_scatter", "fig62_pred_vs_actual"],
}
MIN_FIG_IN = 3.0


def _current_renders():
    import hashlib
    out = {}
    for d in ("slides/tierA", "slides/tierB"):
        for f in (H.ROOT / "results" / "presentation_2026_08" / d).glob("*.png"):
            out[hashlib.md5(f.read_bytes()).hexdigest()] = f
    return out


def _render_index():
    """Every current render, indexed by md5 AND by pixel size.

    Pixel size is what makes the swap content-matched rather than
    title-matched: a figure is re-rendered from the same `figsize` at the same
    dpi, so its raster dimensions are stable across grids while its pixels are
    not. An embedded picture whose dimensions match exactly one current render
    IS that figure, whatever slide it sits on and whatever the title says.
    """
    import hashlib
    from PIL import Image
    by_md5, by_size = {}, {}
    for f in _current_renders().values():
        by_md5[hashlib.md5(f.read_bytes()).hexdigest()] = f
        with Image.open(f) as im:
            by_size.setdefault(im.size, []).append(f)
    return by_md5, by_size


def refresh_figures(prs, n_inherited):
    """Replace stale hand-placed figures with the current render.

    Three ways to identify what an embedded picture is, in order of strength:

    1. its md5 is already a current render -- nothing to do;
    2. its raster dimensions match exactly ONE current render -- that is the
       figure, swap it (content match, title-independent);
    3. the slide's title is in `REFRESH` -- use that mapping (the old
       behaviour, kept for figures whose dimensions are ambiguous).

    Anything none of the three identifies is left untouched and LISTED, with
    its size and slide, because silently leaving a stale figure in a deck named
    for the new grid is the failure this whole pass exists to prevent.
    """
    import hashlib
    from PIL import Image
    from pptx.util import Emu
    by_md5, by_size = _render_index()
    by_name = {f.stem: f for f in _current_renders().values()}
    swapped = stale = fresh_n = 0
    unmatched = []
    for i, sl in enumerate(list(prs.slides)[:n_inherited], 1):
        ph = [sh for sh in sl.shapes if sh.is_placeholder
              and sh.placeholder_format.idx == 0 and sh.has_text_frame]
        title = (ph[0].text_frame.text if ph else "").lower()
        want = next((v for k, v in REFRESH.items() if k in title), None)
        pics = [sh for sh in sl.shapes if "PICTURE" in str(sh.shape_type)
                and sh.width is not None
                and Emu(sh.width).inches >= MIN_FIG_IN]
        pics.sort(key=lambda sh: Emu(sh.left).inches)
        for j, sh in enumerate(pics):
            blob = sh.image.blob
            if hashlib.md5(blob).hexdigest() in by_md5:
                fresh_n += 1
                continue
            target, how = None, ""
            try:
                import io
                with Image.open(io.BytesIO(blob)) as im:
                    size = im.size
                cands = by_size.get(size, [])
                if len(cands) == 1:
                    target, how = cands[0], "content match"
                elif len(cands) > 1 and want and j < len(want):
                    hit = [c for c in cands if c.stem == want[j]]
                    if hit:
                        target, how = hit[0], "content + title"
            except Exception:
                pass
            if target is None and want and j < len(want):
                target, how = by_name.get(want[j]), "title map"
            if target is None:
                stale += 1
                unmatched.append(
                    f"slide {i}: {Emu(sh.width).inches:.1f} in picture, "
                    f"raster {size if 'size' in dir() else '?'}")
                continue
            box = (sh.left, sh.top, sh.width, sh.height)
            el = sh._element
            parent = el.getparent()
            idx = list(parent).index(el)
            parent.remove(el)
            new_pic = sl.shapes.add_picture(str(target), *box)
            parent.remove(new_pic._element)
            parent.insert(idx, new_pic._element)      # keep the z-order
            swapped += 1
            print(f"  ~ slide {i}: {target.stem} ({how})")
    print(f"  figures: {swapped} swapped to the current render, "
          f"{fresh_n} already current, {stale} unmatched")
    for u in unmatched:
        print(f"  ? {u} matches no current render and is not mapped "
              f"-- left as it is")
    return swapped


# ═══════════════════════════════════════════════════════════════════════════
# restating the hand-authored body on the revision grid
# ═══════════════════════════════════════════════════════════════════════════
# `SRC` is the author's own working deck. Its results slides state the
# submission's headline -- 22.8 / 18.5 / 13.5 % routing saving -- as current
# fact, and this revision supersedes exactly those numbers: the same
# routing-optimal plan is 20.4 % routing and MINUS 7.8 % operator at P = 0.
# Appending a withdrawal chapter and shipping the file as `_rev2026-08` while
# its body still argues the old numbers would be worse than not shipping it.
#
# So the body is restated in place, on the grid, and everything that cannot be
# restated is stamped. Two rules:
#
#   * a slide whose numbers HAVE a revision counterpart is rewritten to it and
#     says which plan and which lens it is in, and gets the provisional chip;
#   * a slide whose numbers have NO counterpart yet -- the VROOM validation,
#     whose revision re-run is still being produced -- keeps them and gets a
#     visible stamp, not a footnote.
#
# Slides are matched by a text fragment, never by index: 96_ inserts slides of
# its own and the working deck is edited between builds, so an index silently
# starts pointing at the wrong slide (the same reason `resolve_targets()`
# matches by title).
def _para_clone_format(dst_p, src_p):
    """Copy paragraph-level properties (alignment, spacing) between paragraphs."""
    dst_p.alignment = src_p.alignment
    if src_p.line_spacing is not None:
        dst_p.line_spacing = src_p.line_spacing
    if src_p.space_after is not None:
        dst_p.space_after = src_p.space_after


def _retext(tf, text):
    """Replace a text frame's content, keeping the first run's formatting.

    Hand-made slides carry their type in the runs, so the safe way to restate
    one is to keep run 0 of each paragraph and drop the rest -- not to rebuild
    the frame, which would lose the author's sizes and colours.
    """
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    if not p0.runs:
        tf.text = text
        return
    template = p0.runs[0]
    while len(tf.paragraphs) > 1:
        tf._txBody.remove(tf.paragraphs[-1]._p)
    for r in list(p0.runs)[1:]:
        r._r.getparent().remove(r._r)
    p0.runs[0].text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        _para_clone_format(p, p0)
        r = p.add_run()
        r.text = line
        r.font.name = template.font.name
        r.font.size = template.font.size
        r.font.bold = template.font.bold
        try:
            if template.font.color and template.font.color.type is not None:
                r.font.color.rgb = template.font.color.rgb
        except Exception:
            pass


def _swap(tf, pairs):
    """Run-level string substitution: formatting survives, numbers change."""
    hit = False
    for para in tf.paragraphs:
        for run in para.runs:
            for old, new in pairs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    hit = True
    return hit


def _shape_by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    return None


def _find(prs, fragment, limit=None):
    """The first slide whose text contains `fragment`, searched in order."""
    for i, sl in enumerate(prs.slides, 1):
        if limit and i > limit:
            break
        for sh in sl.shapes:
            if sh.has_text_frame and fragment in sh.text_frame.text:
                return i, sl
            if getattr(sh, "has_table", False) and sh.has_table:
                if any(fragment in c.text for r in sh.table.rows
                       for c in r.cells):
                    return i, sl
    return None, None


def _has(slide, fragment):
    for sh in slide.shapes:
        if sh.has_text_frame and fragment in sh.text_frame.text:
            return True
    return False


def stamp(slide):
    """The submission-era banner. One implementation, in `_revision`."""
    _RV.stamp(slide)
    return slide


def _table_rows(tbl, rows, *, header=None):
    """Rewrite a table's cells in place, keeping each cell's own formatting."""
    if header is not None:
        for ci, txt_ in enumerate(header):
            if ci < len(tbl.columns):
                _retext(tbl.cell(0, ci).text_frame, txt_)
    for ri, row in enumerate(rows, start=1):
        if ri >= len(tbl.rows):
            break
        for ci, txt_ in enumerate(row):
            if ci < len(tbl.columns):
                _retext(tbl.cell(ri, ci).text_frame, txt_)


def restate_body(prs, limit):
    """Put the hand-authored results slides on the revision grid.

    `limit` is how many leading slides belong to the copied working deck; the
    appended backup is generated here and is already on the revision grid.
    Returns the set of slide indices that were restated, so the sweep that
    follows does not stamp them as well.
    """
    f, RV = _F, _RV
    h = f.headline
    done = set()

    def touched(i, sl, note, cite):
        done.add(i)
        RV.notes(sl, note, cite=cite)
        RV.provisional(sl, enabled=_TAG)

    LENS_LINE = ("Revision grid, θ = 100 %. “Routing” = what the tours cost to "
                 "drive; “operator” = kilometres plus 1 134.90 € per peak "
                 "vehicle and hub. Rows marked stage 2 are the "
                 "operator-polished plan.")

    # ── the saving/wait headline ─────────────────────────────────────────
    i, sl = _find(prs, "Saving peaks at 22.8%", limit)
    if sl is not None:
        _retext(_shape_by_id(sl, 7).text_frame,
                f"Routing-optimal plan: {h[0.0]['rout1']:.1f} % routing saving "
                f"at P = 0, wait {h[0.0]['wait1']:.2f} d — but −"
                f"{abs(h[0.0]['op1']):.1f} % in the operator lens\n"
                f"Operator-polished plan: {h[0.0]['op2']:.1f} % operator "
                f"saving, {h[0.0]['rout2']:.1f} % routing, wait "
                f"{h[0.0]['wait2']:.2f} d, peak fleet "
                f"{h[0.0]['peak2_pct']:+.0f} %")
        _retext(_shape_by_id(sl, 4).text_frame,
                f"Revision grid, complete θ grid. Baselines: "
                f"{RV.eur(f.base_routing)} € routing and "
                f"{RV.eur(f.base_operator)} € operator per week. Figure "
                f"rendered from the submission grid — the shape holds, the "
                f"levels are restated here.")
        touched(i, sl, RV.plan_notes(f), RV.cites("§40.12"))

    # ── the efficiency-range table ───────────────────────────────────────
    i, sl = _find(prs, "the cost-optimal extreme", limit)
    if sl is not None:
        tbl = next(sh.table for sh in sl.shapes
                   if getattr(sh, "has_table", False) and sh.has_table)
        _table_rows(tbl, [
            [f"P = {P:g}", f"{h[P]['rout2']:.1f} %", f"{h[P]['op2']:.1f} %",
             f"{h[P]['wait2']:.2f} d"] for P in (0.0, 0.25, 0.5)],
            header=["PENALTY", "ROUTING SAVING", "OPERATOR SAVING",
                    "ADDED WAIT"])
        touched(i, sl,
                "All three rows are the operator-polished plan (stage 2), "
                "priced in both lenses. The submission's single 'cost saving' "
                "column was routing euro on the routing-optimal plan and is "
                "superseded twice over: by the universal tour rule and by the "
                "operator polish.", ["§40.15", "§40.18"])

    # ── the conclusion's "what we found" line (two slides use it) ────────
    found_line = (f"{h[0.5]['op2']:.1f} to {h[0.25]['op2']:.1f} % cheaper "
                  f"per week for an operator, at under half a day of added "
                  f"waiting.")
    for _ in range(2):
        i, sl = _find(prs, "13.5 to 18.5 % cheaper per week", limit)
        if sl is None:
            break
        _retext(_shape_by_id(sl, 17).text_frame, found_line)
        for sid, new in ((20, f"{h[0.25]['peak2_pct']:+.1f} %"),
                         (21, "peak fleet at P = 0.25, and 83 % less weekday "
                              "variation"),
                         (33, f"DHL tops out at {_MAXSAV('DHL'):.1f} %, GLS at "
                              f"{_MAXSAV('GLS'):.1f} % (operator lens)"),
                         (22, "+0.9…2.8 pp"),
                         (23, "the real solver beat every prediction "
                              "(submission grid)")):
            sh = _shape_by_id(sl, sid)
            if sh is not None and sh.has_text_frame:
                _retext(sh.text_frame, new)
        sh = _shape_by_id(sl, 3)
        if sh is not None and sh.has_text_frame:
            _retext(sh.text_frame,
                    f"Revision grid, operator lens on the stage-2 plan · "
                    f"Region Hannover, seven carriers, 1.26 M parcels a week "
                    f"against a {RV.eur(f.base_operator)} € operator "
                    f"daily-delivery baseline.")
        touched(i, sl, "The conclusion now quotes the OPERATOR lens on the "
                       "operator-polished plan, which is the lens an LSP with "
                       "salaried drivers faces. The routing-lens figure for "
                       "the same range is 13.2–17.1 %.",
                ["§40.12", "§40.15"])

    # ── the closing statement ────────────────────────────────────────────
    i, sl = _find(prs, "of a weekly delivery bill, verified against real "
                       "routing", limit)
    if sl is not None:
        _retext(_shape_by_id(sl, 7).text_frame,
                f"{h[0.5]['op2']:.1f} – {h[0.25]['op2']:.1f} % of an "
                f"operator's weekly delivery bill.")
        touched(i, sl,
                "The submission's closing line said '13.5–18.5 %, verified "
                "against real routing'. The range is restated in the operator "
                "lens, and the verification clause is dropped: the revision's "
                "VROOM re-run of both plans is still being produced, so no "
                "revision number is solver-verified yet.",
                ["§40.15", "§40.18"])

    # ── "what a fee buys" ────────────────────────────────────────────────
    i, sl = _find(prs, "SAVING GIVEN UP", limit)
    if sl is not None:
        tbl = next(sh.table for sh in sl.shapes
                   if getattr(sh, "has_table", False) and sh.has_table)
        base_op = h[0.0]["op2"]
        rows = []
        for P in (0.0, 0.25, 0.5, 0.75, 1.0, 2.0):
            x = h[P]
            given = "—" if P == 0 else f"{base_op - x['op2']:.1f} points"
            removed = ("—" if P == 0 else
                       f"{(1 - x['wait2'] / h[0.0]['wait2']) * 100:.0f} %")
            rows.append([("none" if P == 0 else f"P = {P:g}"),
                         f"{x['op2']:.1f} %", f"{x['wait2']:.2f} d", given,
                         removed])
        _table_rows(tbl, rows,
                    header=["FEE", "OPERATOR SAVING", "PEOPLE WAIT",
                            "SAVING GIVEN UP", "WAITING REMOVED"])
        sh = _shape_by_id(sl, 3)
        if sh is not None:
            _retext(sh.text_frame,
                    f"Everyone taking part, operator-polished plan. Saving "
                    f"against the {RV.eur(f.base_operator)} € weekly operator "
                    f"baseline; waiting averaged over all parcels, including "
                    f"those that never wait.")
        sh = _shape_by_id(sl, 6)
        if sh is not None:
            _retext(sh.text_frame,
                    "The first step is still the bargain: give up under two "
                    "points of operator saving and half the waiting "
                    "disappears.\nEvery step after that buys less and costs "
                    "more.")
        touched(i, sl, LENS_LINE, ["§40.15", "§40.18"])

    # ── "why not simply take the biggest saving" ─────────────────────────
    i, sl = _find(prs, "the cheapest week — but a full day of waiting", limit)
    if sl is not None:
        for sid, P in ((7, 0.0), (12, 0.25), (17, 0.5)):
            sh = _shape_by_id(sl, sid)
            if sh is not None:
                _retext(sh.text_frame, f"{h[P]['op2']:.1f} %")
        for sid, P in ((8, 0.0), (13, 0.25), (18, 0.5)):
            sh = _shape_by_id(sl, sid)
            if sh is not None:
                _retext(sh.text_frame, f"+{h[P]['wait2']:.2f} d waiting")
        for sid, new in (
                (9, "the cheapest week for an operator — but the longest wait"),
                (14, "half the waiting, and the same peak-fleet cut"),
                (19, f"peak fleet down {abs(h[0.5]['peak2_pct']):.0f} %, "
                     f"weekday swings down 67 %")):
            sh = _shape_by_id(sl, sid)
            if sh is not None:
                _retext(sh.text_frame, new)
        sh = _shape_by_id(sl, 3)
        if sh is not None:
            _retext(sh.text_frame,
                    "Operator lens, operator-polished plan, revision grid. "
                    "Fleet figures are the sum of hub peaks against the "
                    "daily-delivery baseline of the same grid.")
        touched(i, sl,
                "In the operator lens P = 0 is still the largest number "
                "(24.7 % against 22.8 %), but it costs 0.77 d of waiting "
                "against 0.39 d for the same peak-fleet cut — which is why "
                "the knee stays at P = 0.25.", ["§40.15", "§40.18"])

    # ── the "Validation" slide: current facts, current picture ───────────
    # It used to carry the submission's realised-saving pairs (22.8 -> 23.7 %)
    # with no disclosure at all, beside a stale render of fig62 that the
    # raster match could not recognise. Both are replaced here; the picture is
    # swapped by `refresh_figures()` through the REFRESH map above.
    i, sl = _find(prs, "realized saving above prediction", limit)
    if sl is not None:
        vf = RV.validation_facts()
        for sid, new in (
                (6, f"+{vf['gap_lo']:.1f} … +{vf['gap_hi']:.1f} %"),
                # The author's own label boxes are one line tall at his size,
                # so a replacement has to stay about as short as what it
                # replaces -- roughly 30 characters, not 47.
                (7, "model above solver, same tours"),
                (9, f"{vf.get('mape', float('nan')):.2f} %"),
                (10, "cost error against the solver"),
                (12, f"{vf['peak_pred']} vs {vf['peak_actual']}"
                     if "peak_pred" in vf else f"{vf.get('r2', 0):.3f}"),
                (13, "peak fleet, pred vs actual at P = 0"
                     if "peak_pred" in vf else
                     f"R² over {vf.get('n_clean', 0)} instances")):
            sh = _shape_by_id(sl, sid)
            if sh is not None and sh.has_text_frame:
                _retext(sh.text_frame, new)
        _RV.stamp(sl, text=_RV.VALIDATION_STAMP)
        done.add(i)
        RV.notes(sl, RV.validation_notes(vf), cite=RV.cites())
        print(f"  restated the Validation slide ({i}) on "
              f"{vf['grid']}/validation")

    # ── the backup validation table has no revision counterpart yet ──────
    i, sl = _find(prs, "THE SOLVER DELIVERED", limit)
    if sl is not None:
        stamp(sl)
        done.add(i)
        RV.notes(sl, "Kept as it stands: this is the SUBMISSION grid's VROOM "
                     "validation, and the revision's re-run of both plans is "
                     "still being produced. The direction of the error -- the "
                     "surrogate under-promises -- is what carries over; the "
                     "levels do not.", cite="§40.18")

    return done


# Numbers that can ONLY be the submission's. Deliberately not 22.8 / 18.5 %:
# those are also the revision's operator-lens savings at P = 0.25 and P = 0.5,
# so matching them stamps slides that were just restated. The fingerprints
# below are the submission's routing-lens wait (0.98 d), its baseline, its
# peak-fleet and CV figures, and its four VROOM-validated realised savings.
_SUBMISSION_ONLY = (r"0[.,]98\s?d|0[.,]46\s?d|1[  ]909[  ]748|"
                    r"12[.,]9\s?%|23[.,]7\s?%|19[.,]8\s?%|15[.,]6\s?%|"
                    r"13[.,]0\s?%|10[.,]2\s?%")


def stamp_remaining(prs, limit, done):
    """Stamp every slide still showing a submission-only number.

    A safety net for what the restatement missed, not the mechanism: a slide
    that carries the provisional chip was built or restated on the revision
    grid and is skipped, and `stamp()` is idempotent, so a slide stamped where
    it was made is not stamped twice.
    """
    import re
    pat = re.compile(_SUBMISSION_ONLY)
    n = 0
    for i, sl in enumerate(prs.slides, 1):
        if i > limit or i in done:
            continue
        if _has(sl, _RV.tag_text()) or _RV.stamped(sl):
            continue
        text = []
        for sh in sl.shapes:
            if sh.has_text_frame:
                text.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                text += [c.text for r in sh.table.rows for c in r.cells]
        if pat.search("\n".join(text)):
            stamp(sl)
            n += 1
            print(f"  stamped slide {i}: submission-era number, no revision "
                  f"counterpart applied")
    return n


# ═══════════════════════════════════════════════════════════════════════════
def build(out: Path) -> Path:
    prs = Presentation(str(SRC))
    n_before = len(prs.slides)
    _RESOLVED.update(resolve_targets(prs))
    refresh_figures(prs, n_before)

    # The closing summary belongs in the talk, not in the backup: build it,
    # then move it in front of the contact slide, which is the last one.
    CONC.slide_conclusion(prs)
    CONC.slide_takeaway(prs)
    contact_at = n_before - 1                    # 0-based index of the last
    H.move_slide(prs, len(prs.slides) - 2, contact_at)
    H.move_slide(prs, len(prs.slides) - 1, contact_at + 1)
    print(f"  conclusion + takeaway inserted as slides "
          f"{contact_at + 1}–{contact_at + 2}, before the contact slide")

    # The hand-authored body states the submission headline as current fact.
    # Restate what the revision supersedes, stamp what it cannot, and only
    # then append the backup -- so the deck is never internally inconsistent.
    done = set()
    if _F is not None:
        body = len(prs.slides)
        done = restate_body(prs, body)
        print(f"  restated {len(done)} body slide(s) on {_RV.D.REV.name}")

    divider(prs, "B", "Backup", "Why the results\nlook like this",
            "Seven parts, each answering the questions a close reader asks — "
            "in the order they come up")
    block_contents(prs)
    if _F is None:
        # --no-revision rebuilds the submission-era deck, which is the one
        # legitimate reason to draw the withdrawn illustrations.
        DUM.allow_withdrawn()
        block_mix(prs)     # 1 · the odd thing in the frequency picture
    else:
        block_mix_rev(prs)  # 1 · withdrawn: it explained a pricing artefact
    block_trade(prs)       # 2 · what a fee actually buys
    block_maps(prs)        # 3 · where it happens first
    block_where(prs)       # 4 · why dense areas gain nothing
    block_valid(prs)       # 5 · can the numbers be trusted
    block_providers(prs)   # 6 · the seven carriers, side by side
    block_carrier_full(prs)  # 7 · every figure, one carrier at a time
    fill_contents()
    # The safety net, over the WHOLE deck and after everything is on it: any
    # slide still carrying a submission-era headline number that nothing above
    # restated gets a visible stamp. It is deliberately a net and not the plan
    # -- if it ever fires on a slide nobody expected, that slide needs porting,
    # not stamping.
    if _F is not None:
        n = stamp_remaining(prs, len(prs.slides), done)
        print(f"  {n} slide(s) stamped: submission-era number, no revision "
              f"counterpart")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"  {n_before} original slides kept, "
          f"{len(prs.slides) - n_before} appended")
    return out


def audit_v5() -> int:
    """Recompute every number the RESTATED slides claim; fail if one moved.

    The expectations are not literals in this file: they come from
    `_revision.GRID_EXPECT[<grid>]`, whose entries record what the compendium
    (v5) or the deep dive (v6) states. So this audit compares the deck against
    a written-down source, not against the grid it was built from -- which
    would be circular. `Facts.load()` runs first and carries its own asserts on
    the headline, the knee set and the break-even band.

    The submission-era audit is still here, behind `--submission`.
    """
    import _data as D
    import _revision as RV
    if D.SCHEMA != D.SCHEMA_V2:
        print(f"audit needs a v2 grid; REV is {D.REV} ({D.SCHEMA})",
              file=sys.stderr)
        return 1
    f = RV.Facts.load()          # every assert in here is part of the audit
    e = RV.expect()
    if not e:
        print(f"{D.REV.name} has no recorded expectation in "
              f"_revision.GRID_EXPECT, so there is nothing to audit the "
              f"slides against. Record its values first.", file=sys.stderr)
        return 1
    print(f"  auditing {D.REV.name} against {e['source']}")
    bad = []

    def check(name, got, want, tol):
        ok = abs(got - want) <= tol
        print(f"  {'ok ' if ok else 'FAIL'} {name}: {got:.4g} "
              f"(source says {want:g})")
        if not ok:
            bad.append(name)

    a = e.get("audit", {})
    for P, want in a.get("operator saving, stage-2 plan", {}).items():
        check(f"operator saving at P={P:g}, theta=1", f.headline[P]["op2"],
              want, 0.02)
    for P, want in a.get("wait, stage-2 plan", {}).items():
        check(f"wait at P={P:g}, stage-2 plan", f.headline[P]["wait2"],
              want, 0.006)
    for (P, th), want in a.get("areas consolidating, stage-1 plan", {}).items():
        check(f"areas consolidating at P={P:g}, theta={th:g}",
              f.consolidating[(P, th)]["plan1"], want, 0.06)
    # the headline pair and the peak, from the same source
    check("routing saving at P=0, stage-1 plan", f.headline[0.0]["rout1"],
          e["rout1"], 0.02)
    check("operator saving at P=0, stage-1 plan", f.headline[0.0]["op1"],
          e["op1"], 0.02)
    check("peak fleet at P=0, stage-2 plan [%]", f.headline[0.0]["peak2_pct"],
          e["peak2"], 0.02)
    # the one-area depot count the slides state, derived from the hub table
    check("DHL one-area depots", f.one_cell_hubs, 8, 0)
    check("DHL depots", f.dhl_hubs, 16, 0)
    # the discount line must name a winner only where there is one
    opt = RV.discount_optima(f)
    line = RV.discount_optimum_line(f)
    for lens in ("operator", "routing"):
        P, net, runner, margin = opt[lens]
        tie = margin < RV.DISCOUNT_TIE_PP
        want = (f"level in the {lens} lens" if tie
                else f"P = {P:g} in the {lens} lens")
        ok = want in line
        print(f"  {'ok ' if ok else 'FAIL'} discount line, {lens} lens: "
              f"{'tie' if tie else 'winner'} at {margin:.3f} pp margin")
        if not ok:
            bad.append(f"discount line/{lens}")
    print("\n" + ("AUDIT FAILED: " + ", ".join(bad) if bad else "audit clean"))
    return 1 if bad else 0


def audit_submission() -> int:
    """The SUBMISSION body's numbers, on the submission grid.

    Kept as the record of what `block_mix()` and the pre-revision results
    slides asserted. It only means anything against
    `PRES_REV_DIR=results/revision_2026_07`; on the revision grid every one of
    these is expected to have moved, which is why the deck no longer states
    them and why this is not the default audit any more.
    """
    import _data as D
    import _stage3_common as C
    bad = []

    def check(name, got, want, tol):
        ok = abs(got - want) <= tol
        print(f"  {'ok ' if ok else 'FAIL'} {name}: {got:.4g} "
              f"(slide says {want:g})")
        if not ok:
            bad.append(name)

    s = D.load_chosen_stage3()
    col = "schedule_size_system_smoothed"
    for pen, th, want in ((5.0, 0.1, 49.7), (5.0, 0.2, 42.9), (5.0, 0.3, 14.7),
                          (10.0, 0.1, 41.7), (10.0, 0.2, 10.9),
                          (0.0, 0.1, 87.5)):
        sub = s[np.isclose(s.penalty, pen) & np.isclose(s.share_willing, th)]
        check(f"consolidating share P={pen:g} θ={th:g}",
              100 * (sub[col] < 6).mean(), want, 0.15)

    check("B2B willing at θ=0.1", 100 * C._willing_b2b(0.1), 45.8, 0.1)
    check("B2C willing at θ=0.1", 100 * C._willing_b2c(0.1), 0.1, 0.05)
    check("B2B willing at θ=0.2", 100 * C._willing_b2b(0.2), 72.4, 0.1)
    check("B2C willing at θ=0.2", 100 * C._willing_b2c(0.2), 5.5, 0.1)

    rt = D.load_raumtyp()
    sub = (s[np.isclose(s.penalty, 0.25) & np.isclose(s.share_willing, 1.0)]
           .merge(rt, on="plz", how="left"))
    for cls, want in (("rural", 60.2), ("suburban", 29.0), ("urban", 4.3)):
        m = sub[sub.raumtyp_3 == cls]
        check(f"two-day share, {cls}", 100 * (m[col] == 2).mean(), want, 0.15)

    c = D.load_costs()
    t = (c.groupby(["penalty", "share_willing"], as_index=False)
          .total_stage3_eur.sum())
    t["saved"] = D.BASE_TOTAL - t.total_stage3_eur

    def at(pen, th):
        return float(t[np.isclose(t.penalty, pen)
                       & np.isclose(t.share_willing, th)].saved.iloc[0])

    # the 2026-08-25 metric fix: a no-op at theta = 1, large below it
    w = D.load_wait()

    def wait(pen, th):
        return float(w[np.isclose(w.penalty, pen)
                       & np.isclose(w.share_willing, th)].avg_wait_d_stage3.iloc[0])

    for pen, want in ((0.0, 0.975), (0.25, 0.455), (0.5, 0.228)):
        check(f"wait at P={pen:g}, θ=1 (fix is a no-op here)", wait(pen, 1.0),
              want, 0.001)
    fl = D.fleet_totals()

    def peak(pen, th):
        return float(fl[np.isclose(fl.penalty, pen)
                        & np.isclose(fl.share_willing, th)].peak_s3.iloc[0])

    check("peak fleet cut at P=0.5, θ=1 [%]",
          100 * (1 - peak(0.5, 1.0) / peak(0.5, 0.0)), 12.9, 0.1)

    def vehweek(pen, th):
        return float(fl[np.isclose(fl.penalty, pen)
                        & np.isclose(fl.share_willing, th)].mean_s3.iloc[0]) * 6

    check("van-days saved at P=10, θ=0.1", vehweek(0.0, 0.0) - vehweek(10.0, 0.1),
          180, 2)
    check("van-days saved at P=0.25, θ=1", vehweek(0.0, 0.0) - vehweek(0.25, 1.0),
          524, 2)

    for th, want_prize, want_kept in ((0.1, 145041, 68425), (0.2, 159146, 15109)):
        check(f"prize at θ={th:g} (no fee)", at(0.0, th), want_prize, 60)
        check(f"kept at θ={th:g}, P=10", at(10.0, th), want_kept, 60)
        check(f"the fee costs at θ={th:g}", at(0.0, th) - at(10.0, th),
              want_prize - want_kept, 90)

    g = D.saving_grid().merge(D.load_wait(), on=["penalty", "share_willing"])
    at1 = g[np.isclose(g.share_willing, 1.0)]
    for pen, want in ((0.0, 22.8), (0.25, 18.5), (0.5, 13.5), (0.75, 10.2)):
        check(f"saving at P={pen:g}, θ=1",
              float(at1[np.isclose(at1.penalty, pen)].saving_pct.iloc[0]),
              want, 0.06)
    print("\n" + ("AUDIT FAILED: " + ", ".join(bad) if bad else "audit clean"))
    return 1 if bad else 0


def main() -> int:
    global _F, _RV, _TAG
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=Path, default=DEFAULT_OUT)
    ap.add_argument("--audit", action="store_true",
                    help="recompute every number the built slides claim, "
                         "then stop")
    ap.add_argument("--submission", action="store_true",
                    help="with --audit: run the SUBMISSION body's old checks "
                         "instead, which need PRES_REV_DIR="
                         "results/revision_2026_07")
    ap.add_argument("--rev-dir", type=Path, default=None,
                    help="the revision grid to read (default: $PRES_REV_DIR, "
                         "else results/revision_2026_08_v5)")
    ap.add_argument("--no-provisional", action="store_true",
                    help="drop the 'v5 · provisional' footer tag")
    ap.add_argument("--no-revision", action="store_true",
                    help="keep the submission-era chapter 1 and read no grid")
    G.add_args(ap)
    a = ap.parse_args()
    if a.audit:
        if a.submission:
            import _data as D
            if D.SCHEMA != D.SCHEMA_LEGACY:
                print(f"--audit --submission needs the submission grid; REV "
                      f"is {D.REV} ({D.SCHEMA}). Set PRES_REV_DIR="
                      f"results/revision_2026_07.", file=sys.stderr)
                return 1
            return audit_submission()
        if a.rev_dir is not None:
            import _data as D
            D.set_rev_dir(a.rev_dir)
        return audit_v5()
    if not SRC.exists():
        print(f"source deck not found: {SRC}", file=sys.stderr)
        return 1
    out = G.resolve(a.out, a.out_suffix, overwrite=a.overwrite)
    if not a.no_revision:
        import _data as D
        import _revision as RV
        if a.rev_dir is not None:
            D.set_rev_dir(a.rev_dir)
        if D.SCHEMA != D.SCHEMA_V2:
            raise SystemExit(
                f"{D.REV} is a {D.SCHEMA} grid; the withdrawal chapter needs "
                f"the two-plan tables. Pass --rev-dir or set PRES_REV_DIR.")
        print(f"  revision grid: {D.REV.relative_to(D.ROOT)}")
        _F, _RV, _TAG = RV.Facts.load(), RV, not a.no_provisional
    keep = SRC.with_name(SRC.stem + "_untouched_copy.pptx")
    if not keep.exists():
        shutil.copy2(SRC, keep)
        print(f"kept an untouched copy at {keep.name}")
    p = build(out)
    print(f"wrote {p}")
    print(f"  {len(Presentation(str(p)).slides)} slides total, "
          f"{p.stat().st_size / 1048576:.1f} MB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
