"""The 2026-08 revision, as slide-ready facts with their sources attached.

The revision changed what the talk claims, not just the numbers behind it: the
tours are costed by one universal rule for baseline and scenario alike, there
are now two cost lenses and two weekly plans, and the operating point that wins
depends on which lens the operator is in. Both decks (`91_build_pptx.py`, the
v1 grammar, and `94_build_house_deck.py`, the house grammar) have to tell that
story, so the *evidence* lives here once and each deck renders it in its own
form.

Three rules this module exists to enforce
-----------------------------------------
**Every number is read from the grid, never typed in.** :class:`Facts` pulls
each figure through `_data`'s v2 adapter and asserts it against the value the
compendium records, so a v6 re-run that moves a number fails the build instead
of quietly contradicting the slide it is printed on.

**Every claim carries its compendium section.** Each fact has a `cite`, and
:func:`notes` writes it into the slide's speaker notes. A slide with no
traceable source does not get built.

**Provisional numbers say so.** Everything here is v5 — the head-free grid —
and Task 11's head grid will move it. :func:`provisional` stamps the small
footer tag that Part B removes when it re-runs against v6.

Nothing in this module draws a slide. It returns strings, rows and numbers; the
decks own their own geometry.
"""
from __future__ import annotations

import sys
from dataclasses import dataclass, field
from pathlib import Path

import numpy as np

sys.path.insert(0, str(Path(__file__).resolve().parent))

import _data as D                                              # noqa: E402

# ── the provisional tag ────────────────────────────────────────────────────
# Part B (v6 head grid + Task 12 validation + Task 13 fix round) rebuilds with
# `provisional=False` and this disappears from every slide it stamped.
def tag_text() -> str:
    """What the footer chip says, derived from the grid actually in use.

    The chip is not decoration: it states which grid a number came from and
    what about it is still open. v5 was provisional because a head grid was
    coming; v6's grid is final but its VROOM re-validation is still being
    produced, so the honest chip names that instead of claiming the deck is
    finished. When both are settled the chip disappears on its own.
    """
    short = D.REV.name.rsplit("_", 1)[-1]          # revision_2026_08_v6 -> v6
    if D.vroom_actual_baseline_available():
        return ""
    return f"{short} · validation pending"


TAG_TEXT = "v5 · provisional"      # the literal Part A shipped; see tag_text()
TAG_L, TAG_T, TAG_W, TAG_H = 11.62, 6.68, 1.66, 0.28   # right edge 13.28 in

COMPENDIUM = "docs/PAPER_COMPENDIUM_2026_05_24.md"

# What each grid is EXPECTED to say, and where that expectation is written
# down. The asserts in `Facts` are the point of this table: a grid whose
# numbers have moved fails the build instead of quietly restating every slide,
# and adding a grid means recording its values somewhere citable first.
#
# v5: compendium 40.15 / 40.17 / 40.18.
# v6: results/revision_2026_08_v6/DEEP_DIVE_V6_PAPER_IMPACT.md, which states
#     every one of these against its v5 predecessor.
GRID_EXPECT = {
    "revision_2026_08_v5": dict(
        source="docs/PAPER_COMPENDIUM_2026_05_24.md 40.15/40.17/40.18",
        cite=["§40.15", "§40.17", "§40.18"],
        rout1=23.10, rout2=20.43, op1=-7.79, op2=24.69, peak2=-16.87,
        breakeven_lo=0.77, breakeven_hi=2.24,
        pstar_moved={"Amazon", "FedEx", "Hermes"},
        # further values the source records, checked by 96_ --audit
        audit={
            "operator saving, stage-2 plan": {0.0: 24.69, 0.25: 22.82,
                                              0.5: 18.48, 0.75: 14.79,
                                              1.0: 11.82, 2.0: 6.04},
            "wait, stage-2 plan": {0.0: 0.77, 0.25: 0.39, 0.5: 0.23},
            "areas consolidating, stage-1 plan": {(10.0, 0.1): 2.9,
                                                  (0.0, 0.1): 57.1},
        },
    ),
    "revision_2026_08_v6": dict(
        source="results/revision_2026_08_v6/DEEP_DIVE_V6_PAPER_IMPACT.md",
        cite=["deep dive §2", "deep dive §3", "deep dive §6"],
        rout1=22.64, rout2=19.96, op1=-8.37, op2=24.30, peak2=-16.87,
        # The break-even band comes from the grid's own
        # `_peek/discount_scenarios_v6.csv`, not from the deep dive's prose:
        # the two disagree at the top of the band (grid 2.248, deep dive
        # "etwa 2,37"), and the deck follows the table it can recompute. The
        # same disagreement runs through the net-saving series (grid
        # 8.11/12.89/12.12/10.44/9.05 against the deep dive's
        # 8.05/13.02/12.27/10.58/9.18). Reported, not smoothed over.
        breakeven_lo=0.750, breakeven_hi=2.248,
        pstar_moved={"Amazon", "GLS", "Hermes"},
        # DEEP_DIVE_V6_PAPER_IMPACT.md 2 (savings), 3 (P = 0.25 and the
        # discount), 4 (the corner cells). Every one of these is a number the
        # document states, so the audit compares the deck with the write-up
        # rather than with itself.
        audit={
            "operator saving, stage-2 plan": {0.0: 24.30, 0.25: 22.55,
                                              0.5: 17.77, 0.75: 14.10},
            "wait, stage-2 plan": {0.0: 0.773, 0.25: 0.393},
            "areas consolidating, stage-1 plan": {(10.0, 0.1): 9.62,
                                                  (0.0, 0.1): 58.3},
        },
    ),
}


def expect() -> dict:
    """The recorded expectation for the grid in use, or {} if it has none."""
    return GRID_EXPECT.get(D.REV.name, {})


def cites(*extra) -> list:
    """Methodology sections first, then whatever records THIS grid's numbers.

    A note about the tour rule cites 40.7-40.9 whichever grid drew it -- the
    method did not change. A note carrying a number cites, in addition, the
    document that records that number for the grid in use, which is the
    compendium for v5 and the deep dive for v6.
    """
    out = [x for x in extra if x]
    out += [c for c in expect().get("cite", []) if c not in out]
    return out


def provisional(slide, *, tag: str | None = None, enabled: bool = True):
    """Stamp the small "v5 · provisional" chip into the slide's footer band.

    Any text box already reaching into the chip's rectangle is clipped back to
    just left of it. Both decks put a source line along the bottom, and the
    house deck's is wide enough to run underneath the chip; clipping it here
    keeps `_verify_layout.py` honest instead of teaching it to ignore an
    overlap.
    """
    if not enabled:
        return None
    tag = tag_text() if tag is None else tag
    if not tag:
        return None
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    left_edge = Inches(TAG_L - 0.10)
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if sh.left is None or sh.width is None or sh.top is None:
            continue
        bottom = Emu(sh.top).inches + Emu(sh.height or 0).inches
        if bottom < TAG_T or Emu(sh.top).inches > TAG_T + TAG_H:
            continue
        if sh.left + sh.width > left_edge > sh.left:
            sh.width = int(left_edge - sh.left)

    box = slide.shapes.add_textbox(Inches(TAG_L), Inches(TAG_T),
                                   Inches(TAG_W), Inches(TAG_H))
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = tag
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xBE, 0x1E, 0x3C)
    return box


# What a slide gets when its numbers have no revision counterpart yet -- the
# VROOM validation, whose re-run of both plans is still being produced. A
# visible banner, not a footnote: a reader who does not read the source line
# must still see that the numbers on the slide are the submission's.
STAMP_TEXT = "Stand Einreichung 2026-07 — wird in Teil B aktualisiert"


def stamped(slide) -> bool:
    return any(sh.has_text_frame and STAMP_TEXT in sh.text_frame.text
               for sh in slide.shapes)


def stamp(slide, *, text: str = STAMP_TEXT, top: float = 0.98,
          height: float = 0.40):
    """The submission-era banner, across the top of the body area.

    Idempotent, and it makes room for itself: a full-width box already sitting
    on this row (the backup slides carry an "EXPLAINS SLIDE n" tag there) is
    clipped to the left and the banner takes the right-hand side, rather than
    printing across it and turning a warning into a layout fault.
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    if stamped(slide):
        return None
    left, width = 0.63, 12.21
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if sh.top is None or sh.height is None or sh.width is None:
            continue
        if not (top - 0.08 <= Emu(sh.top).inches <= top + height - 0.10):
            continue
        keep = 5.9
        if Emu(sh.width).inches > keep:
            sh.width = Inches(keep)
        left, width = 0.63 + keep + 0.24, 12.21 - keep - 0.24

    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left),
                                 Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFD, 0xF3, 0xF5)
    box.line.color.rgb = RGBColor(0xBE, 0x1E, 0x3C)
    box.line.width = Pt(1.0)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    run = p0.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xBE, 0x1E, 0x3C)
    return box


def notes(slide, text: str, *, cite: str | list[str] | None = None):
    """Write the speaker notes, ending in the compendium section(s) cited."""
    parts = [text.strip()]
    if cite:
        secs = [cite] if isinstance(cite, str) else list(cite)
        parts.append("Source: " + "; ".join(secs) + f" ({COMPENDIUM}).")
    slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)
    return slide


# ── the facts ──────────────────────────────────────────────────────────────
def eur(x: float) -> str:
    """1909432.42 -> '1 909 432'. A thin space, applied to the NUMBER only.

    A blanket `.replace(",", " ")` on a whole sentence eats its commas, which
    is exactly what happened on the first build of the revision block.
    """
    return f"{x:,.0f}".replace(",", " ")


def _pct(x: float) -> str:
    return f"{x:.1f} %"


def _sign_pct(x: float) -> str:
    return f"{x:+.1f} %"


@dataclass
class Facts:
    """Every revision number the decks print, read from the grid in use.

    Construct once per build (`Facts.load()`); the constructor does the reading
    and the checking, so a slide body is pure formatting.
    """

    rev_dir: Path
    base_routing: float
    base_operator: float
    base_peak: int
    base_vehicle_days: int
    headline: dict = field(default_factory=dict)      # (P) -> row dict
    pstar: list = field(default_factory=list)
    discount: dict = field(default_factory=dict)
    bantorf_after: list = field(default_factory=list)
    consolidating: dict = field(default_factory=dict)

    # The stage-1 profile of a one-cell hub is not in the v2 tables -- the grid
    # keeps only the final plan per hub and day -- so it is quoted from the
    # compendium and labelled as a quote wherever it appears. It is the ONLY
    # typed-in number in this class; everything else, including the one-area
    # depot count below, is derived and asserted.
    BANTORF_BEFORE = [0, 0, 33, 0, 0, 29]
    BANTORF_HUB = "Bantorf"

    # Filled by _load_one_cell_hub() from the hub assignment. They used to be
    # literals, and one of them was wrong (9, against the compendium's and the
    # figure's 8) on two built decks -- which is the argument for deriving even
    # a fact this small.
    one_cell_hubs: int = 0
    dhl_hubs: int = 0

    @classmethod
    def load(cls) -> "Facts":
        b = D.baseline_v2()
        f = cls(rev_dir=D.REV, base_routing=b["routing_eur"],
                base_operator=b["operator_eur"], base_peak=b["hub_peak"],
                base_vehicle_days=b["vehicle_days"])
        f._load_headline()
        f._load_pstar()
        f._load_discount()
        f._load_one_cell_hub()
        f._load_consolidating()
        return f

    # ---- loaders ---------------------------------------------------------
    def _load_headline(self) -> None:
        h = D.load_headline_v2()
        for r in h.itertuples():
            self.headline[float(r.penalty)] = dict(
                rout1=float(r.routing_saving_plan1_pct),
                rout2=float(r.routing_saving_plan2_pct),
                op1=float(r.operator_saving_plan1_pct),
                op2=float(r.operator_saving_plan2_pct),
                peak1=int(r.sum_hub_peak_plan1),
                peak2=int(r.sum_hub_peak_plan2),
                peak1_pct=float(r.hub_peak_plan1_vs_base_pct),
                peak2_pct=float(r.hub_peak_plan2_vs_base_pct),
                wait1=float(r.wait_d_plan1), wait2=float(r.wait_d_plan2),
                days1=float(r.mean_days_plan1), days2=float(r.mean_days_plan2),
            )
        # The headline table is Task 13's; the adapter recomputes the same
        # quantity from the raw grid. If the two ever disagree, one of them is
        # stale and the deck must not pick a winner silently.
        for plan, lens, key in ((D.PLAN_ROUTING, D.LENS_ROUTING, "rout1"),
                                (D.PLAN_OPERATOR, D.LENS_ROUTING, "rout2"),
                                (D.PLAN_ROUTING, D.LENS_OPERATOR, "op1"),
                                (D.PLAN_OPERATOR, D.LENS_OPERATOR, "op2")):
            g = D.saving_grid_v2(plan, lens)
            g = g[np.isclose(g.share_willing, 1.0)]
            for row in g.itertuples():
                want = self.headline[float(row.penalty)][key]
                assert abs(row.saving_pct - want) < 5e-3, (
                    f"{plan}/{lens} at P={row.penalty:g}: adapter says "
                    f"{row.saving_pct:.4f} %, tab_headline_theta1_v2 says "
                    f"{want:.4f} %")
        # The numbers the story slides are written around, against whatever
        # document records this grid. A grid with no recorded expectation is
        # allowed but says so, loudly: nothing downstream can then claim the
        # slides were checked.
        h0 = self.headline[0.0]
        e = expect()
        if not e:
            print(f"  [facts] {D.REV.name} has no recorded expectation in "
                  f"GRID_EXPECT; the headline numbers on the slides are NOT "
                  f"cross-checked against any document")
        else:
            for key, got in (("rout1", h0["rout1"]), ("rout2", h0["rout2"]),
                             ("op1", h0["op1"]), ("op2", h0["op2"]),
                             ("peak2", h0["peak2_pct"])):
                assert abs(got - e[key]) < 0.02, (
                    f"{D.REV.name} {key} = {got:.4f} % but "
                    f"{e['source']} records {e[key]} %; the deck and its "
                    f"source disagree")

    def _load_pstar(self) -> None:
        k = D.load_pstar_v2()
        self.pstar = [dict(
            provider=str(r.provider),
            p_submission=float(r.P_star_submission),
            p_routing=float(r.P_star_routing),
            p_operator=float(r.P_star_operator),
            class_routing=str(r.carrier_class_routing),
            class_operator=str(r.carrier_class_operator),
        ) for r in k.itertuples()]
        moved = {m["provider"] for m in self.pstar_moved()}
        e = expect()
        if e:
            assert moved == e["pstar_moved"], (
                f"{e['source']} names "
                f"{sorted(e['pstar_moved'])} as the providers whose knee moves "
                f"between lenses; {D.REV.name} says {sorted(moved)}")

    def _load_discount(self) -> None:
        d = D.load_discount_v2()
        sub = d[np.isclose(d.th, 1.0) & (d.plan == "operator-plan")]
        for r in sub.itertuples():
            self.discount[float(r.P)] = dict(
                delayed=float(r.delayed_parcels),
                op_shadow=float(r.sav_operator_pct),
                op_perday=float(r.net_operator_perday_pct),
                op_flat=float(r.net_operator_flat_pct),
                rout_flat=float(r.net_routing_flat_pct),
                be_op=float(r.breakeven_flat_operator),
                be_rout=float(r.breakeven_flat_routing),
            )
        lo, hi = self.discount[0.0]["be_op"], self.discount[1.0]["be_op"]
        e = expect()
        if e:
            assert (abs(lo - e["breakeven_lo"]) < 0.02
                    and abs(hi - e["breakeven_hi"]) < 0.02), (
                f"{e['source']} pins the operator-lens break-even band to "
                f"{e['breakeven_lo']}-{e['breakeven_hi']} EUR; "
                f"{D.REV.name} says {lo:.2f}-{hi:.2f}")

    def _load_one_cell_hub(self) -> None:
        hub, profile = D.hub_day_profile_v2(self.BANTORF_HUB, 0.0, 1.0)
        self.bantorf_hub_name = hub
        self.bantorf_after = [int(x) for x in profile]
        assert self.bantorf_after == [10, 11, 13, 12, 10, 8], (
            f"compendium 40.14/40.18 quote Bantorf as 10 11 13 12 10 8 after "
            f"the operator polish; the grid says {self.bantorf_after}")
        self.one_cell_hubs, self.dhl_hubs = D.one_cell_hubs("DHL")
        assert (self.one_cell_hubs, self.dhl_hubs) == (8, 16), (
            f"compendium 40.14 says 8 of DHL's 16 depots serve exactly one "
            f"cell and 40.18's fig 6 (f) counts n = 8 in the single-cell "
            f"bucket; the hub assignment says "
            f"{self.one_cell_hubs} of {self.dhl_hubs}")
        # DHL is the only multi-depot network in the case study, which is what
        # makes this a DHL statement rather than a general law.
        for other in [p for p in D.PROVIDERS if p != "DHL"]:
            assert D.one_cell_hubs(other)[1] == 1, (
                f"{other} is no longer single-depot; the one-area-depot slide "
                f"and fig 6 (f)'s within-DHL caveat both need rewriting")

    def _load_consolidating(self) -> None:
        for P, th in ((10.0, 0.1), (10.0, 0.2), (5.0, 0.1), (0.0, 0.1)):
            self.consolidating[(P, th)] = dict(
                plan1=D.consolidating_share_v2(P, th, D.PLAN_ROUTING),
                plan2=D.consolidating_share_v2(P, th, D.PLAN_OPERATOR),
            )

    # ---- derived views ---------------------------------------------------
    def pstar_moved(self) -> list:
        """The providers whose knee sits at a different P in the two lenses."""
        return [p for p in self.pstar if p["p_routing"] != p["p_operator"]]

    def recommended(self) -> dict:
        """P = 0.25: the point the revision recommends, in both lenses."""
        return self.headline[0.25]

    def partial_adoption(self, penalty: float = 0.0) -> list:
        """(theta, routing %, operator %) for the operator plan below theta=1.

        The finding this supports: at partial adoption the operator lens stays
        positive throughout while the routing lens barely moves.
        """
        gr = D.saving_grid_v2(D.PLAN_OPERATOR, D.LENS_ROUTING)
        go = D.saving_grid_v2(D.PLAN_OPERATOR, D.LENS_OPERATOR)
        m = gr.merge(go, on=["penalty", "share_willing"],
                     suffixes=("_rout", "_op"))
        m = m[np.isclose(m.penalty, penalty) & (m.share_willing > 0)]
        return [(float(r.share_willing), float(r.saving_pct_rout),
                 float(r.saving_pct_op)) for r in m.itertuples()]


# ── slide content: the rows and sentences both decks share ─────────────────
# Each entry pairs the text with the compendium section it comes from. The deck
# renders `rows`; `notes` goes into the speaker notes verbatim.
def lens_rows(f: Facts) -> list:
    """The two cost lenses, with their formulas — for a two-row table."""
    return [
        ["Routing euro",
         D.LENS_FORMULA[D.LENS_ROUTING],
         f"{eur(f.base_routing)} EUR/wk"],
        ["Operator euro",
         D.LENS_FORMULA[D.LENS_OPERATOR],
         f"{eur(f.base_operator)} EUR/wk"],
    ]


LENS_NOTES = (
    "One euro is not one euro. FIXED_COST_EUR = 189.15 is per vehicle and per "
    "DAY and already contains the driver, so a routing-euro saving on a "
    "skipped delivery day is mostly avoided driver cost. An operator with "
    "salaried drivers staffs each hub for its weekly peak: below that peak "
    "only the kilometres are a real outlay, and a vehicle removed from the "
    "peak is worth 6 x 189.15 = 1 134.90 EUR a week. Both baselines are the "
    "same daily-delivery system, priced twice.")


def plan_rows(f: Facts) -> list:
    """Routing plan against operator plan at theta = 1, P = 0."""
    h = f.headline[0.0]
    return [
        ["Routing-optimal (stage 1)", _pct(h["rout1"]), _sign_pct(h["op1"]),
         f"{h['peak1']}", f"{h['wait1']:.2f} d"],
        ["Operator-polished (stage 2)", _pct(h["rout2"]), _pct(h["op2"]),
         f"{h['peak2']}", f"{h['wait2']:.2f} d"],
    ]


def plan_notes(f: Facts) -> str:
    """Why there are two plans, with this grid's own numbers in the sentence."""
    h = f.headline[0.0]
    return (
        f"Two plans, not two ways of reporting one plan. Stage 1 minimises "
        f"routing cost; stage 2 then polishes that plan for operator cost, "
        f"and at theta > 0 it is frequency-free -- it may change HOW OFTEN a "
        f"cell is served, not just on which days. The routing-optimal plan is "
        f"worse than doing nothing in the operator lens "
        f"({h['op1']:.1f} %), because two-day patterns treble the hub peaks. "
        f"The polish turns that into {h['op2']:.1f} % -- "
        f"{h['op2'] - h['op1']:.1f} points -- for "
        f"{h['rout1'] - h['rout2']:.1f} points of routing saving, and it also "
        f"shortens the wait ({h['wait1']:.2f} d to {h['wait2']:.2f} d), "
        f"because serving more days is what lowers a one-cell hub's peak.")


PLAN_NOTES = None      # superseded by plan_notes(f); kept out of use on purpose


def one_cell_rows(f: Facts) -> list:
    return [
        ["Routing-optimal plan", " ".join(str(x) for x in f.BANTORF_BEFORE),
         "peak 33"],
        ["Operator-polished plan", " ".join(str(x) for x in f.bantorf_after),
         f"peak {max(f.bantorf_after)}"],
    ]


ONE_CELL_NOTES = (
    "Eight of DHL's sixteen hubs serve exactly one postal-code area, derived "
    "from the hub assignment and asserted against the compendium; DHL is the "
    "only multi-depot network in the case study. Under any "
    "rotation a one-cell hub on a two-day pattern has the profile 0 0 33 0 0 "
    "29: the whole week's demand lands on two days, and the peak only comes "
    "down if the hub delivers on more days. Stage 2's old frequency lock could "
    "re-time but not re-frequency, so it could not touch this. Freed, it sends "
    "the one-cell hubs back to near-daily service. The message: temporal "
    "consolidation pays where a depot can rotate delivery days across several "
    "areas; at a single-area depot it saves kilometres, not fleet. The "
    "before-profile is quoted from the compendium -- the v2 tables keep only "
    "the final plan per hub and day.")


def discount_rows(f: Facts) -> list:
    """The discount scenario at theta = 1, operator plan, in BOTH lenses.

    The flat-discount optimum is lens-specific -- P = 0.25 in the operator
    lens, P = 0.5 in the routing lens -- so a table that shows only one of them
    cannot be read for the other, and the presenter has nothing to point at.
    """
    out = []
    for P in (0.0, 0.25, 0.5, 0.75, 1.0):
        d = f.discount[P]
        out.append([f"P = {P:g}", f"{d['delayed'] / 1000:.0f} k",
                    _pct(d["op_shadow"]), _pct(d["op_flat"]),
                    _pct(d["rout_flat"]), f"{d['be_op']:.2f} EUR"])
    return out


def discount_optima(f: Facts) -> dict:
    """Where the flat-0.50-EUR discount peaks, per lens: {lens: (P, net %)}."""
    best = {}
    for lens, key in (("operator", "op_flat"), ("routing", "rout_flat")):
        ranked = sorted(f.discount, key=lambda p: -f.discount[p][key])
        P, runner = ranked[0], ranked[1]
        margin = f.discount[P][key] - f.discount[runner][key]
        best[lens] = (P, f.discount[P][key], runner, margin)
    return best


# Below this the winner is not distinguishable from its runner-up at the
# precision this model has, and the slide has to say "a tie" instead of naming
# a winner. On v6 the routing lens sits 0.009 pp apart between P = 0.25 and
# P = 0.5 in `_peek/discount_scenarios_v6.csv`, which is exactly the case this
# guard exists for. (DEEP_DIVE_V6_PAPER_IMPACT.md 3 quotes 0.021 pp for the
# same gap and calls it "practically a tie" too; the margin here is recomputed
# from the table, so it is the table's number that decides.)
DISCOUNT_TIE_PP = 0.20


def discount_optimum_line(f: Facts) -> str:
    """One sentence naming each lens's flat-discount optimum, or calling it."""
    o = discount_optima(f)
    parts = []
    for lens in ("operator", "routing"):
        P, net, runner, margin = o[lens]
        if margin < DISCOUNT_TIE_PP:
            parts.append(f"P = {P:g} and P = {runner:g} are level in the "
                         f"{lens} lens ({net:.1f} % vs "
                         f"{net - margin:.1f} %)")
        else:
            parts.append(f"P = {P:g} in the {lens} lens ({net:.1f} %)")
    return ("Read as a payout the penalty halves the saving — and the optimum "
            "is lens-specific: " + "; ".join(parts) + ".")


def discount_notes(f: Facts) -> str:
    """The discount scenario, read off this grid rather than remembered."""
    d = f.discount
    rout = " / ".join(f"{d[P]['rout_flat']:.1f}" for P in
                      (0.0, 0.25, 0.5, 0.75, 1.0))
    return (
        f"What if the service penalty is not a shadow price but money actually "
        f"paid to the waiting customer? At a flat 0.50 EUR per delayed parcel, "
        f"P = 0 is no longer the best point -- it delays "
        f"{d[0.0]['delayed'] / 1000:.0f} thousand parcels a week -- and the "
        f"optimum moves off the corner. "
        + discount_optimum_line(f)
        + f" The routing-lens series is {rout} % at P = 0 ... 1. The "
        f"break-even discount, what the operator could pay per delayed parcel "
        f"and still be at zero, runs {d[0.0]['be_op']:.2f} EUR at P = 0 to "
        f"{d[1.0]['be_op']:.2f} EUR at P = 1 in the operator lens and "
        f"{d[0.0]['be_rout']:.2f}-{d[1.0]['be_rout']:.2f} EUR in the routing "
        f"lens. Interpreted as a payout the penalty roughly halves the saving; "
        f"it does not remove it.")


DISCOUNT_NOTES = None  # superseded by discount_notes(f)


def pstar_rows(f: Facts) -> list:
    return [[p["provider"], f"{p['p_routing']:g}", f"{p['p_operator']:g}",
             p["class_routing"], p["class_operator"]] for p in f.pstar]


def pstar_headline(f: Facts) -> str:
    """The one-line reading under the knee table, true for THIS grid.

    v5's "three LSPs move up one class" stopped being accurate on v6, where
    GLS lands at P* = 1.0 and outside the paper's class band altogether -- so
    the line is built from the moved set rather than written down.
    """
    moved = f.pstar_moved()
    out_of_band = [m for m in moved if m["p_operator"] > 0.75]
    line = (f"{len(moved)} of {len(f.pstar)} LSPs sit at a different knee in "
            f"the operator lens: peak smoothing only starts to pay at a "
            f"higher penalty.")
    if out_of_band:
        line += ("  " + ", ".join(m["provider"] for m in out_of_band)
                 + " lands outside the paper's [0.25, 0.75] class band, so the "
                   "three-class statement holds in the routing lens only.")
    return line


def pstar_notes(f: Facts) -> str:
    """Which LSPs' knees move between the lenses, on the grid in use."""
    moved = f.pstar_moved()
    who = "; ".join(f"{m['provider']} {m['p_routing']:g} -> "
                    f"{m['p_operator']:g}" for m in moved)
    out = (f"The knee of the cost/wait front is lens-dependent. In the routing "
           f"lens every provider's P* is unchanged from the submission. In the "
           f"operator lens {len(moved)} move -- {who} -- because peak "
           f"smoothing only starts to bite at a higher penalty. The carrier "
           f"classes in the paper therefore need a lens qualifier; they are "
           f"not a property of the carrier alone.")
    outside = [m for m in moved if m["p_operator"] > 0.75]
    if outside:
        out += (" " + ", ".join(m["provider"] for m in outside)
                + " now sits OUTSIDE the [0.25, 0.75] band the paper claims, "
                  "so the three-class statement holds in the routing lens "
                  "only.")
    return out


PSTAR_NOTES = None     # superseded by pstar_notes(f)


TOUR_RULE_NOTES = (
    "The revision made the tour rule universal: one rule prices baseline and "
    "scenario alike, with no code branch that could tell them apart. The "
    "unbounded hub-pooled express tour is gone -- standard parcels of a "
    "non-delivering area ride that area's own tour -- and a minimum tour size "
    "of 230 parcels (one van) stops the optimiser from buying savings with "
    "mini-tours it would never dispatch. The measured consequence: the "
    "theta < 1 savings fall to an honest floor, and the apparent bump at "
    "theta = 10 % disappears, because it was an artefact of the old pooled "
    "express price rather than a behaviour of the system.")


def bulge_rows(f: Facts) -> list:
    """What the (P, theta) cells that produced the old bump do now."""
    out = []
    for (P, th), lbl in (((0.0, 0.1), "P = 0, 10 % join"),
                         ((5.0, 0.1), "P = 5, 10 % join"),
                         ((10.0, 0.1), "P = 10, 10 % join"),
                         ((10.0, 0.2), "P = 10, 20 % join")):
        c = f.consolidating[(P, th)]
        g = D.saving_grid_v2(D.PLAN_ROUTING, D.LENS_ROUTING)
        row = g[np.isclose(g.penalty, P) & np.isclose(g.share_willing, th)]
        out.append([lbl, f"{c['plan1']:.1f} %", f"{float(row.saving_pct.iloc[0]):.2f} %"])
    return out


# ── the VROOM validation, as the slides have to state it ───────────────────
# A realised SAVING needs a solved theta = 0 baseline. The v2 validation solves
# the scenario points first, so until item 0 lands the only honest figures are
# predicted-against-actual COST on the same tours, and the peak-fleet count the
# solver actually produced. `validation_facts()` reads those off whichever
# validation directory `_data.VAL` points at, so a slide can never restate the
# submission's realised savings by accident -- which is precisely what two
# slides of the explainer deck did until 2026-08-28.
def validation_stamp() -> str:
    """The chip a validation slide carries, derived from the run in use.

    It named "v5, v6 mit Baseline folgt" while v6's theta = 0 baseline was
    still being solved. That baseline landed 2026-08-28 (item 0 of
    `results/revision_2026_08_v6/validation/`), so the stamp now names the
    grid it is actually reading and says nothing is outstanding. Derived
    rather than typed, for the same reason `tag_text()` is: a frozen literal
    is a claim nobody re-checks.
    """
    grid = D.VAL.parent.name.rsplit("_", 1)[-1]
    if not D.vroom_actual_baseline_available():
        return f"Validierung {grid} · Baseline folgt"
    return f"Validierung {grid} · inkl. Baseline"


VALIDATION_STAMP = validation_stamp()


def validation_facts() -> dict:
    """What the validation in use actually supports, as numbers for a slide."""
    sv = D.load_savings_validation()
    if "saving_defined" in sv.columns:
        sv = sv[sv.saving_defined]          # see validation_rows()
    dg = D.load_vroom_diagnostics()
    allrow = dg[dg.penalty.astype(str) == "ALL"]
    out = dict(
        grid=D.VAL.parent.name,
        has_baseline=D.vroom_actual_baseline_available(),
        n=int(sv.n.sum()),
        gap_lo=float(sv.gap_pct.min()), gap_hi=float(sv.gap_pct.max()),
        points=sorted({float(x) for x in sv.penalty}),
        plans=sorted({str(x) for x in sv.plan}) if "plan" in sv else [],
    )
    if len(allrow):
        out.update(mape=float(allrow.MAPE_pct.iloc[0]),
                   bias=float(allrow.bias_pct.iloc[0]),
                   r2=float(allrow.R2.iloc[0]),
                   n_clean=int(allrow.n.iloc[0]))
    try:
        v = D.load_vroom_v2()
        sub = v[np.isclose(v.penalty, 0.0) & np.isclose(v.share_willing, 1.0)
                & (v.plan.astype(str) == "balanced")]

        def _peak(col):
            per = sub.groupby(["hub_name", "day"])[col].sum().reset_index()
            return int(per.groupby("hub_name")[col].max().sum())

        out.update(peak_pred=_peak("predicted_n_routes"),
                   peak_actual=_peak("vroom_n_routes"))
    except Exception:                                       # pragma: no cover
        pass
    return out


def validation_rows(vf: dict) -> list:
    """Per validated point: what the surrogate priced against what VROOM cost.

    Restricted to the points a system percentage exists for: not the theta = 0
    baseline (it IS the reference) and not an item the G6 budget rule sampled
    rather than enumerated, whose totals are a subset of the system. Both are
    real solves and both are reported in the validation's own report; they
    just do not belong in a table the speaker note calls "both plans at
    theta = 1".
    """
    sv = D.load_savings_validation().sort_values(["plan", "penalty"])
    if "saving_defined" in sv.columns:
        sv = sv[sv.saving_defined]
    return [[f"P = {r.penalty:g}",
             (str(r.plan).replace("balanced", "operator-polished")
              .replace("stage1", "routing-optimal")
              if "plan" in sv.columns else "-"),
             f"{r.surrogate_total_eur / 1e6:.3f} M EUR",
             f"{r.vroom_actual_total_eur / 1e6:.3f} M EUR",
             f"+{r.gap_pct:.1f} %"] for r in sv.itertuples()]


def validation_notes(vf: dict) -> str:
    """The speaker note for any slide showing the validation."""
    peak = ""
    if "peak_pred" in vf:
        peak = (f" The peak-fleet count survives contact with the solver: "
                f"{vf['peak_pred']} predicted against {vf['peak_actual']} "
                f"actual at (P = 0, theta = 1) on the operator-polished plan, "
                f"which is the number the operator lens is built on.")
    return (
        f"VROOM re-validation on {vf['grid']}, {vf['n']} solved instances "
        f"across both plans at theta = 1. The surrogate prices the SAME tours "
        f"{vf['gap_lo']:.1f}-{vf['gap_hi']:.1f} % ABOVE the solver at every "
        f"point, i.e. it under-promises -- MAPE {vf.get('mape', float('nan')):.2f} %, "
        f"bias +{vf.get('bias', float('nan')):.2f} %, "
        f"R2 {vf.get('r2', float('nan')):.3f} over {vf.get('n_clean', 0)} "
        f"clean instances.{peak} What is NOT shown is a realised SAVING "
        f"percentage: that needs a solved theta = 0 baseline, and the "
        f"validation solves the scenario points first. The submission's "
        f"22.8 -> 23.7 % pairs are a different grid AND a different quantity "
        f"and have been removed from this deck.")


BULGE_OLD_SHARE = 41.7     # what the submission-era deck claimed at (10, 0.1)


def bulge_notes(f: Facts) -> str:
    """Why the theta = 10 % bump was withdrawn, in this grid's numbers.

    The honest form of the statement changed between grids and the note has to
    change with it: on v5 the corner was practically zero, on v6 (where the
    certified bundle head prices supported pools) it is small but clearly
    positive. What is withdrawn either way is the ARTEFACT -- the 3.6-point
    bump produced by an unbounded hub-pooled express tour -- and the reading
    that the effective knob is the product of penalty and adoption.
    """
    c = f.consolidating[(10.0, 0.1)]
    g = D.saving_grid_v2(D.PLAN_ROUTING, D.LENS_ROUTING)
    row = g[np.isclose(g.penalty, 10.0) & np.isclose(g.share_willing, 0.1)]
    sav = float(row.saving_pct.iloc[0])
    return (
        f"This slide replaces the old 'the bump at theta = 10 % survives even "
        f"a punitive penalty' slide. Under the pre-revision pooling, "
        f"{BULGE_OLD_SHARE:.1f} % of areas still gave up daily delivery at "
        f"P = 10, theta = 0.1 and the grid showed a real saving there; that "
        f"bump was priced by a hub-pooled express tour no operator would run. "
        f"With one universal tour rule the same cell consolidates "
        f"{c['plan1']:.1f} % of areas in the routing-optimal plan "
        f"({c['plan2']:.1f} % after the operator polish) and saves "
        f"{sav:.2f} % of routing cost. The asymmetry-driven bump is refuted; "
        f"what remains is small, positive and plan-dependent, and the "
        f"P-times-theta reading of the old slide -- a description of the "
        f"artefact rather than of the mechanism -- is withdrawn.")


BULGE_NOTES = None     # superseded by bulge_notes(f)


# ═══════════════════════════════════════════════════════════════════════════
# Task 20: the temporal-flexibility and express-cost facts of grid v6
# ═══════════════════════════════════════════════════════════════════════════
# `Facts` above answers "what does the revision claim". The flexibility deck
# (`98_deck_v6_flexibility.py`) asks two further questions the existing loaders
# do not cover:
#
#   what does the customer concede -- delayed parcels, waiting days, delivery
#   days per week -- and what does the operator get for it (peak fleet, vehicle
#   days, cost); and
#
#   what does the express pooled tour cost, against a regular delivery tour,
#   over the adoption share.
#
# Both come from the tables 75_/77_/78_ wrote, which no deck read before. Every
# value here is read from those tables and cross-checked against the compendium
# section that records it; a grid whose numbers moved fails the build rather
# than silently restating a slide.
FLEX_EXPECT = {
    # 40.20 (express prices at P = 0.25), 40.21 (headline), 40.22b (fleet),
    # 40.28 (kilometres).
    "express_eur_per_vd_lo": 339.0,      # theta = 0.1
    "express_eur_per_vd_hi": 519.0,      # theta = 0.9
    "regular_eur_per_vd": 298.0,
    "express_share_08": 9.6,
    "express_share_09": 7.4,
    "peaks_P0": (1239, 1666, 1030),      # baseline, routing plan, operator plan
    "peaks_P025": (1239, 1314, 1026),
    "saturday_P025": (781, 944),         # baseline, operator plan
    "freq_changed": 5525,
    "cell_rows": 27456,
    "total_parcels": 1263130,
    "co2_baseline_km": 305310,
    "co2_operator_P0_km": 224782,
}

# What VROOM found, as the compendium records it (40.24 clean basis, 40.25).
VROOM_EXPECT_V6 = {
    "baseline_gap_routing": 4.385,
    "baseline_gap_operator": 4.047,
    "realised_routing_plan_P0": 20.58,     # routing lens
    "realised_operator_plan_P0": 22.08,    # operator lens
    "realised_routing_plan_P0_oplens": -12.09,
}

VROOM_OK_STATUS = ("OK", "CACHED")


def _flex_table(stem: str):
    """A `tables/<stem>.csv` of the grid in use, read through the adapter."""
    return D._read(D.REV / "tables" / f"{stem}.csv")


def clean_vroom(df):
    """The rows a realised number may be computed from.

    A solve that left jobs unassigned or had jobs removed prices a different
    problem than the one the surrogate was asked about, so it cannot enter a
    predicted-against-actual comparison. This is the basis 40.24's addendum I1
    calls "clean", and it is what the paper prints.
    """
    return df[(df.n_unassigned == 0) & (df.jobs_removed == 0)
              & (df.vroom_status.isin(VROOM_OK_STATUS))]


def lens_totals(df, prefix: str, *, fixed_eur: float = None) -> dict:
    """Routing euro, operator euro and peak fleet of one solved point.

    `prefix` is "predicted_" or "vroom_". The operator euro is rebuilt the way
    the grid builds it -- variable cost plus six days of fixed cost for every
    vehicle a hub must keep -- so a validation row and a grid row mean the same
    thing. The peak is the sum over hubs of that hub's busiest day.
    """
    fixed_eur = D.FIXED_COST_EUR if fixed_eur is None else fixed_eur
    cost = float(df[f"{prefix}cost_eur"].sum())
    routes = float(df[f"{prefix}n_routes"].sum())
    per_day = (df.groupby(["hub_name", "day"])[f"{prefix}n_routes"]
               .sum().reset_index())
    peak = float(per_day.groupby("hub_name")[f"{prefix}n_routes"].max().sum())
    variable = cost - fixed_eur * routes
    return dict(routing=cost, operator=variable + 6.0 * fixed_eur * peak,
                peak=peak, variable=variable, routes=routes)


@dataclass
class Vroom:
    """What the v6 VROOM validation supports, on the clean basis.

    Constructed from the validation directory `_data.VAL` points at, which the
    deck sets explicitly -- the module default still names the v5 run, because
    the other decks were audited against it.
    """

    grid: str
    baseline: dict = field(default_factory=dict)     # routing/operator, pred+act
    points: list = field(default_factory=list)       # one dict per solved point
    subset_points: list = field(default_factory=list)
    by_kind: list = field(default_factory=list)

    @classmethod
    def load(cls) -> "Vroom":
        df = clean_vroom(D.load_vroom_v2())
        base = df[df.item == 0]
        if not len(base):
            raise SystemExit(
                f"{D.VAL} has no item-0 rows: without a solved theta = 0 "
                f"baseline there is no realised saving to put on a slide. "
                f"Point --val-dir at the v6 validation.")
        bp = lens_totals(base, "predicted_")
        ba = lens_totals(base, "vroom_")
        v = cls(grid=D.VAL.parent.name,
                baseline=dict(pred=bp, actual=ba, n=int(len(base)),
                              gap_routing=(bp["routing"] / ba["routing"] - 1)
                              * 100.0,
                              gap_operator=(bp["operator"] / ba["operator"] - 1)
                              * 100.0))
        for (item, plan, P, th), g in df[df.item != 0].groupby(
                ["item", "plan", "penalty", "share_willing"]):
            p, a = lens_totals(g, "predicted_"), lens_totals(g, "vroom_")
            row = dict(
                item=int(item), plan=str(plan), penalty=float(P),
                share_willing=float(th), n=int(len(g)),
                gap_routing=(p["routing"] / a["routing"] - 1) * 100.0,
                gap_operator=(p["operator"] / a["operator"] - 1) * 100.0,
                peak_pred=p["peak"], peak_actual=a["peak"],
                # a realised saving needs the SAME population on both sides,
                # so a subset item is kept apart (40.27's report defect).
                #
                # `pred_*` is the surrogate's saving over THIS point's clean
                # population, which is what `gap_*` has to be formed against.
                # It is NOT the predicted saving a slide should print: dropping
                # one partial instance from item 2 at P = 0 moves it from the
                # grid's 22.64 % to 23.29 %, and a slide showing that beside
                # the grid figure would contradict itself. Print
                # `Facts.headline[...]` for prediction and `act_*` for reality,
                # which is the pairing 40.25 records.
                pred_routing=(1 - p["routing"] / bp["routing"]) * 100.0,
                act_routing=(1 - a["routing"] / ba["routing"]) * 100.0,
                pred_operator=(1 - p["operator"] / bp["operator"]) * 100.0,
                act_operator=(1 - a["operator"] / ba["operator"]) * 100.0,
                subset=float(th) < 1.0)
            (v.subset_points if row["subset"] else v.points).append(row)
        v.points.sort(key=lambda r: (r["plan"], r["penalty"]))
        v._load_by_kind(df)
        v._check()
        return v

    def _load_by_kind(self, df) -> None:
        """MAPE and bias per instance kind at the theta < 1 point (40.27)."""
        sub = df[df.item == 3]
        if not len(sub):
            return
        for kind, g in sub.groupby("instance_kind"):
            err = (g.predicted_cost_eur - g.vroom_cost_eur) / g.vroom_cost_eur
            self.by_kind.append(dict(
                kind=str(kind), n=int(len(g)),
                mape=float(np.abs(err).mean() * 100.0),
                bias=float(err.mean() * 100.0),
                pred=float(g.predicted_cost_eur.sum()),
                actual=float(g.vroom_cost_eur.sum())))
        self.by_kind.sort(key=lambda r: -r["bias"])

    def _check(self) -> None:
        e = VROOM_EXPECT_V6
        assert abs(self.baseline["gap_routing"] - e["baseline_gap_routing"]) \
            < 0.02, (f"40.25 records a +{e['baseline_gap_routing']} % routing "
                     f"gap at the daily baseline; this validation says "
                     f"{self.baseline['gap_routing']:.3f} %")
        assert abs(self.baseline["gap_operator"] - e["baseline_gap_operator"]) \
            < 0.02, (f"40.25 records +{e['baseline_gap_operator']} % on the "
                     f"operator lens; this says "
                     f"{self.baseline['gap_operator']:.3f} %")
        r = self.point("stage1", 0.0)
        assert abs(r["act_routing"] - e["realised_routing_plan_P0"]) < 0.05, (
            f"40.25 records a realised {e['realised_routing_plan_P0']} % for "
            f"the routing plan at P = 0; this says {r['act_routing']:.2f} %")
        assert abs(r["act_operator"] - e["realised_routing_plan_P0_oplens"]) \
            < 0.05, (
            f"40.25 records {e['realised_routing_plan_P0_oplens']} % for the "
            f"routing plan in the operator lens; this says "
            f"{r['act_operator']:.2f} %")
        b = self.point("balanced", 0.0)
        assert abs(b["act_operator"] - e["realised_operator_plan_P0"]) < 0.05, (
            f"40.25 records a realised {e['realised_operator_plan_P0']} % for "
            f"the operator plan; this says {b['act_operator']:.2f} %")

    def point(self, plan: str, penalty: float) -> dict:
        for r in self.points:
            if r["plan"] == plan and abs(r["penalty"] - penalty) < 1e-9:
                return r
        raise KeyError(f"no validated point {plan} at P = {penalty:g}")

    def gap_range(self, lens: str = "routing") -> tuple:
        g = [r[f"gap_{lens}"] for r in self.points]
        return min(g), max(g)

    def n_solved(self) -> int:
        return (self.baseline["n"]
                + sum(r["n"] for r in self.points + self.subset_points))


@dataclass
class Flex:
    """Temporal flexibility and express cost, read off grid v6.

    One object per build. Everything a slide of the flexibility deck prints
    comes from here, so the slide bodies stay pure formatting and a moved grid
    number fails an assert instead of quietly contradicting the figure beside
    it.
    """

    mech: object = None           # tab_mechanism_theta_P_v2
    week: object = None           # tab_fleet_week_by_provider_v2
    classes: object = None        # tab_fleet_week_classes_v2
    co2: object = None            # tab_co2_km_v2
    head_kind: object = None      # tab_head_usage_summary_kind_v2
    total_parcels: int = 0
    regular_eur_per_vd: float = 0.0
    freq_changed: int = 0
    cell_rows: int = 0

    SYSTEM = "All seven LSPs"

    @classmethod
    def load(cls) -> "Flex":
        f = cls(mech=_flex_table("tab_mechanism_theta_P_v2"),
                week=_flex_table("tab_fleet_week_by_provider_v2"),
                classes=_flex_table("tab_fleet_week_classes_v2"),
                co2=_flex_table("tab_co2_km_v2"),
                head_kind=_flex_table("tab_head_usage_summary_kind_v2"))
        w = D.load_wait_v2()
        base = w[np.isclose(w.penalty, 0.0) & np.isclose(w.share_willing, 0.0)]
        f.total_parcels = int(base.total_parcels.sum())
        c = D.load_costs_v2()
        fl = D.load_fleet_v2()
        cb = c[np.isclose(c.penalty, 0.0) & np.isclose(c.share_willing, 0.0)]
        fb = fl[np.isclose(fl.penalty, 0.0) & np.isclose(fl.share_willing, 0.0)]
        # what one ORDINARY delivery vehicle-day costs, so the express price
        # below has something to be expensive against (40.20)
        f.regular_eur_per_vd = float(cb.cost_stage1_eur.sum() / fb.fleet.sum())
        f.freq_changed = int(c.cells_freq_changed_vs_stage1.sum())
        f.cell_rows = int(len(D.load_chosen_v2()))
        f._check()
        return f

    def _check(self) -> None:
        e = FLEX_EXPECT
        assert self.total_parcels == e["total_parcels"], (
            f"the case study moves {e['total_parcels']} parcels a week; this "
            f"grid says {self.total_parcels}")
        assert abs(self.regular_eur_per_vd - e["regular_eur_per_vd"]) < 1.0, (
            f"40.20 puts a regular delivery tour at about "
            f"{e['regular_eur_per_vd']} EUR per vehicle-day; this grid says "
            f"{self.regular_eur_per_vd:.2f}")
        lo, hi = self.express_price_band(0.25)
        assert abs(lo - e["express_eur_per_vd_lo"]) < 1.5 and \
            abs(hi - e["express_eur_per_vd_hi"]) < 1.5, (
            f"40.20 puts the express tour at {e['express_eur_per_vd_lo']}-"
            f"{e['express_eur_per_vd_hi']} EUR per vehicle-day at P = 0.25; "
            f"this grid says {lo:.0f}-{hi:.0f}")
        s8, s9 = (self.express_share(0.25, 0.8), self.express_share(0.25, 0.9))
        assert abs(s8 - e["express_share_08"]) < 0.1 and \
            abs(s9 - e["express_share_09"]) < 0.1, (
            f"40.20 records an express cost share of {e['express_share_08']} "
            f"-> {e['express_share_09']} % between theta = 0.8 and 0.9; this "
            f"grid says {s8:.1f} -> {s9:.1f}")
        for P, want in ((0.0, e["peaks_P0"]), (0.25, e["peaks_P025"])):
            got = self.peaks(P)[self.SYSTEM]
            assert got == want, (
                f"40.22b records system hub peaks {want} at P = {P:g}; the "
                f"fleet-week table says {got}")
        assert self.saturday(0.25) == e["saturday_P025"], (
            f"40.22b records Saturday going from {e['saturday_P025'][0]} to "
            f"{e['saturday_P025'][1]} vehicles at P = 0.25; this grid says "
            f"{self.saturday(0.25)}")
        assert (self.freq_changed, self.cell_rows) == (e["freq_changed"],
                                                       e["cell_rows"]), (
            f"the operator polish was measured changing {e['freq_changed']} "
            f"of {e['cell_rows']} cell rows' FREQUENCY; this grid says "
            f"{self.freq_changed} of {self.cell_rows}")
        km = self.km_co2()
        assert abs(km["baseline"][0] - e["co2_baseline_km"]) < 50 and \
            abs(km["operator"][0.0][0] - e["co2_operator_P0_km"]) < 50, (
            f"40.28 records {e['co2_baseline_km']} baseline km and "
            f"{e['co2_operator_P0_km']} at P = 0; this grid says "
            f"{km['baseline'][0]:.0f} / {km['operator'][0.0][0]:.0f}")

    # ---- what the customer concedes --------------------------------------
    def delayed(self, facts: "Facts", penalty: float) -> dict:
        """Delayed parcels, their share, and the wait a delayed parcel takes.

        The grid reports the mean added wait over ALL parcels; a customer who
        is actually held waits longer than that, and the ratio of the two is
        the honest way to say so. Both inputs are table values; the ratio is
        derived, and every slide that prints it says so.
        """
        d = facts.discount[penalty]
        h = facts.headline[penalty]
        share = d["delayed"] / self.total_parcels
        return dict(parcels=d["delayed"], share_pct=share * 100.0,
                    wait_all=h["wait2"],
                    wait_delayed=h["wait2"] / share if share else 0.0,
                    days=h["days2"], days_plan1=h["days1"])

    # ---- the express tour ------------------------------------------------
    def express_row(self, penalty: float, share_willing: float):
        m = self.mech
        r = m[np.isclose(m.penalty, penalty)
              & np.isclose(m.share_willing, share_willing)]
        if not len(r):
            raise KeyError(f"no mechanism row at ({penalty:g}, "
                           f"{share_willing:g})")
        return r.iloc[0]

    def express_price_band(self, penalty: float) -> tuple:
        """Express euro per vehicle-day at theta = 0.1 and 0.9."""
        return (float(self.express_row(penalty, 0.1).express_eur_per_vd),
                float(self.express_row(penalty, 0.9).express_eur_per_vd))

    def express_share(self, penalty: float, share_willing: float) -> float:
        return float(self.express_row(penalty, share_willing).express_share_pct)

    def express_peak_share(self, penalty: float) -> tuple:
        """(theta, share) where the express cost share is highest."""
        m = self.mech[np.isclose(self.mech.penalty, penalty)
                      & (self.mech.share_willing < 1.0)]
        r = m.loc[m.express_share_pct.idxmax()]
        return float(r.share_willing), float(r.express_share_pct)

    def express_curve(self, penalty: float) -> list:
        """(theta, eur/vehicle-day, cost share %, vehicle-days) over theta."""
        m = self.mech[np.isclose(self.mech.penalty, penalty)
                      & (self.mech.share_willing > 0)]
        return [(float(r.share_willing), float(r.express_eur_per_vd),
                 float(r.express_share_pct), float(r.express_veh_days))
                for r in m.itertuples()]

    def days_curve(self, penalty: float, plan: str = "plan2") -> list:
        """(theta, mean weekly delivery days) for one plan."""
        m = self.mech[np.isclose(self.mech.penalty, penalty)]
        col = f"mean_days_{plan}"
        return [(float(r.share_willing), float(getattr(r, col)))
                for r in m.itertuples()]

    def days_rebound(self, penalty: float, plan: str = "plan2") -> tuple:
        """The theta = 0.8 -> 0.9 rebound: the curve turning back upwards."""
        d = dict(self.days_curve(penalty, plan))
        return d[0.8], d[0.9]

    # ---- what the operator gains -----------------------------------------
    def peaks(self, penalty: float) -> dict:
        """{provider: (baseline, routing plan, operator plan)} peak fleet."""
        w = self.week[np.isclose(self.week.penalty, penalty)]
        out = {}
        for prov, g in w.groupby("provider", sort=False):
            r = g.iloc[0]
            out[str(prov)] = (int(r.peak_baseline), int(r.peak_plan1),
                              int(r.peak_plan2))
        return out

    def system_week(self, penalty: float) -> dict:
        """Weekday vehicle counts of the whole system, per plan."""
        w = self.week[np.isclose(self.week.penalty, penalty)
                      & (self.week.provider == self.SYSTEM)].sort_values("day")
        return dict(days=[str(x) for x in w.weekday],
                    baseline=[int(x) for x in w.baseline],
                    plan1=[int(x) for x in w.plan1],
                    plan2=[int(x) for x in w.plan2])

    def vehicle_days(self, penalty: float) -> dict:
        """Weekly vehicle-days of the system, and their change vs baseline."""
        s = self.system_week(penalty)
        b = sum(s["baseline"])
        return dict(baseline=b, plan1=sum(s["plan1"]), plan2=sum(s["plan2"]),
                    plan1_pct=(sum(s["plan1"]) / b - 1) * 100.0,
                    plan2_pct=(sum(s["plan2"]) / b - 1) * 100.0)

    def saturday(self, penalty: float) -> tuple:
        s = self.system_week(penalty)
        i = s["days"].index("Sat")
        return s["baseline"][i], s["plan2"][i]

    def cv(self, penalty: float, group: str = None) -> dict:
        """Mo-Sa coefficient of variation per plan, one class or the system."""
        group = group or self.SYSTEM
        c = self.classes[np.isclose(self.classes.penalty, penalty)
                         & (self.classes.group == group)]
        return {str(r.plan): float(r.cv) for r in c.itertuples()}

    def class_rows(self, penalty: float) -> list:
        """Per carrier class: peak fleet and CV of baseline and both plans."""
        c = self.classes[np.isclose(self.classes.penalty, penalty)]
        out = []
        for grp, g in c.groupby("group", sort=False):
            by = {str(r.plan): r for r in g.itertuples()}
            out.append(dict(
                group=str(grp),
                peaks=(int(by["baseline"].sum_hub_peak),
                       int(by["plan1_routing_optimal"].sum_hub_peak),
                       int(by["plan2_operator_polished"].sum_hub_peak)),
                cv=(float(by["baseline"].cv),
                    float(by["plan1_routing_optimal"].cv),
                    float(by["plan2_operator_polished"].cv))))
        return out

    def km_co2(self) -> dict:
        """Weekly vehicle-km and tonnes of CO2, baseline and both plans."""
        c = self.co2[~self.co2.is_g6_subset.astype(bool)]
        base = c[(c.plan == "stage1") & np.isclose(c.share_willing, 0.0)].iloc[0]
        out = dict(baseline=(float(base.km_week), float(base.co2_t_week)),
                   operator={}, routing={})
        for r in c[np.isclose(c.share_willing, 1.0)].itertuples():
            key = "operator" if r.plan == "balanced" else "routing"
            out[key][float(r.penalty)] = (
                float(r.km_week), float(r.co2_t_week),
                (float(r.km_week) / float(base.km_week) - 1) * 100.0)
        return out

    def head_share(self) -> dict:
        """What fraction of the pooled cost the certified bundle head priced."""
        return {str(r.kind): float(r.head_cost_share)
                for r in self.head_kind.itertuples()}
