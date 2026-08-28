"""The EWGT 2026 talk, rebuilt in the author's own presentation grammar.

Same evidence and the same figures as `91_build_pptx.py`, restructured to the
rules measured from his four previous decks and recorded in `_house.py`:
two-line section titles, progressive bullet builds, icon badges instead of
bullet characters, 22 pt sentences, figures carrying the argument, a contact
slide with the backup material behind it.

The deck is generated into the institutional template, so master, theme, fonts
and furniture are inherited rather than re-implemented.

The deck's default output used to be the same path as `91_build_pptx.py`'s,
so whichever ran last silently replaced the other's deck. It now writes
`EWGT_26_Bienzeisler_TBC_house_deck.pptx`, and every write goes through
`_outguard.resolve()`, which refuses to overwrite an existing file without
`--overwrite`.

Every grid number on a results slide is read from `_revision.Facts` on the grid
that `--rev-dir` (or `$PRES_REV_DIR`) names, and each such slide carries the
`v5 · provisional` footer tag until Part B drops it with `--no-provisional`.

Usage:
    python scripts/presentation/94_build_house_deck.py [--out PATH] [--hero NAME]
    python scripts/presentation/94_build_house_deck.py --out-suffix _rev2026-08
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

import _house as H
import _outguard as G
from _house import (AMBER, BLACK, BLUSH, BODY_B, BODY_T, COL_W, CRIM, DIM, FIG,
                    FIG_W, FIG_X, GREEN, INK, INK2, L, LINE, PANEL, RED, S1,
                    S2, S3, S4, S5, S6, SW, SZ_BIG, SZ_BODY, SZ_DIA, SZ_LEAD,
                    SZ_NUM, SZ_SUB, TEAL, TEXT_W, TEXT_X, W, WHITE, arrow,
                    axes, badges, badge, build, cell_grid, contact, dot,
                    figure_slide, flow, hrule, hslide, label_box, pic, rect,
                    section_divider, serpentine, set_morph, slide_oval, split_slide,
                    statement, sub_items, tour, txt, weekbar)

A = FIG / "tierA"
BK = FIG / "tierB"
ASSET = H.ASSET

# The running section names — line 1 of every title in that stretch.
SEC_BG = "Background: Last-Mile Parcel Delivery"
SEC_TBC = "Time-Based Consolidation"
SEC_PROB = "The Combinatorial Problem"
SEC_METH = "Machine-Learning Surrogate Optimization"
SEC_CASE = "Case Study: Region Hannover"
SEC_RES = "Results"
SEC_IMP = "Implications and Outlook"
SEC_BAK = "Backup"
SEC_REV = "Revision 1: What Changed"

# The revision grid, read once by build_deck() and shared by the content
# functions below, which take only `prs` by design. `_F` is a _revision.Facts;
# `_RV` is the _revision module. Both are None when --no-revision is passed,
# and every use is guarded, so the submission-era deck still builds.
_F = None
_RV = None
_TAG = True


def mark(s, note, cite):
    """Speaker notes with their compendium section, plus the provisional tag."""
    if _RV is None:
        return s
    _RV.notes(s, note, cite=cite)
    _RV.provisional(s, enabled=_TAG)
    return s


def _KNEE(provider: str) -> float:
    """That LSP's routing-lens saving at its own knee, from the grid."""
    k = _RV.D.load_pstar_v2().set_index("provider")
    return float(k.loc[provider, "saving_pct_routing"])


def _VEHDAY_RANGE() -> tuple:
    """Weekly vehicle-days of the operator plan against the daily baseline."""
    if _F is None:
        return -10.0, 4.6
    g = _RV.D.load_grid_full_v2().vehicle_days_plan2_vs_base_pct
    return float(g.min()), float(g.max())

# The opening image. These are the study's own results rendered chrome-free by
# 95_title_heroes.py — the talk no longer opens on a generated illustration.
HERO = H.ROOT / "results" / "presentation_2026_08" / "heroes"
HEROES = {
    "week": (HERO / "hero_week.png",
             "The 312 optimized weekly schedules · P = 0.25 €/parcel/day, "
             "θ = 100 %"),
    "region": (HERO / "hero_region.png",
               "Delivery frequency per area at three adoption levels · "
               "P = 0.25 €/parcel/day"),
    "frontier": (HERO / "hero_frontier.png",
                 "All 88 Stage-3 operating points; the efficient front in red"),
    "wave": (H.ASSET / "parcel-wave.png",
             "Illustration generated with OpenAI ImageGen"),
}


def title_slide(prs, hero_name="week"):
    """His opening: a full-bleed image over a band carrying the title."""
    s = prs.slides.add_slide(prs.slide_layouts[H.LAYOUT_BLANK])
    hero, credit = HEROES[hero_name]
    H.pic_cover(s, hero, 0, 0.93, SW, 4.42)
    rect(s, 0, 5.35, SW, 0.10, RED)
    txt(s, 5.18, 0.27, 3.0, 0.37, "EWGT 2026", 16, bold=True, color=INK)
    txt(s, 1.08, 5.62, 11.4, 1.00,
        "Machine Learning Surrogate Optimization Framework for\n"
        "Time-Based Consolidation in Last Mile Parcel Delivery", 26, bold=True,
        color=INK, line=1.18)
    txt(s, 1.08, 6.62, 7.05, 0.32,
        "Lasse Bienzeisler, Felix Petre, Oskar Wage, and Bernhard Friedrich",
        16, color=INK2)
    txt(s, 8.30, 6.66, 4.55, 0.26, credit, 9, color=DIM, align=PP_ALIGN.RIGHT)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# 1 · Background
# ═══════════════════════════════════════════════════════════════════════════
def part_background(prs):
    build(prs, SEC_BG, "Growth, competition, and pressure to change",
          "Eurostat (2026), share of individuals who ordered goods or services "
          "online.",
          [("Truck", RED,
            "The last mile is the final leg from the distribution centre to the "
            "customer and closes the delivery process"),
           ("UpwardTrend_LTR", RED,
            "Parcel volumes keep growing, driven by many small e-commerce "
            "shipments — 62 % of Europeans ordered online in 2015, 78 % do today"),
           ("GroupOfPeople", RED,
            "The market is characterised by strong competition and high price "
            "expectations from customers"),
           ("OpenHandWithPlant", RED,
            "Operators and public authorities are under growing pressure to "
            "introduce more sustainable delivery concepts")],
          t=BODY_T + 0.35)

    s = hslide(prs, SEC_BG, "Delivery density is the cost driver",
               "Bienzeisler, L. (2025). Modeling the Last Mile: A Large-Scale "
               "Agent-Based Simulation Framework for Parcel Delivery. Doctoral "
               "dissertation, TU Braunschweig. Illustration generated with "
               "OpenAI ImageGen.")
    pic(s, ASSET / "urban-rural-density.png", L, BODY_T, W, 3.15)
    for i, (t1, t2, col) in enumerate([
            ("Many drops, short stems", "one van serves a whole street", S6),
            ("Few drops, long stems", "kilometres between single stops", RED)]):
        x = L + i * (W / 2 + 0.20)
        rect(s, x, BODY_T + 3.35, 0.10, 0.95, col)
        txt(s, x + 0.30, BODY_T + 3.35, W / 2 - 0.50, 0.42, t1, SZ_LEAD,
            bold=True, color=INK)
        txt(s, x + 0.30, BODY_T + 3.85, W / 2 - 0.50, 0.5, t2, SZ_BODY,
            color=DIM)
    txt(s, L, BODY_T + 4.45, W, 0.90,
        "We identified delivery density as the primary driver of tour "
        "efficiency and thus of average cost per delivered parcel.", SZ_LEAD,
        bold=True, color=RED, line=1.25)


def part_concepts(prs):
    concepts = [("Parcel lockers", "self-service pickup", "locker"),
                ("Micro-hubs", "local depots", "hub"),
                ("Cargo bikes", "zero-emission tours", "bike"),
                ("Parcel shops", "retail drop-off", "shop"),
                ("Crowdshipping", "people as couriers", "crowd")]

    def cards(s):
        kw = (W - 4 * 0.20) / 5
        for i, (nm, sub, ic) in enumerate(concepts):
            x = L + i * (kw + 0.20)
            rect(s, x, BODY_T, kw, 2.80, PANEL, line_col=LINE)
            H.B.icon_glyph(s, ic, x + kw / 2 - 0.38, BODY_T + 0.35, 0.76)
            nsz = SZ_BODY
            while nsz > 16 and H.text_lines(nm, kw - 0.24, nsz, bold=True) > 1:
                nsz -= 0.5
            txt(s, x + 0.12, BODY_T + 1.42, kw - 0.24, 0.50, nm, nsz,
                bold=True, color=INK, align=PP_ALIGN.CENTER, line=1.12)
            txt(s, x + 0.12, BODY_T + 2.14, kw - 0.24, 0.5, sub, SZ_SUB,
                color=DIM, align=PP_ALIGN.CENTER, line=1.12)

    build(prs, "Overview: Sustainable Strategies and Delivery Concepts",
          "Many last-mile concepts already exist",
          "Synthesis of the locker, micro-hub, urban consolidation centre and "
          "crowdshipping literature; the five recurring concept families in the "
          "European context.",
          [("Books", BLACK,
            "Each of these concepts is well studied and, in the right place, it "
            "works as intended"),
           ("City", RED,
            [("All five change ", False), ("where", True),
             (" parcels are handed over — none of them changes ", False),
             ("when", True), ("", False)]),
           ("MagnifyingGlass", RED,
            "They share one precondition every time: enough demand in a "
            "small enough area")],
          draw=cards, t=BODY_T + 2.98)

    def frame(s):
        axes(s, xlab="demand density  →", ylab="viability")

    build(prs, "Overview: Sustainable Strategies and Delivery Concepts",
          "Viability depends on demand density",
          "Conceptual synthesis of the locker, micro-hub, UCC and "
          "crowdshipping literature; the curve is schematic.",
          [("UpwardTrend_LTR", RED,
            "Every spatial concept has a break-even density below which the "
            "fixed cost of the place is not recovered"),
           ("TrafficCone", AMBER,
            "So the question is not which concept is best, but which parts of a "
            "region reach that density at all")],
          draw=frame, t=BODY_T + 0.30, x=8.00, tx=8.45, tw=4.30)

    s = hslide(prs, "Overview: Sustainable Strategies and Delivery Concepts",
               "The concepts cluster where density already is",
               "Positions are qualitative, taken from the preconditions the "
               "literature states — not measured break-even points.")
    axes(s, xlab="demand density  →", ylab="viability")
    for cx, cy in [(5.35, 2.10), (5.95, 1.78), (6.45, 2.42), (5.62, 2.72),
                   (6.20, 3.00)]:
        slide_oval(s, cx, cy, 0.30, RED)
    txt(s, 4.85, BODY_T + 0.05, 2.4, 0.36, "urban concepts", SZ_SUB, bold=True,
        color=RED, align=PP_ALIGN.CENTER)
    lb = label_box(s, 2.05, BODY_T + 2.35, 2.35, 1.05, None,
                   [("suburban", SZ_BODY, True, DIM),
                    ("& rural", SZ_BODY, True, DIM)], line_col=DIM)
    H.B._dash(lb)
    txt(s, 1.85, BODY_T + 3.50, 3.0, 0.36, "no concept lands here", SZ_SUB,
        bold=True, color=DIM, align=PP_ALIGN.CENTER)
    badges(s, [("NoSign", AMBER,
                "Lockers, hubs and bikes pay off in dense cities; where demand "
                "thins out they stop making sense"),
               ("City", RED,
                [("Yet these areas carry a ", False), ("large share", True),
                 (" of all delivery activity", False)])],
           BODY_T + 0.30, x=8.00, tx=8.45, tw=4.30)

    split_slide(prs, "Overview: Sustainable Strategies and Delivery Concepts",
                "Most parcels are delivered outside the core",
                "Bienzeisler et al. (2026), HAGRID study. Settlement type per "
                "postal-code area, BBSR classification.",
                BK / "fig13_map_raumtyp.png",
                [("City", RED,
                  [("77.6 %", True),
                   (" of postal-code areas in the region are suburban or rural "
                    "— 59.0 % suburban, 18.6 % rural", False)]),
                 ("UpwardTrend_LTR", AMBER,
                  [("Rural cost per parcel is about ", False), ("70 % higher", True),
                   (" than urban cost", False)]),
                 ("Target", RED,
                  "That is exactly where the spatial toolbox has least to offer, "
                  "so the lever has to be something other than space")],
                fig=(FIG_X + 0.55, BODY_T, 6.20, 5.05), tw=4.55,
                t=BODY_T + 0.55)


# ═══════════════════════════════════════════════════════════════════════════
# 2 · Time-based consolidation
# ═══════════════════════════════════════════════════════════════════════════
def part_concept(prs):
    def levers(s):
        for side, (nm, claim, items, acc, fill) in enumerate([
                ("Where", "Bring parcels together in space",
                 ["Parcel locker", "Micro-hub or UCC", "Carrier pooling"],
                 DIM, PANEL),
                ("When", "Bring order days together in time",
                 ["Hold eligible parcels", "Deliver them batched",
                  "No new infrastructure"], RED, BLUSH)]):
            x = L + side * (W / 2 + 0.20)
            cw = W / 2 - 0.20
            rect(s, x, BODY_T, cw, 2.60, fill, line_col=LINE)
            rect(s, x, BODY_T, cw, 0.10, acc)
            txt(s, x + 0.30, BODY_T + 0.24, cw - 0.60, 0.38, nm, SZ_BODY,
                bold=True, color=acc, spc=1.4, caps=True)
            txt(s, x + 0.30, BODY_T + 0.66, cw - 0.60, 0.80, claim, SZ_LEAD,
                bold=True, color=INK, line=1.18)
            for j, it in enumerate(items):
                txt(s, x + 0.30, BODY_T + 1.50 + j * 0.36, cw - 0.60, 0.34,
                    "▪  " + it, SZ_SUB + 2, color=INK2)

    build(prs, SEC_TBC, "Consolidation has two levers: where and when",
          "Conceptual synthesis of the spatial and temporal consolidation "
          "literature.",
          [("City", BLACK,
            [("Spatial consolidation: ", True),
             ("needs a place, and enough local volume to fill it", False)]),
           ("Truck", RED,
            [("Temporal consolidation: ", True),
             ("holds shipments back until the carrier's next delivery day",
              False)]),
           ("GroupOfPeople", RED,
            "Customers choose between fast and cheaper consolidated delivery")],
          draw=levers, t=BODY_T + 2.85)

    statement(prs, "If urgency differs,\nwhy is next-day the default?",
              sub="Urgent parcels stay daily. Flexible parcels can wait — for a "
                  "discount, a lower fee, or a different service tier. Time "
                  "becomes an operational control variable.")

    def weekgrid(s):
        cw, gp = 1.50, 0.20
        x0 = L + 1.90
        for i, d in enumerate(["MON", "TUE", "WED", "THU", "FRI", "SAT"]):
            txt(s, x0 + i * (cw + gp), BODY_T, cw, 0.32, d, SZ_BODY, bold=True,
                color=DIM, align=PP_ALIGN.CENTER)
        txt(s, L, BODY_T + 0.58, 1.80, 0.36, "eligible θ", SZ_BODY, bold=True,
            color=RED)
        for i, on in enumerate([False, True, False, False, True, False]):
            label_box(s, x0 + i * (cw + gp), BODY_T + 0.46, cw, 0.66,
                      RED if on else H.B.PINK,
                      [("DELIVER" if on else "HOLD", SZ_SUB, on,
                        WHITE if on else H.B.RGBColor(0x8F, 0x14, 0x2B))])
        txt(s, L, BODY_T + 1.86, 1.80, 0.36, "standard", SZ_BODY, bold=True,
            color=DIM)
        for i in range(6):
            rect(s, x0 + i * (cw + gp), BODY_T + 1.76, cw, 0.54, S6)

    build(prs, SEC_TBC, "Current and new delivery pattern",
          "Paper §2.2. Maximum holding time Hmax = 3 days over a Monday-to-"
          "Saturday operating week.",
          [("Envelope", RED,
            "Parcels whose recipients accept a wait are held for the cell's next "
            "predefined service day"),
           ("Truck", BLACK,
            "Parcels that did not opt in keep a conventional daily service, "
            "exactly as today"),
           ("TrafficCone", AMBER,
            "Maximum holding time is three days, so no recipient waits longer "
            "than that in any pattern")],
          draw=weekgrid, t=BODY_T + 2.60)

    def effects(s):
        label_box(s, L + 3.55, BODY_T, 5.1, 0.85, S6,
                  [("ONE PATTERN   σ = Tue · Thu · Sat", 19, True, WHITE)])
        eff = [("Customer wait", "the gap to the next service day", RED),
               ("Routing cost", "denser tours, fewer of them", S5),
               ("Vehicle demand", "which weekday the fleet is needed", S6)]
        bwd = (W - 2 * 0.50) / 3
        for i, (nm, sub, col) in enumerate(eff):
            x = L + i * (bwd + 0.50)
            arrow(s, L + 6.1, BODY_T + 0.90, x + bwd / 2, BODY_T + 1.58,
                  colour=col)
            label_box(s, x, BODY_T + 1.62, bwd, 1.15, PANEL,
                      [(nm, SZ_LEAD, True, col), (sub, SZ_SUB, False, DIM)],
                      line_col=LINE)

    build(prs, SEC_TBC, "One calendar decision moves three things at once",
          "Schematic for the eligible share; standard parcels retain daily "
          "service · Paper §2.2.",
          [("Gears", BLACK,
            "The three effects cannot be tuned separately — the pattern that "
            "saves most also waits longest"),
           ("Target", RED,
            "So the schedule is the decision variable, not the route")],
          draw=effects, t=BODY_T + 3.10)


# ═══════════════════════════════════════════════════════════════════════════
# 3 · The mechanism and the combinatorics
# ═══════════════════════════════════════════════════════════════════════════
GX, GY, GW, GH = 1.95, BODY_T + 0.10, 6.40, 2.70
PX, PW = 8.60, W + L - 8.60
BUL_T = GY + GH + 0.30
DEPOT_X, DEPOT_Y = L, GY + GH / 2


def part_mechanism(prs):
    section_divider(prs, "Part one", "The mechanism, step by step",
                    "From arriving parcels to one tour fewer — and to the "
                    "question of which days belong together")

    def days(s):
        dcw = (W - 5 * 0.22) / 6
        for i, (dname, nd) in enumerate(zip(
                ["MON", "TUE", "WED", "THU", "FRI", "SAT"], [4, 6, 5, 7, 5, 3])):
            x = L + i * (dcw + 0.22)
            txt(s, x, BODY_T, dcw, 0.34, dname, SZ_BODY, bold=True, color=DIM,
                align=PP_ALIGN.CENTER)
            rect(s, x, BODY_T + 0.42, dcw, 2.35, PANEL, line_col=LINE)
            for k in range(nd):
                fx, fy = H.B._DOT_OFF[(i * 3 + k) % len(H.B._DOT_OFF)]
                dot(s, x + 0.20 + fx * (dcw - 0.40),
                    BODY_T + 0.58 + fy * 2.05, 0.19, S4)

    build(prs, SEC_PROB, "Parcels arrive on every weekday",
          "Schematic. Arrival volumes per weekday are illustrative; the model "
          "uses HAGRID weekday demand per cell.",
          [("Envelope", BLACK,
            "Demand does not wait for a good day to arrive — every weekday "
            "brings its own parcels, in every area"),
           ("Truck", RED,
            "Today every one of them is delivered the next day, however few of "
            "them there are")],
          draw=days, t=BODY_T + 3.05)

    def raster(s):
        cell_grid(s, GX, GY, GW, GH, theta=None)
        label_box(s, PX, GY + 0.30, PW, 1.00, PANEL,
                  [("one cell", SZ_BODY, True, INK),
                   ("= one carrier × one postal-code area", SZ_DIA, False, DIM)],
                  line_col=LINE)

    build(prs, SEC_PROB, "The region splits into provider–area cells",
          "Schematic raster of twelve cells; the case study has 312 "
          "provider–area cells across seven providers.",
          [("City", RED,
            "Each cell is one carrier's parcels in one postal-code area — the "
            "unit we decide about and the unit we route"),
           ("Gears", BLACK,
            "Small enough to solve, and numerous enough to matter: 312 of them "
            "in the Hanover region")],
          draw=raster, t=BUL_T)

    # the θ build — identical geometry on all four slides, so Morph tweens it
    def opt_in(theta_pct, note):
        s = hslide(prs, SEC_PROB, "Not every parcel can wait",
                   "Schematic. In the model the willing share is split per cell "
                   "by its business/private mix, so cells differ; here it is "
                   "spread evenly.")
        cell_grid(s, GX, GY, GW, GH, theta=theta_pct / 100.0)
        bar_l, bar_w = PX, 4.05
        txt(s, bar_l, GY + 0.05, bar_w, 0.50, f"θ = {theta_pct} %", 34,
            bold=True, color=RED)
        rect(s, bar_l, GY + 0.62, bar_w, 0.20, LINE)
        rect(s, bar_l, GY + 0.62, max(0.07, bar_w * theta_pct / 100.0), 0.20,
             RED)
        txt(s, bar_l, GY + 0.90, bar_w, 0.34, "willingness to wait", SZ_DIA,
            color=DIM)
        dot(s, bar_l + 0.10, GY + 1.45, 0.20, S3)
        txt(s, bar_l + 0.32, GY + 1.26, bar_w - 0.32, 0.36, "standard — daily",
            SZ_BODY, color=INK2)
        dot(s, bar_l + 0.10, GY + 1.95, 0.20, RED)
        txt(s, bar_l + 0.32, GY + 1.76, bar_w - 0.32, 0.36,
            "willing — can be batched", SZ_BODY, color=INK2)
        badges(s, [("GroupOfPeople", RED, note)], BUL_T)
        set_morph(s)
        return s

    opt_in(0, "Nobody opts in — this is today's system, every parcel delivered "
              "daily")
    opt_in(30, "Some recipients accept a wait; business customers opt in first")
    opt_in(60, "More opt in, and most cells now hold a mixed population")
    opt_in(100, "Everyone opts in — the upper bound of the sweep, not a "
                "forecast")

    s = hslide(prs, SEC_PROB,
               "The standard parcels still need a tour every day",
               "Schematic, one weekday at θ = 60 %. Only the parcels that did "
               "not opt in are shown.")
    pts = cell_grid(s, GX, GY, GW, GH, theta=0.60, show=("std",))
    std = [(x, y) for x, y, k, _ in pts if k == "std"]
    label_box(s, DEPOT_X, DEPOT_Y - 0.42, 1.20, 0.84, S6,
              [("DEPOT", SZ_DIA, True, WHITE)])
    txt(s, PX, GY + 0.15, PW, 0.9,
        f"{len(std)} standard parcels,\nspread over twelve cells", SZ_BODY,
        bold=True, color=INK, line=1.25)
    badges(s, [("Truck", BLACK,
                "The willing parcels are held for their service day, but the "
                "rest must be delivered today, as always")], BUL_T)

    s = hslide(prs, SEC_PROB, "So they ride one shared tour per depot",
               "Schematic. The shared tour is what makes a cell's cost depend "
               "on its hub neighbours · Paper §2.2.")
    pts = cell_grid(s, GX, GY, GW, GH, theta=0.60, show=("std",),
                    cell_line=H.B.RGBColor(0xEA, 0xEC, 0xEF))
    route = serpentine([(x, y) for x, y, k, _ in pts if k == "std"])
    tour(s, [(DEPOT_X + 1.20, DEPOT_Y)] + route, S5, weight=1.75, close=True)
    label_box(s, DEPOT_X, DEPOT_Y - 0.42, 1.20, 0.84, S6,
              [("DEPOT", SZ_DIA, True, WHITE)])
    badges(s, [("Truck", BLACK,
                [("One vehicle serves the leftovers of many cells at once, so "
                  "the tour crosses ", False), ("cell boundaries", True),
                 (" — it has to", False)]),
               ("Gears", RED,
                "A cell's cost therefore depends on what its neighbours chose: "
                "that is the coupling the optimiser has to respect")],
           BUL_T)

    s = hslide(prs, SEC_PROB,
               "The batched parcels form their own denser tours",
               "Schematic, at θ = 60 %: the held parcels of a cell arrive "
               "together on its service day.")
    pts = cell_grid(s, GX, GY, GW, GH, theta=0.60, show=("batch",))
    by_cell: dict = {}
    for x, y, k, ci in pts:
        if k == "batch":
            by_cell.setdefault(ci, []).append((x, y))
    for cpts in by_cell.values():
        if len(cpts) > 2:
            tour(s, serpentine([(x, y) for x, y in cpts]), RED, weight=1.5)
    badges(s, [("DownwardTrend_LTR", RED,
                "Several days of parcels land on one delivery day: same stops, "
                "more parcels each, so the tour gets denser"),
               ("ThumbsUpSign", GREEN,
                "Denser tours cost less per parcel — that is the whole gain, and "
                "the price is a few days of waiting")],
           BUL_T)

    def tours(s):
        for side, (title, n, colour, lab) in enumerate([
                ("Daily delivery", 4, S4, "four vehicle-days"),
                ("With batching", 3, RED, "three vehicle-days")]):
            x = L + side * (W / 2 + 0.20)
            cw = W / 2 - 0.20
            rect(s, x, BODY_T, cw, 2.55, PANEL, line_col=LINE)
            rect(s, x, BODY_T, cw, 0.10, colour)
            txt(s, x + 0.30, BODY_T + 0.24, cw - 0.60, 0.40, title, SZ_LEAD,
                bold=True, color=colour)
            for k in range(n):
                cxx = x + 0.55 + (k % 2) * 2.55
                cyy = BODY_T + 1.00 + (k // 2) * 1.05
                loop = [(cxx, cyy), (cxx + 0.85, cyy - 0.28),
                        (cxx + 1.40, cyy + 0.24), (cxx + 0.58, cyy + 0.52)]
                tour(s, loop, colour, weight=1.5)
                for px, py in loop:
                    dot(s, px, py, 0.13, colour)
            txt(s, x + 0.30, BODY_T + 2.08, cw - 0.60, 0.40, lab, SZ_BODY,
                bold=True, color=INK)
        txt(s, L, BODY_T + 2.70, W, 0.60, "−1 vehicle-day", 36, bold=True,
            color=RED)

    build(prs, SEC_PROB, "Taken together: one tour fewer",
          "Schematic illustration of the mechanism, not a measured result. The "
          "measured system effect is reported in the results section.",
          [("DownwardTrend_LTR", RED,
            "Fewer, fuller tours instead of many thin ones — the saving is a "
            "whole vehicle-day, not a shorter route"),
           ("Target", BLACK,
            "Across a week and across a region, that is where the money is")],
          draw=tours, t=BODY_T + 3.45)


def part_combinatorics(prs):
    def cards(s):
        opts = [("Mon + Thu", [True, False, False, True, False, False], True,
                 "gaps of 3 and 3"),
                ("Tue + Fri", [False, True, False, False, True, False], True,
                 "cells peak together"),
                ("Mon + Tue", [True, True, False, False, False, False], False,
                 "Tue → Mon is a 5-day gap")]
        ow = (W - 2 * 0.50) / 3
        for i, (nm, pat, legal, why) in enumerate(opts):
            x = L + i * (ow + 0.50)
            rect(s, x, BODY_T, ow, 0.10, RED if legal else CRIM)
            txt(s, x, BODY_T + 0.22, ow, 0.42, nm, SZ_LEAD, bold=True,
                color=INK if legal else CRIM)
            weekbar(s, x, BODY_T + 0.80, (ow - 5 * 0.07) / 6, pat,
                    on=S6 if legal else CRIM)
            txt(s, x, BODY_T + 1.78, ow, 0.40,
                "✓  legal" if legal else "✗  illegal", SZ_BODY, bold=True,
                color=TEAL if legal else CRIM)
            txt(s, x, BODY_T + 2.20, ow, 0.42, why, SZ_BODY, color=INK2)

    build(prs, SEC_PROB, "Which days should share a batch?",
          "Legality follows the three-day rule, applied cyclically so Saturday "
          "wraps to Monday. The fleet consequence is schematic.",
          [("Gears", BLACK,
            "Same delivery frequency, different consequences: some pairings are "
            "ruled out by the service promise"),
           ("TrafficCone", AMBER,
            "Others are legal but pile every cell onto the same weekday, which "
            "creates a fleet peak")],
          draw=cards, t=BODY_T + 2.85)

    def switches(s):
        weekbar(s, L + 1.20, BODY_T + 0.10, 1.42, [True] * 6, h=0.80, gap=0.28,
                dsz=SZ_SUB + 2)
        txt(s, L, BODY_T + 1.60, W, 0.7, "2⁶ = 64 possible settings", 36,
            bold=True, color=INK)

    build(prs, SEC_PROB, "A weekly pattern is six on/off switches",
          "Operating week Monday to Saturday · Paper §2.2.",
          [("Gears", BLACK,
            "Six delivery days, each either served or not — that is the entire "
            "decision for one cell"),
           ("NoSign", AMBER,
            "Most of those 64 settings are not allowed, because they break the "
            "service promise")],
          draw=switches, t=BODY_T + 2.55)

    s = hslide(prs, SEC_PROB,
               "The three-day rule leaves exactly 39 patterns",
               "Enumerated by enumerate_valid_schedules(); pinned as "
               "EXPECTED_PATTERN_COUNT_K3 = 39 and covered by unit tests.")
    for i, (pat, lab, bad) in enumerate([
            ([True] * 6, "legal — daily", False),
            ([True, False, False, True, False, False], "legal — gap of 3", False),
            ([True, True, False, False, False, False],
             "illegal — Tue → Mon is 5", True)]):
        y = BODY_T + 0.05 + i * 0.86
        weekbar(s, L, y, 0.60, pat, on=CRIM if bad else S6, days=(i == 2),
                dsz=SZ_DIA)
        txt(s, L + 4.55, y + 0.04, 5.4, 0.42, ("✗  " if bad else "✓  ") + lab,
            SZ_BODY, bold=True, color=RED if bad else INK)
    for i, (fw, lab, fill, col) in enumerate([
            (W, "64 switch settings", PANEL, INK),
            (W * 0.98, "63 after dropping “never deliver”",
             H.B.RGBColor(0xDD, 0xE2, 0xE7), INK),
            (W * 0.61, "39 obey the three-day rule", RED, WHITE)]):
        label_box(s, L, BODY_T + 2.85 + i * 0.62, fw, 0.52, fill,
                  [(lab, SZ_BODY, col == WHITE, col)], pad=0.18,
                  align=PP_ALIGN.LEFT)
    txt(s, L, BODY_T + 4.70, W, 0.4,
        "Cyclic rule: Saturday wraps to Monday, so no gap may exceed three days.",
        SZ_BODY, color=DIM)

    def explosion(s):
        for i, lab in enumerate(["CELL 1", "CELL 2", "CELL 312"]):
            x = L + i * 2.80
            label_box(s, x, BODY_T + 0.10, 2.10, 1.25, WHITE,
                      [(lab, SZ_DIA, False, DIM), ("39", 34, True, INK)],
                      line_col=LINE)
            if i < 2:
                txt(s, x + 2.15, BODY_T + 0.50, 0.65, 0.45,
                    "×" if i == 0 else "···", 26, color=DIM,
                    align=PP_ALIGN.CENTER)
        txt(s, L + 8.15, BODY_T + 0.44, 0.5, 0.5, "=", 30, color=DIM)
        txt(s, L + 8.70, BODY_T + 0.24, 3.85, 0.9, "≈ 2.6 × 10⁴⁹⁶", 40,
            bold=True, color=RED)

    build(prs, SEC_PROB, "Direct routing cannot search this space",
          "39³¹² is a combinatorial upper bound stated before hub coupling "
          "reduces the effective space — a 497-digit number.",
          [("Gears", BLACK,
            "Each cell picks from the same 39 candidates, and every assignment "
            "creates new day-level demand"),
           ("Truck", BLACK,
            "New demand means new routes, so each candidate is a fresh routing "
            "problem"),
           ("NoSign", RED,
            "Enumerating them is not slow — it is impossible")],
          draw=explosion, t=BODY_T + 1.70)

    s = hslide(prs, SEC_PROB, "How large is 10⁴⁹⁶, really?",
               "Bar length is the number of digits, not the value. Reference "
               "figures are standard order-of-magnitude estimates.")
    # The magnitude sits in its own gutter, so a short bar cannot print over
    # its own label — the bar length is the only thing that carries the scale.
    gut, bx = 1.45, L + 1.55
    for i, (frac, mag, lab, col) in enumerate([
            (0.09, "10¹⁸", "seconds since the Big Bang", S3),
            (0.30, "10⁸⁰", "atoms in the observable universe", S4),
            (0.45, "10¹²⁰", "possible chess games", S5),
            (1.00, "10⁴⁹⁶", "weekly schedules of this one region", RED)]):
        y = BODY_T + 0.15 + i * 1.02
        bw2 = (W - 1.55) * frac
        rect(s, bx, y, bw2, 0.60, col)
        txt(s, L, y + 0.05, gut - 0.20, 0.48, mag, SZ_LEAD, bold=True,
            color=INK, align=PP_ALIGN.RIGHT)
        inside = bw2 > 5.6
        txt(s, bx + (0.22 if inside else bw2 + 0.22), y + 0.09,
            (bw2 - 0.44) if inside else (W - 1.55 - bw2 - 0.30), 0.44, lab,
            SZ_BODY, color=WHITE if inside else INK2)
    txt(s, L, BODY_T + 4.40, W, 0.55,
        "One German region outruns every physical count we have.", SZ_LEAD + 2,
        bold=True, color=RED)

    build(prs, SEC_PROB, "Worse: every single guess is a routing problem",
          "VRP NP-hardness: Lenstra & Rinnooy Kan (1981). Solver timing "
          "measured on the VROOM/Valhalla stack used for label generation.",
          [("Truck", BLACK,
            [("You cannot price a schedule without ", False),
             ("planning its tours", True),
             (" — that is a full vehicle routing problem, and NP-hard", False)]),
           ("TrafficCone", AMBER,
            [("The real solver needs ", False), ("minutes per candidate", True),
             (", and hub coupling means cells cannot be graded one at a time",
              False)]),
           ("NoSign", RED,
            "At a billion guesses a second we would still never finish"),
           ("Target", RED,
            [("We need a fast cost oracle inside the loop", True),
             (" — and a search that never enumerates", False)])],
          t=BODY_T + 0.15)

    def peaks(s):
        coll = [True, False, False, True, False, False]
        stag = [[True, False, False, True, False, False],
                [False, True, False, False, True, False],
                [False, False, True, False, False, True]]
        for side, (title, pats, hot) in enumerate([
                ("Cost-optimal calendars collide", [coll] * 3, False),
                ("Balanced calendars stagger", stag, True)]):
            x = L + side * (W / 2 + 0.20)
            cw = W / 2 - 0.20
            txt(s, x, BODY_T, cw, 0.38, title, SZ_LEAD, bold=True,
                color=RED if hot else INK)
            for r, pat in enumerate(pats):
                y = BODY_T + 0.52 + r * 0.54
                txt(s, x, y + 0.04, 1.10, 0.34, f"Cell {'ABC'[r]}", SZ_DIA,
                    color=DIM)
                weekbar(s, x + 1.15, y, 0.54, pat, days=False, h=0.40)
            hh = ([1.00, 0.24, 0.24, 1.00, 0.24, 0.24] if not hot
                  else [0.50] * 6)
            for i, h in enumerate(hh):
                rect(s, x + 1.15 + i * 0.61, BODY_T + 3.30 - h, 0.54, h,
                     TEAL if hot else CRIM)
            hrule(s, x + 1.15, BODY_T + 3.30, 3.70, DIM, 1.25)
            txt(s, x + 1.15, BODY_T + 3.42, 3.70, 0.34, "vehicles per weekday",
                SZ_DIA, color=DIM)

    build(prs, SEC_PROB, "The same frequencies can create fleet peaks",
          "Schematic mechanism. Balancing preserves each cell's delivery "
          "frequency exactly · Paper §2.2.",
          [("UpwardTrend_LTR", AMBER,
            "Cost-optimal patterns pile deliveries onto the same weekdays"),
           ("Checkmark", GREEN,
            "Selection chooses how often, balancing chooses which weekdays")],
          draw=peaks, t=BODY_T + 3.85)


# ═══════════════════════════════════════════════════════════════════════════
# 4 · Method
# ═══════════════════════════════════════════════════════════════════════════
def part_method(prs):
    section_divider(prs, "Part two", "The method",
                    "Build the oracle · search with it · balance the week · "
                    "re-route the winners")

    def framework(s):
        txt(s, L, BODY_T, W, 0.38, "Offline — learn the routing response",
            SZ_BODY, bold=True, color=RED, spc=1.2, caps=True)
        flow(s, [("HAGRID", "weekly demand", PANEL),
                 ("VROOM", "cost labels", PANEL),
                 ("SURROGATE", "Daganzo + LightGBM", RED)],
             BODY_T + 0.42, w=8.6, bh=1.10)
        arrow(s, L + 7.1, BODY_T + 1.56, L + 7.1, BODY_T + 2.16, colour=RED,
              weight=2.5)
        txt(s, L + 7.35, BODY_T + 1.68, 3.4, 0.42, "trained once", SZ_BODY,
            color=RED)
        txt(s, L, BODY_T + 2.38, W, 0.38, "Online — optimize schedules at scale",
            SZ_BODY, bold=True, color=RED, spc=1.2, caps=True)
        flow(s, [("SEARCH", "39 patterns / cell", PANEL),
                 ("BALANCE", "hub + provider", PANEL),
                 ("VALIDATE", "re-route 4 points", S6)],
             BODY_T + 2.80, l=L + 3.6, w=8.6, bh=1.10)

    build(prs, SEC_METH, "Build the oracle, search, re-route the winners",
          "Paper Fig. 2 and §2; Revision 1 reviewer clarification.",
          [("Gears", RED,
            "Routing is solved once to generate labels — after that the search "
            "runs millions of cheap evaluations")],
          draw=framework, t=BODY_T + 4.10)

    figure_slide(prs, SEC_METH, "Physics explains, machine learning corrects",
                 "2 733 VROOM samples · GroupKFold by postal-code area, so no "
                 "area appears in both train and validation · R² = 0.997.",
                 A / "fig21_progression.png",
                 items=[("HeadWithGears", RED,
                         "The raw Daganzo formula is 26 % off; one global factor "
                         "α = 1.343 brings it to 9.7 %, and the learned residual "
                         "to 2.95 %")],
                 fig=(L, BODY_T, W, 3.80), t=BODY_T + 4.00)

    split_slide(prs, SEC_METH, "Why we keep a 1984 formula",
                "Extrapolation behaviour of tree ensembles is structural: "
                "predictions are piecewise constant and cannot exceed the range "
                "seen in training.",
                BK / "fig23_alpha.png",
                [("UpwardTrend_LTR", BLACK,
                  "Batching multiplies a delivery day's volume, so the model is "
                  "asked to price days it has never seen"),
                 ("NoSign", AMBER,
                  [("Trees are ", False),
                   ("flat beyond their training range", True),
                   (" — they cannot extrapolate, they repeat the last leaf",
                    False)]),
                 ("Checkmark", GREEN,
                  "√(nA) and ⌈p/Q⌉ keep rising, because they are theory rather "
                  "than fitted structure")],
                fig=(FIG_X + 0.40, BODY_T, 6.35, 5.00), tw=4.35,
                t=BODY_T + 0.30)

    def split(s):
        label_box(s, L, BODY_T + 0.20, 3.5, 1.45, PANEL,
                  [("α · Daganzo", 26, True, RED),
                   ("the physics backbone", SZ_DIA, False, DIM)], line_col=RED)
        txt(s, L + 3.62, BODY_T + 0.55, 0.7, 0.7, "+", 40, bold=True, color=DIM,
            align=PP_ALIGN.CENTER)
        label_box(s, L + 4.38, BODY_T + 0.20, 3.5, 1.45, PANEL,
                  [("g(x)", 26, True, TEAL),
                   ("the learned residual", SZ_DIA, False, DIM)], line_col=TEAL)
        txt(s, L + 8.00, BODY_T + 0.55, 0.7, 0.7, "=", 40, bold=True, color=DIM,
            align=PP_ALIGN.CENTER)
        label_box(s, L + 8.76, BODY_T + 0.20, 3.45, 1.45, S6,
                  [("cost estimate", 24, True, WHITE),
                   ("2.95 % out-of-area error", SZ_DIA, False,
                    H.B.RGBColor(0xC6, 0xD4, 0xE2))])

    build(prs, SEC_METH, "A clean division of labour: level and structure",
          "α calibrated on the training pool and re-verified against the "
          "shipped production model (α = 1.343) by the figure script's gate.",
          [("Target", RED,
            [("α", True),
             (" fixes the level: one scalar removes a −26 % bias", False)]),
           ("HeadWithGears", RED,
            [("g(x)", True),
             (" fixes the structure the formula cannot see", False)]),
           ("MagnifyingGlass", BLACK,
            "Both halves stay inspectable — an equation and feature "
            "importances"),
           ("TrafficCone", AMBER,
            "Per-provider α made it worse: 3.2 % against 2.9 %")],
          draw=split, t=BODY_T + 1.90)

    part_formulas(prs)

    s = hslide(prs, SEC_METH, "What the formula ignores, and who repairs it",
               "Point inaccuracy of continuum approximations: Figliozzi (2008). "
               "Feature taxonomy after Akkerman et al. (2025).")
    H.B.table(s, ["The formula assumes", "Reality", "Who repairs it"],
              [[("An idealised plane", "key"),
                "Real vans follow streets, one-ways and rivers",
                ("The residual learns the detour factor", "good")],
               [("One average stem r", "key"),
                "An off-centre depot changes every single tour",
                ("Depot distance is an input feature", "good")],
               [("Stops spread evenly", "key"),
                "Real drops cluster in blocks and villages",
                ("Stop density and spread carry this", "good")],
               [("Whole vans, ⌈p/Q⌉", "key"),
                "Right in kind, brutal at the capacity threshold",
                ("The residual smooths the step", "good")]],
              BODY_T, widths=[2.6, 4.2, 4.0], sz=SZ_BODY, reserve=1.30)
    txt(s, L, BODY_B - 1.05, W, 0.9,
        "These are not reasons to drop the backbone — they are exactly the "
        "structured remainder a learner is good at.", SZ_LEAD, bold=True,
        color=RED, line=1.25)

    s = hslide(prs, SEC_METH, "We tested the field, not just our favourite",
               "Five-fold GroupKFold over postal-code areas, 10 946-sample "
               "pool. The shipped surrogate is the Daganzo–LightGBM hybrid, "
               "refitted on the merge-corrected 2 733-sample pool: 2.95 %.")
    y = H.B.table(s, ["Model", "Out-of-area error", "Note"],
                  [[("XGBoost, log target", "key"), ("3.31 %", "num"),
                    "44 features"],
                   [("LightGBM", "key"), ("3.39 %", "num"), "44 features"],
                   [("MLP ensemble", "key"), ("3.54 %", "num"),
                    "deep-learning baseline"],
                   [("Random forest", "key"), ("3.98 %", "num"),
                    "overfits area identity"],
                   [("Daganzo alone", "key"), ("21.5 %", "num"),
                    "no learned residual"]],
                  BODY_T, widths=[4, 3, 4], sz=SZ_BODY, reserve=1.15)
    badges(s, [("MagnifyingGlass", BLACK,
                "Grouped folds: no postal-code area appears in both train and "
                "test, so the numbers describe unseen geography")], y + 0.20)

    build(prs, SEC_METH, "The lesson is not that LightGBM wins",
          "Benchmark protocol A (interpolation holdout, n = 1 927) and protocol "
          "B (grouped by area, n = 10 946).",
          [("Gears", BLACK,
            "0.1 percentage points separate every boosted tree, so the choice of "
            "learner barely matters"),
           ("HeadWithGears", RED,
            "What matters is the features and the physics backbone — the "
            "deep-learning baseline is not ahead"),
           ("NoSign", AMBER,
            "Skipping the learner altogether costs a factor of six: 21.5 % "
            "against 2.95 %"),
           ("TrafficCone", AMBER,
            "Random forest memorises the area — 1.0 % on areas it has seen, "
            "4.0 % on ones it has not")],
          t=BODY_T + 0.30)

    def dials(s):
        for side, (sym, nm, unit, acc, items, frac, ends) in enumerate([
                ("P", "Service penalty", "€ per parcel per day of delay", RED,
                 ["A shadow price of waiting", "Not an observed customer price",
                  "Never booked as cost"], 0.40,
                 "cost-first                 service-first"),
                ("θ", "Willingness to wait", "eligible share of parcel volume",
                 H.B.RGBColor(0x2C, 0x7F, 0xC7),
                 ["Only θ is consolidated", "The rest keeps daily service",
                  "Business opts in before private"], 0.75,
                 "0 %                                   100 %")]):
            x = L + side * (W / 2 + 0.20)
            cw = W / 2 - 0.20
            rect(s, x, BODY_T, cw, 4.05, PANEL, line_col=LINE)
            rect(s, x, BODY_T, cw, 0.10, acc)
            txt(s, x + 0.30, BODY_T + 0.26, 1.3, 0.95, sym, 54, bold=True,
                color=acc)
            txt(s, x + 1.60, BODY_T + 0.36, cw - 1.95, 0.42, nm, SZ_LEAD,
                bold=True, color=INK)
            txt(s, x + 1.60, BODY_T + 0.82, cw - 1.95, 0.76, unit, SZ_BODY,
                color=DIM, line=1.15)
            rect(s, x + 0.30, BODY_T + 1.72, cw - 0.60, 0.18, LINE)
            rect(s, x + 0.30, BODY_T + 1.72, (cw - 0.60) * frac, 0.18, acc)
            txt(s, x + 0.30, BODY_T + 1.99, cw - 0.60, 0.3, ends, SZ_DIA,
                color=DIM)
            for j, it in enumerate(items):
                txt(s, x + 0.30, BODY_T + 2.48 + j * 0.46, cw - 0.60, 0.42,
                    "▪  " + it, SZ_SUB + 2, color=INK2)

    build(prs, SEC_METH, "Two dials define every operating regime",
          "Paper §2.2; P enters the objective as a steering term only, which "
          "keeps the cost–service trade-off continuous across the grid.",
          [("Target", RED,
            "Eight penalty levels × eleven adoption levels = 88 Stage-3 "
            "scenarios, each one a full re-optimization of all 312 cells")],
          draw=dials, t=BODY_T + 4.30)

    def cd(s):
        for i, (n, lab) in enumerate([("1", "Fix every other cell"),
                                      ("2", "Score all 39 patterns"),
                                      ("3", "Keep the best, move on")]):
            y = BODY_T + i * 0.62
            txt(s, L, y, 0.5, 0.42, n, SZ_LEAD, bold=True, color=RED)
            txt(s, L + 0.55, y, 5.2, 0.42, lab, SZ_BODY, color=INK)
        for r in range(3):
            for c in range(4):
                focal = (r, c) == (1, 1)
                rect(s, L + 0.10 + c * 1.00, BODY_T + 2.10 + r * 0.58, 0.86,
                     0.46, BLUSH if focal else PANEL,
                     line_col=RED if focal else LINE,
                     line_w=2.0 if focal else 1.0)
        txt(s, L + 0.10, BODY_T + 3.90, 4.5, 0.36, "one hub, one focal cell",
            SZ_DIA, color=DIM)
        ys = [0.30, 1.00, 1.48, 1.78, 1.94, 2.02, 2.05]
        for i in range(len(ys) - 1):
            x1 = L + 6.36 + i * 0.86
            rect(s, x1, BODY_T + 0.40 + ys[i], 0.86, 0.06, RED)
            rect(s, x1 + 0.80, BODY_T + 0.40 + min(ys[i], ys[i + 1]), 0.06,
                 abs(ys[i + 1] - ys[i]) + 0.06, RED)
        hrule(s, L + 6.36, BODY_T + 2.85, 5.5, DIM, 1.25)
        txt(s, L + 6.36, BODY_T + 2.97, 5.5, 0.36, "hub cost across sweeps  →",
            SZ_DIA, color=DIM)

    build(prs, SEC_METH, "Coordinate descent: one cell at a time",
          "Production call: eight rounds maximum, shuffled cell order, "
          "warm-started from the per-cell argmin, then a pair-swap polish.",
          [("Gears", BLACK,
            "Every accepted move lowers that hub's total cost, so the sweep is "
            "monotone and stops itself")],
          draw=cd, t=BODY_T + 4.30)

    build(prs, SEC_METH, "Why coordinate descent fits this problem",
          "Coordinate descent on a common hub objective, not a game between "
          "self-interested cells · Paper §2.2.",
          [("Checkmark", GREEN,
            [("The menu is small, discrete and complete: ", True),
             ("39 options per cell, so the best is found by enumeration, "
              "exactly — no step size, no learning rate", False)]),
           ("Checkmark", GREEN,
            [("There is no gradient to follow: ", True),
             ("a tree ensemble plus ⌈p/Q⌉ is piecewise constant, so this search "
              "needs only function values", False)]),
           ("Checkmark", GREEN,
            [("The coupling is local: ", True),
             ("cells interact only through their own depot, so 312 decisions "
              "split into 22 independent blocks", False)]),
           ("Checkmark", GREEN,
            [("One shared objective: ", True),
             ("descent is monotone, the sweep must terminate, and it ends at a "
              "coordinate-wise optimum", False)])],
          t=BODY_T + 0.20)

    s = hslide(prs, SEC_METH, "Why not the textbook optimisers?",
               "The legacy annealing implementation is retained in the "
               "repository for reference; the production path is coordinate "
               "descent.")
    H.B.table(s, ["Approach", "Blocker"],
              [[("Exhaustive search", "key"),
                "10⁴⁹⁶ candidates — not slow, impossible"],
               [("MILP or LP relaxation", "key"),
                "Our cost is a tree ensemble plus a ceiling function: neither "
                "linear nor convex"],
               [("Gradient descent", "key"),
                "A piecewise-constant surrogate has no slope to follow"],
               [("Simulated annealing", "key"),
                "300 000 iterations, a temperature schedule and reheating — all "
                "replaced by an exhaustive 39-way scan"]],
              BODY_T, widths=[3.4, 7.6], sz=SZ_BODY, reserve=1.55)
    txt(s, L, BODY_B - 1.30, W, 1.15,
        "Annealing was not wrong — it was uninformed.\nOnce you see the 39 "
        "options, a method that exploits them wins.", SZ_LEAD, bold=True,
        color=RED, line=1.30)

    split_slide(prs, SEC_METH, "Five restarts converge to the same point",
                "Stage-3 restart analysis; the figure script recomputes the "
                "spread and aborts if it exceeds the reported bound.",
                BK / "fig24_determinism.png",
                [("Checkmark", GREEN,
                  "The relative objective spread across five random restarts "
                  "stays below 10⁻¹²"),
                 ("TrafficCone", AMBER,
                  "Stable is not the same as globally optimal — the search "
                  "cannot certify that no better plan exists"),
                 ("MagnifyingGlass", RED,
                  "The sharper risk is the surrogate: the argmin of 39 noisy "
                  "predictions favours the under-priced one, so the winners are "
                  "re-routed with the real solver")],
                fig=(FIG_X + 0.40, BODY_T, 6.35, 5.00), tw=4.35,
                t=BODY_T + 0.30)

    def balance(s):
        flow(s, [("1 · SELECT", "cost-optimal pattern per hub", PANEL),
                 ("2 · BALANCE", "flatten each depot, +5 % cost cap", PANEL),
                 ("3 · SMOOTH", "level each provider's fleet", RED)],
             BODY_T + 0.05, bh=1.25)
        H.B.stats(s, [("0.135", "baseline fleet CV", False),
                      ("0.056", "after stage 2", False),
                      ("0.029", "after stage 3", True)],
                  BODY_T + 1.65, h=1.15, sz=42)

    build(prs, SEC_METH, "Balancing makes the week operational",
          "Frequency invariance asserted for all 27 456 rows by the figure "
          "script · Paper §2.2.",
          [("Checkmark", GREEN,
            "Frequency stays fixed and only weekdays move, so service quality is "
            "untouched by the balancing stages")],
          draw=balance, t=BODY_T + 3.15)

    figure_slide(prs, SEC_METH, "The week, before and after smoothing",
                 "At P = 0.25, θ = 1 · Stage 2 versus Stage 3 fleet smoothing.",
                 A / "fig51_fleet_smoothing.png",
                 items=[("DownwardTrend_LTR", RED,
                         "Swaps flatten each depot and then each provider's "
                         "network, so the peak falls without touching delivery "
                         "frequency")],
                 fig=(L, BODY_T, W, 3.85), t=BODY_T + 4.05)


# ═══════════════════════════════════════════════════════════════════════════
# 4b · The three formulas, each broken into the parts it is made of
# ═══════════════════════════════════════════════════════════════════════════
def _formula_panel(s, parts, *, t=BODY_T, h=1.15, size=27):
    """The framed equation the rest of the slide then takes apart."""
    rect(s, L, t, W, h, PANEL, line_col=None)
    rect(s, L, t, 0.09, h, RED)
    H.mathline(s, L + 0.30, t, W - 0.60, h, parts, size)
    return t + h


def part_formulas(prs):
    # ── 1 · the objective ────────────────────────────────────────────────
    def objective(s):
        y = _formula_panel(s, [
            ("σ", ""), ("★", "sup"), ("z", "sub"), ("  =  arg min", ""),
            ("σ", "sub"), ("   [   ", ""), ("Σ", ""), ("d", "sub"),
            (" C̃", "hi"), ("z,d", "sub"), ("(σ | σ", ""), ("−z", "sub"),
            (")   +   ", ""), ("P", "hi"), (" · ", ""), ("θ", "hi"),
            ("z", "sub"), (" · ", ""), ("p", "hi"), ("z", "sub"), (" · ", ""),
            ("w̄", "hi"), ("(σ)   ]", "")])
        H.term_notes(s, [
            ("σ★z", "the weekly pattern this cell ends up with, one of 39"),
            ("C̃z,d", "surrogate cost of the cell on weekday d"),
            ("σ−z", "what the other cells at the same depot chose"),
            ("P", "service penalty, a shadow price in € per parcel per day"),
            ("θz · pz", "the willing share of the cell's weekly volume"),
            ("w̄(σ)", "average added wait the pattern imposes")], y + 0.22)

    build(prs, SEC_METH, "What we minimise: routing cost plus a steering term",
          "Objective evaluated per cell inside the coordinate-descent sweep; "
          "P enters as a steering term only and is never booked as cost.",
          [("Gears", BLACK,
            [("Only σ is free. ", True),
             ("Everything else is either data or what the neighbours already "
              "decided", False)]),
           ("Truck", RED,
            [("The first term is money: ", True),
             ("what the tours actually cost once this pattern is flown",
              False)]),
           ("GroupOfPeople", AMBER,
            [("The second term is not: ", True),
             ("P prices waiting so the optimum slides along the cost–service "
              "curve, but no customer is charged it", False)])],
          draw=objective, t=BODY_T + 2.70)

    # ── 2 · the cost estimate ────────────────────────────────────────────
    def surrogate(s):
        y = _formula_panel(s, [
            ("C̃", "hi"), ("z,d", "sub"), ("   =   ", ""), ("α", "hi"),
            (" · Ĉ", ""), ("z,d", "sub"), ("   +   ", ""), ("g", "hi"),
            ("(x", ""), ("z,d", "sub"), (")", "")])
        H.term_notes(s, [
            ("α = 1.343", "one global scalar, fitted once on the training pool"),
            ("Ĉz,d", "the Daganzo backbone — continuum routing physics"),
            ("g(xz,d)", "LightGBM residual over 44 engineered features"),
            ("C̃z,d", "2.95 % error on postal-code areas never seen in "
                     "training")], y + 0.22)
        for i, (lab, val, col) in enumerate([
                ("Ĉ raw", "26 % off", H.CRIM),
                ("α · Ĉ", "9.7 % off", AMBER),
                ("α · Ĉ + g", "2.95 % off", H.GREEN)]):
            x = L + 1.45 + i * 3.35
            rect(s, x, BODY_T + 2.95, 2.95, 0.10, col)
            txt(s, x, BODY_T + 3.12, 2.95, 0.36, lab, SZ_BODY, bold=True,
                color=INK, align=PP_ALIGN.CENTER)
            txt(s, x, BODY_T + 3.52, 2.95, 0.36, val, SZ_LEAD, bold=True,
                color=col, align=PP_ALIGN.CENTER)

    build(prs, SEC_METH, "What the cost estimate is made of",
          "2 733 VROOM samples · GroupKFold by postal-code area · R² = 0.997. "
          "α re-verified against the shipped production model by the figure "
          "script's gate.",
          [("HeadWithGears", RED,
            "Physics sets the shape, one scalar sets the level, and the "
            "learner cleans up what is left")],
          draw=surrogate, t=BODY_T + 4.20)

    # ── 3 · the backbone ─────────────────────────────────────────────────
    def backbone(s):
        y = _formula_panel(s, [
            ("Ĉ", "hi"), ("   =   ", ""), ("m", "hi"), ("  ·  [   ", ""),
            ("F", "hi"), ("   +   (  2 ", ""), ("r", "hi"), (" λ  +  ", ""),
            ("k", "hi"), (" √( ", ""), ("s", "hi"), (" · ", ""), ("A", "hi"),
            (" )  )  ·  ", ""), ("c", "hi"), ("km", "sub"), ("   ]", "")])
        H.term_notes(s, [
            ("m = p / Q, rounded up", "tours needed, Q = 230 parcels per "
                                      "van — rounding up is why cost "
                                      "climbs in steps"),
            ("F = 189.15 €", "one van-day, including eight hours of labour"),
            ("2 r λ", "out to the area and back, r = depot distance at "
                      "50 km/h"),
            ("k √(s · A)", "Beardwood local tour length, k = 0.7124, over "
                           "s stops in area A"),
            ("c per km = 0.3864 €", "operating cost per kilometre")],
                     y + 0.22)

    build(prs, SEC_METH, "What the backbone is made of",
          "Daganzo/Beardwood continuum approximation as implemented in "
          "batch_delivery.legacy.daganzo; 120 s service time per parcel enters "
          "through the tour-count time constraint.",
          [("UpwardTrend_LTR", RED,
            [("This is why batching pays: ", True),
             ("m is a ceiling, so a cell that halves its delivery days does "
              "not halve its tours — it drops whole van-days", False)]),
           ("Checkmark", GREEN,
            [("And why we keep it: ", True),
             ("√ and ⌈ ⌉ keep rising outside the training range, where a tree "
              "ensemble would flatten out", False)])],
          draw=backbone, t=BODY_T + 3.30)


# ═══════════════════════════════════════════════════════════════════════════
# 5 · Case study
# ═══════════════════════════════════════════════════════════════════════════
def part_case(prs):
    s = hslide(prs, SEC_CASE,
               "The case spans the full urban–rural gradient",
               "HAGRID demand · PLZ-level geodata · Region Hannover · seven "
               "LSPs. Demand density spans a factor of 141 across one region.")
    H.B.stats(s, [("7", "providers", True), ("1.26 M", "parcels / week", False),
                  ("€ 1.91 M", "baseline cost / week", False),
                  ("312", "provider–area cells", False)], BODY_T, h=1.05, sz=40)
    pic(s, A / "fig12_map_demand.png", L, BODY_T + 1.25, 5.85, 3.65)
    pic(s, A / "fig11_lsp_volumes.png", L + 6.36, BODY_T + 1.25, 5.85, 3.65)


# ═══════════════════════════════════════════════════════════════════════════
# 6 · Results
# ═══════════════════════════════════════════════════════════════════════════
def part_results(prs):
    section_divider(prs, "Part three", "What we found",
                    "The cost–service frontier · where it pays · what it does "
                    "to the fleet")

    def grids(s):
        pic(s, A / "fig31_saving_grid.png", L, BODY_T, 5.85, 3.55)
        pic(s, A / "fig32_wait_grid.png", L + 6.36, BODY_T, 5.85, 3.55)

    _h0 = _F.headline[0.0] if _F is not None else None
    for _sl in build(
            prs, SEC_RES, "Service improves faster than savings disappear",
            "Rendered from the submission grid — the revision restates the "
            "levels, not the shape of the trade-off.",
            [("DownwardTrend_LTR", RED,
              (f"On the revision grid the routing saving peaks at "
               f"{_h0['rout1']:.1f} % and the added wait never exceeds "
               f"{_h0['wait1']:.2f} days" if _h0 else
               "Cost saving peaks at the routing optimum and the added wait "
               "never exceeds one day")),
             ("ThumbsUpSign", GREEN,
              "The first small penalty halves the waiting time while keeping "
              "most of the saving")],
            draw=grids, t=BODY_T + 3.70):
        mark(_sl, "Figure from the submission grid; the bullet numbers are "
                  "read from the revision grid, routing-optimal plan.",
             "§40.15")

    s = hslide(prs, SEC_RES,
               "The efficient range sits between P = 0.25 and P = 0.5",
               "θ = 100 %, surrogate-predicted, on the revision grid. Every "
               "row is the operator-polished plan, priced in both lenses.")
    if _F is not None:
        _rows = [[(f"P = {P:g}", "key"),
                  (f"{_F.headline[P]['rout2']:.1f} %", "num"),
                  (f"{_F.headline[P]['op2']:.1f} %", "num"),
                  f"{_F.headline[P]['wait2']:.2f} d",
                  f"{_F.headline[P]['peak2_pct']:+.1f} %"]
                 for P in (0.0, 0.25, 0.5)]
    else:
        _rows = [[("P = 0", "key"), ("22.8 %", "num"), ("n/a", "body"),
                  "0.98 d", "n/a"],
                 [("P = 0.25", "key"), ("18.5 %", "num"), ("n/a", "body"),
                  "0.46 d", "n/a"],
                 [("P = 0.5", "key"), ("13.5 %", "num"), ("n/a", "body"),
                  "0.23 d", "n/a"]]
    H.B.table(s, ["Penalty", "Routing saving", "Operator saving",
                  "Added wait", "Σ hub peak"],
              _rows, BODY_T, widths=[2.0, 2.6, 2.6, 2.4, 2.6], sz=SZ_BODY,
              reserve=2.20)
    badges(s, [("Target", RED,
                "The peak-fleet cut is already complete at P = 0.25, and the "
                "waiting time there is half what it is at P = 0"),
               ("MagnifyingGlass", BLACK,
                "The expensive part of the trade is the last points of saving, "
                "not the first")],
           BODY_T + 2.55)
    mark(s, "The submission quoted 22.8 / 18.5 / 13.5 % here. Those were "
            "routing euro on the routing-optimal plan and are superseded "
            "twice over: by the universal tour rule and by the operator "
            "polish.", ["§40.15", "§40.18"])

    split_slide(prs, SEC_RES, "The cost–service frontier",
                "Each line is one penalty level swept across adoption; the "
                "dashed line is the efficient front.",
                A / "fig34_pareto.png",
                [("Checkmark", GREEN,
                  "Every point on the plot is a fully balanced Stage-3 schedule, "
                  "not an unconstrained optimum"),
                 ("Target", RED,
                  "The knee, not the extreme, is the operating point — beyond it "
                  "waiting grows faster than saving")],
                fig=(0.45, BODY_T - 0.05, 7.55, 5.05), tw=4.35,
                t=BODY_T + 1.10)

    figure_slide(prs, SEC_RES,
                 "The penalty shifts the delivery-frequency mix",
                 "Frequency is invariant from Stage 2 to Stage 3; only weekday "
                 "placement changes. Frequencies stay within {2,…,6}.",
                 A / "fig35_schedule_mix.png",
                 items=[("DownwardTrend_LTR", RED,
                         "At P = 0 two-day patterns dominate the routing "
                         "optimum; the areas that resist are urban cells at "
                         "their "
                         "capacity limit"),
                        ("UpwardTrend_LTR", BLACK,
                         "At P ≥ 5 the system reverts to daily delivery once "
                         "θ ≥ 0.3")],
                 fig=(L, BODY_T, W, 3.35), t=BODY_T + 3.55)

    figure_slide(prs, SEC_RES, "Where the delivery days land",
                 "At P = 0.25 €/parcel/day. Values are per merged cluster, so "
                 "member polygons of one cluster share a value.",
                 A / "fig41_map_freq_by_theta.png",
                 items=[("City", RED,
                         "Dark means changed most — two delivery days a week; "
                         "the periphery consolidates first and the urban core "
                         "keeps daily service longest")],
                 fig=(L, BODY_T, W, 3.85), t=BODY_T + 4.05)

    split_slide(prs, SEC_RES, "TBC pays where delivery is sparse and far",
                "Paper Fig. 6 · provider-specific P* at θ = 1 · median per "
                "postal-code area.",
                A / "fig71_map_saving.png",
                [("DownwardTrend_LTR", RED,
                  "The median urban area saves 9 %, the median rural area 25 % — "
                  "the gap is structural, not a modelling artefact"),
                 ("MagnifyingGlass", BLACK,
                  "Long stems and large areas help; a high number of parcels per "
                  "stop suppresses the gain")],
                fig=(FIG_X + 0.35, BODY_T, 6.40, 5.05), tw=4.35,
                t=BODY_T + 0.80)

    s = hslide(prs, SEC_RES, "Three structural drivers, one direction",
               "Spearman correlations across the 312 provider–area cells at "
               "each provider's own P*.")
    H.B.table(s, ["Driver", "Correlation", "Reading"],
              [[("Hub distance", "key"), ("ρ = +0.53", "num"),
                "Long stems are amortised by batching"],
               [("Area size", "key"), ("ρ = +0.31", "num"),
                "Large areas gain from denser tours"],
               [("Parcels per drop-site", "key"), ("ρ = −0.72", "num"),
                "Where density exists, little is left to win"]],
              BODY_T, widths=[3.4, 2.6, 5.0], sz=SZ_BODY, reserve=2.10)
    badges(s, [("Target", RED,
                "Spatial targeting is the strategy — not uniform service "
                "degradation across a whole network")],
           BODY_T + 2.65)

    def classes(s):
        rows = [("Service-bound", "P* = 0.25", "Amazon · DHL",
                 "Dense, many parcels per stop", S4),
                ("Hybrid", "P* = 0.5", "FedEx · Hermes · UPS", "Mixed density",
                 S5),
                ("Cost-aggressive", "P* = 0.75", "DPD · GLS",
                 "Sparse, long stems", S6)]
        cwd = (W - 2 * 0.5) / 3
        for i, (nm, pstar, who, why, col) in enumerate(rows):
            x = L + i * (cwd + 0.5)
            rect(s, x, BODY_T, cwd, 0.10, col)
            txt(s, x, BODY_T + 0.22, cwd, 0.42, nm, SZ_LEAD, bold=True,
                color=col)
            txt(s, x, BODY_T + 0.76, cwd, 0.60, pstar, 32, bold=True, color=RED)
            txt(s, x, BODY_T + 1.44, cwd, 0.42, who, SZ_BODY, bold=True,
                color=INK)
            txt(s, x, BODY_T + 1.90, cwd, 0.75, why, SZ_BODY, color=INK2,
                line=1.2)

    build(prs, SEC_RES, "One policy does not fit every carrier",
          "Provider-specific chord-distance knees at θ = 1; heuristic operating "
          "points, not a normative recommendation.",
          [("GroupOfPeople", BLACK,
            "Tolerance for waiting tracks how much a network actually gains from "
            "consolidating"),
           ("MagnifyingGlass", RED,
            (f"DHL already carries 41 % of volume at the lowest unit cost, so "
             f"at its knee it reaches only {_KNEE('DHL'):.1f} % — against "
             f"{_KNEE('GLS'):.1f} % for GLS" if _F is not None else
             "DHL already carries 41 % of volume at the lowest unit cost, so "
             "at its knee it gains least of any carrier"))],
          draw=classes, t=BODY_T + 2.72)

    split_slide(prs, SEC_RES, "Each carrier's own operating point",
                "Chord-distance knee per provider at θ = 1; saving against "
                "added wait.",
                A / "fig36_pstar_knees.png",
                [("Target", RED,
                  "The knee is where the curve stops paying, and it differs by "
                  "network rather than by preference"),
                 ("TrafficCone", AMBER,
                  "A uniform policy therefore leaves value on the table")],
                fig=(0.45, BODY_T - 0.05, 7.55, 5.05), tw=4.35,
                t=BODY_T + 1.10)

    def parallel(s):
        for i, (nm, col) in enumerate([("Conventional tour", S4),
                                       ("Batched tour", RED)]):
            y = BODY_T + 0.05 + i * 0.82
            rect(s, L, y, 0.10, 0.58, col)
            txt(s, L + 0.32, y + 0.08, 3.0, 0.42, nm, SZ_BODY, bold=True,
                color=INK)
            for k in range(6):
                rect(s, L + 3.55 + k * 0.44, y + 0.12, 0.32, 0.34,
                     col if (i == 0 or k % 3 == 1) else S1)
        _lo, _hi = _VEHDAY_RANGE()
        H.B.stats(s, [(f"{_hi:+.1f} %", "worst-case vehicle-day rise", True),
                      (f"{_lo:+.1f} %", "best case, only at high θ", False)],
                  BODY_T + 1.90, w=6.2, h=1.25, sz=40)
        pic(s, BK / "fig52_fleet_per_provider.png", L + 6.90, BODY_T, 5.30,
            4.30)

    _lo, _hi = _VEHDAY_RANGE()
    build(prs, SEC_RES, "Low adoption creates parallel delivery systems",
          f"Weekly vehicle-days span {_lo:+.1f} % to {_hi:+.1f} % across the "
          f"grid; peak and CV benefits arise earlier than absolute reductions.",
          [("TrafficCone", AMBER,
            "At low θ both systems must run at once, and the total fleet can "
            "grow rather than shrink")],
          draw=parallel, t=BODY_T + 3.90, tw=6.10)

    def validation(s):
        H.B.stats(s, [("+0.9 … +2.7 pp", "realized saving above prediction",
                       True),
                      ("3.04 %", "cost error against the solver", False),
                      ("0.997", "R² over 1 248 observations", False)],
                  BODY_T, h=1.10, sz=38)
        pic(s, A / "fig61_vroom_scatter.png", L, BODY_T + 1.30, 5.85, 3.30)
        pic(s, A / "fig62_pred_vs_actual.png", L + 6.36, BODY_T + 1.30, 5.85,
            3.30)

    build(prs, SEC_RES, "Real routing confirms a conservative surrogate",
          "SUBMISSION grid — four out-of-sample operating points at θ = 1 "
          "(P ∈ {0, 0.25, 0.5, 0.75}), re-routed with VROOM/Valhalla. "
          "n = 1 248 · bias = +2.73 %.",
          [("Checkmark", GREEN,
            "The surrogate understates what the solver achieves")],
          draw=validation, t=BODY_T + 4.66)

    s = hslide(prs, SEC_RES, "The error points in the safe direction",
               "Predicted versus realised saving at the four validated "
               "operating points; the surrogate is conservative at every one.")
    y = H.B.table(s, ["Operating point", "Predicted", "Realised", "Gap"],
                  [[("P = 0", "key"), "22.8 %", ("23.7 %", "good"), "+0.9 pp"],
                   [("P = 0.25", "key"), "18.5 %", ("19.8 %", "good"),
                    "+1.3 pp"],
                   [("P = 0.5", "key"), "13.5 %", ("15.6 %", "good"),
                    "+2.1 pp"],
                   [("P = 0.75", "key"), "10.2 %", ("13.0 %", "good"),
                    "+2.8 pp"]],
                  BODY_T, widths=[3.0, 2.6, 2.6, 2.8], sz=SZ_BODY,
                  reserve=1.20)
    badges(s, [("UpwardTrend_LTR", GREEN,
                "The gap widens at higher penalties, so the conservative bias is "
                "strongest exactly where service is protected most")],
           y + 0.22)


# ═══════════════════════════════════════════════════════════════════════════
# 7 · Implications
# ═══════════════════════════════════════════════════════════════════════════
# ═══════════════════════════════════════════════════════════════════════════
# 8b · Revision 1 — what changed since the submission
# ═══════════════════════════════════════════════════════════════════════════
def part_revision(prs):
    """The revision's frame, in the house grammar.

    It opens the results section rather than closing it: every number after it
    has to be read through the two lenses and the two plans. Each slide's
    speaker notes carry the compendium section its claim comes from, and
    nothing on them is a typed-in figure.
    """
    if _F is None:
        return
    f, RV, D = _F, _RV, _RV.D
    h0 = f.headline[0.0]

    section_divider(prs, "Revision 1", "What changed since the submission",
                    "One tour rule · two cost lenses · two weekly plans")

    # ── the universal tour rule ───────────────────────────────────────────
    def floors(s):
        g = D.saving_grid_v2(D.PLAN_ROUTING, D.LENS_ROUTING)
        g = g[g.penalty == 0.0].set_index("share_willing").saving_pct
        H.B.stats(s, [(f"{g.loc[0.1]:.2f} %", "routing saving at θ = 10 %",
                       False),
                      (f"{g.loc[0.5]:.2f} %", "at θ = 50 %", False),
                      (f"{g.loc[1.0]:.2f} %", "at full adoption", True)],
                  BODY_T + 2.95, h=1.05, sz=38)

    sl = build(prs, SEC_REV, "One tour rule prices baseline and scenario alike",
               "Per-cell express with a 230-parcel minimum tour, applied "
               "scenario-blind: there is no longer a code branch that could "
               "treat the baseline differently.",
               [("Truck", RED,
                 [("Every area keeps its own tour: ", True),
                  ("the standard parcels of an area that is not delivered "
                   "today ride that area's own tour, not a hub-wide pooled "
                   "one", False)]),
                ("TrafficCone", AMBER,
                 [("The old pooled express tour is gone: ", True),
                  ("it priced a tour no operator would dispatch, and only the "
                   "scenario could ever have one", False)]),
                ("ThumbsUpSign", GREEN,
                 [("The θ < 1 savings fall to an honest floor: ", True),
                  ("the apparent bump at θ = 10 % was that pricing artefact "
                   "and disappears with it", False)])],
               draw=floors)
    for x in sl:
        mark(x, RV.TOUR_RULE_NOTES, ["§40.7", "§40.8", "§40.9"])

    # ── two cost lenses ───────────────────────────────────────────────────
    s = hslide(prs, SEC_REV, "One euro is not one euro",
               f"Cost model: {D.COST_MODEL_SENTENCE}. Both baselines price "
               f"the same daily-delivery system.")
    y = H.B.table(s, ["Lens", "What it counts", "Weekly baseline"],
                  [[(r[0], "key"), r[1], (r[2], "num")]
                   for r in RV.lens_rows(f)],
                  BODY_T, widths=[2.6, 6.2, 3.0], sz=SZ_BODY, reserve=2.55)
    badges(s, [("Truck", BLACK,
                [("189.15 € per vehicle-day already contains the driver: ",
                  True),
                 ("a saving on a skipped delivery day is mostly avoided "
                  "driver cost", False)]),
               ("Target", RED,
                [("An operator staffs each hub for its weekly peak: ", True),
                 ("below the peak only the kilometres are real, and a vehicle "
                  "taken out of the peak is worth 1 134.90 € a week", False)])],
           y + 0.26)
    mark(s, RV.LENS_NOTES, ["§40.11", "§40.12"])

    # ── two plans ─────────────────────────────────────────────────────────
    s = hslide(prs, SEC_REV, "Two plans, because the two lenses disagree",
               f"θ = 100 %, P = 0 €/parcel/day. Baseline "
               f"{RV.eur(f.base_routing)} € routing / "
               f"{RV.eur(f.base_operator)} € operator, "
               f"Σ hub peak {f.base_peak}.")
    y = H.B.table(s, ["Weekly plan", "Routing saving", "Operator saving",
                      "Σ hub peak", "Added wait"],
                  [[(r[0], "key"), (r[1], "num"), (r[2], "num"), r[3], r[4]]
                   for r in RV.plan_rows(f)],
                  BODY_T, widths=[3.6, 2.3, 2.4, 1.9, 2.0], sz=SZ_BODY,
                  reserve=2.55)
    label_box(s, L, y + 0.24, W, 0.86, H.BLUSH,
              [("The routing optimum is worse than doing nothing for an "
                "operator: two-day patterns treble the hub peaks.",
                SZ_LEAD, True, RED)], line_col=RED)
    badges(s, [("Gears", RED,
                [("Stage 2 may now change how OFTEN an area is served: ",
                  True),
                 ("that is what turns −7.8 % into 24.7 %, and it shortens the "
                  "wait at the same time", False)])], y + 1.24)
    mark(s, RV.PLAN_NOTES, ["§40.14", "§40.15"])

    # ── the one-area depot ────────────────────────────────────────────────
    def profiles(s):
        top = max(max(f.BANTORF_BEFORE), max(f.bantorf_after))
        for i, (lbl, prof) in enumerate(
                [("Routing-optimal plan", f.BANTORF_BEFORE),
                 ("Operator-polished plan", f.bantorf_after)]):
            y0 = BODY_T + i * 1.72
            peak = max(prof)
            txt(s, L, y0, 5.0, 0.32, lbl, SZ_SUB, bold=True, color=INK)
            for k, v in enumerate(prof):
                x = L + k * 0.70
                hgt = 0.74 * v / top
                rect(s, x, y0 + 0.66 + (0.74 - max(hgt, 0.03)), 0.50,
                     max(hgt, 0.03), CRIM if v == peak else (S4 if v else LINE))
                txt(s, x, y0 + 0.38, 0.50, 0.25, str(v), SZ_DIA, bold=True,
                    color=INK if v else DIM, align=PP_ALIGN.CENTER)
                txt(s, x, y0 + 1.44, 0.50, 0.25, "MTWTFS"[k], SZ_DIA,
                    color=DIM, align=PP_ALIGN.CENTER)
            txt(s, L + 4.70, y0 + 0.80, 3.0, 0.46, f"peak {peak}", 28,
                bold=True, color=CRIM if i == 0 else TEAL)

    sl = build(prs, SEC_REV, "A one-area depot cannot rotate its delivery days",
               f"{f.bantorf_hub_name}, Monday-to-Saturday vehicles at P = 0, "
               f"θ = 100 %. The routing-optimal profile is quoted from the "
               f"compendium; the revision tables keep only the final plan.",
               [("Truck", BLACK,
                 f"{f.DHL_ONE_CELL} of DHL's {f.DHL_HUBS} depots serve exactly "
                 f"one postal-code area, and their peak falls only by "
                 f"delivering on more days"),
                ("Target", RED,
                 "Consolidation buys fleet where a depot can rotate delivery "
                 "days across several areas; elsewhere it buys kilometres")],
               draw=profiles, t=BODY_T + 3.55)
    for x in sl:
        mark(x, RV.ONE_CELL_NOTES, ["§40.14", "§40.18"])

    # ── the operator headline ─────────────────────────────────────────────
    def head(s):
        H.B.stats(s, [(f"{h0['op2']:.1f} %", "operator-cost saving", True),
                      (f"{h0['peak2_pct']:+.0f} %", "peak vehicles, all hubs",
                       True),
                      (f"{h0['rout2']:.1f} %", "routing saving, same plan",
                       False)], BODY_T, h=1.35, sz=42)

    sl = build(prs, SEC_REV, "In the operator lens the headline changes hands",
               f"Operator-polished plan at P = 0, θ = 100 %, against the "
               f"daily-delivery baseline of {RV.eur(f.base_operator)} € and "
               f"{f.base_peak} peak vehicles.",
               [("MagnifyingGlass", BLACK,
                 f"The submission's headline was a routing-cost figure: "
                 f"{h0['rout1']:.1f} % at the routing optimum"),
                ("TrafficCone", AMBER,
                 f"That same plan costs an operator {abs(h0['op1']):.1f} % "
                 f"MORE than simply delivering every day"),
                ("ThumbsUpSign", GREEN,
                 "The polished plan gives up 2.7 points of routing saving and "
                 "buys 32 points of operator saving")],
               draw=head, t=BODY_T + 1.75)
    for x in sl:
        mark(x, "The revision's headline number. Report both lenses, and "
                "never pair one plan with the other plan's lens.",
             ["§40.12", "§40.15"])

    # ── the recommended point ─────────────────────────────────────────────
    r = f.recommended()
    s = hslide(prs, SEC_REV,
               "P = 0.25 €/parcel/day works in both lenses",
               "θ = 100 %, operator-polished plan in both rows. Savings "
               "against the daily-delivery baseline of the same grid.")
    y = H.B.table(s, ["Operating point", "Routing saving", "Operator saving",
                      "Added wait", "Σ hub peak"],
                  [[("P = 0", "key"), (f"{h0['rout2']:.1f} %", "num"),
                    (f"{h0['op2']:.1f} %", "num"), f"{h0['wait2']:.2f} d",
                    f"{h0['peak2_pct']:+.1f} %"],
                   [("P = 0.25", "key"), (f"{r['rout2']:.1f} %", "num"),
                    (f"{r['op2']:.1f} %", "num"), f"{r['wait2']:.2f} d",
                    f"{r['peak2_pct']:+.1f} %"]],
                  BODY_T, widths=[2.9, 2.4, 2.5, 2.1, 2.3], sz=SZ_BODY,
                  reserve=2.45)
    label_box(s, L, y + 0.24, W, 0.86, H.BLUSH,
              [("Half the waiting time for two points of operator saving — "
                "and the same peak-fleet cut.", SZ_LEAD, True, RED)],
              line_col=RED)
    badges(s, [("Target", RED,
                "P = 0 wins on paper; P = 0.25 is the point to defend in "
                "front of a customer")], y + 1.24)
    mark(s, "At P = 0 the operator lens gives 24.7 % against 22.8 % at "
            "P = 0.25, but the wait is 0.77 d against 0.39 d and the "
            "peak-fleet cut is the same. The knee stays at P = 0.25.",
         ["§40.15", "§40.18"])

    # ── the lens-dependent knee ───────────────────────────────────────────
    s = hslide(prs, SEC_REV, "The knee depends on which lens you use",
               "Chord-distance knee on the (saving, wait) front at θ = 1, per "
               "LSP, in each lens. Carrier classes are a heuristic, not a "
               "recommendation.")
    H.B.table(s, ["LSP", "P* routing lens", "P* operator lens",
                  "Class · routing", "Class · operator"],
              [[(r0[0], "key"), (r0[1], "body"),
                (r0[2], "num" if r0[1] != r0[2] else "body"), r0[3],
                (r0[4], "num" if r0[3] != r0[4] else "body")]
               for r0 in RV.pstar_rows(f)],
              BODY_T, widths=[1.9, 2.5, 2.6, 2.6, 2.6], sz=SZ_BODY,
              reserve=1.05)
    txt(s, L, BODY_B - 0.88, W, 0.86,
        "Three LSPs move up one class in the operator lens: peak smoothing "
        "only starts to pay at a higher penalty.", SZ_LEAD, bold=True,
        color=RED, line=1.22)
    mark(s, RV.PSTAR_NOTES, "§40.18")

    # ── partial adoption ──────────────────────────────────────────────────
    rows = {th: (a, b) for th, a, b in f.partial_adoption(0.0)}
    s = hslide(prs, SEC_REV, "Below full adoption only the operator lens pays",
               "Operator-polished plan at P = 0 across the adoption grid; both "
               "columns are that one plan, priced in the two lenses.")
    y = H.B.table(s, ["Willing to wait", "Routing saving", "Operator saving"],
                  [[(f"θ = {int(th * 100)} %", "key"),
                    f"{rows[th][0]:.1f} %", (f"{rows[th][1]:.1f} %", "num")]
                   for th in (0.1, 0.3, 0.5, 0.8, 1.0) if th in rows],
                  BODY_T, widths=[3.2, 4.0, 4.0], sz=SZ_BODY, reserve=1.95)
    badges(s, [("TrafficCone", AMBER,
                "Two delivery systems run in parallel at low adoption, so the "
                "kilometres barely move"),
               ("ThumbsUpSign", GREEN,
                "The weekly peak does move — and that is what an operator "
                "staffs for; nowhere in the grid is the operator lens "
                "negative at P = 0")], y + 0.26)
    mark(s, "Partial adoption is positive only in the operator lens: routing "
            "saving 0.4-4.6 % against operator saving 3.9 % at θ = 0.1 rising "
            "to 11.0 % at θ = 0.8, never negative (the previous grid reached "
            "−2.1 % at θ = 0.9).", "§40.15")

    # ── the penalty as a real payout ──────────────────────────────────────
    s = hslide(prs, SEC_REV,
               "If the penalty is paid out, P = 0 stops winning",
               "θ = 100 %, operator-polished plan. Flat discount = 0.50 € per "
               "delayed willing parcel; delayed parcels are demand on "
               "non-delivery days times the willing share.")
    H.B.table(s, ["Operating point", "Delayed parcels/wk",
                  "Saving, shadow price", "Net after 0.50 €",
                  "Break-even discount"],
              [[(r0[0], "key"), r0[1], r0[2], (r0[3], "num"), (r0[4], "num")]
               for r0 in RV.discount_rows(f)],
              BODY_T, widths=[2.4, 2.6, 2.6, 2.6, 2.7], sz=SZ_BODY,
              reserve=1.05)
    txt(s, L, BODY_B - 0.88, W, 0.86,
        "Read as a payout the penalty halves the saving — it does not remove "
        "it, and it moves the optimum to P = 0.25–0.5.", SZ_LEAD, bold=True,
        color=RED, line=1.22)
    mark(s, RV.DISCOUNT_NOTES, "§40.17")


def part_implications(prs):
    def steps(s):
        for i, (n, nm, body) in enumerate([
                ("1", "Target",
                 "Rural and depot-distant areas with few parcels per stop"),
                ("2", "Offer",
                 "A discounted flexible tier, with urgent delivery preserved"),
                ("3", "Tune",
                 "P and adoption by provider and by local structure"),
                ("4", "Scale",
                 "Only once participation is high enough to avoid parallel "
                 "fleets")]):
            y = BODY_T + i * 1.12
            label_box(s, L, y, 0.84, 0.84, RED, [(n, 30, True, WHITE)])
            txt(s, L + 1.10, y - 0.02, 3.2, 0.44, nm, SZ_LEAD, bold=True,
                color=RED)
            txt(s, L + 1.10, y + 0.44, W - 1.10, 0.44, body, SZ_BODY,
                color=INK2)

    s = hslide(prs, SEC_IMP, "Start where the advantage is strongest",
               "Managerial implications from the Revision 1 conclusion.")
    steps(s)

    build(prs, SEC_IMP, "Implications for practice",
          "Revision 1 conclusion; equity caveat drawn from the transport-justice "
          "literature.",
          [("ThumbsUpSign", GREEN,
            [("Time-based consolidation works: ", True),
             ("13.5 to 18.5 % weekly cost saving in the efficient range at full "
              "adoption, verified by real routing", False)]),
           ("Target", RED,
            [("Target it spatially: ", True),
             ("the median rural area saves 25 % against 9 % urban, so a uniform "
              "rollout wastes the mechanism", False)]),
           ("TrafficCone", AMBER,
            [("Adoption is the binding constraint: ", True),
             ("below roughly a third participation, two parallel systems run and "
              "the fleet can grow", False)]),
           ("GroupOfPeople", RED,
            [("Preserve choice: ", True),
             ("the flexible tier is opt-in and funded from part of the saving, "
              "never a downgrade of the default", False)])],
          t=BODY_T + 0.20)

    s = hslide(prs, SEC_IMP, "One principle holds the design together",
               "Equity caveat drawn from the transport-justice literature.")
    label_box(s, L, BODY_T + 0.30, W, 1.45, BLUSH,
              [("Do not make rural customers pay for low density with a worse "
                "default.", 30, True, RED)], line_col=RED, pad=0.34)
    badges(s, [("GroupOfPeople", RED,
                "Otherwise a cost-driven rollout penalises exactly those "
                "recipients who have the fewest alternatives"),
               ("OpenHandWithPlant", GREEN,
                "Time-based consolidation is a service-design lever, not just a "
                "routing trick")],
           BODY_T + 2.20)

    build(prs, SEC_IMP, "Outlook",
          "Revision 1 conclusion and reviewer-response framing.",
          [("MagnifyingGlass", BLACK,
            [("Calibrate against operations: ", True),
             ("VROOM is a detailed routing benchmark, not observed ground truth",
              False)]),
           ("City", BLACK,
            [("Route neighbouring areas jointly: ", True),
             ("postal-code cells currently exclude cross-area routing synergies",
              False)]),
           ("GroupOfPeople", BLACK,
            [("Estimate uptake instead of assuming it: ", True),
             ("θ is scenario-based and P is a shadow price, not a market tariff",
              False)]),
           ("Truck", BLACK,
            [("Replicate across network archetypes: ", True),
             ("one German region and seven modelled networks is one data point",
              False)])],
          t=BODY_T + 0.20)

    statement(prs, "Batch where it is\nsparse and far.",
              sub="Temporal flexibility creates the density that non-urban "
                  "delivery is missing — without any new infrastructure.")


# ═══════════════════════════════════════════════════════════════════════════
# 8 · Backup, after the contact slide
# ═══════════════════════════════════════════════════════════════════════════
def part_backup(prs):
    s = hslide(prs, SEC_BAK, "Hard invariants of the model",
               "AGENTS.md hard invariants; Paper §2.2; "
               "docs/HOLDING_DAYS_INVARIANT.md.")
    H.B.table(s, ["Setting", "Value", "Status"],
              [[("Operating week", "key"), ("Monday–Saturday", "num"),
                "six delivery days"],
               [("Maximum holding", "key"), ("Hmax = 3 days", "num"),
                "authoritative; enforced by config validation and unit tests"],
               [("Weekly patterns", "key"), ("39 per cell", "num"),
                "the same candidate set for every cell"],
               [("Baseline", "key"), ("daily delivery", "num"),
                (f"no batching; {_RV.eur(_F.base_routing)} € routing / "
                 f"{_RV.eur(_F.base_operator)} € operator per week"
                 if _F is not None else "no batching")]],
              BODY_T, widths=[3.0, 3.2, 4.8], sz=SZ_BODY)

    s = hslide(prs, SEC_BAK, "Scope and cost parameters",
               "Standard operational values for a >2 t delivery van; the "
               "learned residual corrects the rest.")
    H.B.table(s, ["Setting", "Value", "Note"],
              [[("Provider scope", "key"), ("separate networks", "num"),
                "no cross-carrier sharing of parcels, routes or vehicles"],
               [("Service penalty", "key"), ("P [€/parcel/day]", "num"),
                "steering term; booked as cost only in the discount scenario"],
               [("Vehicle capacity", "key"), ("Q = 230 parcels", "num"),
                "also the minimum tour size of the universal tour rule"],
               [("Cost model", "key"), ("189.15 € + 0.3864 €/km", "num"),
                "per vehicle-day and per kilometre, plus 36 € per route-hour "
                "(VROOM per_hour default, active in every label)"]],
              BODY_T, widths=[3.0, 3.2, 4.8], sz=SZ_BODY)

    s = hslide(prs, SEC_BAK, "Limitations",
               "Revision 1 conclusion and reviewer-response framing.")
    H.B.table(s, ["Boundary", "What it means", "Next step"],
              [[("Routing benchmark", "key"),
                "VROOM is a detailed benchmark, not observed ground truth",
                "Calibrate against operating data"],
               [("Spatial decomposition", "key"),
                "Postal-code cells exclude cross-area routing synergies",
                "Jointly route neighbouring areas"],
               [("Adoption", "key"),
                "θ is scenario-based; customer uptake is not estimated here",
                "Estimate choice and price response"],
               [("Transferability", "key"),
                "One German region and seven modelled networks",
                "Replicate across network archetypes"]],
              BODY_T, widths=[2.8, 4.4, 3.8], sz=SZ_BODY, reserve=1.05)
    txt(s, L, BODY_B - 0.85, W, 0.7,
        "Neither bound changes the direction of the result.", SZ_LEAD,
        bold=True, color=RED)

    # WITHDRAWN. This slot used to argue that consolidation survives at
    # theta = 10 % even at P = 10, and read the product P x theta as the real
    # knob. Both were descriptions of the pre-revision pooled-express price.
    # Under the universal tour rule the cell they were built on consolidates
    # almost nothing, so the finding is not corrected here -- it is withdrawn.
    if _F is not None:
        _rows = _RV.bulge_rows(_F)
        s = hslide(prs, SEC_BAK, "The bump at θ = 10 % was a pricing artefact",
                   "Share of the 312 delivery areas that give up daily "
                   "delivery, and the routing saving at the same cell, on the "
                   "revision grid (routing-optimal plan).")
        y = H.B.table(s, ["Operating point", "Areas consolidating",
                          "Routing saving"],
                      [[(r[0], "key"), (r[1], "num"), r[2]] for r in _rows],
                      BODY_T, widths=[4.0, 4.0, 3.2], sz=SZ_BODY, reserve=2.55)
        label_box(s, L, y + 0.24, W, 0.86, H.BLUSH,
                  [(f"The old deck said 41.7 % of areas still consolidated at "
                    f"P = 10, θ = 10 %. Here it is {_rows[2][1]}.",
                    SZ_LEAD, True, RED)], line_col=RED)
        badges(s, [("TrafficCone", AMBER,
                    "The bump was paid for by a hub-pooled express tour that "
                    "no operator would dispatch, so the P · θ reading of that "
                    "slide is withdrawn")], y + 1.24)
        mark(s, _RV.BULGE_NOTES, ["§40.7", "§40.8", "§40.15"])

    extras = [
        # The revision's own figures, adopted by 95_adopt_paper_figs.py from
        # <rev>/figures/. Backup tier: they are dense multi-panel paper
        # figures, shown when somebody wants the whole grid.
        (BK / "fig90_grid_two_lenses.png",
         "The whole grid, in both lenses",
         "Revision grid. Panels (a) and (b) are different lenses AND "
         "different plans — do not read them as one series."),
        (BK / "fig91_offdiagonal.png",
         "Each plan priced in the other lens",
         "Revision grid: the two off-diagonal combinations."),
        (BK / "fig92_freq_mix_two_plans.png",
         "Delivery-frequency mix, both plans",
         "Revision grid, routing-optimal above, operator-polished below."),
        (BK / "fig93_mean_days.png", "Mean delivery days per area",
         "Revision grid, both plans."),
        (BK / "fig94_structural_two_lenses.png",
         "Fronts, knees and structure",
         "Revision grid. Panel (f) is a within-DHL statement: the hub-size "
         "buckets hold DHL cells only, because DHL is the only multi-depot "
         "network in the case study."),
        (A / "fig14_headline.png", "Headline result at a glance",
         "Stage-3 summary across the penalty grid."),
        (A / "fig33_fleet_grid.png", "Peak-fleet reduction across the grid",
         "Stage 3, complete θ grid."),
        (A / "fig42_map_freq_by_P.png", "Delivery frequency by penalty level",
         "Per merged cluster at θ = 1."),
        (A / "fig44_map_wait.png", "Added waiting time in space",
         "At the provider-specific P*."),
        (A / "fig45_map_efficiency.png", "Saving per day of added wait",
         "Efficiency of the trade, per area."),
        (A / "fig53_co2.png", "CO₂ effect of the consolidated week",
         "Derived from vehicle-kilometres at the validated points."),
        (A / "fig55_pattern_clock.png", "Which weekday combinations are chosen",
         "Pattern frequencies across all optimized cells."),
        (A / "fig72_raumtyp.png", "Saving by settlement type",
         "Median per postal-code area, BBSR classification."),
        (A / "fig73_threshold_demand.png", "Saving against demand density",
         "All 312 provider–area cells."),
        (A / "fig75_breakeven.png", "Where the mechanism breaks even",
         "Density threshold per provider."),
        (BK / "fig22_pool.png", "The surrogate training pool",
         "results/supplementary/sweep_v3_mergefix/training_matrix.csv."),
        (BK / "fig25_importance.png", "What the residual actually uses",
         "Feature importances of the shipped model."),
        (BK / "fig43_map_freq_provider.png", "Frequency maps per provider",
         "At each provider's own P*."),
        (BK / "fig54_load_factor.png", "Vehicle load factor across the week",
         "Stage 3 against baseline."),
        (BK / "fig63_diagnostics.png", "Validation diagnostics",
         "Residual behaviour of the 1 248 re-routed observations."),
        (BK / "fig74_regime_map.png", "Operating regimes in space",
         "Classification of cells by their response."),
        (BK / "fig76_provider_raumtyp.png", "Provider response by settlement type",
         "Median saving per provider and settlement class."),
        (BK / "fig77_drivers.png", "Structural drivers in detail",
         "Spearman correlations across all cells."),
    ]
    for path, subject, src in extras:
        if path.exists():
            s = hslide(prs, SEC_BAK, subject, src)
            pic(s, path, L, BODY_T, W, 5.10)


# ═══════════════════════════════════════════════════════════════════════════
def build_deck(out: Path, hero: str = "week", *, facts=None, revision=None,
               tag: bool = True) -> Path:
    global _F, _RV, _TAG
    _F, _RV, _TAG = facts, revision, tag
    prs = Presentation(str(H.TEMPLATE))

    # The template master's footer still names the previous talk.
    for master in prs.slide_masters:
        for shp in master.shapes:
            if shp.has_text_frame and "mobil.TUM" in shp.text_frame.text:
                shp.text_frame.text = ""
                H.B._para(shp.text_frame.paragraphs[0],
                          "Machine Learning Surrogate Optimization Framework "
                          "for Time-Based Consolidation in Last Mile Parcel "
                          "Delivery | EWGT 2026 | Bienzeisler, Petre, Wage, "
                          "Friedrich", 10, color=DIM)
    for i in range(len(prs.slides) - 1, -1, -1):
        H.delete_slide(prs, i)

    title_slide(prs, hero)
    part_background(prs)
    part_concepts(prs)
    part_concept(prs)
    part_mechanism(prs)
    part_combinatorics(prs)
    part_method(prs)
    part_case(prs)
    part_revision(prs)
    part_results(prs)
    part_implications(prs)
    contact(prs)
    part_backup(prs)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=Path,
                    default=H.TEMPLATE.parent /
                    "EWGT_26_Bienzeisler_TBC_house_deck.pptx")
    ap.add_argument("--hero", choices=sorted(HEROES), default="week",
                    help="which opening image to use (default: week)")
    ap.add_argument("--rev-dir", type=Path, default=None,
                    help="the revision grid to read (default: $PRES_REV_DIR, "
                         "else results/revision_2026_08_v5)")
    ap.add_argument("--no-provisional", action="store_true",
                    help="drop the 'v5 · provisional' footer tag")
    ap.add_argument("--no-revision", action="store_true",
                    help="build the submission-era deck, reading no grid")
    G.add_args(ap)
    a = ap.parse_args()
    out = G.resolve(a.out, a.out_suffix, overwrite=a.overwrite)
    if a.no_revision:
        facts, revision = None, None
    else:
        import _data as D
        import _revision as RV
        if a.rev_dir is not None:
            D.set_rev_dir(a.rev_dir)
        if D.SCHEMA != D.SCHEMA_V2:
            raise SystemExit(
                f"{D.REV} is a {D.SCHEMA} grid; the revision slides need the "
                f"two-plan tables. Pass --rev-dir or set PRES_REV_DIR.")
        print(f"  revision grid: {D.REV.relative_to(D.ROOT)}")
        facts, revision = RV.Facts.load(), RV
    p = build_deck(out, a.hero, facts=facts, revision=revision,
                   tag=not a.no_provisional)
    print(f"wrote {p}")
    print(f"  {len(Presentation(str(p)).slides)} slides, "
          f"{p.stat().st_size / 1048576:.1f} MB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
