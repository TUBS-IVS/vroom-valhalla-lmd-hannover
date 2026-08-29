"""Build the combined EWGT talk deck as an editable PowerPoint file.

The deck is generated *into the institutional template*
`EWGT_26_Bienzeisler_new.pptx`, so it inherits that file's master, theme, fonts
and furniture: the grey header band, the bold Arial title in it, the TU-red
hairline above the footer, and the automatic slide number. Nothing here
re-implements that chrome; the slides only fill the body area beneath it.

Presentation rules this file enforces
-------------------------------------
* **20 pt floor.** Anything the audience reads from a seat is at least
  ``SZ_BODY`` = 20 pt. Only true marginalia go smaller: the provenance line,
  axis-style labels inside schematics, and table column heads.
* **Bullets, not prose.** Body copy is short statements in real PowerPoint
  bulleted lists, so it stays editable as a list.
* **One idea per slide.** Where a topic did not fit at 20 pt it was split
  across two slides rather than shrunk to fit.
* **Diagrams over paragraphs.** Processes, funnels and mechanisms are drawn
  from native autoshapes -- editable, and legible from the back of a room.

Everything is emitted as native shapes, so every element stays selectable in
PowerPoint. Figures come from `results/presentation_2026_08/slides/`, rendered
by 10_*..70_* through the shared colour system in `_style.py`.

Usage:
    python scripts/presentation/91_build_pptx.py [--out PATH]
    python scripts/presentation/91_build_pptx.py --out-suffix _rev2026-08

The template is opened read-only and never written to, and so is every deck
that already exists: the output path goes through `_outguard.resolve()`, which
refuses to overwrite a file unless `--overwrite` was passed on purpose.

Revision numbers
----------------
Every grid figure on a results slide is read from `_revision.Facts`, which
pulls it out of the grid `--rev-dir` (or `$PRES_REV_DIR`) points at and asserts
it against the value the compendium records. Nothing numeric on those slides is
a literal, so re-running against the v6 head grid restates the deck rather than
requiring an edit. While the grid is provisional every such slide carries the
small `v5 · provisional` footer tag; `--no-provisional` removes it.
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

import _outguard as G                                             # noqa: E402

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "results" / "presentation_2026_08" / "slides"
ASSET = ROOT / "results" / "presentation_2026_08" / "assets_generated"
TEMPLATE = Path(
    r"C:/Users/bienzeisler/Documents/Präsentationen/EWGT/2026/EWGT_26_Bienzeisler_new.pptx"
)
DEFAULT_OUT = TEMPLATE.parent / "EWGT_26_Bienzeisler_TBC_deck.pptx"

# ── palette: the template's theme red, plus the figure set's ramps ──────────
RED = RGBColor(0xBE, 0x1E, 0x3C)
INK = RGBColor(0x15, 0x18, 0x1D)
INK2 = RGBColor(0x33, 0x38, 0x40)
DIM = RGBColor(0x66, 0x6D, 0x78)
LINE = RGBColor(0xD5, 0xD9, 0xDE)
PANEL = RGBColor(0xF3, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# Schematic chrome for the hand-drawn diagrams. These are deck furniture, not
# figure colours: the rendered figures follow the paper's palette via _style.py.
S1, S2, S3, S4, S5, S6 = (
    RGBColor(0xE8, 0xEE, 0xF4), RGBColor(0xC3, 0xD5, 0xE6),
    RGBColor(0x8F, 0xB4, 0xD4), RGBColor(0x56, 0x8F, 0xB8),
    RGBColor(0x2C, 0x63, 0x96), RGBColor(0x12, 0x3F, 0x66),
)
TEAL = RGBColor(0x12, 0xA1, 0x88)
CRIM = RGBColor(0xC9, 0x30, 0x4F)
PINK = RGBColor(0xF2, 0xDB, 0xE1)
BLUSH = RGBColor(0xFD, 0xF3, 0xF5)
FONT = "Arial"

# ── type scale.  SZ_BODY is the floor for anything the room must read. ──────
SZ_TITLE = 26.0
SZ_KICK = 12.0
SZ_LEAD = 22.0
SZ_BODY = 20.0      # the floor
SZ_STAT = 46.0
SZ_STATL = 15.0
SZ_CHIP = 15.0
SZ_DIA = 15.0       # labels inside schematics
SZ_DIAB = 19.0
SZ_HEAD = 12.5      # table column heads
SZ_SRC = 9.0        # provenance line

# ── geometry, in inches, matching the master ────────────────────────────────
L = 0.63
W = 12.21
KICK_T = 1.02
BODY_T = 1.46
BODY_B = 6.60
SRC_T = 6.70
SRC_W = 11.0
COL_W = (W - 0.55) / 2
COL2 = L + COL_W + 0.55

LAYOUT_TITLE_ONLY = 4
LAYOUT_BLANK = 5


# ── primitives ─────────────────────────────────────────────────────────────
def _para(p, text, size, *, bold=False, color=INK, font=FONT, spc=None,
          italic=False, caps=False):
    run = p.add_run()
    run.text = text.upper() if caps else text
    f = run.font
    f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
    f.color.rgb = color
    if spc is not None:
        run.font._rPr.set("spc", str(int(spc * 100)))
    return run


def _bullet(p, *, colour=RED, indent=0.32):
    """Turn a paragraph into a real, editable PowerPoint bullet."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(indent))))
    pPr.set("indent", str(int(-Inches(indent))))
    clr = pPr.makeelement(qn("a:buClr"), {})
    clr.append(clr.makeelement(qn("a:srgbClr"), {"val": f"{colour}"}))
    pPr.append(clr)
    pPr.append(pPr.makeelement(qn("a:buSzPct"), {"val": "70000"}))
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "\u25aa"}))


def _frame(slide, l, t, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def txt(slide, l, t, w, h, text, size, *, bold=False, color=INK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spc=None, caps=False,
        line=None):
    box, tf = _frame(slide, l, t, w, h, anchor=anchor)
    for i, chunk in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        _para(p, chunk, size, bold=bold, color=color, spc=spc, caps=caps)
    return box


def bullets(slide, items, t, *, l=L, w=W, size=SZ_BODY, h=None, line=1.22,
            gap=12, color=INK2):
    """A real bulleted list.  `items` = [str] or [[(text, bold), ...]]."""
    box, tf = _frame(slide, l, t, w, h or (BODY_B - t))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line
        p.space_after = Pt(gap)
        _bullet(p)
        for text, bold in ([(item, False)] if isinstance(item, str) else item):
            _para(p, text, size, bold=bold, color=INK if bold else color)
    return box


def rect(slide, l, t, w, h, fill, *, line_col=None, line_w=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_col is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_col
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def label_box(slide, l, t, w, h, fill, lines, *, line_col=None, pad=0.16,
              align=PP_ALIGN.CENTER):
    """A filled box carrying (text, size, bold, colour) lines."""
    sh = rect(slide, l, t, w, h, fill, line_col=line_col)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (text, size, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.18
        _para(p, text, size, bold=bold, color=col)
    return sh


def hrule(slide, l, t, w, colour=LINE, weight=1.0):
    ln = slide.shapes.add_connector(1, Inches(l), Inches(t), Inches(l + w),
                                    Inches(t))
    ln.line.color.rgb = colour
    ln.line.width = Pt(weight)
    return ln


def arrow(slide, x1, y1, x2, y2, colour=DIM, weight=1.75, head=True):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2),
                                    Inches(y2))
    ln.line.color.rgb = colour
    ln.line.width = Pt(weight)
    if head:
        el = ln.line._get_or_add_ln()
        tail = el.find(qn("a:tailEnd"))
        if tail is None:
            tail = el.makeelement(qn("a:tailEnd"), {})
            el.append(tail)
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return ln


# ── the mechanism build-up ─────────────────────────────────────────────────
# Dot placement is fixed, not random: the four opt-in slides must carry the
# *same* dots in the *same* places so PowerPoint's Morph transition can match
# them shape for shape and animate only the fill colour. Any randomness here
# would break the animation and make the deck non-reproducible.
_DOT_OFF = [(0.24, 0.28), (0.58, 0.20), (0.79, 0.50), (0.34, 0.61),
            (0.64, 0.74), (0.17, 0.70), (0.47, 0.44), (0.86, 0.31),
            (0.30, 0.84), (0.71, 0.86)]
_CELL_N = [5, 4, 6, 5, 7, 4, 5, 6, 4, 5, 6, 5]      # dots per cell, 12 cells
_OPT_STRIDE = 23                                     # coprime with sum(_CELL_N)


def dot(slide, cx, cy, d, fill):
    """A filled circle centred on absolute (cx, cy)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2),
                                Inches(cy - d / 2), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def cell_grid(slide, l, t, w, h, *, cols=4, rows=3, theta=None, gap=0.14,
              dot_d=0.17, show=("std", "batch"), grey=S3, red=RED,
              cell_fill=WHITE, cell_line=LINE):
    """Draw the cell raster with its parcels and return the dot geometry.

    `theta` is the system-wide opt-in share in [0, 1]; a dot is a batch parcel
    when its fixed rank falls below it, so the red set grows monotonically as
    theta rises. `show` filters which kinds get drawn, which is how the later
    slides isolate the standard parcels from the batched ones.

    Returns [(cx, cy, kind, cell_index)] for every dot drawn.
    """
    cw = (w - gap * (cols - 1)) / cols
    chh = (h - gap * (rows - 1)) / rows
    total = sum(_CELL_N)
    out, i = [], 0
    for r in range(rows):
        for c in range(cols):
            ci = r * cols + c
            x0, y0 = l + c * (cw + gap), t + r * (chh + gap)
            rect(slide, x0, y0, cw, chh, cell_fill, line_col=cell_line)
            n = _CELL_N[ci % len(_CELL_N)]
            for k in range(n):
                fx, fy = _DOT_OFF[(ci * 3 + k) % len(_DOT_OFF)]
                cx = x0 + 0.16 + fx * (cw - 0.32)
                cy = y0 + 0.16 + fy * (chh - 0.32)
                rank = ((i * _OPT_STRIDE) % total) / total
                kind = "batch" if (theta is not None and rank < theta) else "std"
                if kind in show:
                    dot(slide, cx, cy, dot_d, red if kind == "batch" else grey)
                out.append((cx, cy, kind, ci))
                i += 1
    return out


def tour(slide, pts, colour, *, weight=1.5, close=True, dash=False):
    """Connect a list of (x, y) points into a route."""
    seq = list(pts)
    if close and len(seq) > 2:
        seq = seq + [seq[0]]
    for (x1, y1), (x2, y2) in zip(seq, seq[1:]):
        ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2),
                                        Inches(y2))
        ln.line.color.rgb = colour
        ln.line.width = Pt(weight)
        if dash:
            _dash(ln)
    return slide


def serpentine(pts):
    """Order points row by row, alternating direction — reads as a route."""
    rows: dict = {}
    for x, y, *_ in pts:
        rows.setdefault(round(y * 1.6), []).append((x, y))
    out = []
    for j, key in enumerate(sorted(rows)):
        row = sorted(rows[key], key=lambda p: p[0], reverse=bool(j % 2))
        out.extend(row)
    return out


_NS_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
# Morph was introduced with PowerPoint 2016 and lives in the 2015/09 namespace.
# It is NOT part of p14 (2010): writing <p14:morph> there produces a transition
# that PowerPoint parses for its duration and then silently discards the effect
# from -- verified by round-tripping the file through PowerPoint and seeing the
# element vanish. Keep this namespace.
_NS_P159 = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"


def set_morph(slide, dur=1400):
    """Give the slide a Morph transition, with a fade fallback.

    python-pptx has no transition API, so the element is injected directly.
    Morph is what turns the four opt-in slides into one continuous animation:
    identical shapes, only the fill colour differs, so PowerPoint tweens it.
    """
    from pptx.oxml import parse_xml
    trans = (f'<p:transition xmlns:p="{_NS_P}" xmlns:p14="{_NS_P14}" '
             f'spd="slow" p14:dur="{dur}">')
    xml = (
        f'<mc:AlternateContent xmlns:mc="{_NS_MC}">'
        f'<mc:Choice xmlns:p159="{_NS_P159}" Requires="p159">'
        f'{trans}<p159:morph option="byObject"/></p:transition>'
        f'</mc:Choice>'
        f'<mc:Fallback>{trans}<p:fade/></p:transition></mc:Fallback>'
        f'</mc:AlternateContent>'
    )
    slide._element.append(parse_xml(xml))
    return slide


def legend(slide, items, t, *, l=L, d=0.20, gap=0.42, size=SZ_BODY):
    """A colour key: swatch plus label, laid out in a row."""
    x = l
    for label, colour in items:
        dot(slide, x + d / 2, t + d / 2, d, colour)
        txt(slide, x + d + 0.16, t - 0.06, 4.2, 0.36, label, size, color=INK2)
        x += d + 0.16 + 0.098 * (size / 15) * len(label) + gap
    return t + 0.36


def slide_oval(slide, cx, cy_rel, d, fill):
    """A filled circle, centred on (cx, BODY_T + cy_rel)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2),
                                Inches(BODY_T + cy_rel - d / 2), Inches(d),
                                Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _dash(shape, pattern="dash"):
    """Make an autoshape's outline dashed."""
    ln = shape.line._get_or_add_ln()
    d = ln.find(qn("a:prstDash"))
    if d is None:
        d = ln.makeelement(qn("a:prstDash"), {})
        ln.append(d)
    d.set("val", pattern)
    return shape


def icon_glyph(slide, kind, l, t, sz):
    """A simple line-art icon built from autoshapes, sz inches square."""
    u = sz / 6.0                      # grid unit
    def r(dx, dy, dw, dh, fill=None):
        return rect(slide, l + dx * u, t + dy * u, dw * u, dh * u,
                    fill or RED, line_col=None)
    def o(dx, dy, dd):
        return slide_oval(slide, l + (dx + dd / 2) * u,
                          (t - BODY_T) + (dy + dd / 2) * u, dd * u, RED)
    if kind == "locker":              # a 2x2 bank of compartments
        for dx in (0.4, 3.2):
            for dy in (0.4, 3.2):
                r(dx, dy, 2.4, 2.4)
    elif kind == "hub":               # a depot: roof over a body
        slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(l),
                               Inches(t + 0.3 * u), Inches(6 * u),
                               Inches(2.4 * u)).fill.solid()
        tri = slide.shapes[-1]
        tri.fill.fore_color.rgb = RED
        tri.line.fill.background()
        tri.shadow.inherit = False
        r(1.1, 2.9, 3.8, 2.8)
    elif kind == "bike":              # two wheels and a load box
        o(0.2, 3.4, 2.2)
        o(3.6, 3.4, 2.2)
        r(1.6, 1.4, 2.8, 2.0)
    elif kind == "shop":              # a shopfront with an awning
        r(0.6, 1.1, 4.8, 0.9)
        r(1.0, 2.2, 4.0, 3.4)
    elif kind == "crowd":             # three couriers, linked
        o(2.1, 0.2, 1.8)
        o(0.2, 3.6, 1.8)
        o(4.0, 3.6, 1.8)
        r(1.1, 2.75, 3.8, 0.30)
    return slide


def axes(slide, *, xlab, ylab, ax_l=1.62, ax_r=6.95, top=0.20, bottom=4.05):
    """A bare quadrant frame: y axis, x axis and their labels."""
    y0, y1 = BODY_T + top, BODY_T + bottom
    v = slide.shapes.add_connector(1, Inches(ax_l), Inches(y0), Inches(ax_l),
                                   Inches(y1))
    v.line.color.rgb = LINE
    v.line.width = Pt(1.75)
    h = slide.shapes.add_connector(1, Inches(ax_l), Inches(y1), Inches(ax_r),
                                   Inches(y1))
    h.line.color.rgb = LINE
    h.line.width = Pt(1.75)
    txt(slide, ax_l, y1 + 0.16, ax_r - ax_l, 0.36, xlab, SZ_BODY, color=DIM,
        align=PP_ALIGN.CENTER)
    box = slide.shapes.add_textbox(Inches(ax_l - 2.05), Inches((y0 + y1) / 2 - 0.2),
                                   Inches(1.85), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _para(p, ylab, SZ_BODY, color=DIM)
    box.rotation = 270
    box.left = Inches(ax_l - 1.28)
    return ax_l, ax_r, y0, y1


def pic(slide, path, l, t, w, h):
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / (iw / 96), h / (ih / 96))
    pw, ph = (iw / 96) * scale, (ih / 96) * scale
    return slide.shapes.add_picture(str(path), Inches(l + (w - pw) / 2),
                                    Inches(t + (h - ph) / 2), Inches(pw),
                                    Inches(ph))


# ── slide scaffolding ──────────────────────────────────────────────────────
def new_slide(prs, kicker, headline, source):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
    s.shapes.title.text_frame.text = ""
    _para(s.shapes.title.text_frame.paragraphs[0], headline, SZ_TITLE, bold=True,
          color=INK)
    if kicker:
        txt(s, L, KICK_T, W, 0.3, kicker, SZ_KICK, bold=True, color=RED, spc=1.5,
            caps=True)
    if source:
        txt(s, 0.54, SRC_T, SRC_W, 0.5, source, SZ_SRC, color=DIM, line=1.2)
    return s


def stats(slide, items, t, *, l=L, w=W, h=1.5, gap=0.34, sz=SZ_STAT):
    n = len(items)
    cw = (w - gap * (n - 1)) / n
    for i, (val, label, hot) in enumerate(items):
        x = l + i * (cw + gap)
        hrule(slide, x, t, cw, RED if hot else LINE, 3.0 if hot else 1.25)
        txt(slide, x, t + 0.10, cw, sz / 58 + 0.22, val, sz, bold=True,
            color=RED if hot else INK)
        txt(slide, x, t + sz / 58 + 0.36, cw, max(0.4, h - sz / 58 - 0.36), label,
            SZ_STATL, color=DIM, line=1.2)


def chips(slide, items, t, *, l=L, maxw=W, size=SZ_CHIP):
    x, y = l, t
    for text, hot in items:
        cw = 0.098 * (size / 15) * len(text) + 0.34
        if x > l and x + cw > l + maxw:
            x, y = l, y + 0.54
        label_box(slide, x, y, cw, 0.44, PANEL,
                  [(text, size, hot, RED if hot else INK2)], line_col=LINE,
                  pad=0.10)
        x += cw + 0.13
    return y + 0.44


def table(slide, header, rows, t, *, widths=None, h=None, sz=SZ_BODY, l=L, w=W,
          reserve=0.0):
    """Draw a table and return the y coordinate of its bottom edge."""
    nr = len(rows) + (1 if header else 0)
    nc = len(rows[0])
    th = h or min(BODY_B - reserve - t, (0.46 + sz / 52) * nr)
    shp = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w),
                                 Inches(th))
    tbl = shp.table
    tbl.first_row = bool(header)
    if widths:
        tot = sum(widths)
        for i, wd in enumerate(widths):
            tbl.columns[i].width = Emu(int(Inches(w) * wd / tot))
    for ri, row in enumerate(([header] if header else []) + rows):
        tbl.rows[ri].height = Inches(th / nr)
        for ci, cell in enumerate(row):
            c = tbl.cell(ri, ci)
            c.fill.background()
            c.margin_left = Inches(0.0)
            c.margin_right = Inches(0.20)
            c.margin_top = c.margin_bottom = Inches(0.08)
            c.vertical_anchor = MSO_ANCHOR.TOP
            c.text_frame.word_wrap = True
            p = c.text_frame.paragraphs[0]
            p.line_spacing = 1.15
            text, kind = cell if isinstance(cell, tuple) else (cell, "body")
            if header and ri == 0:
                _para(p, text, SZ_HEAD, bold=True, color=DIM, spc=1.2, caps=True)
            elif kind == "key":
                _para(p, text, sz, bold=True, color=INK)
            elif kind == "num":
                _para(p, text, sz, bold=True, color=RED)
            elif kind == "good":
                _para(p, text, sz, bold=True, color=TEAL)
            else:
                _para(p, text, sz, color=INK2)
    return t + th


def weekbar(slide, l, t, cw, pattern, *, on=S6, off=S1, gap=0.07, h=0.44,
            days=True, dsz=SZ_DIA):
    for i, v in enumerate(pattern):
        col = v if isinstance(v, RGBColor) else (on if v else off)
        rect(slide, l + i * (cw + gap), t, cw, h, col)
    if days:
        for i, d in enumerate("MTWTFS"):
            txt(slide, l + i * (cw + gap), t + h + 0.05, cw, 0.28, d, dsz,
                color=DIM, align=PP_ALIGN.CENTER)


def flow(slide, boxes, t, *, l=L, w=W, bh=1.15, gap=0.36, sz=SZ_DIAB):
    """A left-to-right process chain of labelled boxes with arrows between."""
    n = len(boxes)
    bw = (w - gap * (n - 1)) / n
    for i, (title, sub, fill) in enumerate(boxes):
        x = l + i * (bw + gap)
        dark = fill in (RED, S6)
        sub_col = (RGBColor(0xF3, 0xD6, 0xDC) if fill == RED else
                   RGBColor(0xC6, 0xD4, 0xE2) if fill == S6 else DIM)
        label_box(slide, x, t, bw, bh, fill,
                  [(title, sz, True, WHITE if dark else INK)]
                  + ([(sub, SZ_DIA, False, sub_col)] if sub else []),
                  line_col=None if dark else LINE)
        if i < n - 1:
            arrow(slide, x + bw + 0.04, t + bh / 2, x + bw + gap - 0.04,
                  t + bh / 2)
    return t + bh


def full_bleed(prs, kicker, headline, body, *, closing=None):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    rect(s, 0, 0, 13.333, 7.5, RED)
    y = 1.95
    txt(s, 1.05, y, 10.5, 0.32, kicker, SZ_KICK, bold=True,
        color=RGBColor(0xF2, 0xC8, 0xD1), spc=1.6, caps=True)
    nl = headline.count("\n") + 1
    txt(s, 1.05, y + 0.42, 11.2, 0.80 * nl, headline, 40, bold=True, color=WHITE,
        line=1.06)
    yy = y + 0.52 + 0.78 * nl
    if body:
        txt(s, 1.05, yy, 9.8, 1.4, body, SZ_LEAD,
            color=RGBColor(0xF7, 0xE0, 0xE5), line=1.35)
        yy += 1.25
    if closing:
        hrule(s, 1.05, yy, 1.8, RGBColor(0xE4, 0x9A, 0xA8), 2.0)
        txt(s, 1.05, yy + 0.24, 10.8, 0.8, closing, SZ_LEAD, bold=True,
            color=WHITE, line=1.3)
    return s


def divider(prs, num, kicker, headline, body):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    rect(s, 0, 0, 13.333, 7.5, RED)
    txt(s, 8.2, 1.30, 4.5, 4.6, num, 190, bold=True,
        color=RGBColor(0xCE, 0x50, 0x68), align=PP_ALIGN.RIGHT)
    txt(s, 1.05, 3.00, 8.0, 0.32, kicker, SZ_KICK, bold=True,
        color=RGBColor(0xF2, 0xC8, 0xD1), spc=1.6, caps=True)
    txt(s, 1.05, 3.42, 8.4, 1.6, headline, 42, bold=True, color=WHITE, line=1.06)
    txt(s, 1.05, 5.05, 8.4, 0.9, body, SZ_LEAD,
        color=RGBColor(0xF7, 0xE0, 0xE5), line=1.3)
    return s


def delete_slide(prs, index):
    lst = prs.slides._sldIdLst
    items = list(lst)
    rid = items[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rid)
    lst.remove(items[index])


# ── the 2026-08 revision ───────────────────────────────────────────────────
def part_revision(prs, f, RV, *, tag: bool = True):
    """What the revision changed: one tour rule, two lenses, two plans.

    This block sits at the head of the results act, because every number after
    it has to be read through it. Nothing here is a literal: `f` is a
    `_revision.Facts` read from the grid in use, and each slide's speaker notes
    carry the compendium section its claim comes from.
    """
    def mark(s, note, cite):
        RV.notes(s, note, cite=cite)
        RV.provisional(s, enabled=tag)
        return s

    divider(prs, "R", "Revision 1", "What changed\nsince the submission",
            "One tour rule · two cost lenses · two weekly plans")

    # ── 1 · the universal tour rule ───────────────────────────────────────
    s = new_slide(prs, "Revision · 1",
                  "One tour rule prices the baseline and the scenario",
                  "Per-cell express with a 230-parcel minimum tour, applied "
                  "scenario-blind. There is no longer a code branch that could "
                  "treat the baseline differently.")
    flow(s, [("Own tour per area", "every delivery day", S1),
             ("Standard parcels ride it", "of a non-delivering area", S2),
             ("Pooled below one van", "minimum tour 230 parcels", S3),
             ("Same rule both sides", "baseline and scenario", S6)],
         BODY_T, bh=1.55)
    bullets(s, ["The old model pooled every non-delivering area of a hub into "
                "one express tour.",
                "That priced a tour no operator would dispatch — and only the "
                "scenario could ever have one.",
                "Removing it lowers the θ < 1 savings to an honest floor and "
                "dissolves the θ = 10 % bump."],
             BODY_T + 1.90, h=1.75)
    _fl = RV.D.saving_grid_v2(RV.D.PLAN_ROUTING, RV.D.LENS_ROUTING)
    _fl = _fl[_fl.penalty == 0.0].set_index("share_willing").saving_pct
    stats(s, [(f"{_fl.loc[0.1]:.2f}%", "routing saving at θ = 10 %", False),
              (f"{_fl.loc[0.5]:.2f}%", "at θ = 50 %", False),
              (f"{_fl.loc[1.0]:.2f}%", "at full adoption", True)],
          BODY_T + 3.70, h=1.00, sz=36)
    mark(s, RV.TOUR_RULE_NOTES, RV.cites("§40.7", "§40.8", "§40.9"))

    # ── 2 · two cost lenses ───────────────────────────────────────────────
    s = new_slide(prs, "Revision · 2", "One euro is not one euro",
                  f"Cost model: {RV.D.COST_MODEL_SENTENCE}. Both baselines "
                  f"price the same daily-delivery system.")
    _y = table(s, ["Lens", "What it counts", "Weekly baseline"],
               [[(r[0], "key"), r[1], (r[2], "num")] for r in RV.lens_rows(f)],
               BODY_T, widths=[2.6, 6.2, 3.0], reserve=2.15)
    bullets(s, ["189.15 € per vehicle-day already contains the driver.",
                "An operator staffs each hub for its weekly peak, so below "
                "the peak only the kilometres are a real outlay.",
                "A vehicle taken out of a hub's peak is worth 1 134.90 € a "
                "week."], _y + 0.28, h=2.0)
    mark(s, RV.LENS_NOTES, RV.cites("§40.11", "§40.12"))

    # ── 3 · two plans ─────────────────────────────────────────────────────
    h0 = f.headline[0.0]
    s = new_slide(prs, "Revision · 3",
                  "Two plans, because the two lenses disagree",
                  f"θ = 100 %, P = 0 €/parcel/day. Baseline "
                  f"{RV.eur(f.base_routing)} € routing / "
                  f"{RV.eur(f.base_operator)} € operator, "
                  f"Σ hub peak {f.base_peak}.")
    _y = table(s, ["Weekly plan", "Routing saving", "Operator saving",
                   "Σ hub peak", "Added wait"],
               [[(r[0], "key"), (r[1], "num"), (r[2], "num"), r[3], r[4]]
                for r in RV.plan_rows(f)],
               BODY_T, widths=[3.6, 2.3, 2.4, 1.9, 2.0], reserve=2.30)
    label_box(s, L, _y + 0.26, W, 0.92, BLUSH,
              [("The routing optimum is worse than doing nothing for an "
                "operator: two-day patterns treble the hub peaks.",
                SZ_LEAD, True, RED)], line_col=RED)
    bullets(s, ["Stage 2 may now change how OFTEN an area is served, not "
                "only on which days.",
                f"That is what turns {h0['op1']:.1f} % into "
                f"{h0['op2']:.1f} % — and it shortens the wait at the same "
                f"time."], _y + 1.38, h=1.18)
    mark(s, RV.plan_notes(f), RV.cites("§40.14"))

    # ── 4 · the one-cell depot ────────────────────────────────────────────
    s = new_slide(prs, "Revision · 4",
                  "A one-area depot cannot rotate its delivery days",
                  f"{f.bantorf_hub_name}, Mon–Sat vehicles at P = 0, "
                  f"θ = 100 %. The routing-optimal profile is quoted from the "
                  f"compendium; the v2 tables keep only the final plan.")
    top = max(max(f.BANTORF_BEFORE), max(f.bantorf_after))
    ROW, BAR = 1.80, 0.80        # row pitch and bar-well height, in inches
    for i, (lbl, prof, peak) in enumerate(
            [("Routing-optimal plan", f.BANTORF_BEFORE, max(f.BANTORF_BEFORE)),
             ("Operator-polished plan", f.bantorf_after, max(f.bantorf_after))]):
        y = BODY_T + i * ROW
        txt(s, L, y, 5.0, 0.34, lbl, SZ_BODY, bold=True, color=INK)
        for k, v in enumerate(prof):
            x = L + k * 0.72
            hgt = BAR * v / top
            # A zero day is a hairline on the axis, not a stub: the profile
            # "0 0 33 0 0 29" is the whole point of the slide.
            rect(s, x, y + 0.70 + (BAR - max(hgt, 0.03)), 0.52,
                 max(hgt, 0.03), CRIM if v == peak else (S4 if v else LINE))
            txt(s, x, y + 0.40, 0.52, 0.26, str(v), SZ_DIA, bold=True,
                color=INK if v else DIM, align=PP_ALIGN.CENTER)
            txt(s, x, y + 1.54, 0.52, 0.26, "MTWTFS"[k], SZ_DIA, color=DIM,
                align=PP_ALIGN.CENTER)
        txt(s, L + 4.85, y + 0.86, 3.0, 0.50, f"peak {peak}", 30, bold=True,
            color=CRIM if i == 0 else TEAL)
    bullets(s, [f"{f.one_cell_hubs} of DHL's {f.dhl_hubs} depots serve one "
                f"area only — and DHL is the only multi-depot network here.",
                "Their peak falls only by delivering on more days.",
                "Consolidation buys fleet where a depot rotates days across "
                "areas; elsewhere it buys kilometres."],
             BODY_T + 2 * ROW + 0.10, h=1.45)
    mark(s, RV.ONE_CELL_NOTES, RV.cites("§40.14"))

    # ── 5 · the operator headline ─────────────────────────────────────────
    s = new_slide(prs, "Revision · 5",
                  "In the operator lens the headline changes hands",
                  f"Operator-polished plan at P = 0, θ = 100 %, against the "
                  f"daily-delivery baseline of {RV.eur(f.base_operator)} € "
                  f"and {f.base_peak} peak vehicles.")
    stats(s, [(f"{h0['op2']:.1f}%", "operator-cost saving", True),
              (f"{h0['peak2_pct']:+.0f}%", "peak vehicles across all hubs",
               True),
              (f"{h0['rout2']:.1f}%", "routing-cost saving of the same plan",
               False)], BODY_T, h=1.5)
    bullets(s, [f"The submission's headline was a routing-cost figure: "
                f"{h0['rout1']:.1f} % at the routing optimum.",
                f"That same plan costs an operator {abs(h0['op1']):.1f} % MORE "
                f"than delivering daily.",
                f"The polished plan gives up "
                f"{h0['rout1'] - h0['rout2']:.1f} points of routing saving "
                f"and buys {h0['op2'] - h0['op1']:.1f} points of operator "
                f"saving."],
             BODY_T + 1.85, h=2.4)
    mark(s, "The revision's headline number. The routing lens is the one "
            "comparable with the submission; the operator lens is the one an "
            "LSP with salaried drivers actually faces. Report both, and never "
            "pair a plan with the other plan's lens.", ["§40.12", "§40.15"])

    # ── 6 · the recommended point ─────────────────────────────────────────
    r = f.recommended()
    s = new_slide(prs, "Revision · 6",
                  "P = 0.25 €/parcel/day is the point that works in both lenses",
                  "θ = 100 %. Savings against the daily-delivery baseline of "
                  "the same grid; wait is parcel-weighted over all parcels.")
    _y = table(s, ["Operating point", "Routing saving", "Operator saving",
                   "Added wait", "Σ hub peak"],
               [[("P = 0", "key"), (f"{h0['rout2']:.1f} %", "num"),
                 (f"{h0['op2']:.1f} %", "num"), f"{h0['wait2']:.2f} d",
                 f"{h0['peak2_pct']:+.1f} %"],
                [("P = 0.25", "key"), (f"{r['rout2']:.1f} %", "num"),
                 (f"{r['op2']:.1f} %", "num"), f"{r['wait2']:.2f} d",
                 f"{r['peak2_pct']:+.1f} %"]],
               BODY_T, widths=[2.9, 2.4, 2.5, 2.1, 2.3], reserve=2.30)
    label_box(s, L, _y + 0.26, W, 0.92, BLUSH,
              [("Half the waiting time for two points of operator saving — "
                "and the same peak-fleet cut.", SZ_LEAD, True, RED)],
              line_col=RED)
    bullets(s, ["Both rows are the operator-polished plan, so the comparison "
                "is one plan family, not two.",
                "P = 0 wins on paper; P = 0.25 is the point to defend in "
                "front of a customer."], _y + 1.38, h=1.18)
    mark(s, f"At P = 0 the operator lens is {h0['op2']:.1f} % against "
            f"{r['op2']:.1f} % at P = 0.25, but the wait is "
            f"{h0['wait2']:.2f} d against {r['wait2']:.2f} d "
            f"({(1 - r['wait2'] / h0['wait2']) * 100:.0f} % less) and the peak "
            f"cut is the same. The knee of the front stays at P = 0.25.",
         RV.cites())

    # ── 7 · the knee is lens-dependent ────────────────────────────────────
    s = new_slide(prs, "Revision · 7", "The knee depends on which lens you use",
                  "Chord-distance knee on the (saving, wait) front at θ = 1, "
                  "per LSP, in each lens. Carrier classes are a heuristic, "
                  "not a recommendation.")
    table(s, ["LSP", "P* routing lens", "P* operator lens",
              "Class · routing", "Class · operator"],
          [[(r0[0], "key"), (r0[1], "body"),
            (r0[2], "num" if r0[1] != r0[2] else "body"), r0[3],
            (r0[4], "num" if r0[3] != r0[4] else "body")]
           for r0 in RV.pstar_rows(f)],
          BODY_T, widths=[1.9, 2.5, 2.6, 2.6, 2.6], reserve=1.15)
    txt(s, L, BODY_B - 1.28, W, 1.22, RV.pstar_headline(f), SZ_LEAD,
        bold=True, color=RED, line=1.25)
    mark(s, RV.pstar_notes(f), RV.cites())

    # ── 8 · partial adoption ──────────────────────────────────────────────
    rows = {th: (a, b) for th, a, b in f.partial_adoption(0.0)}
    s = new_slide(prs, "Revision · 8",
                  "Below full adoption only the operator lens pays",
                  "Operator-polished plan at P = 0 across the adoption grid. "
                  "Both columns are the same plan, priced in the two lenses.")
    _y = table(s, ["Willing to wait", "Routing saving", "Operator saving"],
               [[(f"θ = {int(th * 100)} %", "key"),
                 f"{rows[th][0]:.1f} %", (f"{rows[th][1]:.1f} %", "num")]
                for th in (0.1, 0.3, 0.5, 0.8, 1.0) if th in rows],
               BODY_T, widths=[3.2, 4.0, 4.0], reserve=1.95)
    bullets(s, ["Two delivery systems run in parallel at low adoption: the "
                "batched one and the daily one.",
                "The kilometres barely move — but the weekly peak does, and "
                "that is what an operator staffs for.",
                "Nowhere in the grid is the operator lens negative at P = 0."],
             _y + 0.26, h=1.70)
    _rv = [v[0] for v in rows.values()]
    _ov = [v[1] for v in rows.values()]
    mark(s, f"Partial adoption is positive only in the operator lens: routing "
            f"saving {min(_rv):.1f}-{max(_rv):.1f} % against operator saving "
            f"{min(_ov):.1f}-{max(_ov):.1f} % across theta, and never negative "
            f"at P = 0.", RV.cites())

    # ── 9 · the penalty as a real payout ──────────────────────────────────
    s = new_slide(prs, "Revision · 9",
                  "If the penalty is actually paid out, P = 0 stops winning",
                  "θ = 100 %, operator-polished plan. Flat discount = 0.50 € "
                  "per delayed willing parcel; delayed parcels are demand on "
                  "non-delivery days times the willing share.")
    table(s, ["Operating point", "Delayed parcels/wk", "Saving, shadow price",
              "Net after 0.50 €, operator lens",
              "Net after 0.50 €, routing lens", "Break-even discount"],
          [[(r0[0], "key"), r0[1], r0[2], (r0[3], "num"), (r0[4], "num"),
            r0[5]] for r0 in RV.discount_rows(f)],
          BODY_T, widths=[1.9, 2.1, 2.1, 2.2, 2.2, 2.0], reserve=1.15)
    txt(s, L, BODY_B - 1.28, W, 1.22, RV.discount_optimum_line(f), SZ_LEAD,
        bold=True, color=RED, line=1.25)
    mark(s, RV.discount_notes(f), RV.cites("§40.17"))
    return prs


def slide_bulge_replacement(prs, f, RV, *, tag: bool = True):
    """Replaces the old "the θ = 10 % bump survives a punitive fee" slide.

    That slide described an artefact of the pre-revision pooled express price.
    Under the universal tour rule the cell it was built on consolidates almost
    nothing and saves almost nothing, so the slide is not corrected — it is
    withdrawn, and this one says why.
    """
    s = new_slide(prs, "Appendix · withdrawn finding",
                  "The bump at θ = 10 % was a pricing artefact",
                  "Share of the 312 delivery areas that give up daily "
                  "delivery, and the routing saving at the same cell, on the "
                  "revision grid (routing-optimal plan).")
    _rows = RV.bulge_rows(f)
    _old, _new = 41.7, float(_rows[2][1].rstrip(" %"))
    _y = table(s, ["Operating point", "Areas consolidating", "Routing saving"],
               [[(r[0], "key"), (r[1], "num"), r[2]] for r in _rows],
               BODY_T, widths=[4.0, 4.0, 3.2], reserve=2.55)
    label_box(s, L, _y + 0.24, W, 0.92, BLUSH,
              [(f"The old deck said {_old:.1f} % of areas still consolidated "
                f"at P = 10, θ = 10 %. Here it is {_new:.1f} %.",
                SZ_LEAD, True, RED)], line_col=RED)
    bullets(s, ["The bump was paid for by a hub-pooled express tour no "
                "operator would dispatch.",
                "The P · θ reading of that slide is withdrawn."],
             _y + 1.34, h=1.05)
    RV.notes(s, RV.bulge_notes(f), cite=RV.cites("§40.7", "§40.8"))
    RV.provisional(s, enabled=tag)
    return s


# ── the deck ───────────────────────────────────────────────────────────────
def build(out: Path, keep_template_slides: bool, *, facts=None,
          revision=None, tag: bool = True) -> Path:
    prs = Presentation(str(TEMPLATE))

    # The template master's footer still names the previous talk. Retitle it
    # once here rather than per slide, so it stays a single source of truth.
    for master in prs.slide_masters:
        for shp in master.shapes:
            if shp.has_text_frame and "mobil.TUM" in shp.text_frame.text:
                shp.text_frame.text = ""
                _para(shp.text_frame.paragraphs[0],
                      "Machine-Learning Surrogate Optimization for Time-Based "
                      "Consolidation in Last-Mile Parcel Delivery | EWGT 2026 | "
                      "Bienzeisler, Petre, Wage, Friedrich", 10, color=DIM)

    if not keep_template_slides:
        for i in range(len(prs.slides) - 1, -1, -1):
            delete_slide(prs, i)

    A, B = FIG / "tierA", FIG / "tierB"

    # ═══ title ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    rect(s, L, 1.75, 0.09, 3.35, RED)
    txt(s, L + 0.40, 1.82, 10.5, 0.34, "Machine-learning surrogate optimization",
        SZ_KICK, bold=True, color=RED, spc=1.6, caps=True)
    txt(s, L + 0.40, 2.28, 11.5, 1.7,
        "Time-Based Consolidation\nin Last-Mile Delivery", 44, bold=True,
        color=INK, line=1.06)
    txt(s, L + 0.40, 4.05, 11.5, 1.3,
        "Lasse Bienzeisler · Felix Petre · Oskar Wage · Bernhard Friedrich\n"
        "TU Braunschweig · Leibniz University Hannover\n"
        "EWGT 2026 · Revision 1 · manuscript under review",
        SZ_STATL + 1, color=INK2, line=1.7)

    # ═══ the pressure ═════════════════════════════════════════════════════
    s = new_slide(prs, "The pressure", "Every click becomes a delivery task",
                  "Eurostat (2026). Illustration generated with OpenAI ImageGen.")
    stats(s, [("62%", "shopped online in 2015", False),
              ("78%", "do so in 2025", True)], BODY_T, w=5.6, h=1.35)
    bullets(s, ["More demand is not only more parcels.",
                "It is more stops, more routes, more vehicles.",
                "All of it lands on the last mile."],
            BODY_T + 1.70, w=5.6, h=2.5)
    pic(s, ASSET / "parcel-wave.png", 6.55, BODY_T, 6.25, 5.0)

    # ═══ the bottleneck ═══════════════════════════════════════════════════
    s = new_slide(prs, "The bottleneck", "The last mile is a density problem",
                  "Bienzeisler et al. (2026). Illustration generated with "
                  "OpenAI ImageGen.")
    pic(s, ASSET / "urban-rural-density.png", L, BODY_T, W, 3.30)
    y = BODY_T + 3.50
    for i, (t1, t2, col) in enumerate([
            ("Many drops, short stems", "one van serves a whole street", S6),
            ("Few drops, long stems", "kilometres between single stops", RED)]):
        x = L + i * (COL_W + 0.55)
        rect(s, x, y, 0.10, 0.95, col)
        txt(s, x + 0.30, y, COL_W - 0.30, 0.42, t1, SZ_LEAD, bold=True, color=INK)
        txt(s, x + 0.30, y + 0.48, COL_W - 0.30, 0.5, t2, SZ_BODY, color=DIM)

    # ═══ the toolbox ══════════════════════════════════════════════════════
    s = new_slide(prs, "Motivation", "The toolbox mostly consolidates in space",
                  "Literature synthesis. Illustration generated with OpenAI "
                  "ImageGen.")
    pic(s, ASSET / "urban-logistics-concepts.png", L, BODY_T, W, 3.80)
    chips(s, [("Parcel locker", False), ("Micro-hub + cargo bike", False),
              ("Urban consolidation centre", False),
              ("White-label pooling", True)], BODY_T + 4.00)

    # ═══ the five concepts, named ═════════════════════════════════════════
    s = new_slide(prs, "The toolbox", "Many last-mile concepts already exist",
                  "Literature synthesis; the five recurring families of last-mile "
                  "concepts in the European context.")
    concepts = [("Parcel lockers", "self-service pickup", "locker"),
                ("Micro-hubs", "local depots", "hub"),
                ("Cargo bikes", "zero-emission tours", "bike"),
                ("Parcel shops", "retail drop-off", "shop"),
                ("Crowdshipping", "people as couriers", "crowd")]
    kw = (W - 4 * 0.20) / 5
    for i, (nm, sub, icon) in enumerate(concepts):
        x = L + i * (kw + 0.20)
        rect(s, x, BODY_T + 0.15, kw, 2.85, PANEL, line_col=LINE)
        icon_glyph(s, icon, x + kw / 2 - 0.42, BODY_T + 0.50, 0.84)
        txt(s, x + 0.12, BODY_T + 1.70, kw - 0.24, 0.5, nm, SZ_BODY, bold=True,
            color=INK, align=PP_ALIGN.CENTER, line=1.15)
        txt(s, x + 0.12, BODY_T + 2.28, kw - 0.24, 0.5, sub, SZ_STATL, color=DIM,
            align=PP_ALIGN.CENTER, line=1.15)
    bullets(s, ["Each one is well studied and, in the right place, it works.",
                [("All five change ", False), ("where", True),
                 (" parcels are handed over.", False)],
                [("None of them changes ", False), ("when", True), (".", False)]],
            BODY_T + 3.30, h=2.1)

    # ═══ viability rises with density ═════════════════════════════════════
    s = new_slide(prs, "One shared assumption",
                  "Viability depends on demand density",
                  "Conceptual synthesis of the locker, micro-hub, UCC and "
                  "crowdshipping literature.")
    axes(s, xlab="demand density  →", ylab="viability")
    bullets(s, ["Every concept has a break-even density.",
                "Below it, the fixed cost of the place is not recovered.",
                "So the question is not which concept is best, but which parts "
                "of a region reach that density at all."],
            BODY_T + 0.35, l=7.35, w=5.3, h=3.6)
    label_box(s, 2.05, BODY_T + 2.55, 2.35, 1.05, None,
              [("suburban", SZ_BODY, True, DIM), ("& rural?", SZ_BODY, True, DIM)],
              line_col=DIM)
    _dash(s.shapes[-1])

    # ═══ the concepts, positioned ═════════════════════════════════════════
    s = new_slide(prs, "The gap", "The concepts cluster where density already is",
                  "Positions are qualitative, from the literature's stated "
                  "preconditions — not measured break-even points.")
    axes(s, xlab="demand density  →", ylab="viability")
    for cx, cy in [(5.35, 2.30), (5.95, 1.98), (6.45, 2.62), (5.62, 2.92),
                   (6.20, 3.20)]:
        slide_oval(s, cx, cy, 0.30, RED)
    txt(s, 4.85, BODY_T + 0.18, 2.4, 0.36, "urban concepts", SZ_STATL, bold=True,
        color=RED, align=PP_ALIGN.CENTER)
    label_box(s, 2.05, BODY_T + 2.55, 2.35, 1.05, None,
              [("suburban", SZ_BODY, True, DIM), ("& rural", SZ_BODY, True, DIM)],
              line_col=DIM)
    _dash(s.shapes[-1])
    txt(s, 1.85, BODY_T + 3.70, 3.0, 0.36, "no concept lands here", SZ_STATL,
        bold=True, color=DIM, align=PP_ALIGN.CENTER)
    bullets(s, ["Lockers, hubs and bikes pay off in dense cities.",
                "Where demand thins out, they stop making sense.",
                [("Yet these areas carry a ", False), ("large share", True),
                 (" of all delivery activity.", False)],
                "They need a lever that works at low density."],
            BODY_T + 0.35, l=7.35, w=5.3, h=3.6)

    # ═══ density is the condition ═════════════════════════════════════════
    s = new_slide(prs, "Motivation", "Every concept needs the same thing: density",
                  "UCC, locker and horizontal-collaboration literature; EU "
                  "horizontal-cooperation guidance.")
    bw, gp = 2.55, 0.62
    x0 = L + 0.15
    label_box(s, x0, BODY_T + 0.30, bw, 1.55, S6,
              [("IMPACT", 26, True, WHITE),
               ("of any concept", SZ_DIA, False, RGBColor(0xC6, 0xD4, 0xE2))])
    for i, (nm, sub) in enumerate([("volume", "enough parcels"),
                                   ("proximity", "short access"),
                                   ("participation", "users or carriers")]):
        x = x0 + (i + 1) * (bw + gp)
        txt(s, x - gp, BODY_T + 0.85, gp, 0.55, "=" if i == 0 else "×", 30,
            bold=True, color=DIM, align=PP_ALIGN.CENTER)
        label_box(s, x, BODY_T + 0.30, bw, 1.55, PANEL,
                  [(nm, 24, True, RED), (sub, SZ_DIA, False, DIM)], line_col=LINE)
    txt(s, L, BODY_T + 2.35, W, 0.6,
        "Urbanity is not the objective. Density is the performance condition.",
        SZ_LEAD + 2, bold=True, color=RED)
    bullets(s, ["Spatial concepts move parcels or people closer together.",
                "They all need a place — and enough local volume.",
                "Outside the city that volume simply is not there."],
            BODY_T + 3.15, h=2.0)

    # ═══ what each concept needs ══════════════════════════════════════════
    s = new_slide(prs, "Motivation", "What each spatial concept actually requires",
                  "UCC, locker and collaboration literature; EU horizontal-"
                  "cooperation guidance.")
    _y = table(s, ["Concept", "Precondition"],
          [[("Parcel locker", "key"), "Proximity plus a high fill rate"],
           [("Micro-hub", "key"), "A dense catchment and a short final leg"],
           [("Urban consolidation centre", "key"),
            "Critical volume and a viable business model"],
           [("White-label pooling", "key"),
            "Data sharing and neutral governance between competitors"]],
          BODY_T, widths=[3, 5], reserve=1.15)
    label_box(s, L, _y + 0.28, W, 0.90, PANEL,
              [("Each precondition is hardest to meet exactly where cost is "
                "highest.", SZ_LEAD, True, RED)], line_col=RED,
              align=PP_ALIGN.LEFT)

    # ═══ outside the core ═════════════════════════════════════════════════
    s = new_slide(prs, "Motivation", "Most parcels are delivered outside the core",
                  "Bienzeisler et al. (2026), HAGRID study. Settlement type per "
                  "postal-code area, BBSR classification.")
    stats(s, [("77.6%", "suburban or rural", True), ("59.0%", "suburban", False),
              ("18.6%", "rural", False)], BODY_T, w=6.1, h=1.4, sz=40)
    bullets(s, [[("Rural cost per parcel is about ", False), ("70% higher", True),
                 (".", False)],
                "That is where the spatial toolbox has least to offer.",
                "So the lever must be something other than space."],
            BODY_T + 1.75, w=6.1, h=2.6)
    pic(s, B / "fig13_map_raumtyp.png", 6.85, BODY_T, 5.95, 4.6)

    # ═══ two levers ═══════════════════════════════════════════════════════
    s = new_slide(prs, "Concept", "Consolidation has two levers: where and when",
                  "Conceptual synthesis of the spatial and temporal consolidation "
                  "literature.")
    for side, (nm, claim, items, acc, fill) in enumerate([
            ("Where", "Bring parcels together in space",
             ["Parcel locker", "Micro-hub or UCC", "Carrier pooling"], DIM, PANEL),
            ("When", "Bring order days together in time",
             ["Hold eligible parcels", "Deliver them batched",
              "No new infrastructure"], RED, BLUSH)]):
        x = L + side * (COL_W + 0.55)
        rect(s, x, BODY_T, COL_W, 4.50, fill, line_col=LINE)
        rect(s, x, BODY_T, COL_W, 0.10, acc)
        txt(s, x + 0.32, BODY_T + 0.30, COL_W - 0.64, 0.42, nm, SZ_BODY,
            bold=True, color=acc, spc=1.4, caps=True)
        txt(s, x + 0.32, BODY_T + 0.82, COL_W - 0.64, 0.95, claim, SZ_LEAD,
            bold=True, color=INK, line=1.2)
        bullets(s, items, BODY_T + 1.95, l=x + 0.32, w=COL_W - 0.64, h=1.7)
        if side:
            weekbar(s, x + 0.32, BODY_T + 3.50, 0.66,
                    [True, False, False, True, False, False], on=RED,
                    off=RGBColor(0xE9, 0xCF, 0xD6), days=False)
        else:
            txt(s, x + 0.32, BODY_T + 3.55, COL_W - 0.64, 0.6,
                "Needs a place and local volume.", SZ_BODY, color=DIM)

    # ═══ statement ════════════════════════════════════════════════════════
    full_bleed(prs, "The service question",
               "If urgency differs,\nwhy is next-day the default?",
               "Urgent parcels stay daily. Flexible parcels can wait — for a "
               "discount, a lower fee, or a different service tier.",
               closing="Time becomes an operational control variable.")

    # ═══ hold and batch ══════════════════════════════════════════════════
    s = new_slide(prs, "Concept", "Hold eligible parcels. Batch them later.",
                  "Paper §2.2. Hmax = 3 days yields 39 feasible weekly patterns "
                  "per cell.")
    cw, gp = 1.62, 0.24
    x0 = L + 1.85
    for i, d in enumerate(["MON", "TUE", "WED", "THU", "FRI", "SAT"]):
        txt(s, x0 + i * (cw + gp), BODY_T, cw, 0.32, d, SZ_BODY, bold=True,
            color=DIM, align=PP_ALIGN.CENTER)
    txt(s, L, BODY_T + 0.62, 1.80, 0.36, "eligible θ", SZ_BODY, bold=True,
        color=RED)
    for i, on in enumerate([False, True, False, False, True, False]):
        label_box(s, x0 + i * (cw + gp), BODY_T + 0.48, cw, 0.70,
                  RED if on else PINK,
                  [("DELIVER" if on else "HOLD", SZ_STATL + 1, on,
                    WHITE if on else RGBColor(0x8F, 0x14, 0x2B))])
    txt(s, L, BODY_T + 1.98, 1.80, 0.36, "standard", SZ_BODY, bold=True, color=DIM)
    for i in range(6):
        rect(s, x0 + i * (cw + gp), BODY_T + 1.88, cw, 0.58, S6)
    bullets(s, ["Willing parcels wait for the next service day.",
                "Non-willing parcels keep a conventional daily service.",
                "Maximum holding time: three days, Monday to Saturday."],
            BODY_T + 2.80, h=2.1)

    # ═══ one calendar, three effects ══════════════════════════════════════
    s = new_slide(prs, "From concept to decision",
                  "One calendar decision moves three things at once",
                  "Schematic for the eligible share; standard parcels retain daily "
                  "service · Paper §2.2.")
    label_box(s, L + 3.55, BODY_T, 5.1, 0.90, S6,
              [("ONE PATTERN   σ = Tue · Thu · Sat", SZ_DIAB, True, WHITE)])
    eff = [("Customer wait", "the gap to the next service day", RED),
           ("Routing cost", "denser tours, fewer of them", S5),
           ("Vehicle demand", "which weekday the fleet is needed", S6)]
    bwd = (W - 2 * 0.50) / 3
    for i, (nm, sub, col) in enumerate(eff):
        x = L + i * (bwd + 0.50)
        arrow(s, L + 6.1, BODY_T + 0.95, x + bwd / 2, BODY_T + 1.68, colour=col)
        label_box(s, x, BODY_T + 1.72, bwd, 1.25, PANEL,
                  [(nm, SZ_LEAD, True, col), (sub, SZ_STATL + 1, False, DIM)],
                  line_col=LINE)
    bullets(s, ["The three effects cannot be tuned separately.",
                "The pattern that saves most also waits longest.",
                "So the schedule is the decision variable — not the route."],
            BODY_T + 3.25, h=2.1)

    # ═══ no single calendar ═══════════════════════════════════════════════
    s = new_slide(prs, "The combinatorial decision",
                  "No single calendar fits every area",
                  "Illustrative cells — not an empirical classification. 39 "
                  "admissible patterns per cell, Paper §2.2.")
    cells = [("Dense urban", [S6] * 6, "6 days / week",
              "Existing density already pays", False),
             ("Suburban", [S4, S1, S4, S1, S4, S4], "4 days / week",
              "Cost and waiting balance out", False),
             ("Sparse rural", [S1, S6, S1, S1, S6, S1], "2 days / week",
              "Batching amortises the long stem", True)]
    ccw = (W - 2 * 0.55) / 3
    for i, (nm, pat, freq, expl, hot) in enumerate(cells):
        x = L + i * (ccw + 0.55)
        hrule(s, x, BODY_T, ccw, RED if hot else LINE, 3.0 if hot else 1.25)
        txt(s, x, BODY_T + 0.14, ccw, 0.42, nm, SZ_LEAD, bold=True,
            color=RED if hot else INK)
        weekbar(s, x, BODY_T + 0.80, (ccw - 5 * 0.07) / 6, pat, h=0.52)
        txt(s, x, BODY_T + 1.75, ccw, 0.48, freq, 26, bold=True, color=INK)
        txt(s, x, BODY_T + 2.30, ccw, 0.85, expl, SZ_BODY, color=INK2, line=1.2)
    hrule(s, L, BODY_T + 3.35, W, RED, 2.0)
    txt(s, L, BODY_T + 3.55, W, 0.6,
        "312 cells × 39 calendars × shared hub effects", 30, bold=True, color=RED)
    txt(s, L, BODY_T + 4.22, W, 0.5,
        "Same menu everywhere. Different right answer everywhere.", SZ_LEAD,
        color=INK2)

    # ═══ divider 01 · the mechanism ═══════════════════════════════════════
    divider(prs, "01", "Part one", "The mechanism,\nstep by step",
            "From arriving parcels to one tour fewer — and to the question of "
            "which days belong together")

    # geometry shared by the whole build-up, so the raster never jumps between
    # slides.  The left gutter stays free for the depot, the right column for
    # text; both are reserved on every slide even where they are unused.
    GX, GY, GW, GH = 1.95, BODY_T + 0.15, 6.40, 2.70
    PX, PW = 8.60, W + L - 8.60          # right-hand text column
    BUL_T = GY + GH + 0.24               # bullets start below the raster
    DEPOT_X, DEPOT_Y = L, GY + GH / 2

    # ═══ M1 · parcels arrive every day ════════════════════════════════════
    s = new_slide(prs, "The mechanism · 1", "Parcels arrive every day",
                  "Schematic. Arrival volumes per weekday are illustrative; the "
                  "model uses HAGRID weekday demand per cell.")
    dcw = (W - 5 * 0.22) / 6
    per_day = [4, 6, 5, 7, 5, 3]
    for i, (dname, nd) in enumerate(zip(["MON", "TUE", "WED", "THU", "FRI",
                                         "SAT"], per_day)):
        x = L + i * (dcw + 0.22)
        txt(s, x, BODY_T, dcw, 0.34, dname, SZ_BODY, bold=True, color=DIM,
            align=PP_ALIGN.CENTER)
        rect(s, x, BODY_T + 0.44, dcw, 2.55, PANEL, line_col=LINE)
        for k in range(nd):
            fx, fy = _DOT_OFF[(i * 3 + k) % len(_DOT_OFF)]
            dot(s, x + 0.20 + fx * (dcw - 0.40),
                BODY_T + 0.62 + fy * 2.20, 0.19, S4)
    bullets(s, ["Demand does not wait for a good day to arrive.",
                "Every weekday brings its own parcels, in every area.",
                "Today, every one of them is delivered the next day."],
            BODY_T + 3.25, h=2.2)

    # ═══ M2 · the region splits into cells ════════════════════════════════
    s = new_slide(prs, "The mechanism · 2",
                  "The region splits into provider–area cells",
                  "Schematic raster of twelve cells; the case study has 312 "
                  "provider–area cells across seven providers.")
    cell_grid(s, GX, GY, GW, GH, theta=None)
    label_box(s, PX, GY + 0.30, PW, 1.05, PANEL,
              [("one cell", SZ_BODY, True, INK),
               ("= one carrier × one postal-code area", SZ_DIA, False, DIM)],
              line_col=LINE)
    bullets(s, ["Each cell is one carrier's parcels in one postal-code area.",
                "Cells are the unit we decide about — and the unit we route.",
                "Small enough to solve, many enough to matter: 312 of them."],
            BUL_T, h=1.95)

    # ═══ M3–M6 · who opts in?  the Morph build ════════════════════════════
    # One function, four calls: identical geometry is what lets Morph tween
    # the colour change instead of cutting between slides.
    def opt_in_slide(theta_pct, note):
        sl = new_slide(prs, "The mechanism · 3",
                       "Not every parcel can wait",
                       "Schematic. In the model the willing share is split per "
                       "cell by its business/private mix, so cells differ; here "
                       "it is spread evenly.")
        cell_grid(sl, GX, GY, GW, GH, theta=theta_pct / 100.0)
        bar_l, bar_w = PX, 4.05
        txt(sl, bar_l, GY + 0.10, bar_w, 0.46, f"θ = {theta_pct}%", 34,
            bold=True, color=RED)
        rect(sl, bar_l, GY + 0.62, bar_w, 0.20, LINE)
        # The filled part is always drawn, never conditionally: Morph matches
        # shapes by position in the slide, so a shape that exists on one slide
        # and not the next shifts everything after it and breaks the tween.
        rect(sl, bar_l, GY + 0.62, max(0.07, bar_w * theta_pct / 100.0), 0.20,
             RED)
        txt(sl, bar_l, GY + 0.90, bar_w, 0.34, "willingness to wait", SZ_DIA,
            color=DIM)
        dot(sl, bar_l + 0.10, GY + 1.45, 0.20, S3)
        txt(sl, bar_l + 0.32, GY + 1.28, bar_w - 0.32, 0.34, "standard — daily",
            SZ_BODY, color=INK2)
        dot(sl, bar_l + 0.10, GY + 1.90, 0.20, RED)
        txt(sl, bar_l + 0.32, GY + 1.73, bar_w - 0.32, 0.34,
            "willing — can be batched", SZ_BODY, color=INK2)
        txt(sl, bar_l, GY + 2.25, bar_w, 0.9, note, SZ_BODY, color=INK2,
            line=1.25)
        set_morph(sl)
        return sl

    opt_in_slide(0, "Nobody opts in. This is today's system: every parcel "
                    "delivered daily.")
    opt_in_slide(30, "Some recipients accept a wait. Business customers opt in "
                     "first.")
    opt_in_slide(60, "More opt in. Now most cells hold a mixed population.")
    opt_in_slide(100, "Everyone opts in — the upper bound of the sweep, not a "
                      "forecast.")

    # ═══ M7 · the standard parcels still need a daily tour ════════════════
    s = new_slide(prs, "The mechanism · 4",
                  "The standard parcels still need a tour every day",
                  "Schematic, one weekday at θ = 60%. Only the parcels that did "
                  "not opt in are shown.")
    pts = cell_grid(s, GX, GY, GW, GH, theta=0.60, show=("std",))
    std = [(x, y) for x, y, k, _ in pts if k == "std"]
    label_box(s, DEPOT_X, DEPOT_Y - 0.42, 1.20, 0.84, S6,
              [("DEPOT", SZ_DIA, True, WHITE)])
    bullets(s, ["The willing parcels are held for their service day.",
                "The rest must be delivered today, as always.",
                "They are scattered thinly across every cell."],
            BUL_T, h=1.95)
    txt(s, PX, GY + 0.20, PW, 0.9,
        f"{len(std)} standard parcels, spread over twelve cells", SZ_BODY,
        bold=True, color=INK, line=1.25)

    # ═══ M8 · so they ride one shared tour ════════════════════════════════
    s = new_slide(prs, "The mechanism · 5",
                  "So they ride one shared tour per depot",
                  "Schematic. The shared tour is what makes a cell's cost depend "
                  "on its hub neighbours · Paper §2.2.")
    pts = cell_grid(s, GX, GY, GW, GH, theta=0.60, show=("std",),
                    cell_line=RGBColor(0xEA, 0xEC, 0xEF))
    std = [(x, y) for x, y, k, _ in pts if k == "std"]
    route = serpentine([(x, y) for x, y in std])
    tour(s, [(DEPOT_X + 1.20, DEPOT_Y)] + route, S5, weight=1.75, close=True)
    label_box(s, DEPOT_X, DEPOT_Y - 0.42, 1.20, 0.84, S6,
              [("DEPOT", SZ_DIA, True, WHITE)])
    bullets(s, ["One vehicle serves the leftovers of many cells at once.",
                [("The tour crosses ", False), ("cell boundaries", True),
                 (" — it has to.", False)],
                "So a cell's cost depends on what its neighbours chose.",
                "That is the coupling the optimiser has to respect."],
            BUL_T, h=1.95)

    # ═══ M9 · the batched parcels form their own tours ════════════════════
    s = new_slide(prs, "The mechanism · 6",
                  "The batched parcels form their own denser tours",
                  "Schematic, at θ = 60%: the held parcels of a cell arrive "
                  "together on its service day.")
    pts = cell_grid(s, GX, GY, GW, GH, theta=0.60, show=("batch",))
    by_cell: dict = {}
    for x, y, k, ci in pts:
        if k == "batch":
            by_cell.setdefault(ci, []).append((x, y))
    for cpts in by_cell.values():
        if len(cpts) > 2:
            tour(s, serpentine([(x, y, ) for x, y in cpts]), RED, weight=1.5)
    bullets(s, ["Several days of parcels land on one delivery day.",
                "Same stops, more parcels each — the tour gets denser.",
                "Denser tours cost less per parcel. That is the whole gain.",
                "The price is a few days of waiting."],
            BUL_T, h=1.95)

    # ═══ M10 · together: one tour fewer ═══════════════════════════════════
    s = new_slide(prs, "The mechanism · 7", "Taken together: one tour fewer",
                  "Schematic illustration of the mechanism, not a measured "
                  "result. The measured system effect is in Part three.")
    for side, (title, n, colour, lab) in enumerate([
            ("Daily delivery", 4, S4, "four vehicle-days"),
            ("With batching", 3, RED, "three vehicle-days")]):
        x = L + side * (COL_W + 0.55)
        rect(s, x, BODY_T, COL_W, 2.70, PANEL, line_col=LINE)
        rect(s, x, BODY_T, COL_W, 0.10, colour)
        txt(s, x + 0.30, BODY_T + 0.28, COL_W - 0.60, 0.40, title, SZ_LEAD,
            bold=True, color=colour)
        for k in range(n):
            cxx = x + 0.55 + (k % 2) * 2.55
            cyy = BODY_T + 1.05 + (k // 2) * 1.15
            loop = [(cxx, cyy), (cxx + 0.85, cyy - 0.30),
                    (cxx + 1.45, cyy + 0.25), (cxx + 0.60, cyy + 0.55)]
            tour(s, loop, colour, weight=1.5)
            for px, py in loop:
                dot(s, px, py, 0.13, colour)
        txt(s, x + 0.30, BODY_T + 2.22, COL_W - 0.60, 0.40, lab, SZ_BODY,
            bold=True, color=INK)
    txt(s, L, BODY_T + 2.92, W, 0.65, "−1 vehicle-day", 38, bold=True, color=RED)
    bullets(s, ["Fewer, fuller tours instead of many thin ones.",
                "The saving is a whole vehicle-day, not a shorter route.",
                "Across a week and a region, that is where the money is."],
            BODY_T + 4.05, h=1.7)

    # ═══ M11 · which days should share a batch? ═══════════════════════════
    s = new_slide(prs, "The mechanism · 8",
                  "Which days should share a batch?",
                  "Legality follows the three-day rule (cyclic, Saturday wraps "
                  "to Monday). The fleet consequence is schematic.")
    # one short line per card; the argument itself lives in the bullets
    opts = [("Mon + Thu", [True, False, False, True, False, False], True,
             "gaps of 3 and 3"),
            ("Tue + Fri", [False, True, False, False, True, False], True,
             "legal, but cells peak together"),
            ("Mon + Tue", [True, True, False, False, False, False], False,
             "Tue → Mon is a 5-day gap")]
    ow = (W - 2 * 0.50) / 3
    for i, (nm, pat, legal, why) in enumerate(opts):
        x = L + i * (ow + 0.50)
        rect(s, x, BODY_T, ow, 0.10, RED if legal else CRIM)
        txt(s, x, BODY_T + 0.24, ow, 0.42, nm, SZ_LEAD, bold=True,
            color=INK if legal else CRIM)
        weekbar(s, x, BODY_T + 0.85, (ow - 5 * 0.07) / 6, pat,
                on=S6 if legal else CRIM)
        txt(s, x, BODY_T + 1.85, ow, 0.40, "✓ legal" if legal else "✗ illegal",
            SZ_BODY, bold=True, color=TEAL if legal else CRIM)
        txt(s, x, BODY_T + 2.26, ow, 0.42, why, SZ_BODY, color=INK2)
    bullets(s, ["Same delivery frequency, different consequences.",
                "Some pairings are ruled out by the service promise.",
                "Others are legal but pile every cell onto the same weekday.",
                "So: how many pairings are there, and how do we choose?"],
            BODY_T + 2.80, h=2.4)

    # ═══ divider 02 · the problem ═════════════════════════════════════════
    divider(prs, "02", "Part two", "The problem,\nprecisely",
            "What we optimise · why it is coupled · why it cannot be enumerated")

    # ═══ local decision, coupled effects ══════════════════════════════════
    s = new_slide(prs, "Problem", "The decision is local; the effects are coupled",
                  "Paper §2.2; reviewer clarification on per-LSP versus system "
                  "aggregation. 22 depots across seven providers (DHL 16, one each "
                  "for the others).")
    label_box(s, L, BODY_T + 1.30, 1.85, 1.30, S6,
              [("ONE", SZ_DIAB, True, WHITE), ("HUB", SZ_DIAB, True, WHITE)])
    for i in range(3):
        y = BODY_T + 0.25 + i * 1.35
        label_box(s, L + 2.55, y, 3.30, 1.05, WHITE,
                  [(f"CELL {i + 1}", SZ_BODY, True, INK),
                   ("one carrier × one area", SZ_DIA, False, DIM)],
                  line_col=LINE, align=PP_ALIGN.LEFT)
        arrow(s, L + 1.90, BODY_T + 1.95, L + 2.50, y + 0.52, colour=RED,
              head=False)
    label_box(s, COL2 + 0.35, BODY_T + 1.30, COL_W - 0.35, 1.30, PANEL,
              [("SHARED DAILY TOUR", SZ_BODY, True, INK),
               ("carries every cell's non-willing parcels", SZ_DIA, False, DIM)],
              line_col=LINE)
    bullets(s, ["A cell is one carrier in one postal-code area.",
                "Its pattern is chosen locally.",
                [("But its cost moves when ", False), ("hub neighbours", True),
                 (" change days.", False)],
                "312 decisions become 22 coupled blocks."],
            BODY_T + 2.95, h=2.5)

    # ═══ six switches ═════════════════════════════════════════════════════
    s = new_slide(prs, "The combinatorics · 1",
                  "A weekly pattern is six on/off switches",
                  "Operating week Monday–Saturday · Paper §2.2.")
    weekbar(s, L + 1.20, BODY_T + 0.20, 1.42, [True] * 6, h=0.85, gap=0.28,
            dsz=SZ_BODY)
    txt(s, L, BODY_T + 1.85, W, 0.7, "2⁶ = 64 possible settings", 36, bold=True,
        color=INK)
    bullets(s, ["Six delivery days, each either served or not.",
                "That is the entire decision for one cell.",
                "Most of those 64 settings are not allowed."],
            BODY_T + 2.75, h=2.2)

    # ═══ the rule leaves 39 ═══════════════════════════════════════════════
    s = new_slide(prs, "The combinatorics · 2",
                  "The three-day rule leaves exactly 39 patterns",
                  "Enumerated by enumerate_valid_schedules(); pinned as "
                  "EXPECTED_PATTERN_COUNT_K3 = 39 and covered by unit tests.")
    for i, (pat, lab, bad) in enumerate([
            ([True] * 6, "legal — daily", False),
            ([True, False, False, True, False, False], "legal — gap of 3", False),
            ([True, True, False, False, False, False],
             "illegal — Tue→Mon is 5", True)]):
        y = BODY_T + 0.10 + i * 0.92
        weekbar(s, L, y, 0.62, pat, on=CRIM if bad else S6, days=(i == 2),
                dsz=SZ_DIA)
        txt(s, L + 4.65, y + 0.06, 5.0, 0.42, ("✗  " if bad else "✓  ") + lab,
            SZ_BODY, bold=True, color=RED if bad else INK)
    for i, (fw, lab, fill, col) in enumerate([
            (W, "64 switch settings", PANEL, INK),
            (W * 0.98, "63 after dropping “never”", RGBColor(0xDD, 0xE2, 0xE7), INK),
            (W * 0.61, "39 obey the three-day rule", RED, WHITE)]):
        label_box(s, L, BODY_T + 3.05 + i * 0.62, fw, 0.52, fill,
                  [(lab, SZ_BODY, col == WHITE, col)], pad=0.18,
                  align=PP_ALIGN.LEFT)
    txt(s, L, BODY_T + 4.82, W, 0.4,
        "Cyclic rule: Saturday wraps to Monday, so no gap may exceed three days.",
        SZ_BODY, color=DIM)

    # ═══ the explosion ════════════════════════════════════════════════════
    s = new_slide(prs, "The combinatorics · 3",
                  "Direct routing cannot search this space",
                  "39³¹² is a combinatorial upper bound stated before hub coupling "
                  "reduces the effective space. A 497-digit number.")
    for i, lab in enumerate(["CELL 1", "CELL 2", "CELL 312"]):
        x = L + i * 2.80
        label_box(s, x, BODY_T + 0.15, 2.10, 1.30, WHITE,
                  [(lab, SZ_DIA, False, DIM), ("39", 34, True, INK)],
                  line_col=LINE)
        if i < 2:
            txt(s, x + 2.15, BODY_T + 0.58, 0.65, 0.45,
                "×" if i == 0 else "···", 26, color=DIM, align=PP_ALIGN.CENTER)
    txt(s, L + 5.60, BODY_T + 0.58, 0.70, 0.45, "×", 26, color=DIM,
        align=PP_ALIGN.CENTER)
    txt(s, L + 8.15, BODY_T + 0.52, 0.5, 0.5, "=", 30, color=DIM)
    txt(s, L + 8.75, BODY_T + 0.32, 4.0, 0.9, "≈ 2.6 × 10⁴⁹⁶", 40, bold=True,
        color=RED)
    bullets(s, ["Each cell picks from the same 39 candidates.",
                "Every assignment creates new day-level demand.",
                "New demand means new routes — a fresh routing problem.",
                "Enumerating them is not slow. It is impossible."],
            BODY_T + 1.80, h=2.7)

    # ═══ for scale ════════════════════════════════════════════════════════
    s = new_slide(prs, "The combinatorics · 3", "How large is 10⁴⁹⁶, really?",
                  "Bar length is the number of digits, not the value. Reference "
                  "figures are standard order-of-magnitude estimates.")
    for i, (bw2, mag, lab, col) in enumerate([
            (0.95, "10¹⁸", "seconds since the Big Bang", S3),
            (3.20, "10⁸⁰", "atoms in the observable universe", S4),
            (4.75, "10¹²⁰", "possible chess games", S5),
            (W, "10⁴⁹⁶", "weekly schedules of this one region", RED)]):
        y = BODY_T + 0.20 + i * 1.05
        rect(s, L, y, bw2, 0.60, col)
        dark = col == RED
        txt(s, L + 0.18, y + 0.06, 2.2, 0.48, mag, SZ_LEAD, bold=True,
            color=WHITE if dark else INK)
        txt(s, L + (2.40 if dark else bw2 + 0.22), y + 0.10, 8.4, 0.44, lab,
            SZ_BODY, color=WHITE if dark else INK2)
    txt(s, L, BODY_T + 4.55, W, 0.55,
        "One German region outruns every physical count we have.", SZ_LEAD + 2,
        bold=True, color=RED)

    # ═══ every guess is a routing problem ═════════════════════════════════
    s = new_slide(prs, "The combinatorics · 4",
                  "Worse: every single guess is a routing problem",
                  "VRP NP-hardness: Lenstra & Rinnooy Kan (1981). Solver timing "
                  "measured on the VROOM/Valhalla stack used for label generation.")
    bullets(s, [[("You cannot price a schedule without ", False),
                 ("planning its tours", True), (".", False)],
                "That is a full vehicle routing problem — NP-hard.",
                [("The real solver needs ", False), ("minutes per candidate", True),
                 (".", False)],
                "Hub coupling means cells cannot be graded one at a time.",
                "At a billion guesses a second we would still never finish."],
            BODY_T, h=3.6)
    label_box(s, L, BODY_T + 3.80, W, 1.00, BLUSH,
              [("We need a fast cost oracle inside the loop — and a search that "
                "never enumerates.", SZ_LEAD + 2, True, RED)], line_col=RED)

    # ═══ fleet peaks ══════════════════════════════════════════════════════
    s = new_slide(prs, "Why vehicle deployment must be optimized",
                  "The same frequencies can create fleet peaks",
                  "Schematic mechanism. Balancing preserves each cell's delivery "
                  "frequency exactly · Paper §2.2.")
    coll = [True, False, False, True, False, False]
    stag = [[True, False, False, True, False, False],
            [False, True, False, False, True, False],
            [False, False, True, False, False, True]]
    for side, (title, pats, hot) in enumerate([
            ("Cost-optimal calendars collide", [coll] * 3, False),
            ("Balanced calendars stagger", stag, True)]):
        x = L + side * (COL_W + 0.55)
        txt(s, x, BODY_T, COL_W, 0.38, title, SZ_LEAD, bold=True,
            color=RED if hot else INK)
        for r, pat in enumerate(pats):
            y = BODY_T + 0.55 + r * 0.58
            txt(s, x, y + 0.06, 1.10, 0.34, f"Cell {'ABC'[r]}", SZ_DIA, color=DIM)
            weekbar(s, x + 1.15, y, 0.58, pat, days=False, h=0.42)
        hh = [1.05, 0.26, 0.26, 1.05, 0.26, 0.26] if not hot else [0.52] * 6
        for i, h in enumerate(hh):
            rect(s, x + 1.15 + i * 0.65, BODY_T + 3.60 - h, 0.58, h,
                 TEAL if hot else CRIM)
        hrule(s, x + 1.15, BODY_T + 3.60, 3.95, DIM, 1.25)
        txt(s, x + 1.15, BODY_T + 3.72, 3.95, 0.34, "vehicles per weekday",
            SZ_DIA, color=DIM)
    txt(s, L, BODY_T + 4.30, W, 0.55,
        "Selection chooses how often. Balancing chooses which weekdays.",
        SZ_LEAD + 2, bold=True, color=RED)

    # ═══ divider 02 ═══════════════════════════════════════════════════════
    divider(prs, "03", "Part three", "The method",
            "Build the oracle · search with it · balance the week · re-route")

    # ═══ framework ════════════════════════════════════════════════════════
    s = new_slide(prs, "Method", "Build the oracle. Search. Re-route the winners.",
                  "Paper Fig. 2 and §2; Revision 1 reviewer clarification.")
    txt(s, L, BODY_T, W, 0.40, "Offline — learn the routing response", SZ_BODY,
        bold=True, color=RED, spc=1.2, caps=True)
    flow(s, [("HAGRID", "weekly demand", PANEL), ("VROOM", "cost labels", PANEL),
             ("SURROGATE", "Daganzo + LightGBM", RED)],
         BODY_T + 0.44, w=8.6, bh=1.20)
    arrow(s, L + 7.1, BODY_T + 1.68, L + 7.1, BODY_T + 2.34, colour=RED,
          weight=2.5)
    txt(s, L + 7.35, BODY_T + 1.82, 3.4, 0.42, "trained once", SZ_BODY, color=RED)
    txt(s, L, BODY_T + 2.58, W, 0.40, "Online — optimize schedules at scale",
        SZ_BODY, bold=True, color=RED, spc=1.2, caps=True)
    flow(s, [("SEARCH", "39 patterns / cell", PANEL),
             ("BALANCE", "hub + provider", PANEL),
             ("VALIDATE", "re-route 4 points", S6)],
         BODY_T + 3.02, l=L + 3.6, w=8.6, bh=1.20)
    txt(s, L, BODY_T + 4.45, W, 0.5,
        "Routing is solved once for labels — then millions of cheap evaluations.",
        SZ_LEAD, bold=True, color=INK)

    # ═══ physics + ML ═════════════════════════════════════════════════════
    s = new_slide(prs, "Method", "Physics explains; machine learning corrects",
                  "2 733 VROOM samples · GroupKFold by postal-code area, so no area "
                  "appears in both train and validation · R² = 0.997.")
    stats(s, [("26%", "raw Daganzo error", False),
              ("9.7%", "after one global factor α = 1.343", False),
              ("2.95%", "after the learned residual", True)], BODY_T, h=1.15)
    pic(s, A / "fig21_progression.png", L, BODY_T + 1.35, W, 3.7)

    # ═══ why keep the formula ═════════════════════════════════════════════
    s = new_slide(prs, "Why not pure ML · 1",
                  "The formula earns its keep at the edges",
                  "Extrapolation behaviour of tree ensembles is structural: "
                  "predictions are piecewise constant and cannot exceed the range "
                  "seen in training.")
    bullets(s, ["We have solver labels — so why keep a 1984 formula?",
                "Because batching multiplies a delivery day's volume.",
                [("Trees are ", False), ("flat beyond their training range", True),
                 (".", False)],
                "They cannot extrapolate; they repeat the last leaf.",
                "√(nA) and ⌈p/Q⌉ keep rising, because they are theory."],
            BODY_T, w=6.3, h=4.0)
    pic(s, B / "fig23_alpha.png", COL2 - 0.15, BODY_T + 0.15, 5.9, 4.1)

    # ═══ division of labour ═══════════════════════════════════════════════
    s = new_slide(prs, "Why not pure ML · 2",
                  "A clean division of labour: level and structure",
                  "α calibrated on the training pool and re-verified against the "
                  "shipped production model (α = 1.343) by the figure script's gate.")
    label_box(s, L, BODY_T + 0.25, 3.5, 1.50, PANEL,
              [("α · Daganzo", 26, True, RED),
               ("the physics backbone", SZ_DIA, False, DIM)], line_col=RED)
    txt(s, L + 3.62, BODY_T + 0.60, 0.7, 0.7, "+", 40, bold=True, color=DIM,
        align=PP_ALIGN.CENTER)
    label_box(s, L + 4.38, BODY_T + 0.25, 3.5, 1.50, PANEL,
              [("g(x)", 26, True, TEAL),
               ("the learned residual", SZ_DIA, False, DIM)], line_col=TEAL)
    txt(s, L + 8.00, BODY_T + 0.60, 0.7, 0.7, "=", 40, bold=True, color=DIM,
        align=PP_ALIGN.CENTER)
    label_box(s, L + 8.76, BODY_T + 0.25, 3.45, 1.50, S6,
              [("cost estimate", 24, True, WHITE),
               ("2.95% out-of-area error", SZ_DIA, False,
                RGBColor(0xC6, 0xD4, 0xE2))])
    bullets(s, [[("α", True), (" fixes the level: one scalar removes a −26% bias.",
                               False)],
                [("g(x)", True),
                 (" fixes the structure the formula cannot see.", False)],
                "Both halves stay inspectable — equation and feature importances.",
                "Per-provider α made it worse: 3.2% against 2.9%."],
            BODY_T + 2.10, h=2.9)

    # ═══ what the formula ignores, part 1 ═════════════════════════════════
    s = new_slide(prs, "Why not pure ML · 3",
                  "What the formula ignores — and who repairs it",
                  "Point-inaccuracy of continuum approximations: Figliozzi (2008). "
                  "Feature taxonomy after Akkerman et al. (2025).")
    _y = table(s, ["The formula assumes", "Reality", "Who repairs it"],
          [[("An idealised plane", "key"),
            "Real vans follow streets, one-ways and rivers",
            ("The residual learns the detour factor", "good")],
           [("One average stem r", "key"),
            "An off-centre depot changes every single tour",
            ("Depot distance is an input feature", "good")]],
          BODY_T, widths=[2.6, 4.2, 4.0], reserve=2.0)
    bullets(s, ["These are not reasons to drop the backbone.",
                "They are exactly the structured remainder a learner is good at."],
            _y + 0.30, h=1.8)

    # ═══ what the formula ignores, part 2 ═════════════════════════════════
    s = new_slide(prs, "Why not pure ML · 3", "Two more gaps the residual closes",
                  "The visible cost bands in the scatter come from the integer tour "
                  "count, not from a fitting artefact.")
    _y = table(s, ["The formula assumes", "Reality", "Who repairs it"],
          [[("Stops spread evenly", "key"),
            "Real drops cluster in blocks and villages",
            ("Stop density and spread carry this", "good")],
           [("Whole vans, ⌈p/Q⌉", "key"),
            "Right in kind, brutal at the capacity threshold",
            ("The residual smooths the step", "good")]],
          BODY_T, widths=[2.6, 4.2, 4.0], reserve=2.4)
    label_box(s, L, _y + 0.26, W, 0.90, BLUSH,
              [("The hybrid beats either half alone.", SZ_LEAD + 2, True, RED)],
              line_col=RED)
    bullets(s, ["The backbone extrapolates where data thins out.",
                "The residual is accurate where data is dense."],
            _y + 1.30, h=1.6)

    # ═══ benchmark ════════════════════════════════════════════════════════
    s = new_slide(prs, "Model choice", "We tested the field, not just our favourite",
                  "Five-fold GroupKFold over postal-code areas, 10 946-sample pool. "
                  "The shipped surrogate is the Daganzo–LightGBM hybrid, refitted on "
                  "the merge-corrected 2 733-sample pool: 2.95%.")
    _y = table(s, ["Model", "Out-of-area error", "Note"],
          [[("XGBoost, log target", "key"), ("3.31%", "num"), "44 features"],
           [("LightGBM", "key"), ("3.39%", "num"), "44 features"],
           [("MLP ensemble", "key"), ("3.54%", "num"), "deep-learning baseline"],
           [("Random forest", "key"), ("3.98%", "num"), "overfits area identity"],
           [("Daganzo alone", "key"), ("21.5%", "num"), "no learned residual"]],
          BODY_T, widths=[4, 3, 4], reserve=0.75)
    txt(s, L, _y + 0.22, W, 0.5,
        "Grouped folds: no postal-code area appears in both train and test.",
        SZ_BODY, color=DIM)

    # ═══ the benchmark lesson ═════════════════════════════════════════════
    s = new_slide(prs, "Model choice", "The lesson is not that LightGBM wins",
                  "Benchmark protocol A (interpolation holdout, n = 1 927) and "
                  "protocol B (grouped by area, n = 10 946).")
    stats(s, [("0.1 pp", "separates every boosted tree", False),
              ("6×", "penalty for skipping the learner", True)], BODY_T, w=7.2,
          h=1.25)
    bullets(s, ["The choice of learner barely matters.",
                "What matters is the features and the backbone.",
                "The deep-learning baseline is not ahead.",
                "The analytical proxy alone costs a factor of six."],
            BODY_T + 1.65, h=3.0)

    # ═══ rejected, part 1 ═════════════════════════════════════════════════
    s = new_slide(prs, "Model choice", "Every rejection has a reason",
                  "Protocol A versus protocol B comparison; see the benchmark "
                  "report for the full grid.")
    _y = table(s, ["Rejected", "Why"],
          [[("Random forest", "key"),
            "Memorises the area: 1.0% on areas it has seen, 4.0% on ones it has not"],
           [("Stacked ensemble", "key"),
            "The meta-learner over-trusted a leaky member; an honest fix needs 25 "
            "refits"],
           [("Linear and ridge", "key"),
            "Cost spans 200 € to 32 000 € per cell-day; no linear form holds that"]],
          BODY_T, widths=[3, 8], reserve=1.15)
    label_box(s, L, _y + 0.26, W, 0.90, BLUSH,
              [("Models that memorise the area collapse on unseen geography.",
                SZ_LEAD, True, RED)], line_col=RED)

    # ═══ rejected, part 2 ═════════════════════════════════════════════════
    s = new_slide(prs, "Model choice", "Two we kept, for a reason",
                  "The Daganzo baseline is reported throughout as the classical "
                  "lower bound.")
    _y = table(s, ["Model", "Verdict"],
          [[("Deep tabular nets", "key"),
            "Not pursued — the boosted trees were already on the accuracy plateau, "
            "so depth buys effort, not error"],
           [("Daganzo alone", "key"),
            "Kept as the classical baseline; it under-prices every cell by about 21%"]],
          BODY_T, widths=[3, 8], reserve=1.9)
    bullets(s, ["Models that are too rigid cannot span the cost range.",
                "The hybrid avoids both failure modes by construction."],
            _y + 0.28, h=1.8)

    # ═══ training pool ════════════════════════════════════════════════════
    s = new_slide(prs, "Method", "The training pool covers the consolidation range",
                  "Pool: results/supplementary/sweep_v3_mergefix/training_matrix.csv. "
                  "Every variant carries a true VROOM/Valhalla cost label.")
    pic(s, B / "fig22_pool.png", L, BODY_T, W, 3.05)
    chips(s, [("2 733 VROOM samples", True), ("48 areas", False),
              ("all seven providers", False), ("25 → 44 features", True)],
          BODY_T + 3.20)
    bullets(s, ["Baseline operations alone would only show daily delivery.",
                "Augmentation adds scaling, dropout, jitter and multi-day batches."],
            BODY_T + 3.90, h=1.6)

    # ═══ two dials ════════════════════════════════════════════════════════
    s = new_slide(prs, "Method", "Two dials define every operating regime",
                  "Paper §2.2; P enters the objective as a steering term only, which "
                  "keeps the cost–service trade-off continuous across the grid.")
    for side, (sym, nm, unit, acc, items, frac, ends) in enumerate([
            ("P", "Service penalty", "€ per parcel per day of delay", RED,
             ["A shadow price of waiting", "Not an observed customer price",
              "Never booked as cost"], 0.40, "cost-first                 service-first"),
            ("θ", "Willingness to wait", "eligible share of parcel volume",
             RGBColor(0x2C, 0x7F, 0xC7),
             ["Only θ is consolidated", "The rest keeps daily service",
              "Business opts in before private"], 0.75,
             "0%                                    100%")]):
        x = L + side * (COL_W + 0.55)
        rect(s, x, BODY_T, COL_W, 4.50, PANEL, line_col=LINE)
        rect(s, x, BODY_T, COL_W, 0.10, acc)
        txt(s, x + 0.30, BODY_T + 0.30, 1.3, 0.95, sym, 54, bold=True, color=acc)
        txt(s, x + 1.60, BODY_T + 0.40, COL_W - 1.95, 0.42, nm, SZ_LEAD,
            bold=True, color=INK)
        txt(s, x + 1.60, BODY_T + 0.88, COL_W - 1.95, 0.42, unit, SZ_BODY,
            color=DIM)
        rect(s, x + 0.30, BODY_T + 1.62, COL_W - 0.60, 0.18, LINE)
        rect(s, x + 0.30, BODY_T + 1.62, (COL_W - 0.60) * frac, 0.18, acc)
        txt(s, x + 0.30, BODY_T + 1.90, COL_W - 0.60, 0.3, ends, SZ_DIA,
            color=DIM)
        bullets(s, items, BODY_T + 2.40, l=x + 0.30, w=COL_W - 0.60, h=1.9)

    # ═══ the grid ═════════════════════════════════════════════════════════
    s = new_slide(prs, "Method", "Eight penalties × eleven adoption levels",
                  "Every cell of the grid is a full Stage-3 optimization with "
                  "balancing and smoothing applied.")
    # One flat colour: this grid counts scenarios, it does not carry values.
    # Shading it would invent a result the slide has no data for.
    gx, gy, cw2, ch = L + 0.62, BODY_T + 0.22, 0.88, 0.28
    for r in range(8):
        for c in range(11):
            rect(s, gx + c * (cw2 + 0.055), gy + r * (ch + 0.055), cw2, ch, S3)
    txt(s, gx, gy + 8 * (ch + 0.055) + 0.10, 10.4, 0.36,
        "adoption θ   →   0 … 100%   (eleven levels)", SZ_BODY, color=DIM)
    txt(s, L - 0.02, gy + 0.95, 0.56, 1.2, "P", SZ_LEAD, bold=True, color=DIM)
    txt(s, L - 0.02, gy + 1.45, 0.56, 1.2, "↓", SZ_BODY, bold=True, color=DIM)
    stats(s, [("88", "Stage-3 scenarios", True),
              ("312", "cells re-optimized in each", False),
              ("0.25–0.5", "the efficient penalty range", False)],
          gy + 8 * (ch + 0.055) + 0.80, h=1.0, sz=38)

    # ═══ greedy is exact ══════════════════════════════════════════════════
    s = new_slide(prs, "The search · 1",
                  "Sometimes the greedy answer is the exact one",
                  "For batch-only scenarios optimize_cd_ml() returns argmin over the "
                  "39 directly, which is the global optimum by separability.")
    txt(s, L, BODY_T, COL_W, 0.38, "Without the residual — separable", SZ_LEAD,
        bold=True, color=INK)
    for i in range(4):
        label_box(s, L + i * 1.45, BODY_T + 0.52, 1.22, 0.72, WHITE,
                  [("min", SZ_BODY, False, RGBColor(0x2C, 0x7F, 0xC7))],
                  line_col=LINE)
    label_box(s, L, BODY_T + 1.45, COL_W, 0.66, PANEL,
              [("Per-cell argmin is the global optimum", SZ_BODY, True, TEAL)],
              line_col=TEAL)
    txt(s, COL2, BODY_T, COL_W, 0.38, "With the residual — coupled", SZ_LEAD,
        bold=True, color=RED)
    rect(s, COL2, BODY_T + 0.52, COL_W, 1.59, BLUSH, line_col=RED)
    for dx, dy in [(0.25, 0.20), (2.05, 0.20), (3.85, 0.20), (1.15, 0.88),
                   (2.95, 0.88)]:
        rect(s, COL2 + dx, BODY_T + 0.52 + dy, 1.60, 0.58, WHITE, line_col=S6)
    bullets(s, ["Non-willing parcels still need a daily service.",
                "They ride one shared tour per depot.",
                "So a cell's cost depends on its hub neighbours.",
                "312 independent choices become 22 coupled blocks."],
            BODY_T + 2.35, h=2.8)

    # ═══ how CD runs ══════════════════════════════════════════════════════
    s = new_slide(prs, "The search · 2",
                  "One cell at a time, until nothing improves",
                  "Production call: eight rounds maximum, shuffled cell order, "
                  "warm-started from the per-cell argmin, then a pair-swap polish.")
    for i, (n, lab) in enumerate([("1", "Fix every other cell"),
                                  ("2", "Score all 39 patterns"),
                                  ("3", "Keep the best, move on")]):
        y = BODY_T + i * 0.68
        txt(s, L, y, 0.5, 0.42, n, SZ_LEAD, bold=True, color=RED)
        txt(s, L + 0.55, y, 5.2, 0.42, lab, SZ_BODY, color=INK)
    for r in range(3):
        for c in range(4):
            focal = (r, c) == (1, 1)
            rect(s, L + 0.10 + c * 1.00, BODY_T + 2.30 + r * 0.62, 0.86, 0.50,
                 BLUSH if focal else PANEL, line_col=RED if focal else LINE,
                 line_w=2.0 if focal else 1.0)
    txt(s, L + 0.10, BODY_T + 4.22, 4.5, 0.36, "one hub, one focal cell", SZ_DIA,
        color=DIM)
    ys = [0.30, 1.00, 1.48, 1.78, 1.94, 2.02, 2.05]
    for i in range(len(ys) - 1):
        x1 = COL2 + i * 0.86
        rect(s, x1, BODY_T + 0.45 + ys[i], 0.86, 0.06, RED)
        rect(s, x1 + 0.80, BODY_T + 0.45 + min(ys[i], ys[i + 1]), 0.06,
             abs(ys[i + 1] - ys[i]) + 0.06, RED)
    hrule(s, COL2, BODY_T + 3.00, 5.5, DIM, 1.25)
    txt(s, COL2, BODY_T + 3.12, 5.5, 0.36, "hub cost across sweeps   →", SZ_DIA,
        color=DIM)
    txt(s, COL2, BODY_T + 3.65, 5.7, 1.0,
        "Monotone: cost never rises.\nThe sweep stops itself.", SZ_LEAD,
        bold=True, color=INK, line=1.3)

    # ═══ why CD fits, part 1 ══════════════════════════════════════════════
    s = new_slide(prs, "The search · 3 — why this method",
                  "The problem picks the method: two structural facts",
                  "Coordinate descent on a common hub objective, not a game between "
                  "self-interested cells · Paper §2.2.")
    for i, (num, head, items) in enumerate([
            ("1", "The menu is small, discrete and complete",
             ["39 options per cell",
              "So the best is found by enumeration, exactly",
              "No step size, no learning rate, nothing to tune"]),
            ("2", "There is no gradient to follow",
             ["A tree ensemble plus ⌈p/Q⌉ is piecewise constant",
              "Gradient methods and LP relaxations have nothing to grip",
              "This search needs only function values"])]):
        y = BODY_T + i * 2.45
        txt(s, L, y, 0.65, 0.55, num, 34, bold=True, color=RED)
        txt(s, L + 0.72, y + 0.04, W - 0.72, 0.46, head, SZ_LEAD, bold=True,
            color=INK)
        bullets(s, items, y + 0.62, l=L + 0.72, w=W - 0.72, h=1.7)

    # ═══ why CD fits, part 2 ══════════════════════════════════════════════
    s = new_slide(prs, "The search · 3 — why this method",
                  "And two facts that make it converge",
                  "The fixed point is a coordinate-wise optimum — equivalently a "
                  "Nash equilibrium of an identical-interest hub objective.")
    for i, (num, head, items) in enumerate([
            ("3", "The coupling is local",
             ["Cells interact only through their own depot",
              "So the problem splits into 22 independent blocks",
              "One flip re-scores only the weekdays that changed"]),
            ("4", "One shared objective, so descent is monotone",
             ["Every accepted move lowers that hub's total cost",
              "The sweep must therefore terminate",
              "It ends at a coordinate-wise optimum"])]):
        y = BODY_T + i * 2.45
        txt(s, L, y, 0.65, 0.55, num, 34, bold=True, color=RED)
        txt(s, L + 0.72, y + 0.04, W - 0.72, 0.46, head, SZ_LEAD, bold=True,
            color=INK)
        bullets(s, items, y + 0.62, l=L + 0.72, w=W - 0.72, h=1.7)

    # ═══ alternatives, part 1 ═════════════════════════════════════════════
    s = new_slide(prs, "The search · 4", "Why not the textbook optimisers?",
                  "The legacy annealing implementation is retained in the "
                  "repository for reference; the production path is coordinate "
                  "descent.")
    _y = table(s, ["Approach", "Blocker"],
          [[("Exhaustive search", "key"),
            "10⁴⁹⁶ candidates — not slow, impossible"],
           [("MILP or LP relaxation", "key"),
            "Our cost is a tree ensemble plus a ceiling function: neither linear "
            "nor convex"],
           [("Gradient descent", "key"),
            "A piecewise-constant surrogate has no slope to follow"]],
          BODY_T, widths=[3.4, 7.6])
    label_box(s, L, BODY_T + 3.25, W, 0.90, BLUSH,
              [("None of them can even read our objective function.", SZ_LEAD,
                True, RED)], line_col=RED)

    # ═══ alternatives, part 2 ═════════════════════════════════════════════
    s = new_slide(prs, "The search · 4",
                  "We started with annealing. Then we looked closer.",
                  "300 000 iterations, a temperature schedule, reheating and a "
                  "Boltzmann proposal — all replaced by an exhaustive 39-way scan.")
    for side, (nm, items, acc) in enumerate([
            ("Simulated annealing",
             ["300 000 iterations", "A temperature schedule and reheating",
              "Stochastic — needs tuning",
              "Most evaluations land on rejected moves"], CRIM),
            ("Coordinate descent",
             ["Exhaustive within a cell", "Local across cells",
              "Deterministic at a fixed seed", "Nothing to tune"], TEAL)]):
        x = L + side * (COL_W + 0.55)
        rect(s, x, BODY_T, COL_W, 3.40, PANEL, line_col=LINE)
        rect(s, x, BODY_T, COL_W, 0.10, acc)
        txt(s, x + 0.30, BODY_T + 0.32, COL_W - 0.60, 0.46, nm, SZ_LEAD,
            bold=True, color=acc)
        bullets(s, items, BODY_T + 0.95, l=x + 0.30, w=COL_W - 0.60, h=2.3)
    txt(s, L, BODY_T + 3.65, W, 1.0,
        "Annealing was not wrong — it was uninformed.\nOnce you see the 39 options, "
        "a method that exploits them wins.", SZ_LEAD, bold=True, color=RED,
        line=1.3)

    # ═══ restarts ═════════════════════════════════════════════════════════
    s = new_slide(prs, "The search · 5", "Five restarts converge to the same point",
                  "Stage-3 restart analysis; the figure script recomputes the "
                  "spread and aborts if it exceeds the reported bound.")
    stats(s, [("< 10⁻¹²", "relative objective spread across five random restarts",
               True)], BODY_T, w=6.2, h=1.5, sz=42)
    bullets(s, ["The search reaches the same point from different starts.",
                "Stable is not the same as globally optimal.",
                "It cannot certify that no better plan exists."],
            BODY_T + 1.80, w=6.2, h=2.5)
    pic(s, B / "fig24_determinism.png", COL2 - 0.15, BODY_T + 0.10, 5.9, 4.2)

    # ═══ what it gives up ═════════════════════════════════════════════════
    s = new_slide(prs, "The search · 6 — honesty",
                  "What the search gives up, and how we cover it",
                  "Best-of-K selection bias in surrogate-based optimization; "
                  "addressed here by out-of-sample solver verification.")
    bullets(s, ["Local, not global — no certificate of optimality.",
                "Guards: five restarts, shuffled order, a pair-swap polish.",
                "The sharper risk is the surrogate, not the search.",
                "Argmin of 39 noisy predictions favours the under-priced one.",
                "So the winners get re-routed with the real solver."],
            BODY_T, w=6.4, h=4.0)
    x = COL2 + 0.15
    hs2 = [1.10, 1.30, 1.00, 1.40, 1.17, 0.93, 1.27, 0.53, 1.07, 1.33, 0.97, 1.20]
    for i, h in enumerate(hs2):
        rect(s, x + i * 0.46, BODY_T + 2.35 - h, 0.34, h, RED if i == 7 else S2)
    hrule(s, x, BODY_T + 1.35, 5.5, S6, 1.75)
    txt(s, x + 3.4, BODY_T + 0.95, 2.2, 0.36, "true cost", SZ_DIA, color=S6)
    hrule(s, x, BODY_T + 2.35, 5.5, DIM, 1.25)
    txt(s, x, BODY_T + 2.52, 5.7, 0.42, "picked because under-priced", SZ_BODY,
        bold=True, color=RED)

    # ═══ balancing process ════════════════════════════════════════════════
    s = new_slide(prs, "Method", "Balancing makes the week operational",
                  "Frequency invariance asserted for all 27 456 rows by the figure "
                  "script · Paper §2.2.")
    flow(s, [("1 · SELECT", "cost-optimal pattern per hub", PANEL),
             ("2 · BALANCE", "flatten each depot, +5% cost cap", PANEL),
             ("3 · SMOOTH", "level each provider's fleet", RED)],
         BODY_T + 0.10, bh=1.35)
    stats(s, [("0.135", "baseline fleet CV", False),
              ("0.056", "after stage 2", False),
              ("0.029", "after stage 3", True)], BODY_T + 1.80, h=1.20, sz=42)
    label_box(s, L, BODY_T + 3.35, W, 0.90, PANEL,
              [("Frequency stays fixed. Only weekdays move — so service quality "
                "is untouched.", SZ_LEAD, True, INK)], line_col=LINE)

    # ═══ balancing result ═════════════════════════════════════════════════
    s = new_slide(prs, "Method", "The week, before and after smoothing",
                  "At P = 0.25, θ = 1 · Stage 2 versus Stage 3 fleet smoothing.")
    pic(s, A / "fig51_fleet_smoothing.png", L, BODY_T, W, 3.65)
    bullets(s, ["Cost-optimal patterns pile deliveries onto a few weekdays.",
                "Swaps flatten each depot, then each provider's network.",
                "The peak falls without touching delivery frequency."],
            BODY_T + 3.80, h=2.1)

    # ═══ case study ═══════════════════════════════════════════════════════
    s = new_slide(prs, "Case study",
                  "The Hanover case spans the full urban–rural gradient",
                  "HAGRID demand · PLZ-level geodata · Region Hannover · seven "
                  "LSPs. Demand density spans a factor of 141 across one region.")
    stats(s, [("7", "providers", True), ("1.26 M", "parcels / week", False),
              ("€1.91 M", "baseline cost / week", False),
              ("312", "provider–area cells", False)], BODY_T, h=1.10, sz=40)
    pic(s, A / "fig12_map_demand.png", L, BODY_T + 1.30, 5.85, 3.7)
    pic(s, A / "fig11_lsp_volumes.png", COL2, BODY_T + 1.30, 5.85, 3.7)

    # ═══ divider 03 ═══════════════════════════════════════════════════════
    divider(prs, "04", "Part four", "What we\nfound",
            "The cost–service frontier · where it pays · what it does to the fleet")

    # ═══ the revision block ═══════════════════════════════════════════════
    # Everything after this reads through the revision's frame, so the frame
    # comes first.
    if facts is not None:
        part_revision(prs, facts, revision, tag=tag)

    # ═══ the trade-off in numbers ═════════════════════════════════════════
    # Numbers from the grid in use, both plans, both lenses. The submission's
    # single "cost saving" column no longer exists: it was the routing lens of
    # the routing-optimal plan, which is now one of four cells.
    s = new_slide(prs, "Results", "The efficient range sits between 0.25 and 0.5",
                  "θ = 100 %, surrogate-predicted, on the revision grid. "
                  "Routing = what the tours cost to drive; operator = "
                  "kilometres plus 1 134.90 € per peak vehicle and hub.")
    _rows = []
    if facts is not None:
        for _P in (0.0, 0.25, 0.5):
            _h = facts.headline[_P]
            _rows.append([(f"P = {_P:g}", "key"),
                          (f"{_h['rout2']:.1f} %", "num"),
                          (f"{_h['op2']:.1f} %", "num"),
                          f"{_h['wait2']:.2f} d",
                          f"{_h['peak2_pct']:+.1f} %"])
    else:
        # --no-revision: the submission's own figures, with no second lens
        _rows = [[("P = 0", "key"), ("22.8 %", "num"), ("n/a", "body"),
                  "0.98 d", "n/a"],
                 [("P = 0.25", "key"), ("18.5 %", "num"), ("n/a", "body"),
                  "0.46 d", "n/a"],
                 [("P = 0.5", "key"), ("13.5 %", "num"), ("n/a", "body"),
                  "0.23 d", "n/a"]]
    _y = table(s, ["Penalty", "Routing saving", "Operator saving",
                   "Added wait", "Σ hub peak"],
               _rows, BODY_T, widths=[2.0, 2.6, 2.6, 2.4, 2.6], reserve=1.45)
    txt(s, L, _y + 0.30, W, 1.0,
        "Every row is the operator-polished plan.\nThe expensive part of the "
        "trade is the last points of saving, not the first.", SZ_LEAD,
        bold=True, color=RED, line=1.3)
    if facts is not None:
        revision.notes(
            s, "The submission quoted 22.8 / 18.5 / 13.5 % here. Those were "
               "routing euro on the routing-optimal plan and are superseded "
               "twice over: by the universal tour rule and by the operator "
               "polish. This table is one plan (operator-polished) in both "
               "lenses.", cite=["§40.15", "§40.18"])
        revision.provisional(s, enabled=tag)

    # ═══ pareto ═══════════════════════════════════════════════════════════
    # The shape of the front is what this slide argues, and it is unchanged;
    # the levels on it are not, so the source line says which grid drew it.
    s = new_slide(prs, "Results", "The cost–service frontier",
                  "Each line is one penalty level swept across adoption; the dashed "
                  "line is the efficient front. Rendered from the submission grid "
                  "— the revision restates the levels, not the shape.")
    pic(s, A / "fig34_pareto.png", L, BODY_T, 7.5, 4.85)
    bullets(s, ["Every point is a fully balanced weekly schedule.",
                "The knee, not the extreme, is the operating point.",
                "Beyond it, waiting grows faster than saving."],
            BODY_T + 0.7, l=8.4, w=4.4, h=3.2)

    # ═══ frequency mix ════════════════════════════════════════════════════
    s = new_slide(prs, "Results", "The penalty shifts the delivery-frequency mix",
                  "Shares are the revision grid at θ = 100 %; the figure is the "
                  "submission render of the same quantity. Frequencies stay "
                  "within {2,…,6}.")
    pic(s, A / "fig35_schedule_mix.png", L, BODY_T, W, 3.55)
    if facts is not None:
        _mix = {p: revision.D.chosen_sizes_v2(p) for p in revision.D.PLANS}
        def _two_day(plan):
            m = _mix[plan]
            m = m[(m.penalty == 0.0) & (m.share_willing == 1.0)]
            return 100.0 * (m.schedule_size == 2).mean()
        bullets(s, [f"At P = 0 the routing optimum is almost all two-day "
                    f"patterns — {_two_day(revision.D.PLAN_ROUTING):.1f} % of "
                    f"areas.",
                    f"The operator polish pulls that back to "
                    f"{_two_day(revision.D.PLAN_OPERATOR):.1f} %, because a "
                    f"two-day pattern is what makes a hub peak.",
                    "At P ≥ 5 the system reverts to daily delivery."],
                BODY_T + 3.70, h=2.1)
        revision.notes(s, "The frequency mix is now a two-plan quantity: the "
                          "routing optimum consolidates hard, the operator "
                          "polish trades some of that back for a flatter week.",
                       cite=["§40.14", "§40.15"])
        revision.provisional(s, enabled=tag)
    else:
        bullets(s, ["At P = 0, two-day patterns dominate the routing optimum.",
                    "At P ≥ 5 the system reverts to daily delivery."],
                BODY_T + 3.70, h=2.1)

    # ═══ maps ═════════════════════════════════════════════════════════════
    s = new_slide(prs, "Results", "Where the delivery days land",
                  "At P = 0.25 €/p/d, submission grid. Values are per merged "
                  "cluster, so member polygons of one cluster share a value. "
                  "The revision has no per-area euro map: its per-cell plan "
                  "costs are still being produced.")
    pic(s, A / "fig41_map_freq_by_theta.png", L, BODY_T, W, 3.45)
    bullets(s, ["Dark means changed most: two delivery days a week.",
                "The periphery consolidates first.",
                "The urban core keeps daily service longest."],
            BODY_T + 3.60, h=2.2)

    # ═══ where it pays ════════════════════════════════════════════════════
    # Per-area EURO only exists on the submission grid: the revision's per-cell
    # plan costs (Task 11) are not in the tables yet, so this slide says which
    # grid it is on rather than implying the revision restated it.
    s = new_slide(prs, "Results", "TBC pays where delivery is sparse and far",
                  "Median per postal-code area at each provider's P*, "
                  "submission grid — the only grid with a per-area euro "
                  "decomposition. On the revision grid the same structure is "
                  "visible in delivery frequency (next slides).")
    stats(s, [("9%", "urban median saving", False),
              ("25%", "rural median saving", True)], BODY_T, w=5.9, h=1.35)
    bullets(s, ["The gap is structural, not a modelling artefact.",
                "Long stems and large areas help.",
                "High parcels-per-stop suppresses the gain."],
            BODY_T + 1.70, w=5.9, h=2.5)
    pic(s, A / "fig71_map_saving.png", COL2 - 0.15, BODY_T, 6.05, 4.7)

    # ═══ drivers ══════════════════════════════════════════════════════════
    s = new_slide(prs, "Results", "Three structural drivers, one direction",
                  "Spearman correlations across the 312 provider–area cells at each "
                  "provider's own P*, submission grid. The revision reproduces the "
                  "same three directions in delivery frequency (Fig. 6 d–f).")
    _y = table(s, ["Driver", "Correlation", "Reading"],
          [[("Hub distance", "key"), ("ρ = +0.53", "num"),
            "Long stems are amortised by batching"],
           [("Area size", "key"), ("ρ = +0.31", "num"),
            "Large areas gain from denser tours"],
           [("Parcels per drop-site", "key"), ("ρ = −0.72", "num"),
            "Where density exists, little is left to win"]],
          BODY_T, widths=[3.4, 2.6, 5.0], reserve=1.20)
    label_box(s, L, _y + 0.28, W, 0.95, BLUSH,
              [("Spatial targeting is the strategy — not uniform service "
                "degradation.", SZ_LEAD + 2, True, RED)], line_col=RED)

    # ═══ carrier classes ══════════════════════════════════════════════════
    # The classes are now lens-qualified: three LSPs move up one class in the
    # operator lens, so the card grid is built from the grid's own table.
    s = new_slide(prs, "Results", "One policy does not fit every carrier",
                  "Chord-distance knees at θ = 1 in the ROUTING lens; heuristic "
                  "operating points, not a normative recommendation. Three LSPs "
                  "sit one class higher in the operator lens — see Revision · 7.")
    if facts is not None:
        _by_class = {}
        for _p in facts.pstar:
            _by_class.setdefault(_p["class_routing"], []).append(_p)
        classes = []
        for _nm, _why, _col in (
                ("Service-bound", "Dense, high parcels-per-stop", S4),
                ("Hybrid", "Mixed density", S5),
                ("Cost-aggressive", "Sparse, long stems", S6)):
            _members = _by_class.get(_nm, [])
            _ps = sorted({m["p_routing"] for m in _members})
            classes.append((
                _nm, "P* = " + " / ".join(f"{v:g}" for v in _ps),
                " · ".join(m["provider"] for m in _members), _why, _col))
    else:
        classes = [("Service-bound", "P* = 0.25", "Amazon · DHL",
                    "Dense, high parcels-per-stop", S4),
                   ("Hybrid", "P* = 0.5", "FedEx · Hermes · UPS",
                    "Mixed density", S5),
                   ("Cost-aggressive", "P* = 0.75", "DPD · GLS",
                    "Sparse, long stems", S6)]
    cwd = (W - 2 * 0.5) / 3
    for i, (nm, pstar, who, why, col) in enumerate(classes):
        x = L + i * (cwd + 0.5)
        rect(s, x, BODY_T, cwd, 0.10, col)
        txt(s, x, BODY_T + 0.24, cwd, 0.42, nm, SZ_LEAD, bold=True, color=col)
        txt(s, x, BODY_T + 0.80, cwd, 0.60, pstar, 32, bold=True, color=RED)
        txt(s, x, BODY_T + 1.50, cwd, 0.42, who, SZ_BODY, bold=True, color=INK)
        txt(s, x, BODY_T + 1.98, cwd, 0.75, why, SZ_BODY, color=INK2, line=1.2)
    if facts is not None:
        _sv = {p["provider"]: p for p in facts.pstar}
        _k = revision.D.load_pstar_v2().set_index("provider")
        bullets(s, ["Tolerance for waiting tracks how much a network gains.",
                    "DHL already carries 41% of volume at the lowest unit cost.",
                    f"So at its knee it reaches only "
                    f"{_k.loc['DHL', 'saving_pct_routing']:.1f} % — against "
                    f"{_k.loc['GLS', 'saving_pct_routing']:.1f} % for GLS."],
                BODY_T + 2.90, h=2.2)
        revision.notes(s, revision.pstar_notes(facts), cite=revision.cites())
        revision.provisional(s, enabled=tag)
    else:
        bullets(s, ["Tolerance for waiting tracks how much a network gains.",
                    "DHL already carries 41% of volume at the lowest unit cost.",
                    "So at its knee it reaches the least of any carrier."],
                BODY_T + 2.90, h=2.2)

    # ═══ knees figure ═════════════════════════════════════════════════════
    s = new_slide(prs, "Results", "Each carrier's own operating point",
                  "Chord-distance knee per provider at θ = 1, submission render; "
                  "the revision's knees and their second lens are on "
                  "Revision · 7.")
    pic(s, A / "fig36_pstar_knees.png", L, BODY_T, 7.5, 4.85)
    bullets(s, ["The knee is where the curve stops paying.",
                "It differs by network, not by preference.",
                "Uniform policy leaves value on the table."],
            BODY_T + 0.8, l=8.4, w=4.4, h=3.2)

    # ═══ parallel systems ═════════════════════════════════════════════════
    if facts is not None:
        _g = revision.D.load_grid_full_v2()
        _vd_lo = float(_g.vehicle_days_plan2_vs_base_pct.min())
        _vd_hi = float(_g.vehicle_days_plan2_vs_base_pct.max())
    else:
        _vd_lo, _vd_hi = -10.0, 4.6
    s = new_slide(prs, "Results", "Low adoption creates parallel delivery systems",
                  f"Weekly vehicle-days span {_vd_lo:+.1f} % to {_vd_hi:+.1f} % "
                  f"across the grid; peak and CV benefits arise earlier than "
                  f"absolute reductions.")
    for i, (nm, col) in enumerate([("Conventional tour", S4),
                                   ("Batched tour", RED)]):
        y = BODY_T + 0.10 + i * 0.88
        rect(s, L, y, 0.10, 0.62, col)
        txt(s, L + 0.32, y + 0.10, 3.0, 0.42, nm, SZ_BODY, bold=True, color=INK)
        for k in range(6):
            rect(s, L + 3.55 + k * 0.44, y + 0.14, 0.32, 0.36,
                 col if (i == 0 or k % 3 == 1) else S1)
    label_box(s, L, BODY_T + 1.98, 6.2, 0.85, BLUSH,
              [("At low θ, both systems must run.", SZ_LEAD, True, RED)],
              line_col=RED)
    stats(s, [(f"{_vd_hi:+.1f}%", "worst-case vehicle-day rise", True),
              (f"{_vd_lo:+.1f}%", "best case, only at high θ", False)],
          BODY_T + 3.10, w=6.2, h=1.35, sz=40)
    pic(s, B / "fig52_fleet_per_provider.png", COL2 + 0.35, BODY_T, 5.5, 4.4)
    if facts is not None:
        revision.notes(s, "Weekly vehicle-days of the operator-polished plan "
                          "against the daily baseline. The submission's +4.6 % "
                          "worst case came from the old pooled-express fleet "
                          "count and is gone: on the revision grid no cell "
                          "raises vehicle-days by more than "
                          f"{_vd_hi:+.2f} %.", cite=["§40.9", "§40.15"])
        revision.provisional(s, enabled=tag)

    # ═══ validation ═══════════════════════════════════════════════════════
    # The revision's VROOM validation (both plans, both lenses) is still being
    # produced; until it is, these two slides stay on the submission's
    # validation and say so.
    s = new_slide(prs, "Validation",
                  "Real routing confirms a conservative surrogate",
                  "Four out-of-sample operating points at θ = 1 "
                  "(P ∈ {0, 0.25, 0.5, 0.75}), re-routed with VROOM/Valhalla, "
                  "SUBMISSION grid. n = 1 248 · bias = +2.73%. The revision's "
                  "re-validation of both plans is in progress.")
    stats(s, [("+0.9 … +2.7 pp", "realized saving above prediction", True),
              ("3.04%", "cost error against the solver", False),
              ("0.997", "R² over 1 248 observations", False)], BODY_T, h=1.15,
          sz=38)
    pic(s, A / "fig61_vroom_scatter.png", L, BODY_T + 1.35, 5.85, 3.4)
    pic(s, A / "fig62_pred_vs_actual.png", COL2, BODY_T + 1.35, 5.85, 3.4)

    # ═══ the error direction ══════════════════════════════════════════════
    s = new_slide(prs, "Validation", "The error points in the safe direction",
                  "Predicted versus realised saving at the four validated operating "
                  "points of the SUBMISSION grid; the surrogate is conservative at "
                  "every one. The direction, not the level, is what carries over.")
    _y = table(s, ["Operating point", "Predicted", "Realised", "Gap"],
          [[("P = 0", "key"), "22.8%", ("23.7%", "good"), "+0.9 pp"],
           [("P = 0.25", "key"), "18.5%", ("19.8%", "good"), "+1.3 pp"],
           [("P = 0.5", "key"), "13.5%", ("15.6%", "good"), "+2.1 pp"],
           [("P = 0.75", "key"), "10.2%", ("13.0%", "good"), "+2.8 pp"]],
          BODY_T, widths=[3.0, 2.6, 2.6, 2.8], reserve=2.15)
    bullets(s, ["The surrogate understates what the solver achieves.",
                "The gap widens at higher penalties.",
                "A schedule chosen on the surrogate does at least as well."],
            _y + 0.26, h=2.1)
    if facts is not None:
        # No revision counterpart yet, so the numbers stay and the slide says
        # so where nobody can miss it -- a banner, not a source line.
        revision.stamp(s)
        revision.notes(s, "SUBMISSION grid. The direction of the error -- the "
                          "surrogate under-promises -- is what carries over "
                          "to the revision; the levels do not. The revision's "
                          "re-run of both plans is still being produced.",
                       cite="§40.18")

    # ═══ implications ═════════════════════════════════════════════════════
    s = new_slide(prs, "Implications", "Start where the advantage is strongest",
                  "Managerial implications from the Revision 1 conclusion.")
    for i, (n, nm, body) in enumerate([
            ("1", "Target",
             "Rural and depot-distant areas with few parcels per stop"),
            ("2", "Offer",
             "A discounted flexible tier, urgent delivery preserved"),
            ("3", "Tune", "P and adoption by provider and local structure"),
            ("4", "Scale", "Only once participation avoids parallel fleets")]):
        y = BODY_T + i * 1.22
        label_box(s, L, y, 0.88, 0.88, RED, [(n, 30, True, WHITE)])
        txt(s, L + 1.15, y + 0.02, 3.2, 0.44, nm, SZ_LEAD, bold=True, color=RED)
        txt(s, L + 1.15, y + 0.50, W - 1.15, 0.44, body, SZ_BODY, color=INK2)

    # ═══ design principle ═════════════════════════════════════════════════
    s = new_slide(prs, "Implications",
                  "One principle holds the whole design together",
                  "Equity caveat drawn from the transport-justice literature.")
    label_box(s, L, BODY_T + 0.15, W, 1.50, BLUSH,
              [("Do not make rural customers pay for low density with a worse "
                "default.", 30, True, RED)], line_col=RED, pad=0.34)
    bullets(s, ["Preserve choice: the flexible tier is opt-in, not a downgrade.",
                "Use part of the saving to fund that option.",
                "Otherwise a cost-driven rollout penalises those with fewest "
                "alternatives.",
                "TBC is a service-design lever, not just a routing trick."],
            BODY_T + 2.00, h=3.3)

    # ═══ takeaway ═════════════════════════════════════════════════════════
    if facts is not None:
        _r, _h = facts.recommended(), facts.headline[0.5]
        _close = (f"{_h['op2']:.1f}–{_r['op2']:.1f}% operator saving in the "
                  f"efficient range at full adoption, with "
                  f"{abs(_r['peak2_pct']):.0f}% fewer peak vehicles.")
    else:
        _close = "Saving in the efficient range at full adoption."
    full_bleed(prs, "The takeaway", "Batch where it is\nsparse and far.",
               "Temporal flexibility creates the density that non-urban delivery "
               "is missing — without new infrastructure.",
               closing=_close)

    # ═══ appendix: invariants 1 ═══════════════════════════════════════════
    s = new_slide(prs, "Appendix", "Hard invariants of the model",
                  "AGENTS.md hard invariants; Paper §2.2; "
                  "docs/HOLDING_DAYS_INVARIANT.md.")
    _y = table(s, ["Setting", "Value", "Status"],
          [[("Operating week", "key"), ("Monday–Saturday", "num"),
            "six delivery days"],
           [("Maximum holding", "key"), ("Hmax = 3 days", "num"),
            "authoritative; enforced by config validation and unit tests"],
           [("Weekly patterns", "key"), ("39 per cell", "num"),
            "the same candidate set for every cell"],
           [("Baseline", "key"), ("daily delivery", "num"),
            (f"no batching; {revision.eur(facts.base_routing)} € routing / "
             f"{revision.eur(facts.base_operator)} € operator per week"
             if facts is not None else "no batching")]],
          BODY_T, widths=[3.0, 3.2, 4.8])

    # ═══ appendix: invariants 2 ═══════════════════════════════════════════
    s = new_slide(prs, "Appendix", "Scope and cost parameters",
                  "Standard operational values for a >2 t delivery van; the learned "
                  "residual corrects the rest.")
    _y = table(s, ["Setting", "Value", "Note"],
          [[("Provider scope", "key"), ("separate networks", "num"),
            "no cross-carrier sharing of parcels, routes or vehicles"],
           [("Service penalty", "key"), ("P [€/parcel/day]", "num"),
            "steering term; booked as cost only in the discount scenario"],
           [("Vehicle capacity", "key"), ("Q = 230 parcels", "num"),
            "also the minimum tour size of the universal tour rule"],
           [("Cost model", "key"), ("189.15 € + 0.3864 €/km", "num"),
            "per vehicle-day and per kilometre, plus 36 € per route-hour "
            "(VROOM per_hour default, active in every label)"]],
          BODY_T, widths=[3.0, 3.2, 4.8])
    if facts is not None:
        revision.notes(
            s, "The effective cost model measured on the training pool is "
               f"{revision.D.COST_MODEL_SENTENCE} — 72 / 6 / 22 % of the "
               "total. The route-hour term is VROOM's per_hour default, "
               "active in every label because requests.py does not set "
               "per_hour when it is 0. The pipeline is internally consistent "
               "(labels = surrogate = validation = bundle pool, same "
               "builder); the direction is conservative, because service time "
               "is nearly constant between baseline and scenario and a "
               "near-constant term in numerator and denominator depresses the "
               "percentage.", cite="§40.12")
        revision.provisional(s, enabled=tag)
        # the withdrawn theta = 10 % finding, immediately after the parameters
        slide_bulge_replacement(prs, facts, revision, tag=tag)
        # the revision's own multi-panel figures, adopted by
        # 95_adopt_paper_figs.py. Backup tier: too dense to project as a
        # statement, exactly right when somebody asks for the whole grid.
        for _nm, _subj, _src in (
                ("fig90_grid_two_lenses", "The whole grid, in both lenses",
                 "Revision grid. Panels (a) and (b) are different lenses AND "
                 "different plans — do not read them as one series."),
                ("fig91_offdiagonal", "Each plan priced in the other lens",
                 "Revision grid: the two off-diagonal combinations."),
                ("fig92_freq_mix_two_plans",
                 "Delivery-frequency mix, both plans",
                 "Revision grid, routing-optimal above, operator-polished "
                 "below."),
                ("fig93_mean_days", "Mean delivery days per area",
                 "Revision grid, both plans."),
                ("fig95_operator_lens", "The operator lens on its own",
                 "Revision grid: the operator-cost view without the routing "
                 "lens beside it."),
                ("fig94_structural_two_lenses",
                 "Fronts, knees and structure",
                 "Revision grid. Panel (f) is a within-DHL statement: the "
                 "hub-size buckets hold DHL cells only, because DHL is the "
                 "only multi-depot network in the case study.")):
            _p = B / f"{_nm}.png"
            if not _p.exists():
                print(f"  ! {_nm}.png missing — run 95_adopt_paper_figs.py")
                continue
            _s = new_slide(prs, "Appendix · revision figures", _subj, _src)
            pic(_s, _p, L, BODY_T, W, 5.05)
            revision.provisional(_s, enabled=tag)

    # ═══ appendix: limitations 1 ══════════════════════════════════════════
    s = new_slide(prs, "Appendix", "Limitations · what bounds the result",
                  "Revision 1 conclusion and reviewer-response framing.")
    _y = table(s, ["Boundary", "What it means", "Next step"],
          [[("Routing benchmark", "key"),
            "VROOM is a detailed benchmark, not observed ground truth",
            "Calibrate against operating data"],
           [("Spatial decomposition", "key"),
            "Postal-code cells exclude cross-area routing synergies",
            "Jointly route neighbouring areas"]],
          BODY_T, widths=[2.8, 4.4, 3.8], reserve=1.05)
    bullets(s, ["Neither bound changes the direction of the result."],
            _y + 0.28, h=1.0)

    # ═══ appendix: limitations 2 ══════════════════════════════════════════
    s = new_slide(prs, "Appendix", "Limitations · what needs new data",
                  "Revision 1 conclusion and reviewer-response framing.")
    _y = table(s, ["Boundary", "What it means", "Next step"],
          [[("Adoption", "key"),
            "θ is scenario-based; customer uptake is not estimated here",
            "Estimate choice and price response"],
           [("Pricing", "key"), "P is a shadow price, not a market tariff",
            "Connect service design to willingness-to-pay"],
           [("Transferability", "key"),
            "One German region and seven modelled networks",
            "Replicate across network archetypes"]],
          BODY_T, widths=[2.8, 4.4, 3.8], reserve=1.15)
    label_box(s, L, _y + 0.26, W, 0.90, PANEL,
              [("The framework is reusable; the numerical operating point is "
                "context-specific.", SZ_LEAD, True, RED)], line_col=LINE)

    # ═══ close ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    rect(s, L, 2.20, 0.09, 2.35, RED)
    txt(s, L + 0.40, 2.28, 10.5, 0.34, "Thank you", SZ_KICK, bold=True, color=RED,
        spc=1.6, caps=True)
    txt(s, L + 0.40, 2.72, 11.5, 0.95, "Questions?", 44, bold=True, color=INK)
    txt(s, L + 0.40, 3.85, 11.5, 1.3,
        "l.bienzeisler@tu-braunschweig.de\n"
        "github.com/TUBS-IVS/vroom-valhalla-lmd-hannover\n"
        "Figures, tables and the full pipeline are reproducible from the repository.",
        SZ_STATL + 1, color=INK2, line=1.7)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def load_facts(rev_dir: Path | None):
    """Read the revision grid and its facts, or explain why it cannot be."""
    import _data as D
    import _revision as RV
    if rev_dir is not None:
        D.set_rev_dir(rev_dir)
    if D.SCHEMA != D.SCHEMA_V2:
        raise SystemExit(
            f"{D.REV} is a {D.SCHEMA} grid; the revision slides need the "
            f"two-plan tables. Pass --rev-dir or set PRES_REV_DIR.")
    print(f"  revision grid: {D.REV.relative_to(D.ROOT)}")
    return RV.Facts.load(), RV


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=Path, default=DEFAULT_OUT)
    ap.add_argument("--keep-template-slides", action="store_true",
                    help="append after the template's own slides instead of "
                         "starting from an empty deck")
    ap.add_argument("--rev-dir", type=Path, default=None,
                    help="the revision grid to read (default: $PRES_REV_DIR, "
                         "else results/revision_2026_08_v5)")
    ap.add_argument("--no-provisional", action="store_true",
                    help="drop the 'v5 · provisional' footer tag — for Part B, "
                         "once the head grid is final")
    ap.add_argument("--no-revision", action="store_true",
                    help="build the submission-era deck without the revision "
                         "block and without reading any grid")
    G.add_args(ap)
    a = ap.parse_args()
    if not TEMPLATE.exists():
        print(f"template not found: {TEMPLATE}", file=sys.stderr)
        return 1
    missing = [p for p in (FIG / "tierA", FIG / "tierB", ASSET) if not p.exists()]
    if missing:
        print(f"figure directories missing: {missing}", file=sys.stderr)
        return 1
    out = G.resolve(a.out, a.out_suffix, overwrite=a.overwrite)
    facts, revision = (None, None) if a.no_revision else load_facts(a.rev_dir)
    p = build(out, a.keep_template_slides, facts=facts, revision=revision,
              tag=not a.no_provisional)
    print(f"wrote {p}")
    print(f"  {len(Presentation(str(p)).slides)} slides, "
          f"{p.stat().st_size / 1048576:.1f} MB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
