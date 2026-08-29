"""Export an HTML slide deck to PowerPoint, one rendered image per slide.

This is a *faithful* export, not a conversion: each `<section class="slide">` of
the source deck is rendered by headless Chrome at 16:9 and placed full-bleed on
its own PowerPoint slide. The result looks exactly like the HTML deck and can be
presented from PowerPoint, but the slides are images -- they cannot be edited
shape by shape.

Use it when the HTML deck itself is the artefact worth carrying into PowerPoint.
For an editable deck, use 91_build_pptx.py, which emits native shapes.

The output goes through `_outguard.resolve()`, so a rebuild lands beside the
existing export instead of on top of it unless `--overwrite` says otherwise.

Usage:
    python scripts/presentation/92_export_html_deck_pptx.py
    python scripts/presentation/92_export_html_deck_pptx.py --src DECK.html --out OUT.pptx
    python scripts/presentation/92_export_html_deck_pptx.py --out-suffix _rev2026-08
    python scripts/presentation/92_export_html_deck_pptx.py --scale 2   # 2x raster
"""
from __future__ import annotations

import argparse
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))

import _outguard as G                                             # noqa: E402

ROOT = Path(__file__).resolve().parents[2]
DEFAULT_SRC = Path(r"C:/Users/bienzeisler/.agent/diagrams/tbc-ewgt2026-slides.html")
DEFAULT_OUT = Path(
    r"C:/Users/bienzeisler/Documents/Präsentationen/EWGT/2026/"
    r"EWGT_26_Bienzeisler_TBC_slides_from_html.pptx"
)
CHROME_CANDIDATES = (
    r"C:/Program Files/Google/Chrome/Application/chrome.exe",
    r"C:/Program Files (x86)/Google/Chrome/Application/chrome.exe",
    r"C:/Program Files (x86)/Microsoft/Edge/Application/msedge.exe",
)
W_IN, H_IN = 13.333, 7.5          # 16:9, matching the other deck


def find_chrome() -> str:
    for c in CHROME_CANDIDATES:
        if Path(c).exists():
            return c
    raise SystemExit("no Chrome or Edge found for headless rendering")


DECK_OPEN = '<div class="deck">'


def split_slides(src: Path) -> tuple[str, list[str]]:
    """Return (document head, [slide section html]) for the source deck."""
    html = src.read_text(encoding="utf-8")
    marker = DECK_OPEN
    if marker not in html:
        raise SystemExit(f"{src.name}: no <div class=\"deck\"> found")
    head = html[: html.index(marker)]
    sections = re.findall(r'(<section class="slide.*?</section>)', html, re.S)
    if not sections:
        raise SystemExit(f"{src.name}: no slide sections found")
    return head, sections


def render(src: Path, out: Path, scale: int, keep: bool) -> Path:
    chrome = find_chrome()
    head, sections = split_slides(src)
    # The deck reveals slides on scroll, so force every element visible and
    # disable the entry transitions before shooting a single slide.
    force = ("<style>.slide,.slide .reveal,.slide .bullets li{opacity:1!important;"
             "transform:none!important;animation:none!important}"
             ".slide .heading::after{transform:scaleX(1)!important}"
             ".deck-progress,.deck-dots,.deck-foot,.deck-hints,.sect{display:none!important}"
             "</style>")
    # Relative asset paths in the source resolve against its own folder.
    base = f'<base href="{src.parent.as_uri()}/">'
    head = head.replace("<meta charset=\"utf-8\">",
                        "<meta charset=\"utf-8\">" + base, 1)

    work = Path(tempfile.mkdtemp(prefix="deck_export_"))
    # The viewport stays at 1920x1080 CSS pixels whatever the scale: the deck's
    # type and padding are in fixed px, so widening the window to 3840 shrinks
    # every element relative to the canvas instead of sharpening it. Resolution
    # comes from the device scale factor, which multiplies the output raster
    # while leaving layout untouched.
    pngs: list[Path] = []
    try:
        for i, sec in enumerate(sections, 1):
            page = work / f"s{i:03d}.html"
            sec_vis = (sec.replace('class="slide ', 'class="slide visible ')
                          .replace('class="slide"', 'class="slide visible"'))
            page.write_text(head + DECK_OPEN + sec_vis + "</div>" + force
                            + "</body></html>", encoding="utf-8")
            png = work / f"s{i:03d}.png"
            subprocess.run(
                [chrome, "--headless=new", "--disable-gpu", "--no-sandbox",
                 "--hide-scrollbars", "--window-size=1920,1080",
                 f"--force-device-scale-factor={scale}",
                 f"--virtual-time-budget={4000 + 1500 * scale}",
                 f"--screenshot={png}", page.as_uri()],
                check=False, capture_output=True, timeout=180,
            )
            if not png.exists():
                raise SystemExit(f"slide {i}: Chrome produced no image")
            pngs.append(png)
            if i % 10 == 0 or i == len(sections):
                print(f"  rendered {i}/{len(sections)}")

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(W_IN), Inches(H_IN)
        blank = prs.slide_layouts[6]
        for png in pngs:
            s = prs.slides.add_slide(blank)
            s.shapes.add_picture(str(png), 0, 0, Inches(W_IN), Inches(H_IN))
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
    finally:
        if keep:
            print(f"  renders kept in {work}")
        else:
            shutil.rmtree(work, ignore_errors=True)
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--src", type=Path, default=DEFAULT_SRC)
    ap.add_argument("--out", type=Path, default=DEFAULT_OUT)
    ap.add_argument("--scale", type=int, default=2,
                    help="device pixel ratio: 1 = 1920x1080 per slide, "
                         "2 = 3840x2160. The CSS viewport stays 1920x1080 "
                         "either way, so raising this sharpens the image "
                         "without shrinking the deck's type.")
    ap.add_argument("--keep-renders", action="store_true")
    G.add_args(ap)
    a = ap.parse_args()
    if not a.src.exists():
        print(f"source deck not found: {a.src}", file=sys.stderr)
        return 1
    out = G.resolve(a.out, a.out_suffix, overwrite=a.overwrite)
    print(f"exporting {a.src.name} at {1920 * a.scale}x{1080 * a.scale} "
          f"per slide (1920x1080 CSS px at {a.scale}x)")
    p = render(a.src, out, a.scale, a.keep_renders)
    print(f"wrote {p}")
    print(f"  {len(Presentation(str(p)).slides)} slides, "
          f"{p.stat().st_size / 1048576:.1f} MB")
    print("  note: slides are images -- use 91_build_pptx.py for an editable deck")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
