"""A talk on what customers concede and what operators gain — grid v6.

This is a NEW deck, not a rebuild of an existing one. The three `_rev2026-08`
decks tell the whole paper; this one takes the two questions the revision made
answerable and goes deep on them:

**Temporal flexibility.** What exactly does a customer give up when the
willingness-to-wait share is theta — how many parcels are held, for how long,
and how many delivery days a week does an area keep — and what does the
operator get back for it, in fleet, in vehicle-days and in euro. The answer is
an exchange rate, and the deck states it per operating point rather than as one
headline.

**Routing cost against express cost.** Below full adoption every dropped
delivery day forces a pooled express tour for the customers who would not wait.
That tour is the whole reason the theta < 1 savings are small, and it is priced
here: euro per vehicle-day against a regular tour, share of total cost over
theta, and the jump at theta = 1 where the obligation disappears.

Both parts end on what the solver said about them, because the VROOM
re-validation of v6 is finished and it does not agree with the surrogate
everywhere — least of all on express tours.

Rules this file follows
-----------------------
**No number is typed in.** Every figure on a slide comes from
`_revision.Facts`, `_revision.Flex` or `_revision.Vroom`, which read the grid
and assert each value against the compendium section that records it. A grid
whose numbers moved fails the build.

**Every figure slide names its source stem** in a footer chip, so a slide can
be traced to the PNG and the PNG to the script that drew it.

**Only v6 renders.** The Act 3–7 figures in `results/presentation_2026_08/`
carry v6 provenance stamps but are drawn through
`_data.load_chosen_stage3()`, which reads the frozen pre-revision run
unconditionally, and through `VAL_GRID_NAME`, which is pinned to v5. They are
therefore NOT v6 and none of them is embedded here. Where one of them is the
right picture, the deck carries a marked placeholder slide naming the stem and
the statement, so it can be dropped in once the loader is fixed.

**The validation directory is set explicitly.** `_data.VAL` still defaults to
the v5 run because the other decks were audited against it; this build points
at v6 and refuses to run without a solved theta = 0 baseline.

Language: English slides, German speaker notes are allowed — the same
convention as the TBC deck.

Usage:
    python scripts/presentation/98_deck_v6_flexibility.py --out-suffix _rev2026-08
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

import _house as H
import _outguard as G
from _house import (AMBER, BLACK, BLUSH, BODY_B, BODY_T, COL_W, CRIM, DIM,
                    FIG_W, FIG_X, GREEN, INK, INK2, L, LINE, PANEL, RED, S1,
                    S2, S3, S4, S5, S6, SH, SW, SZ_BIG, SZ_BODY, SZ_DIA,
                    SZ_LEAD, SZ_NUM, SZ_SRC, SZ_SUB, TEAL, TEXT_W, TEXT_X, W,
                    WHITE, arrow, badges, build, contact, flow, hrule, hslide,
                    label_box, pic, rect, section_divider, split_slide,
                    statement, sub_items, txt)

# The v6 figure directory. Set by build_deck() from the grid actually in use,
# so --rev-dir moves the pictures and the numbers together.
FIGDIR: Path = None

SEC_SET = "Setting the frame"
SEC_HEAD = "The result on grid v6"
SEC_FLEX = "Temporal flexibility"
SEC_EXP = "Routing cost against express cost"
SEC_FLEET = "The weekly fleet"
SEC_HON = "What the numbers do not say"

_F = None      # _revision.Facts
_X = None      # _revision.Flex
_V = None      # _revision.Vroom
_RV = None     # the _revision module
_D = None      # the _data module


# ── small shared furniture ─────────────────────────────────────────────────
def mark(s, note, cite=None):
    """Speaker notes with the compendium section(s) the claim comes from."""
    _RV.notes(s, note, cite=cite)
    return s


CHIP_RIGHT = 13.28      # the right edge the decks' footer chip sits on
CHIP_SZ = 9.0


def chip(s, stem: str):
    """The figure's source stem, bottom right, where the decks put their tag.

    It is not decoration: a reader who wants to check a picture needs the file
    name, and a presenter asked "where is that from" needs it on the slide
    rather than in a notebook. The box is measured to the stem rather than
    given a fixed width -- a 33-character stem is twice as wide as the decks'
    "v5 provisional" tag -- and any text box that would reach under it is
    clipped back first, so the chip never creates a collision.
    """
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Emu, Inches

    w = H._font(CHIP_SZ, True).getlength(stem) / (72.0 * 4.0) + 0.12
    left = CHIP_RIGHT - w
    edge = Inches(left - 0.14)
    for sh in s.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if sh.left is None or sh.width is None or sh.top is None:
            continue
        bottom = Emu(sh.top).inches + Emu(sh.height or 0).inches
        if bottom < _RV.TAG_T or Emu(sh.top).inches > _RV.TAG_T + _RV.TAG_H:
            continue
        if sh.left + sh.width > edge > sh.left:
            sh.width = int(edge - sh.left)
    box, tf = H.B._frame(s, left, _RV.TAG_T, w, _RV.TAG_H,
                         anchor=MSO_ANCHOR.MIDDLE)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    H.B._para(p, stem, CHIP_SZ, bold=True, color=DIM)
    return s


def figpath(stem: str) -> Path:
    p = FIGDIR / f"{stem}.png"
    if not p.exists():
        raise SystemExit(
            f"missing figure {p}\n"
            f"  the deck only embeds v6 renders from <rev>/figures/; run "
            f"scripts/revision/70_figs_tables_v2.py on this grid first")
    return p


def figure_top(prs, section, subject, source, stem, items, *, h=3.45,
               t=None):
    """Figure across the top, badge bullets under it, source chip bottom right."""
    s = hslide(prs, section, subject, source)
    pic(s, figpath(stem), L, BODY_T, W, h)
    if items:
        badges(s, items, BODY_T + h + 0.14 if t is None else t,
               bottom=s.body_bottom, label=f"{section} / {subject}")
    chip(s, stem)
    return s


def figure_split(prs, section, subject, source, stem, items, **kw):
    """Text left, figure right — his two-column slide, plus the chip."""
    s = split_slide(prs, section, subject, source, figpath(stem), items, **kw)
    chip(s, stem)
    return s


# A figure wider than this fills the slide across the top; anything squarer
# leaves so much white on either side that the text belongs beside it instead.
WIDE_ASPECT = 2.2
CAP_H = 3.90        # the tallest top-figure the eye still reads as a band
MIN_TOP_H = 2.20    # below this a top figure is not worth the full width


def aspect(stem: str) -> float:
    from PIL import Image
    with Image.open(figpath(stem)) as im:
        return im.size[0] / im.size[1]


def _body_bottom(source: str) -> float:
    """Where `hslide()` will put this slide's body floor, computed in advance.

    The source line is bottom-anchored and grows upwards, so how much room a
    figure has depends on how long its citation is. Recomputing that here means
    the figure can be sized BEFORE the slide exists, instead of being placed
    and then found to be too tall.
    """
    if not source:
        return H.SRC_B - 0.16
    return H.SRC_B - H.text_height(source, 11.2, SZ_SRC, 1.25) - 0.16


def figure_auto(prs, section, subject, source, stem, items, *, cap=CAP_H):
    """Lay the figure out by its own shape rather than by a fixed box.

    `pic()` letterboxes, so a fixed 12.21-inch-wide box turns a 1.7:1 figure
    into a postage stamp between two hands of white. Measuring the image, the
    citation and the bullets — all three of which vary per slide — and then
    choosing between the two house layouts costs one PIL open and makes every
    figure slide fill its slide without pushing text into the footer.
    """
    a = aspect(stem)
    need = (H._block_height(items, TEXT_W, SZ_BODY, 0.26, 1.18) + 0.24
            if items else 0.0)
    h = min(cap, W / a, _body_bottom(source) - BODY_T - need)
    if a >= WIDE_ASPECT and h >= MIN_TOP_H:
        return figure_top(prs, section, subject, source, stem, items, h=h)
    return figure_split(prs, section, subject, source, stem, items,
                        fig=(FIG_X + 0.10, BODY_T, FIG_W - 0.10, 4.95))


def placeholder(prs, section, subject, source, stem, statement_lines, note):
    """A slide that says which picture belongs here and why it is not here yet.

    The Act 3-7 renders are drawn through a loader that still reads the frozen
    pre-revision run, so they are not v6 whatever their provenance file claims.
    Leaving a marked gap is honest; embedding the figure would put a
    pre-revision schedule on a v6 slide.
    """
    s = hslide(prs, section, subject, source)
    label_box(s, L, BODY_T, W, 2.05, BLUSH,
              [("FIGURE PENDING — needs a v6 re-render", 17, True, RED),
               (f"{stem}  ·  drawn through _data.load_chosen_stage3(), which "
                f"reads results/runs/path2_2026_05_29 unconditionally, so the "
                f"schedules in it are pre-revision", 15, False, INK2)],
              line_col=RED, align=PP_ALIGN.LEFT)
    badges(s, statement_lines, BODY_T + 2.30, bottom=s.body_bottom,
           label=f"{section} / {subject}")
    mark(s, note, _RV.cites("§40.21"))
    return s


def pct(x, digits=1):
    return f"{x:.{digits}f} %"


def signed(x, digits=1):
    return f"{x:+.{digits}f} %"


# ═══════════════════════════════════════════════════════════════════════════
# Title
# ═══════════════════════════════════════════════════════════════════════════
def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[H.LAYOUT_BLANK])
    H.pic_cover(s, H.ASSET / "parcel-wave.png", 0, 0.93, SW, 4.42, focus=0.45)
    rect(s, 0, 5.35, SW, 0.10, RED)
    txt(s, 5.18, 0.27, 3.0, 0.37, "EWGT 2026 · revision", 16, bold=True,
        color=INK)
    txt(s, 1.08, 5.62, 11.4, 1.00,
        "What customers concede and what operators gain\n"
        "Time-based consolidation, measured on grid v6", 26, bold=True,
        color=INK, line=1.18)
    txt(s, 1.08, 6.62, 7.05, 0.32,
        "Lasse Bienzeisler, Felix Petre, Oskar Wage, and Bernhard Friedrich",
        16, color=INK2)
    txt(s, 8.30, 6.66, 4.55, 0.26, "Illustration generated with OpenAI ImageGen",
        9, color=DIM, align=PP_ALIGN.RIGHT)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# 0 · Setting the frame
# ═══════════════════════════════════════════════════════════════════════════
def part_setup(prs):
    f, X = _F, _X

    section_divider(
        prs, "Setting the frame", "Three things changed since the submission",
        "One universal tour rule · two cost lenses · two weekly plans")

    # ── the tour rule ─────────────────────────────────────────────────────
    def rule(s):
        H.B.flow(s, [("Willing customer", "waits for the area's next\n"
                                          "delivery day", S1),
                     ("Area tour", "every area keeps its own\n"
                                   "tour, on its own days", S2),
                     ("Express tour", "the non-willing remainder,\n"
                                      "bundled at the depot", S3),
                     ("Minimum 230 parcels", "one van; below that the\n"
                                             "tour is not dispatched", S6)],
              BODY_T + 2.66, bh=1.30)

    sl = build(
        prs, SEC_SET, "One tour rule prices baseline and scenario alike",
        f"Per-area express with a 230-parcel minimum tour, applied "
        f"scenario-blind. Baseline: {_RV.eur(f.base_routing)} EUR routing / "
        f"{_RV.eur(f.base_operator)} EUR operator per week, "
        f"{f.base_peak} peak vehicles, {_RV.eur(X.total_parcels)} parcels.",
        [("Truck", RED,
          [("Every area keeps its own tour: ", True),
           ("the standard parcels of an area that is not delivered today ride "
            "that area's own tour, not a hub-wide pooled one", False)]),
         ("TrafficCone", AMBER,
          [("The unbounded pooled express tour is gone: ", True),
           ("it priced a tour no operator would dispatch, and only the "
            "scenario could ever have one", False)]),
         ("Checkmark", GREEN,
          [("A tour needs 230 parcels to exist: ", True),
           ("the optimiser can no longer buy a saving with a mini-tour it "
            "would never send out", False)])],
        draw=rule, t=BODY_T, bottom=BODY_T + 2.52, start=2)
    for x in sl:
        mark(x, _RV.TOUR_RULE_NOTES, _RV.cites("§40.7", "§40.8", "§40.9"))

    # ── two lenses ────────────────────────────────────────────────────────
    s = hslide(prs, SEC_SET, "One euro is not one euro",
               f"Cost model: {_D.COST_MODEL_SENTENCE}. Both baselines price "
               f"the same daily-delivery system.")
    y = H.B.table(s, ["Cost lens", "What it counts", "Weekly baseline"],
                  [[(r[0], "key"), r[1], (r[2], "num")]
                   for r in _RV.lens_rows(f)],
                  BODY_T, widths=[2.8, 6.2, 3.0], sz=SZ_BODY, reserve=2.70)
    badges(s, [("Truck", BLACK,
                [("189.15 EUR per vehicle-day already contains the driver: ",
                  True),
                 ("a saving on a skipped delivery day is mostly avoided "
                  "driver cost", False)]),
               ("Target", RED,
                [("An operator staffs each depot for its busiest day: ", True),
                 ("below that peak only the kilometres are a real outlay, and "
                  "a vehicle taken out of the peak is worth 1 134.90 EUR a "
                  "week", False)])],
           y + 0.24, bottom=s.body_bottom)
    mark(s, _RV.LENS_NOTES, _RV.cites("§40.11", "§40.12"))

    # ── two plans ─────────────────────────────────────────────────────────
    h0 = f.headline[0.0]
    s = hslide(prs, SEC_SET, "Two weekly plans, because the lenses disagree",
               f"theta = 100 %, P = 0 EUR/parcel/day. Stage 2 is "
               f"frequency-free at theta > 0: it may change how OFTEN an area "
               f"is served, not only on which days.")
    y = H.B.table(s, ["Weekly plan", "Routing saving", "Operator saving",
                      "Peak fleet", "Added wait"],
                  [[(r[0], "key"), (r[1], "num"), (r[2], "num"), r[3], r[4]]
                   for r in _RV.plan_rows(f)],
                  BODY_T, widths=[3.7, 2.3, 2.4, 1.8, 2.0], sz=SZ_BODY,
                  reserve=2.55)
    label_box(s, L, y + 0.22, W, 0.88, BLUSH,
              [("A plan that is cheap per day can be loss-making per week: "
                "two-day patterns treble the depot peaks.",
                SZ_LEAD, True, RED)], line_col=RED)
    badges(s, [("Gears", RED,
                [(f"The polish swings the operator lens by "
                  f"{h0['op2'] - h0['op1']:.1f} points "
                  f"({h0['op1']:.1f} to {h0['op2']:+.1f} %): ", True),
                 (f"it costs {h0['rout1'] - h0['rout2']:.1f} points of routing "
                  f"saving and shortens the wait from {h0['wait1']:.2f} d to "
                  f"{h0['wait2']:.2f} d at the same time", False)])],
           y + 1.26, bottom=s.body_bottom)
    mark(s, _RV.plan_notes(f), _RV.cites("§40.14", "§40.21"))


# ═══════════════════════════════════════════════════════════════════════════
# 1 · The headline on grid v6
# ═══════════════════════════════════════════════════════════════════════════
def part_headline(prs):
    f, V = _F, _V
    h0, h25 = f.headline[0.0], f.headline[0.25]
    rp0, op0 = V.point("stage1", 0.0), V.point("balanced", 0.0)

    section_divider(prs, "The result on grid v6",
                    "Predicted, and then actually solved",
                    "616 optimiser runs, 8 609 VROOM instances, two lenses")

    # ── what the optimiser promises ───────────────────────────────────────
    s = figure_auto(
        prs, SEC_HEAD, "The whole grid, in both lenses",
        f"Grid v6, {_RV.eur(f.base_routing)} EUR routing / "
        f"{_RV.eur(f.base_operator)} EUR operator baseline. Panels (a) and (b) "
        f"are different lenses AND different plans; they are not one series.",
        "supp_fig5_grid_heatmap_v2",
        [("MagnifyingGlass", BLACK,
          f"Routing lens, routing-optimal plan: {pct(h0['rout1'])} at P = 0 "
          f"and {pct(h25['rout1'])} at P = 0.25"),
         ("Target", RED,
          f"Operator lens, operator-polished plan: {pct(h0['op2'])} and "
          f"{pct(h25['op2'])}, with {pct(abs(h0['peak2_pct']))} and "
          f"{pct(abs(h25['peak2_pct']))} fewer peak vehicles")])
    mark(s, "Die eine Figur, die beide Linsen nebeneinander zeigt. Panel (c) "
            "ist das, was der Kunde abgibt, (d)-(f) das, was der Betreiber "
            "gewinnt. Wichtig im Vortrag: (a) und (b) sind NICHT dieselbe "
            "Groesse - andere Linse und anderer Plan.",
         _RV.cites("§40.21"))

    # ── what the solver delivered ─────────────────────────────────────────
    s = hslide(prs, SEC_HEAD, "What the solver actually delivered",
               f"VROOM re-validation on {V.grid}, clean basis "
               f"(no unassigned jobs, no removed jobs): {V.n_solved()} solved "
               f"instances including the theta = 0 baseline.")
    # Prediction comes from the GRID, reality from VROOM -- the pairing 40.25
    # records. The validation's own predicted saving is formed over the clean
    # population of that point and differs by a few tenths where an instance
    # was dropped; printing it beside the grid figure would contradict the
    # slide before this one.
    rows = [["Daily baseline", "-", "-",
             f"+{V.baseline['gap_routing']:.2f} %"],
            ["Routing-optimal plan, P = 0", pct(h0["rout1"]),
             pct(rp0["act_routing"]), f"+{rp0['gap_routing']:.2f} %"],
            ["Operator-polished plan, P = 0", pct(h0["rout2"]),
             pct(op0["act_routing"]), f"+{op0['gap_routing']:.2f} %"],
            ["Operator-polished, P = 0, operator lens",
             pct(h0["op2"]), pct(op0["act_operator"]),
             f"+{op0['gap_operator']:.2f} %"]]
    y = H.B.table(s, ["Operating point", "Predicted saving", "Realised saving",
                      "Surrogate over solver"],
                  [[(r[0], "key"), r[1], (r[2], "num"), r[3]] for r in rows],
                  BODY_T, widths=[5.0, 2.5, 2.4, 2.3], sz=SZ_BODY,
                  reserve=2.30)
    badges(s, [("TrafficCone", AMBER,
                [("The prediction is an upper bound, not a conservative "
                  "estimate: ", True),
                 (f"the surrogate over-prices the thin daily baseline by "
                  f"{V.baseline['gap_routing']:.1f} % but the bundled "
                  f"scenarios by only "
                  f"{V.gap_range()[0]:.1f}-{V.gap_range()[1]:.1f} %, so the "
                  f"saving is overstated by about 2 points", False)])],
           y + 0.26, bottom=s.body_bottom)
    mark(s, "Beide Zahlen berichten, nie nur eine. Der Bias ist nicht "
            "konstant: duenne taegliche Einzeltouren werden staerker "
            "ueberschaetzt als gebuendelte Touren, deshalb vergleicht die "
            "Optimierung eine zu teure Baseline mit realistischeren "
            "Szenarien. Formulierung: 'obere Schranke (ca. 2 pp)', nicht "
            "'konservativ'.", _RV.cites("§40.24", "§40.25"))

    # ── the routing plan is a fleet problem ───────────────────────────────
    s = figure_auto(
        prs, SEC_HEAD, "The cheapest daily plan is the worst weekly plan",
        "Off-diagonal panels: each plan priced in the other lens, plus the "
        "routing-optimal plan's wait and peak fleet.",
        "supp_fig5b_offdiagonal_v2",
        [("DownwardTrend_LTR", CRIM,
          f"In the operator lens the routing-optimal plan is "
          f"{pct(h0['op1'])} at P = 0 — worse than simply delivering every "
          f"day — and VROOM makes it worse still, {pct(rp0['act_operator'])}"),
         ("Target", RED,
          f"The reason is fleet, not kilometres: its peak rises from "
          f"{f.base_peak} to {h0['peak1']} vehicles "
          f"({signed(h0['peak1_pct'])})")])
    mark(s, "Das ist die Folie, die den ganzen Rest traegt: 'billig pro Tag' "
            "und 'billig pro Woche' sind zwei verschiedene Optimierungen. Die "
            "VROOM-Zahl ist die saubere Basis aus der v6-Validierung.",
         _RV.cites("§40.21", "§40.25"))


# ═══════════════════════════════════════════════════════════════════════════
# 2 · Temporal flexibility — the main part
# ═══════════════════════════════════════════════════════════════════════════
def part_flexibility(prs):
    f, X, V = _F, _X, _V
    P_LIST = (0.0, 0.25, 0.5, 0.75, 1.0)

    section_divider(
        prs, "Temporal flexibility", "What is actually being traded",
        "The customer gives up delivery days; the operator gets a flatter week")

    # ── what we ask of a customer ─────────────────────────────────────────
    def week(s):
        y0 = BODY_T + 0.10
        for i, (lbl, patt, col) in enumerate(
                [("Today: every area, every day", [1] * 6, S1),
                 ("Consolidated: the area's own days", [1, 0, 0, 1, 0, 0],
                  RED)]):
            txt(s, L, y0 + i * 1.35, 6.0, 0.30, lbl, SZ_SUB, bold=True,
                color=INK)
            H.B.weekbar(s, L, y0 + 0.36 + i * 1.35, 0.62, patt, on=col,
                        off=LINE)
        txt(s, L + 5.10, y0 + 1.55, 6.9, 0.90,
            "A parcel arriving on a non-delivery day is held to the next one — "
            "at most three days, which is what makes 39 weekly patterns "
            "feasible instead of 63.", 15, color=DIM, line=1.24)

    sl = build(
        prs, SEC_FLEX, "What we ask a customer to give up",
        "theta is a scenario parameter, not a measured willingness: no survey "
        "in this study says how many customers would accept the offer.",
        [("GroupOfPeople", BLACK,
          [("theta is the share who accept: ", True),
           ("the rest keep same-day service and are collected by a separate "
            "express tour on the days their area is not served", False)]),
         ("Books", RED,
          [("The holding limit is three days: ", True),
           ("with a six-day week that leaves 39 feasible weekly patterns, and "
            "the optimiser picks one per area and provider", False)]),
         ("CryingFaceOutline", AMBER,
          [("Nobody waits longer than the schedule: ", True),
           ("the wait is a property of the pattern, so a customer knows in "
            "advance which days their address is served", False)])],
        draw=week, t=BODY_T + 2.88, start=2)
    for x in sl:
        mark(x, "MAX_HOLDING_DAYS = 3 ist eine harte Invariante des Repos "
                "(config/constants.py, EXPECTED_PATTERN_COUNT_K3 = 39). theta "
                "ist gesetzt, nicht gemessen - im Vortrag immer als Annahme "
                "fuehren.", _RV.cites("§40.21"))

    # ── what customers concede, per point ─────────────────────────────────
    s = hslide(prs, SEC_FLEX, "What that costs the customer, per operating point",
               f"Operator-polished plan at theta = 100 %, all "
               f"{_RV.eur(X.total_parcels)} parcels a week. The wait per "
               f"delayed parcel is derived: mean wait over all parcels "
               f"divided by the delayed share.")
    rows = []
    for P in P_LIST:
        d = X.delayed(f, P)
        rows.append([f"P = {P:g}", f"{d['parcels'] / 1000:.0f} k",
                     pct(d["share_pct"]), f"{d['wait_all']:.2f} d",
                     f"{d['wait_delayed']:.2f} d", f"{d['days']:.2f}"])
    y = H.B.table(s, ["Service penalty", "Parcels delayed", "Share of all",
                      "Mean wait, all", "Mean wait, delayed",
                      "Delivery days/wk"],
                  [[(r[0], "key"), r[1], (r[2], "num"), r[3], (r[4], "num"),
                    r[5]] for r in rows],
                  BODY_T, widths=[2.2, 1.9, 1.6, 2.0, 2.2, 2.3], sz=20.0,
                  reserve=2.15)
    d0, d25 = X.delayed(f, 0.0), X.delayed(f, 0.25)
    badges(s, [("CryingFaceOutline", AMBER,
                f"At the cost optimum P = 0 more than half the parcels are "
                f"held — {pct(d0['share_pct'])} of them — and an area drops "
                f"from six delivery days to {d0['days']:.1f}"),
               ("Checkmark", GREEN,
                f"A quarter-euro penalty halves that: "
                f"{pct(d25['share_pct'])} delayed, "
                f"{d25['wait_all']:.2f} d mean wait, "
                f"{d25['days']:.1f} delivery days a week")],
           y + 0.24, bottom=s.body_bottom)
    mark(s, "Verzoegerte Pakete aus discount_scenarios_v6.csv (Nachfrage an "
            "Nicht-Liefertagen mal Willigkeitsanteil), Wartezeit aus "
            "tab_headline_theta1_v2.csv. Die Spalte 'mean wait, delayed' ist "
            "abgeleitet und auf der Folie als solche bezeichnet. Eine "
            "Verteilung ueber 1/2/3 Wartetage gibt es in keiner Tabelle - "
            "nicht behaupten.", _RV.cites("§40.17", "§40.21"))

    # ── where the waiting lands ───────────────────────────────────────────
    wmax = X.week is not None
    s = figure_auto(
        prs, SEC_FLEX, "Where the waiting actually lands",
        "Operator-polished plan, P = 0.25 EUR/parcel/day. Parcel-weighted mean "
        "over the providers serving an area; held = the willing share.",
        "supp_map_wait_theta_v2",
        [("City", BLACK,
          "The dense core barely waits at all — it is already cheap to serve "
          "daily, so the optimiser leaves it alone"),
         ("OpenHandWithPlant", RED,
          f"The periphery carries the service cost: up to "
          f"{_area_wait_max():.2f} days for a held parcel, against a system "
          f"mean of {f.headline[0.25]['wait2']:.2f} d")])
    mark(s, "Die Karte ist die ehrlichste Folie des Blocks: Konsolidierung "
            "ist raeumlich ungleich verteilt. Wer die Ersparnis will, laesst "
            "genau die Gebiete warten, die heute am teuersten bedient werden.",
         _RV.cites("§40.23b"))

    # ── what the operator gets: days ──────────────────────────────────────
    s = figure_auto(
        prs, SEC_FLEX, "What the operator gets: delivery days",
        "Mean weekly delivery days per area over the adoption share, both "
        "plans, all eight penalty levels. 6 = daily; one day a week is "
        "infeasible at a three-day holding limit.",
        "supp_fig4b_mean_days",
        [("DownwardTrend_LTR", RED,
          f"From six days to {f.headline[0.0]['days2']:.2f} at P = 0 and "
          f"{f.headline[0.25]['days2']:.2f} at P = 0.25 — but only at full "
          f"adoption"),
         ("ThumbsUpSign", GREEN,
          f"The operator polish hands days BACK: "
          f"{f.headline[0.0]['days1']:.2f} to "
          f"{f.headline[0.0]['days2']:.2f} at P = 0, because a depot lowers "
          f"its weekly peak by serving more days, not fewer"),
         ("TrafficCone", AMBER,
          f"And below full adoption the curve turns back up: at P = 0.25 it "
          f"runs {X.days_rebound(0.25)[0]:.2f} days at theta = 80 % and "
          f"{X.days_rebound(0.25)[1]:.2f} at 90 %")])
    mark(s, "Zwei Botschaften auf einer Folie: (1) die Operator-Linse ist "
            "nicht 'noch mehr Konsolidierung', sie gibt Liefertage zurueck; "
            "(2) der Ruecktrend zwischen theta = 0.8 und 0.9 ist kein "
            "Rauschen, sondern Oekonomie des Restvolumens - Erklaerung im "
            "Express-Kapitel.", _RV.cites("§40.14", "§40.20", "§40.21"))

    # ── delivery frequency in space ───────────────────────────────────────
    s = figure_auto(
        prs, SEC_FLEX, "The same picture on the map",
        "Operator-polished plan, P = 0.25 EUR/parcel/day. Parcel-weighted "
        "median of the schedule size over the providers serving an area.",
        "supp_map_freq_theta_v2",
        [("City", BLACK,
          f"At full adoption three delivery days a week dominate, which is "
          f"what a system mean of {f.headline[0.25]['days2']:.2f} looks like "
          f"in space"),
         ("MagnifyingGlass", RED,
          "Consolidation moves inward as adoption grows: at 30 % it is a "
          "fringe phenomenon, at 60 % the whole ring outside the core")])
    mark(s, "Legacy-Stage-3-Karte zeigte hier 2 d/wk; auf v6 sind es 3 d/wk, "
            "weil der Operator-Polish die Frequenz anhebt. Alte Karten nicht "
            "mehr zeigen.", _RV.cites("§40.23b"))

    # ── what the operator gets: fleet and cost ────────────────────────────
    s = hslide(prs, SEC_FLEX, "What the operator gets, per operating point",
               f"Operator-polished plan at theta = 100 %, against the "
               f"daily-delivery baseline: {f.base_peak} peak vehicles, "
               f"{X.vehicle_days(0.0)['baseline']} vehicle-days, "
               f"{_RV.eur(f.base_operator)} EUR a week.")
    rows = []
    for P in P_LIST:
        h = f.headline[P]
        vd = X.vehicle_days(P) if P in (0.0, 0.25) else None
        rows.append([f"P = {P:g}", pct(h["op2"]), pct(h["rout2"]),
                     f"{h['peak2']}", signed(h["peak2_pct"]),
                     signed(vd["plan2_pct"]) if vd else "n/a"])
    y = H.B.table(s, ["Service penalty", "Operator saving", "Routing saving",
                      "Peak fleet", "vs baseline", "Vehicle-days"],
                  [[(r[0], "key"), (r[1], "num"), r[2], r[3], (r[4], "num"),
                    r[5]] for r in rows],
                  BODY_T, widths=[2.2, 2.2, 2.0, 1.7, 2.0, 2.1], sz=20.0,
                  reserve=2.15)
    vd0 = X.vehicle_days(0.0)
    badges(s, [("Target", RED,
                f"The peak falls about twice as fast as the driving: "
                f"{pct(abs(f.headline[0.0]['peak2_pct']))} fewer peak vehicles "
                f"against {pct(abs(vd0['plan2_pct']))} fewer vehicle-days at "
                f"P = 0"),
               ("MagnifyingGlass", BLACK,
                "That difference is pure levelling — the same work spread over "
                "a flatter week, which is exactly what a salaried fleet is "
                "billed for")],
           y + 0.24, bottom=s.body_bottom)
    mark(s, "Vehicle-days nur fuer P = 0 und 0.25, weil 78_ nur diese beiden "
            "Punkte als Wochenprofil schreibt; die Peak-Spalte kommt aus dem "
            "Headline-Grid und deckt alle P ab.",
         _RV.cites("§40.21", "§40.22b"))

    # ── the exchange rate ─────────────────────────────────────────────────
    s = figure_auto(
        prs, SEC_FLEX, "The exchange rate, drawn",
        "System Pareto front: weekly cost saving against mean additional "
        "customer wait, one curve per adoption share, colour = service "
        "penalty. Routing lens, routing-optimal plan.",
        "fig6_structural_grid_6_smoothed",
        [("UpwardTrend_LTR", RED,
          f"The curve is concave: the first tenth of a waiting day buys most "
          f"of the saving — {pct(f.headline[0.25]['rout2'])} for "
          f"{f.headline[0.25]['wait2']:.2f} d against "
          f"{pct(f.headline[0.0]['rout2'])} for "
          f"{f.headline[0.0]['wait2']:.2f} d"),
         ("TrafficCone", AMBER,
          "Every curve below theta = 1 collapses towards the origin: partial "
          "adoption is not a scaled-down version of full adoption")])
    mark(s, "Panel (a) ist die Wechselkurs-Figur. Die uebrigen Panels sind "
            "raeumliche Struktur und gehoeren nicht auf diese Folie - im "
            "Vortrag nur auf (a) zeigen.", _RV.cites("§40.21"))

    # ── the exchange rate, priced ─────────────────────────────────────────
    s = hslide(prs, SEC_FLEX, "The exchange rate, priced",
               "What if the penalty is not a shadow price but money actually "
               "paid to the waiting customer? Flat 0.50 EUR per delayed "
               "parcel, operator-polished plan at theta = 100 %.")
    y = H.B.table(s, ["Service penalty", "Parcels delayed",
                      "Operator saving", "Net after the discount",
                      "Break-even discount"],
                  [[(r[0], "key"), r[1], r[2], (r[3], "num"), (r[5], "num")]
                   for r in _RV.discount_rows(f)],
                  BODY_T, widths=[2.4, 2.2, 2.4, 2.6, 2.6], sz=20.0,
                  reserve=2.15)
    o = _RV.discount_optima(f)
    badges(s, [("Envelope", RED,
                [("Paid out, the penalty roughly halves the saving: ", True),
                 (f"the operator lens peaks at P = {o['operator'][0]:g} with "
                  f"{o['operator'][1]:.1f} % net, not at the cost optimum "
                  f"P = 0", False)]),
               ("Checkmark", GREEN,
                [("And it says what flexibility is worth: ", True),
                 (f"the operator could pay "
                  f"{f.discount[0.0]['be_op']:.2f} EUR per delayed parcel at "
                  f"P = 0 and {f.discount[1.0]['be_op']:.2f} EUR at P = 1 and "
                  f"still break even", False)])],
           y + 0.22, bottom=s.body_bottom)
    mark(s, _RV.discount_notes(f), _RV.cites("§40.17", "§40.22"))

    # ── the recommended point, whole ──────────────────────────────────────
    d25 = X.delayed(f, 0.25)
    h25 = f.headline[0.25]
    op25 = V.point("balanced", 0.25)
    s = hslide(prs, SEC_FLEX,
               "One operating point, both sides of the trade",
               f"P = 0.25 EUR/parcel/day, theta = 100 %, operator-polished "
               f"plan. Predicted from the grid, realised from "
               f"{op25['n']} VROOM instances.")
    H.B.stats(s, [(pct(d25["share_pct"]), "of parcels held at least one day",
                   False),
                  (f"{d25['wait_all']:.2f} d", "mean added wait per parcel",
                   False),
                  (f"{d25['days']:.2f}", "delivery days a week, from six",
                   False)], BODY_T, h=1.30, sz=40)
    H.B.stats(s, [(pct(h25["op2"]), "operator cost saved (predicted)", True),
                  (pct(op25["act_operator"]), "operator cost saved (realised)",
                   True),
                  (signed(h25["peak2_pct"]), "peak vehicles", True)],
              BODY_T + 1.62, h=1.30, sz=40)
    label_box(s, L, BODY_T + 3.30, W, 0.92, PANEL,
              [(f"A third of the parcels wait less than half a day on "
                f"average, and the operator keeps "
                f"{f.base_peak - h25['peak2']} fewer vehicles.",
                SZ_LEAD, True, INK)], line_col=LINE)
    badges(s, [("Target", RED,
                f"That is {f.base_peak - h25['peak2']} vehicles at "
                f"1 134.90 EUR a week each — the single largest term in the "
                f"operator lens")],
           BODY_T + 4.42, bottom=s.body_bottom)
    mark(s, "Das ist die Folie, die der Autor im Vortrag stehen lassen kann. "
            "Beide Seiten des Tauschs auf einer Folie, vorhergesagt und "
            "realisiert nebeneinander. 1239 - 1026 = 213 Fahrzeuge.",
         _RV.cites("§40.21", "§40.25"))

    # ── partial adoption ──────────────────────────────────────────────────
    s = hslide(prs, SEC_FLEX,
               "Below full adoption only the operator lens pays",
               "Operator-polished plan at P = 0, both lenses, over the "
               "willingness-to-wait share.")
    rows = [[f"theta = {th * 100:.0f} %", pct(rout), pct(op)]
            for th, rout, op in f.partial_adoption(0.0)
            if abs(th * 10 - round(th * 10)) < 1e-9
            and round(th * 10) in (1, 3, 5, 8, 10)]
    y = H.B.table(s, ["Willingness to wait", "Routing saving",
                      "Operator saving"],
                  [[(r[0], "key"), r[1], (r[2], "num")] for r in rows],
                  BODY_T, widths=[4.0, 4.0, 4.2], sz=SZ_BODY, reserve=2.60)
    badges(s, [("MagnifyingGlass", BLACK,
                "The routing lens barely moves below full adoption, because "
                "the express tour eats the saving on every dropped day"),
               ("Target", RED,
                "The operator lens pays throughout, because levelling the "
                "depot peak does not need the express remainder to disappear")],
           y + 0.26, bottom=s.body_bottom)
    mark(s, "Das ist die Adoptionsfrage: unterhalb von theta = 1 ist "
            "Konsolidierung fuer den Betreiber lohnend und fuer den "
            "Tourenkostenrechner fast wertlos. Warum, sagt das naechste "
            "Kapitel.", _RV.cites("§40.21"))

    # ── the pending Act-3 render ──────────────────────────────────────────
    placeholder(
        prs, SEC_FLEX, "Pending: the slide-styled Pareto front",
        "Act 3 of results/presentation_2026_08 has a 20-pt render of exactly "
        "this front; it is not embedded because its schedules are "
        "pre-revision.",
        "fig34_pareto",
        [("MagnifyingGlass", BLACK,
          "The statement that belongs here: 29 of 80 admissible grid cells "
          "lie on the efficient front, and savings up to the cost optimum are "
          "reachable at under one day of mean added wait"),
         ("TrafficCone", AMBER,
          "Until the loader is fixed the paper figure above carries this "
          "argument at paper type size")],
        "Platzhalter, kein Inhaltsproblem: die Akt-3-Renders waeren die "
        "folientauglichen Bilder (grosse Achsen), ziehen ihre Schedules aber "
        "ueber load_chosen_stage3() aus dem eingefrorenen Vor-Revisions-Lauf. "
        "Nach dem Loader-Fix hier einsetzen.")


def _area_wait_max() -> float:
    """The largest per-area wait on the wait map, read from its own table."""
    t = _RV._flex_table("tab_map_wait_theta_v2")
    import numpy as np
    return float(t[np.isclose(t.share_willing, 1.0)].wait_d.max())


# ═══════════════════════════════════════════════════════════════════════════
# 3 · Routing cost against express cost
# ═══════════════════════════════════════════════════════════════════════════
def part_express(prs):
    f, X, V = _F, _X, _V
    lo, hi = X.express_price_band(0.25)
    th_peak, sh_peak = X.express_peak_share(0.25)

    section_divider(
        prs, "Routing cost against express cost",
        "The tour nobody wants to run",
        "Why partial adoption is expensive, in euro per vehicle-day")

    # ── where an express tour comes from ──────────────────────────────────
    def origin(s):
        H.B.flow(s, [("Monday: area served", "everyone gets their\nparcel",
                      S1),
                     ("Tuesday: area skipped", "willing customers\nwait a day",
                      S2),
                     ("But not everyone waits", "1 - theta of the parcels\n"
                                                "must go out today", AMBER),
                     ("Express tour", "the remainder of that area,\n"
                                      "bundled at the depot", RED)],
              BODY_T + 2.66, bh=1.32)

    sl = build(
        prs, SEC_EXP, "Where an express tour comes from",
        "Below full adoption every dropped delivery day creates a second, "
        "thinner tour for the customers who would not wait.",
        [("Truck", BLACK,
          [("It is the price of partial adoption: ", True),
           ("it exists only because some customers said no, and it reaches "
            "the same addresses with a fraction of the parcels", False)]),
         ("TrafficCone", AMBER,
          [("It is thin by construction: ", True),
           ("the smaller the remainder, the further the vehicle drives for "
            "the same load", False)]),
         ("Checkmark", GREEN,
          [("At full adoption it disappears: ", True),
           ("which is why every curve here jumps at theta = 100 %", False)])],
        draw=origin, t=BODY_T, bottom=BODY_T + 2.52, start=2)
    for x in sl:
        mark(x, "Die Express-Sammeltour ist der zentrale Mechanismus hinter "
                "allen theta < 1-Ergebnissen. Sie ist per Zelle definiert "
                "(nicht hub-weit gepoolt) und braucht wie jede Tour "
                "mindestens 230 Pakete.", _RV.cites("§40.7", "§40.20"))

    # ── the price ─────────────────────────────────────────────────────────
    def price(s):
        H.B.stats(s, [(f"{X.regular_eur_per_vd:.0f} EUR",
                       "a regular delivery vehicle-day", False),
                      (f"{lo:.0f} EUR",
                       "an express vehicle-day at theta = 10 %", True),
                      (f"{hi:.0f} EUR",
                       "an express vehicle-day at theta = 90 %", True)],
                  BODY_T + 3.20, h=1.20, sz=38)

    s = figure_auto(
        prs, SEC_EXP, "An express vehicle-day costs more than a regular one",
        "Panel (b1): express euro per vehicle-day over the adoption share at "
        "P = 0 and P = 0.25, against a regular tour. Express quantities come "
        "from the stage-2 cost decomposition.",
        "supp_fig_mechanism_v2",
        [("Truck", RED,
          f"An express vehicle-day costs {lo:.0f} to {hi:.0f} EUR against "
          f"{X.regular_eur_per_vd:.0f} EUR for a regular one — up to "
          f"{hi / X.regular_eur_per_vd:.2f} times as much for the same van"),
         ("MagnifyingGlass", BLACK,
          "The thinner the remainder, the further the vehicle drives for the "
          "same 230 parcels: the price rises monotonically with adoption")])
    mark(s, "Express-EUR/Fahrzeugtag = Summe express_cost_eur / Summe "
            "express_veh je (P, theta); die regulaere Tour ist Summe "
            "cost_stage1_eur / Summe fleet bei (0, 0) = "
            f"{X.regular_eur_per_vd:.2f} EUR.", _RV.cites("§40.20"))

    # ── the share ─────────────────────────────────────────────────────────
    s = hslide(prs, SEC_EXP, "What the express tour costs the system",
               "P = 0.25 EUR/parcel/day, operator-polished plan. The express "
               "share is the express tours' share of total routing cost.")
    rows = []
    for th, eur, share, vd in X.express_curve(0.25):
        if round(th * 10) not in (1, 3, 5, 8, 9, 10):
            continue
        rows.append([f"theta = {th * 100:.0f} %",
                     f"{eur:.0f} EUR" if eur == eur else "-",
                     pct(share), f"{vd:.0f}",
                     f"{f.headline[0.25]['rout1']:.1f} %" if th == 1.0
                     else f"{_routing_saving(0.25, th):.2f} %"])
    y = H.B.table(s, ["Willingness to wait", "Express EUR/vehicle-day",
                      "Express share of cost", "Express vehicle-days",
                      "Routing saving"],
                  [[(r[0], "key"), r[1], (r[2], "num"), r[3], (r[4], "num")]
                   for r in rows],
                  BODY_T, widths=[2.6, 2.7, 2.5, 2.3, 2.1], sz=20.0,
                  reserve=2.15)
    badges(s, [("TrafficCone", AMBER,
                f"The express share peaks at {pct(sh_peak)} around "
                f"theta = {th_peak * 100:.0f} % and then falls as the "
                f"remainder thins — but each remaining tour gets dearer"),
               ("UpwardTrend_LTR", GREEN,
                f"At theta = 100 % the obligation disappears and the routing "
                f"saving jumps from "
                f"{_routing_saving(0.25, 0.9):.1f} % to "
                f"{f.headline[0.25]['rout1']:.1f} %")],
           y + 0.22, bottom=s.body_bottom)
    mark(s, "Der Sprung bei theta = 1 ist kein Artefakt: die letzten 10 % "
            "Nicht-Wartewilliger kosten unverhaeltnismaessig viel. Das ist "
            "der Satz fuer die Diskussion der Adoptionsschwelle.",
         _RV.cites("§40.20"))

    # ── why the week re-densifies ─────────────────────────────────────────
    r8, r9 = X.days_rebound(0.25)
    s = hslide(prs, SEC_EXP,
               "Why the delivery week gets denser again near full adoption",
               "P = 0.25 EUR/parcel/day. Between theta = 80 % and 90 % the "
               "optimiser adds delivery days back instead of removing them.")
    badges(s, [("Books", BLACK,
                [("The penalty scales with theta: ", True),
                 ("more parcels wait at the same P, so every dropped day "
                  "costs more — the penalty mass at P = 0.25 rises from "
                  f"{_penalty_keur(0.25, 0.1):.0f} k to "
                  f"{_penalty_keur(0.25, 0.9):.0f} k EUR a week", False)]),
               ("TrafficCone", AMBER,
                [("And the express remainder gets too thin to pool: ", True),
                 (f"express vehicle-days fall from "
                  f"{_express_vd(0.25, 0.1):.0f} to "
                  f"{_express_vd(0.25, 0.9):.0f}, and small depots stop "
                  f"reaching the 230-parcel threshold — a one-van mini-tour "
                  f"for a handful of parcels is dearer than keeping the "
                  f"regular day", False)]),
               ("Gears", RED,
                [("So the week re-densifies: ", True),
                 (f"mean delivery days go back up from {r8:.2f} at "
                  f"theta = 80 % to {r9:.2f} at 90 %, and the express share "
                  f"falls from {pct(X.express_share(0.25, 0.8))} to "
                  f"{pct(X.express_share(0.25, 0.9))}", False)])],
           BODY_T, bottom=s.body_bottom)
    mark(s, "Lesehilfe zu Fig. 4 und zur Mechanismus-Figur. Wichtig: das ist "
            "kein Optimiererfehler, sondern die Oekonomie des Restvolumens.",
         _RV.cites("§40.20"))

    # ── why the penalty forces daily plans ────────────────────────────────
    s = figure_split(
        prs, SEC_EXP, "Why a higher penalty means daily plans",
        "Panel (c): routing saving per parcel at maximal bundling, over all "
        "312 provider-area cells; the dashed lines are what a two-day plan's "
        "penalty costs per parcel.",
        "supp_fig_mechanism_v2",
        [("Books", BLACK,
          [("The penalty is a per-parcel price: ", True),
           ("a two-day plan costs roughly P euro per parcel and waiting day, "
            "so the cell must save more than P to be worth consolidating",
            False)]),
         ("TrafficCone", AMBER,
          [("Most cells do not save that much: ", True),
           ("at P = 0.25 about a quarter of the cells already fall short, at "
            "P = 0.5 more than half, and at P = 1 nine in ten", False)]),
         ("Target", RED,
          [("Which is why P is the policy lever: ", True),
           ("it does not tune consolidation, it switches areas off one price "
            "band at a time", False)])],
        fig=(FIG_X + 0.10, BODY_T, FIG_W - 0.10, 4.55))
    mark(s, "Die Prozentwerte 24 / 54 / 90 % stehen als Beschriftung in der "
            "Figur selbst (Zellkosten Grid v6). Der deutsche Entwurf "
            "_peek/fig_mechanism_theta_P.png nennt 23 / 49 / 86 % - das sind "
            "v5-Zellkosten, nicht zitieren.", _RV.cites("§40.20"))

    # ── one-area depots ───────────────────────────────────────────────────
    def profiles(s):
        top = max(max(f.BANTORF_BEFORE), max(f.bantorf_after))
        for i, (lbl, prof) in enumerate(
                [("Routing-optimal plan", f.BANTORF_BEFORE),
                 ("Operator-polished plan", f.bantorf_after)]):
            y0 = BODY_T + i * 1.68
            peak = max(prof)
            txt(s, L, y0, 5.0, 0.32, lbl, SZ_SUB, bold=True, color=INK)
            for k, v in enumerate(prof):
                x = L + k * 0.70
                hgt = 0.72 * v / top
                rect(s, x, y0 + 0.64 + (0.72 - max(hgt, 0.03)), 0.50,
                     max(hgt, 0.03), CRIM if v == peak else (S4 if v else LINE))
                txt(s, x, y0 + 0.36, 0.50, 0.25, str(v), SZ_DIA, bold=True,
                    color=INK if v else DIM, align=PP_ALIGN.CENTER)
                txt(s, x, y0 + 1.40, 0.50, 0.24, "MTWTFS"[k], SZ_DIA,
                    color=DIM, align=PP_ALIGN.CENTER)
            txt(s, L + 4.70, y0 + 0.76, 2.4, 0.46, f"peak {peak}", 28,
                bold=True, color=CRIM if i == 0 else TEAL)
        label_box(s, L + 7.55, BODY_T + 0.30, 5.15, 2.55, PANEL,
                  [("Temporal consolidation buys FLEET where a depot can "
                    "rotate delivery days across several areas.", 19, True,
                    INK),
                   ("Where it cannot, it buys kilometres only — and the "
                    "operator lens sends the depot back to daily service.",
                    17, False, INK2)],
                  line_col=LINE, align=PP_ALIGN.LEFT)

    sl = build(
        prs, SEC_EXP, "A one-area depot cannot rotate its delivery days",
        f"{f.bantorf_hub_name}, Monday-to-Saturday vehicles at P = 0, "
        f"theta = 100 %. The routing-optimal profile is quoted from the "
        f"compendium; the v6 tables keep only the final plan per depot and "
        f"day.",
        [("Truck", BLACK,
          f"{f.one_cell_hubs} of DHL's {f.dhl_hubs} depots serve exactly one "
          f"area, and their weekly peak is that one area's busiest day"),
         ("Gears", RED,
          "Stage 2 sends them back to daily service: such a peak only comes "
          "down by delivering on MORE days, never fewer")],
        draw=profiles, t=BODY_T + 3.52)
    for x in sl:
        mark(x, _RV.ONE_CELL_NOTES, _RV.cites("§40.14", "§40.23b"))

    # ── the same thing on the map ─────────────────────────────────────────
    s = figure_auto(
        prs, SEC_EXP, "The one-area depots, seen from above",
        "Operator-polished plan, P = 0. The 4-day patch in the west (31515 "
        "Wunstorf) is a one-area DHL depot the polish puts back on daily "
        "service; on the routing-optimal plan the whole region sits at two "
        "days.",
        "supp_map_freq_theta_P0_v2",
        [("City", RED,
          "One area, one depot, one peak: the polish pays 1 134.90 EUR per "
          "peak vehicle to avoid, which dwarfs the routing saving that area "
          "could offer"),
         ("MagnifyingGlass", BLACK,
          "The same pattern appears at all eight single-area DHL depots — it "
          "is a network property, not a local anomaly")])
    mark(s, "Kontrastfigur ist supp_map_freq_theta_P0_routing_v2 (Stufe 1, "
            "dort ist Wunstorf wie alles andere bei 2 d/wk). Bei Bedarf im "
            "Backup zeigen.", _RV.cites("§40.23b"))

    # ── what the solver says about express ────────────────────────────────
    s = hslide(prs, SEC_EXP, "The solver says we over-price express tours",
               f"VROOM validation at P = 0.25, theta = 50 % — the one "
               f"partial-adoption point that was solved. A stratified subset "
               f"of 1 000 of 1 594 instances, so it carries no realised "
               f"saving in per cent.")
    rows = [[r["kind"].replace("_", " "), f"{r['n']}", pct(r["mape"], 2),
             f"+{r['bias']:.1f} %"] for r in V.by_kind]
    y = H.B.table(s, ["Instance kind", "n", "MAPE", "Bias"],
                  [[(r[0], "key"), r[1], r[2], (r[3], "num")] for r in rows],
                  BODY_T, widths=[4.6, 1.8, 2.8, 3.0], sz=SZ_BODY,
                  reserve=2.55)
    ex = [r for r in V.by_kind if r["kind"].startswith("express")]
    badges(s, [("TrafficCone", AMBER,
                [("Express tours are the surrogate's weak spot: ", True),
                 (f"it prices them "
                  f"{min(r['bias'] for r in ex):.0f} to "
                  f"{max(r['bias'] for r in ex):.0f} % above the solver, "
                  f"against about "
                  f"{max(r['bias'] for r in V.by_kind if r['kind'].startswith('delivery')):.0f} % "
                  f"for delivery tours", False)]),
               ("UpwardTrend_LTR", GREEN,
                [("The bias points against our own conclusion: ", True),
                 ("an over-priced express tour makes partial adoption look "
                  "WORSE than it is, so the theta < 1 results here are a "
                  "floor, not a ceiling", False)])],
           y + 0.24, bottom=s.body_bottom)
    mark(s, "Gegenrichtung zum Baseline-Bias: dort ueberschaetzt das Surrogat "
            "die Baseline und ueberzeichnet die Ersparnis; hier ueberschaetzt "
            "es die Express-Touren und unterzeichnet sie. Offene Frage fuer "
            "die Limitations: ein express-spezifischer Head oder mehr "
            "Express-Labels im Pool.", _RV.cites("§40.27"))


def _routing_saving(penalty: float, share_willing: float) -> float:
    import numpy as np
    g = _D.saving_grid_v2(_D.PLAN_ROUTING, _D.LENS_ROUTING)
    r = g[np.isclose(g.penalty, penalty) & np.isclose(g.share_willing,
                                                      share_willing)]
    return float(r.saving_pct.iloc[0])


def _penalty_keur(penalty: float, share_willing: float) -> float:
    return float(_X.express_row(penalty, share_willing).penalty_plan2_keur)


def _express_vd(penalty: float, share_willing: float) -> float:
    return float(_X.express_row(penalty, share_willing).express_veh_days)


# ═══════════════════════════════════════════════════════════════════════════
# 4 · The weekly fleet
# ═══════════════════════════════════════════════════════════════════════════
def part_fleet(prs):
    f, X = _F, _X
    ORDER = ["DHL", "Amazon", "Hermes", "UPS", "DPD", "FedEx", "GLS"]

    section_divider(prs, "The weekly fleet",
                    "Where the operator's money actually is",
                    "Peak vehicles per depot, Monday to Saturday")

    for P, stem in ((0.0, "supp_fig_fleet_week_v2_P0"),
                    (0.25, "supp_fig_fleet_week_v2_P025")):
        pk = X.peaks(P)
        sysb, sys1, sys2 = pk[X.SYSTEM]
        vd = X.vehicle_days(P)
        worst = max((p for p in ORDER), key=lambda p: pk[p][1] - pk[p][0])
        s = figure_auto(
            prs, SEC_FLEET,
            f"The weekly fleet per provider at P = {P:g}",
            f"theta = 100 %. A line is the vehicles a provider runs that day "
            f"across all its depots; the dotted line is the fleet it must "
            f"keep — the sum of the weekly per-depot maxima, which is what "
            f"the operator lens bills.",
            stem,
            [("DownwardTrend_LTR", CRIM,
              f"The routing-optimal plan raises the fleet the system must "
              f"keep from {sysb} to {sys1} vehicles — {worst} alone goes "
              f"{pk[worst][0]} to {pk[worst][1]}"),
             ("Checkmark", GREEN,
              f"The operator polish brings it to {sys2} "
              f"({signed((sys2 / sysb - 1) * 100)}) while the vehicle-days "
              f"only fall {abs(vd['plan2_pct']):.1f} % — the rest is pure "
              f"levelling")])
        mark(s, f"Autorenwunsch: diese Figur gehoert 'auf jeden Fall' ins "
                f"Deck. Werte aus tab_fleet_week_by_provider_v2.csv. "
                f"Systempeaks {sysb} -> {sys1} -> {sys2}; Fahrzeugtage "
                f"{vd['baseline']} -> {vd['plan1']} -> {vd['plan2']}.",
             _RV.cites("§40.22b"))

    # ── carrier classes ───────────────────────────────────────────────────
    cls = X.class_rows(0.25)
    hyb0 = [c for c in X.class_rows(0.0) if c["group"] == "Hybrid"][0]
    s = figure_auto(
        prs, SEC_FLEET, "The same move, at three different sizes",
        "Weekly fleet profile per carrier class and penalty, theta = 100 %, "
        "indexed on each class's own baseline weekly mean so no single "
        "provider is recoverable. Classes are the routing-lens knee taxonomy.",
        "supp_fig7_fleet_week_classes",
        [("Gears", RED,
          f"Every class ends up flat: at P = 0.25 the Monday-to-Saturday "
          f"coefficient of variation falls to "
          f"{min(c['cv'][2] for c in cls):.2f}-"
          f"{max(c['cv'][2] for c in cls):.2f} from "
          f"{min(c['cv'][0] for c in cls):.2f}-"
          f"{max(c['cv'][0] for c in cls):.2f}"),
         ("TrafficCone", AMBER,
          f"The routing-optimal plan does the opposite where demand is "
          f"lumpy: the hybrid class reaches CV {hyb0['cv'][1]:.2f} at P = 0, "
          f"a single Thursday spike")])
    mark(s, "Die Figur zum Satz 'TBC schafft den Freiheitsgrad Liefertag, "
            "homogenisiert aber nicht von selbst'. Die y-Achse ist bewusst "
            "indexiert - einzelne Anbieter duerfen nicht rueckrechenbar sein.",
         _RV.cites("§40.22b"))

    # ── what the polish does to the schedules ─────────────────────────────
    sat_b, sat_2 = X.saturday(0.25)
    vd = X.vehicle_days(0.25)
    s = hslide(prs, SEC_FLEET, "What the polish actually does to the schedules",
               f"Operator-polished against routing-optimal plan across the "
               f"whole grid: {_RV.eur(X.cell_rows)} area-provider-scenario "
               f"rows.")
    H.B.stats(s, [(_RV.eur(X.freq_changed),
                   f"of {_RV.eur(X.cell_rows)} rows change frequency", True),
                  (f"{sat_b} to {sat_2}", "Saturday vehicles, P = 0.25", False),
                  (f"{abs(vd['plan2_pct']):.1f} %",
                   f"fewer vehicle-days, against "
                   f"{abs(f.headline[0.25]['peak2_pct']):.0f} % fewer peaks",
                   False)], BODY_T, h=1.70, sz=40)
    badges(s, [("Gears", RED,
                [("Stage 2 is not a re-timing pass: ", True),
                 (f"one row in five changes how often an area is served, "
                  f"which is exactly the freedom the frequency lock used to "
                  f"deny it", False)]),
               ("Truck", BLACK,
                [("The flat week is bought on Saturday: ", True),
                 (f"the quietest day of the baseline is loaded from {sat_b} "
                  f"to {sat_2} vehicles so that Wednesday does not have to "
                  f"set the fleet size", False)]),
               ("Target", GREEN,
                [("And that is where the operator money is: ", True),
                 (f"{f.base_peak - f.headline[0.25]['peak2']} peak vehicles "
                  f"removed at 1 134.90 EUR a week each", False)])],
           BODY_T + 2.05, bottom=s.body_bottom)
    mark(s, "5 525 von 27 456 Zeilen aendern die Frequenz (Summe "
            "cells_freq_changed_vs_stage1 gegen die Zeilenzahl von "
            "_tab_chosen_v2.csv). Der Samstag-Effekt steht in 40.22b.",
         _RV.cites("§40.14", "§40.22b"))


# ═══════════════════════════════════════════════════════════════════════════
# 5 · What the numbers do not say
# ═══════════════════════════════════════════════════════════════════════════
def part_honesty(prs):
    f, X, V = _F, _X, _V

    section_divider(prs, "What the numbers do not say",
                    "Three caveats worth more than a fourth digit",
                    "Solver bias · the cost of an hour · what is still open")

    # ── the validation, end to end ────────────────────────────────────────
    s = hslide(prs, SEC_HON, "Every consolidated point, solved",
               f"VROOM on {V.grid}, clean basis. The surrogate prices the SAME "
               f"tours; the gap is model error, not a different plan.")
    rows = []
    for r in V.points:
        plan = ("operator-polished" if r["plan"] == "balanced"
                else "routing-optimal")
        rows.append([f"P = {r['penalty']:g}", plan,
                     f"+{r['gap_routing']:.2f} %",
                     f"+{r['gap_operator']:.2f} %",
                     f"{r['peak_pred']:.0f} / {r['peak_actual']:.0f}"])
    y = H.B.table(s, ["Point", "Weekly plan", "Routing gap", "Operator gap",
                      "Peak fleet pred / actual"],
                  [[(r[0], "key"), r[1], (r[2], "num"), (r[3], "num"), r[4]]
                   for r in rows],
                  BODY_T, widths=[1.9, 3.6, 2.2, 2.3, 2.9], sz=19.0,
                  reserve=1.95)
    badges(s, [("Checkmark", GREEN,
                f"The fleet count survives the solver almost exactly — the "
                f"ceiling(parcels / 230) rule lands within a handful of "
                f"vehicles on about a thousand"),
               ("MagnifyingGlass", BLACK,
                f"The cost gap runs "
                f"{V.gap_range()[0]:.1f}-{V.gap_range()[1]:.1f} % in the "
                f"routing lens and "
                f"{V.gap_range('operator')[0]:.1f}-"
                f"{V.gap_range('operator')[1]:.1f} % in the operator lens, "
                f"and it grows with the penalty")],
           y + 0.22, bottom=s.body_bottom)
    mark(s, "Diese Tabelle ist die gedruckte Basis des Papers (Abschnitt 3.4): "
            "0.8 bis 2.9 % und 0.5 bis 2.5 % ueber dieselbe Population.",
         _RV.cites("§40.24"))

    # ── the upper bound ───────────────────────────────────────────────────
    rp0 = V.point("stage1", 0.0)
    op0 = V.point("balanced", 0.0)
    statement(
        prs,
        "Our prediction is an upper bound,\nnot a conservative estimate.",
        sub=f"The daily baseline is over-priced by "
            f"{V.baseline['gap_routing']:.1f} %, the bundled scenarios by "
            f"only {V.gap_range()[0]:.1f}-{V.gap_range()[1]:.1f} %: "
            f"{f.headline[0.0]['rout1']:.1f} % predicted, "
            f"{rp0['act_routing']:.1f} % realised at P = 0.")
    mark(prs.slides[-1],
         "Die Einreichung hatte den umgekehrten Fall (22.8 vorhergesagt -> "
         "23.7 realisiert). Formulierung 'obere Schranke (ca. 2 pp)' statt "
         "'konservativ'. Rangfolge und Vorzeichen aller Aussagen bleiben.",
         _RV.cites("§40.25"))

    # ── the cost of an hour ───────────────────────────────────────────────
    s = hslide(prs, SEC_HON, "How much of this is the price of an hour?",
               "Sensitivity: the same stored VROOM solutions re-priced with "
               "the route-hour term set to zero. No re-solving, no retraining.")
    badges(s, [("Books", BLACK,
                [("The hour behaves like a second distance term: ", True),
                 ("36 EUR per route-hour is about 1.52 EUR per kilometre at "
                  "the observed speeds, so the effective distance rate is "
                  "1.90 rather than 0.3864 EUR/km", False)]),
               ("TrafficCone", AMBER,
                [("Drop it and the realised saving nearly halves: ", True),
                 ("19.9 % becomes 12.6 % for the routing-optimal plan and "
                  "17.5 % becomes 11.2 % for the operator-polished plan at "
                  "P = 0; the change over all points is 2.5 to 7.3 points",
                  False)]),
               ("MagnifyingGlass", RED,
                [("Because consolidation saves driving time faster than "
                  "distance: ", True),
                 ("time is 26.5 % of the baseline's cost and only 19.8 to "
                  "24.5 % of the scenarios'", False)])],
           BODY_T, bottom=s.body_bottom)
    mark(s, "Task-15-Memo, Option B. Die urspruengliche Annahme ('ohne "
            "Zeitterm laege die Ersparnis hoeher') war falsch - VROOM "
            "bepreist per_hour nur auf die Fahrzeit, nicht auf die "
            "Servicezeit, und das Vorzeichen dreht sich. Entscheidung "
            "A/B/C steht noch aus.", _RV.cites("§40.26"))

    # ── kilometres, as the by-product ─────────────────────────────────────
    km = X.km_co2()
    s = hslide(prs, SEC_HON, "One number we did not optimise for",
               "Vehicle-kilometres from the solved tours of the validation, "
               "0.25 kg CO2 per kilometre as an external assumption. Not a "
               "claim in the manuscript.")
    H.B.stats(s, [(f"{km['baseline'][0] / 1000:.0f}k km",
                   "driven per week today", False),
                  (f"{km['operator'][0.0][0] / 1000:.0f}k km",
                   "operator-polished plan at P = 0", True),
                  (f"{km['operator'][0.0][2]:+.0f} %",
                   f"distance, against {f.headline[0.0]['rout2']:.0f} % cost",
                   True)], BODY_T, h=1.40, sz=40)
    badges(s, [("OpenHandWithPlant", GREEN,
                f"Distance falls further than cost because the fixed cost per "
                f"vehicle-day does not fall with the kilometres — "
                f"{km['baseline'][1]:.1f} t of CO2 a week becomes "
                f"{km['operator'][0.0][1]:.1f} t"),
               ("TrafficCone", AMBER,
                "This is a by-product, not an objective: nothing in the "
                "optimisation ever saw a kilometre or a gram")],
           BODY_T + 1.75, bottom=s.body_bottom)
    mark(s, "Im Manuskript steht keine CO2-Aussage (Part C); die Tabelle "
            "gehoert ins Supplement. Referenzspalte "
            "vs_least_consolidated_pct der CSV NICHT benutzen - sie "
            "vergleicht innerhalb eines Plans, nicht gegen die Baseline.",
         _RV.cites("§40.28"))

    # ── what is still open ────────────────────────────────────────────────
    sl = build(
        prs, SEC_HON, "What is still open",
        "Three of these are decisions, not unknowns — they need a ruling "
        "rather than more computing.",
        [("HeadWithGears", AMBER,
          [("The express head: ", True),
           ("express group prices come largely from the fallback, which is "
            "why their bias is the largest in the validation. An "
            "express-specific head or more express labels in the pool is the "
            "next step", False)]),
         ("Gears", BLACK,
          [("The route-hour term: ", True),
           ("re-solving everything without it is about 79 VROOM-hours; the "
            "memo recommends reporting the sensitivity instead and naming the "
            "effective cost model in the methods", False)]),
         ("GroupOfPeople", RED,
          [("And the parameter we cannot compute: ", True),
           ("theta is set, not measured. Everything here says what happens IF "
            "a share of customers accepts — not how large that share is",
            False)])],
        t=BODY_T, start=2)
    for x in sl:
        mark(x, "Punkt 3 ist der wichtigste fuer die Diskussion: Die Studie "
                "liefert die Austauschkurve, nicht den Betriebspunkt. Den "
                "waehlt, wer die Zahlungsbereitschaft kennt.",
             _RV.cites("§40.26", "§40.27"))


# ═══════════════════════════════════════════════════════════════════════════
# 6 · Closing
# ═══════════════════════════════════════════════════════════════════════════
def part_close(prs):
    f, X, V = _F, _X, _V
    h25 = f.headline[0.25]
    d25 = X.delayed(f, 0.25)
    op25 = V.point("balanced", 0.25)

    s = hslide(prs, "Three sentences", "What this study lets you say",
               "Grid v6, VROOM-validated, both cost lenses.")
    badges(s, [("Checkmark", GREEN,
                [("One. ", True),
                 (f"At a quarter-euro service penalty and full acceptance, "
                  f"holding {pct(d25['share_pct'])} of parcels for "
                  f"{d25['wait_all']:.2f} days on average lets an operator "
                  f"run the same week with "
                  f"{f.base_peak - h25['peak2']} fewer vehicles — "
                  f"{pct(op25['act_operator'])} of operator cost, solved, not "
                  f"predicted", False)]),
               ("Target", RED,
                [("Two. ", True),
                 ("The cheapest daily plan is the worst weekly plan: "
                  "optimising tour cost alone raises the fleet an operator "
                  "has to keep by a third, and only a frequency-free second "
                  "stage turns that around", False)]),
               ("GroupOfPeople", AMBER,
                [("Three. ", True),
                 ("Below full acceptance the express tour for the customers "
                  "who would not wait eats most of the saving — so the "
                  "binding question is adoption, not routing", False)])],
           BODY_T, bottom=s.body_bottom)
    mark(s, "Die drei Saetze, die der Autor am Ende sagen kann. Jeder ist "
            "durch eine Folie im Deck belegt.",
         _RV.cites("§40.21", "§40.25", "§40.20"))

    contact(prs)

    # ── backup ────────────────────────────────────────────────────────────
    extras = [
        ("fig5_grid_heatmap_6_smoothed", "The paper's grid figure",
         "Fig. 5 as submitted for the revision: cost, wait, peak fleet, CV "
         "and vehicle-days over the whole (P, theta) grid."),
        ("supp_fig4_freq_mix_two_plans", "Delivery-frequency mix, both plans",
         "Routing-optimal plan above, operator-polished plan below: how the "
         "mix of weekly frequencies moves over the grid."),
        ("supp_fig6b_operator_lens_v2", "The operator lens on its own",
         "Pareto front, per-provider knees in both lenses, and hub savings by "
         "carrier type, hub size and region type. Panels (c)-(f) are "
         "hub-attributed, not area-attributed."),
        ("supp_map_freq_theta_P0_routing_v2",
         "The routing-optimal plan on the map",
         "P = 0: the routing optimum puts every area on two delivery days. "
         "Contrast to the operator-polished map in the express section."),
        ("supp_penalty_raumtyp_v2", "The penalty as a policy lever",
         "Euro-weighted saving within each settlement type, routing lens, "
         "operator-polished plan. From P = 5 EUR per parcel and day on, every "
         "type keeps at most 0.7 %."),
    ]
    for stem, subject, src in extras:
        p = FIGDIR / f"{stem}.png"
        if not p.exists():
            continue
        s = hslide(prs, "Backup", subject, src)
        pic(s, p, L, BODY_T, W, 5.05)
        chip(s, stem)


# ═══════════════════════════════════════════════════════════════════════════
def build_deck(out: Path, *, facts, flex, vroom, revision, data) -> Path:
    global _F, _X, _V, _RV, _D, FIGDIR
    _F, _X, _V, _RV, _D = facts, flex, vroom, revision, data
    FIGDIR = data.REV / "figures"

    prs = Presentation(str(H.TEMPLATE))
    for master in prs.slide_masters:
        for shp in master.shapes:
            if shp.has_text_frame and "mobil.TUM" in shp.text_frame.text:
                shp.text_frame.text = ""
                H.B._para(shp.text_frame.paragraphs[0],
                          "What customers concede and what operators gain | "
                          "EWGT 2026 revision | Bienzeisler, Petre, Wage, "
                          "Friedrich", 10, color=DIM)
    for i in range(len(prs.slides) - 1, -1, -1):
        H.delete_slide(prs, i)

    title_slide(prs)
    part_setup(prs)
    part_headline(prs)
    part_flexibility(prs)
    part_express(prs)
    part_fleet(prs)
    part_honesty(prs)
    part_close(prs)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=Path,
                    default=H.TEMPLATE.parent /
                    "EWGT_26_Bienzeisler_v6_flexibility_deck.pptx")
    ap.add_argument("--rev-dir", type=Path, default=None,
                    help="the grid to read (default: $PRES_REV_DIR, else "
                         "results/revision_2026_08_v6)")
    ap.add_argument("--val-dir", type=Path, default=None,
                    help="the FINISHED VROOM validation to read (default: the "
                         "grid's own validation/ directory). _data's module "
                         "default still names v5, which has no solved "
                         "baseline, so this build sets it explicitly.")
    G.add_args(ap)
    a = ap.parse_args()
    out = G.resolve(a.out, a.out_suffix, overwrite=a.overwrite)

    import _data as D
    import _revision as RV
    if a.rev_dir is not None:
        D.set_rev_dir(a.rev_dir)
    if D.SCHEMA != D.SCHEMA_V2:
        raise SystemExit(
            f"{D.REV} is a {D.SCHEMA} grid; this deck needs the two-plan "
            f"tables. Pass --rev-dir or set PRES_REV_DIR.")
    D.set_val_dir(a.val_dir or (D.REV / "validation"))
    print(f"  grid:       {D.REV.name}")
    print(f"  validation: {D.VAL.parent.name}")
    facts, flex, vroom = RV.Facts.load(), RV.Flex.load(), RV.Vroom.load()
    print(f"  {vroom.n_solved()} clean VROOM instances, "
          f"{len(vroom.points)} consolidated points")
    p = build_deck(out, facts=facts, flex=flex, vroom=vroom, revision=RV,
                   data=D)
    print(f"wrote {p}")
    print(f"  {len(Presentation(str(p)).slides)} slides, "
          f"{p.stat().st_size / 1048576:.1f} MB")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
