"""Draw a generated slide as a PNG, from the shape tree, without PowerPoint.

`_verify_layout.py` measures a deck; this one shows it. Neither PowerPoint COM
(`pywin32`/`comtypes` are not installed) nor LibreOffice is available on this
machine, so a build cannot be proof-read by exporting slide images the usual
way. What python-pptx does give is every shape's rectangle, fill, line and text
runs with their real point sizes — which is enough to redraw the slide to scale
and see whether it reads.

This is a GEOMETRY preview, not a PowerPoint render: gradients, the master's
own chrome, autoshape geometry beyond rectangles and ovals, tables' internal
borders and any real font metrics beyond Arial are approximations. Use it to
answer "does this slide hold together and is anything colliding or empty", not
"is this pixel-accurate".

    python scripts/presentation/_preview_slide.py DECK.pptx 62 63 64 --out DIR
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

DPI = 96
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
_FONTS: dict = {}


def _font(size_pt: float, bold: bool):
    key = (round(size_pt, 1), bold)
    if key not in _FONTS:
        name = "arialbd.ttf" if bold else "arial.ttf"
        try:
            _FONTS[key] = ImageFont.truetype(rf"C:\Windows\Fonts\{name}",
                                             max(6, int(round(size_pt * DPI / 72))))
        except OSError:                                   # pragma: no cover
            _FONTS[key] = ImageFont.load_default()
    return _FONTS[key]


def _rgb(colour, default=None):
    try:
        if colour is None or colour.type is None:
            return default
        return "#%02X%02X%02X" % tuple(colour.rgb)
    except Exception:
        return default


def _fill(shape):
    try:
        f = shape.fill
        if f.type is None or f.type == 5:            # inherited / background
            return None
        if f.type == 1:                              # solid
            return _rgb(f.fore_color)
    except Exception:
        pass
    return None


def _line(shape):
    try:
        if shape.line.fill.type in (None, 5):
            return None
        return _rgb(shape.line.color)
    except Exception:
        return None


def _wrap(draw, text, font, width_px):
    out, cur = [], ""
    for word in text.split(" "):
        trial = (cur + " " + word).strip()
        if cur and draw.textlength(trial, font=font) > width_px:
            out.append(cur)
            cur = word
        else:
            cur = trial
    out.append(cur)
    return out


def _draw_text_frame(draw, tf, box, scale):
    l, t, w, h = box
    y = t + 2
    for para in tf.paragraphs:
        runs = [r for r in para.runs if r.text]
        if not runs:
            y += 6
            continue
        size = max((r.font.size.pt if r.font.size else 18.0) for r in runs)
        bold = any(bool(r.font.bold) for r in runs)
        col = next((_rgb(r.font.color) for r in runs
                    if _rgb(r.font.color)), "#15181D")
        font = _font(size * scale, bold)
        text = "".join(r.text for r in runs)
        for line in _wrap(draw, text, font, max(8, w - 6)):
            draw.text((l + 3, y), line, font=font, fill=col)
            y += font.size * 1.20
        y += 3


def render(deck: Path, index: int, out: Path, scale: float = 1.0) -> Path:
    prs = Presentation(str(deck))
    slide = list(prs.slides)[index - 1]
    W = int(SLIDE_W_IN * DPI * scale)
    H = int(SLIDE_H_IN * DPI * scale)
    img = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(img)

    def box(sh):
        return (Emu(sh.left).inches * DPI * scale,
                Emu(sh.top).inches * DPI * scale,
                Emu(sh.width).inches * DPI * scale,
                Emu(sh.height).inches * DPI * scale)

    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        l, t, w, h = box(sh)
        if sh.shape_type is not None and sh.has_text_frame is False \
                and sh.shape_type == 13:                       # picture
            try:
                im = Image.open(sh.image.blob and __import__("io").BytesIO(
                    sh.image.blob))
                im.thumbnail((max(1, int(w)), max(1, int(h))))
                img.paste(im.convert("RGB"), (int(l), int(t)))
                continue
            except Exception:
                pass
        fill, line = _fill(sh), _line(sh)
        if fill or line:
            d.rectangle([l, t, l + w, t + h], fill=fill, outline=line or None)
        if sh.has_table:
            tbl = sh.table
            ys, x0 = t, l
            for row in tbl.rows:
                rh = Emu(row.height).inches * DPI * scale
                d.line([x0, ys, x0 + w, ys], fill="#D5D9DE")
                xs = x0
                for ci, cell in enumerate(row.cells):
                    cw = Emu(tbl.columns[ci].width).inches * DPI * scale
                    _draw_text_frame(d, cell.text_frame, (xs, ys, cw, rh), scale)
                    xs += cw
                ys += rh
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            _draw_text_frame(d, sh.text_frame, (l, t, w, h), scale)

    d.rectangle([0, 0, W - 1, H - 1], outline="#BBBBBB")
    out.parent.mkdir(parents=True, exist_ok=True)
    img.save(out)
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck", type=Path)
    ap.add_argument("slides", type=int, nargs="+")
    ap.add_argument("--out", type=Path, required=True,
                    help="directory to write slide-NN.png into")
    ap.add_argument("--scale", type=float, default=1.0)
    a = ap.parse_args()
    for i in a.slides:
        p = render(a.deck, i, a.out / f"slide-{i:03d}.png", a.scale)
        print(f"  wrote {p}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
