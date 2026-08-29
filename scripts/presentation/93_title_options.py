"""Five alternative opening slides for the EWGT talk, to pick one from.

Each is a complete title slide in the institutional template, built from the
same helpers as the main deck (91_build_pptx.py), so whichever one wins can be
copied straight into it — or this file's slide can be dragged to position 1.

Every number on these slides is a real figure from the case study:

    1 263 130 parcels per week   results/.../tab_baseline_per_provider.csv
    1 909 748 EUR baseline/week  D.BASE_TOTAL, pinned
    312 provider-area cells      the model's unit count
    7 providers, 39 patterns     hard invariants
    22.8 % saving, 0.98 d wait   Stage-3 optimum at (P, theta) = (0, 1)

Usage:
    python scripts/presentation/93_title_options.py [--out PATH]
"""
from __future__ import annotations

import argparse
import importlib.util
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[1]


def _load_builder():
    """Import 91_build_pptx.py, whose module name is not a valid identifier."""
    spec = importlib.util.spec_from_file_location(
        "deck_builder", HERE / "91_build_pptx.py")
    mod = importlib.util.module_from_spec(spec)
    sys.modules["deck_builder"] = mod
    spec.loader.exec_module(mod)
    return mod


B = _load_builder()
TITLE = "Time-Based Consolidation\nin Last-Mile Delivery"
AUTHORS = "Lasse Bienzeisler · Felix Petre · Oskar Wage · Bernhard Friedrich"
AFFIL = "TU Braunschweig · Leibniz University Hannover · EWGT 2026"


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[B.LAYOUT_BLANK])


def _byline(s, y, *, colour=None, size=None):
    B.txt(s, B.L + 0.40, y, 11.6, 0.9, f"{AUTHORS}\n{AFFIL}",
          size or (B.SZ_STATL + 1), color=colour or B.INK2, line=1.65)


def opt_a_volume(prs):
    """A · The number that stops the room."""
    s = _blank(prs)
    B.txt(s, B.L + 0.40, 1.55, 11.6, 0.40, "One region · one week",
          B.SZ_KICK, bold=True, color=B.RED, spc=1.6, caps=True)
    B.txt(s, B.L + 0.30, 1.95, 12.0, 1.5, "1 263 130", 116, bold=True,
          color=B.INK)
    B.txt(s, B.L + 0.40, 3.45, 11.6, 0.5,
          "parcels, delivered by seven carriers, every single week.",
          B.SZ_LEAD + 2, color=B.INK2)
    B.hrule(s, B.L + 0.40, 4.15, 3.2, B.RED, 3.0)
    B.txt(s, B.L + 0.40, 4.35, 11.6, 1.1, TITLE.replace("\n", " "), 30,
          bold=True, color=B.INK, line=1.15)
    _byline(s, 5.55)
    return s


def opt_b_map(prs):
    """B · The region as the stage."""
    s = _blank(prs)
    B.rect(s, B.L, 1.70, 0.09, 3.30, B.RED)
    B.txt(s, B.L + 0.40, 1.78, 6.4, 0.36, "Region Hannover", B.SZ_KICK,
          bold=True, color=B.RED, spc=1.6, caps=True)
    B.txt(s, B.L + 0.40, 2.20, 6.6, 1.6, TITLE, 38, bold=True, color=B.INK,
          line=1.08)
    B.txt(s, B.L + 0.40, 3.95, 6.4, 0.9,
          "312 provider–area cells · seven carriers\n"
          "1.26 million parcels a week", B.SZ_LEAD, color=B.INK2, line=1.4)
    _byline(s, 5.30)
    B.pic(s, B.FIG / "tierA" / "fig12_map_demand.png", 7.30, 1.35, 5.55, 4.60)
    return s


def opt_c_week(prs):
    """C · The mechanism, before a word is said."""
    s = _blank(prs)
    B.txt(s, B.L + 0.40, 1.35, 11.6, 0.36, "Six days · two deliveries",
          B.SZ_KICK, bold=True, color=B.RED, spc=1.6, caps=True)
    cw, gap = 1.72, 0.26
    x0 = B.L + 0.40
    for i, on in enumerate([False, True, False, False, True, False]):
        B.rect(s, x0 + i * (cw + gap), 1.85, cw, 1.15,
               B.RED if on else B.PANEL,
               line_col=None if on else B.LINE)
        B.txt(s, x0 + i * (cw + gap), 3.08, cw, 0.32, "MTWTFS"[i],
              B.SZ_BODY, color=B.DIM, align=PP_ALIGN.CENTER)
    B.txt(s, x0, 3.65, 11.6, 1.2, TITLE, 38, bold=True, color=B.INK, line=1.08)
    B.txt(s, x0, 5.05, 11.6, 0.5,
          "Hold what can wait. Deliver it together. Nothing else changes.",
          B.SZ_LEAD, color=B.INK2)
    _byline(s, 5.75)
    return s


def opt_d_trade(prs):
    """D · The trade, stated up front."""
    s = _blank(prs)
    B.txt(s, B.L + 0.40, 1.45, 11.6, 0.36, "The trade in two numbers",
          B.SZ_KICK, bold=True, color=B.RED, spc=1.6, caps=True)
    for i, (val, lab, hot) in enumerate([
            ("22.8%", "cheaper to run the week", True),
            ("0.98 d", "longer for those who opted in", False)]):
        x = B.L + 0.40 + i * 6.10
        B.hrule(s, x, 1.95, 5.30, B.RED if hot else B.LINE, 3.0 if hot else 1.5)
        B.txt(s, x, 2.10, 5.30, 1.3, val, 88, bold=True,
              color=B.RED if hot else B.INK)
        B.txt(s, x, 3.45, 5.30, 0.5, lab, B.SZ_LEAD, color=B.INK2)
    B.txt(s, B.L + 0.40, 4.25, 11.6, 1.2, TITLE, 34, bold=True, color=B.INK,
          line=1.10)
    B.txt(s, B.L + 0.40, 5.45, 11.6, 0.44,
          "Region Hannover · 312 cells · seven carriers", B.SZ_BODY,
          color=B.DIM)
    _byline(s, 5.95)
    return s


def opt_e_statement(prs):
    """E · The claim first, the title second."""
    s = _blank(prs)
    B.rect(s, 0, 0, 13.333, 7.5, B.RED)
    B.txt(s, 1.05, 1.65, 11.0, 0.36, "EWGT 2026", B.SZ_KICK, bold=True,
          color=B.RGBColor(0xF2, 0xC8, 0xD1),
          spc=1.6, caps=True)
    B.txt(s, 1.05, 2.10, 11.2, 1.9, "Batch where it is\nsparse and far.", 62,
          bold=True, color=B.WHITE, line=1.04)
    B.hrule(s, 1.05, 4.35, 2.0, B.RGBColor(0xE4, 0x9A, 0xA8), 2.5)
    B.txt(s, 1.05, 4.60, 10.8, 1.0, TITLE.replace("\n", " "), 26, bold=True,
          color=B.WHITE, line=1.2)
    B.txt(s, 1.05, 5.55, 11.0, 0.9, f"{AUTHORS}\n{AFFIL}", B.SZ_STATL + 1,
          color=B.RGBColor(0xF7, 0xE0, 0xE5), line=1.65)
    return s


OPTIONS = [("A · The volume", opt_a_volume),
           ("B · The region", opt_b_map),
           ("C · The week", opt_c_week),
           ("D · The trade", opt_d_trade),
           ("E · The claim", opt_e_statement)]


def build(out: Path) -> Path:
    prs = Presentation(str(B.TEMPLATE))
    for i in range(len(prs.slides) - 1, -1, -1):
        B.delete_slide(prs, i)
    for label, fn in OPTIONS:
        s = fn(prs)
        # a small marker so the options stay identifiable once separated
        box = s.shapes.add_textbox(Inches(11.0), Inches(0.16), Inches(2.0),
                                   Inches(0.34))
        box.text_frame.margin_left = box.text_frame.margin_right = 0
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = label
        run.font.name, run.font.size, run.font.bold = B.FONT, Pt(11), True
        run.font.color.rgb = (B.RGBColor(0xF2, 0xC8, 0xD1)
                              if fn is opt_e_statement else B.DIM)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", type=Path,
                    default=B.TEMPLATE.parent /
                    "EWGT_26_Bienzeisler_TBC_title_options.pptx")
    a = ap.parse_args()
    p = build(a.out)
    print(f"wrote {p}")
    print(f"  {len(Presentation(str(p)).slides)} title options")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
