"""The author's own presentation grammar, measured from four of his decks.

This module is not a taste statement. Every constant in it was extracted from
`FreightSim_Bienzeisler.pptx` (78 slides), `mobilTUM_25_Bienzeisler.pptx` (42),
`Universitätstagung_2023_Bienzeisler.pptx` (46) and `Disputation_Bienzeisler.pptx`
(112) — 278 slides in total — so a deck built on it looks like the decks he
already gives.

What the measurement found
--------------------------
**Two-line titles.** 73–89 % of his titled slides put two lines in the single
title placeholder at (0.63, 0.12) 12.21 x 0.774, separated by ``<a:br/>``:
line 1 is the *section*, repeated verbatim on every slide of that section, and
line 2 is what this particular slide shows. Both lines inherit the master's
22 pt bold. That is why his average title runs 60–76 characters — it is two
titles, not one long one.

**Progressive builds.** The same slide is emitted two to five times in a row
with one more bullet each time (mobilTUM 16–20, 28–29, 32–34, 40–42;
Disputation 10–12; FreightSim 9–12). ``build()`` reproduces this.

**Icon badges instead of bullet characters.** 124 white circles of 0.69 in with
a 1 pt coloured ring and a stock icon inside, text starting at x = 1.91. The
ring colour is semantic: TU red for the subject, green for a result that holds,
amber for a caveat, black for neutral machinery.

**Sentences, not fragments.** Body copy is 22–24 pt (his most frequent run
sizes are 22 pt and 24 pt), two to four items per slide, each a complete
declarative sentence of roughly 90–140 characters, no terminal period. Key
messages open with a bold lead-in and a colon — "Time-based consolidation
works: …", "Trade off: …", "Need for adaptive strategies: …".

**Figures carry the argument.** 4–5 pictures per slide on average, no native
charts and almost no tables; a text column at x = 0.63 (width ~5.3) with the
figure at x = 5.95 (width ~6.9), frequently bleeding off the bottom edge.

**The end of the deck.** A contact slide with name, institute, phone, e-mail
and repository — and then the backup slides *after* it.

The schematic primitives (week bars, cell rasters, tours, flows, Morph) are
imported from the existing generator rather than rewritten, so the two decks
cannot drift apart.
"""
from __future__ import annotations

import importlib.util
import sys

if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")
from pathlib import Path

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[1]
FIG = ROOT / "results" / "presentation_2026_08" / "slides"
ASSET = ROOT / "results" / "presentation_2026_08" / "assets_generated"
ICONS = ROOT / "results" / "presentation_2026_08" / "icons"
ICON_CACHE = ICONS / "_tinted"


def _load_primitives():
    """Import 91_build_pptx.py, whose module name is not a valid identifier."""
    if "deck_builder" in sys.modules:
        return sys.modules["deck_builder"]
    spec = importlib.util.spec_from_file_location(
        "deck_builder", HERE / "91_build_pptx.py")
    mod = importlib.util.module_from_spec(spec)
    sys.modules["deck_builder"] = mod
    spec.loader.exec_module(mod)
    return mod


B = _load_primitives()
TEMPLATE = B.TEMPLATE

# ── colour: the template's red, plus his three semantic ring colours ────────
RED = B.RED                              # BE1E3C, the subject
GREEN = RGBColor(0x00, 0xB0, 0x50)       # a result that holds
AMBER = RGBColor(0xE8, 0x9A, 0x00)       # a caveat or a trade-off
BLACK = RGBColor(0x00, 0x00, 0x00)       # neutral machinery
INK = B.INK
INK2 = B.INK2
DIM = B.DIM
LINE = B.LINE
PANEL = B.PANEL
WHITE = B.WHITE
BLUSH = B.BLUSH
TEAL = B.TEAL
CRIM = B.CRIM
S1, S2, S3, S4, S5, S6 = B.S1, B.S2, B.S3, B.S4, B.S5, B.S6

# ── type scale: his measured run sizes ──────────────────────────────────────
SZ_BODY = 22.0      # the dominant body size in every deck
SZ_LEAD = 24.0      # the second most frequent; used for a slide's key line
SZ_SUB = 16.0       # sub-items under a bullet
SZ_SRC = 10.0       # sources and image credits
SZ_DIA = 15.0       # labels inside a schematic
SZ_BIG = 44.0       # a headline number
SZ_NUM = 32.0       # the numeral of a numbered finding

# ── geometry, in inches, measured off his slides ────────────────────────────
SW, SH = 13.333, 7.5
L, W = B.L, B.W                 # 0.63 / 12.21 — the master's title box
BODY_T = 1.26                   # first content row, as on his slides
BODY_B = 6.62
BADGE_X = 0.95                  # centre-left of the icon circle
BADGE_D = 0.69                  # his circle diameter
TEXT_X = 1.91                   # where the sentence starts
TEXT_W = 10.76
ROW_H = 1.02                    # row pitch between badge bullets
COL_W = 5.05                    # text column when a figure sits on the right
FIG_X = 5.95                    # his figure column
FIG_W = 6.75
SRC_B = 6.96                    # sources sit on this line, growing upwards
SRC_T = SRC_B - 0.20            # kept for callers that place a one-line source

LAYOUT_TITLE_ONLY = B.LAYOUT_TITLE_ONLY
LAYOUT_BLANK = B.LAYOUT_BLANK

# re-export the primitives the content file needs, so it imports one module
txt, rect, label_box, hrule, arrow, dot, pic = (
    B.txt, B.rect, B.label_box, B.hrule, B.arrow, B.dot, B.pic)
weekbar, flow, cell_grid, tour, serpentine = (
    B.weekbar, B.flow, B.cell_grid, B.tour, B.serpentine)
set_morph, delete_slide, slide_oval, axes = (
    B.set_morph, B.delete_slide, B.slide_oval, B.axes)
# raw units, for callers that draw connectors themselves
Inches, Pt = B.Inches, B.Pt


# ── icons ───────────────────────────────────────────────────────────────────
# The PNG fallbacks of the stock icons he uses, lifted from his own files by
# scripts/presentation/_extract_icons.py. Four of the extracted images are
# fully opaque rasters rather than line art and are excluded.
_SOLID = {"Box", "Database", "HighVoltage", "Receiver"}


def icon_path(name: str, colour: RGBColor = RED) -> Path:
    """Return a tinted copy of a stock icon, rendering it once and caching it.

    The icons are flat line art on transparency, so tinting is a fill through
    the alpha channel — it keeps the anti-aliased edge and drops the original
    two-tone (his Truck is black-and-amber) in favour of one semantic colour.
    """
    src = ICONS / f"{name}.png"
    if not src.exists():
        raise FileNotFoundError(f"no icon {name!r} in {ICONS}")
    ICON_CACHE.mkdir(parents=True, exist_ok=True)
    dst = ICON_CACHE / f"{name}_{colour}.png"
    if not dst.exists():
        with Image.open(src) as im:
            alpha = im.convert("RGBA").getchannel("A")
            tint = Image.new("RGBA", im.size,
                             (colour[0], colour[1], colour[2], 0))
            tint.putalpha(alpha)
            tint.save(dst)
    return dst


def pic_cover(slide, path, l, t, w, h, *, focus=0.5):
    """Fill the box completely, cropping the overflowing side.

    `pic()` letterboxes, which leaves the master's chrome showing through on a
    full-bleed slide. `focus` is the fraction of the cropped axis to keep
    centred on: 0 is top/left, 1 is bottom/right.
    """
    sh = slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w),
                                  Inches(h))
    with Image.open(path) as im:
        iw, ih = im.size
    want, have = (w / h), (iw / ih)
    if have > want:                       # image too wide: crop left/right
        keep = want / have
        off = (1 - keep) * focus
        sh.crop_left, sh.crop_right = off, 1 - keep - off
    elif have < want:                     # image too tall: crop top/bottom
        keep = have / want
        off = (1 - keep) * focus
        sh.crop_top, sh.crop_bottom = off, 1 - keep - off
    return sh


def badge(slide, cx, cy, name, colour=RED, *, d=BADGE_D):
    """His bullet marker: a white circle with a coloured ring and an icon."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2),
                                Inches(cy - d / 2), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = colour
    sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    if name:
        g = d * 0.60
        slide.shapes.add_picture(str(icon_path(name, colour)),
                                 Inches(cx - g / 2), Inches(cy - g / 2),
                                 Inches(g), Inches(g))
    return sh


def _runs(p, item, size, colour):
    """Render a bullet body: a plain string, or [(text, bold), ...]."""
    if isinstance(item, str):
        item = [(item, False)]
    for text, bold in item:
        B._para(p, text, size, bold=bold, color=INK if bold else colour)


# Text height is measured, never guessed. A fixed row pitch is what makes
# generated slides overlap: a 120-character sentence is one line in a 10.8 in
# column and five lines in a 4.3 in one, and no single pitch is right for both.
_FONTS = {}


def _font(size_pt, bold):
    from PIL import ImageFont
    key = (round(size_pt, 1), bold)
    if key not in _FONTS:
        name = "arialbd.ttf" if bold else "arial.ttf"
        try:
            f = ImageFont.truetype(rf"C:\Windows\Fonts\{name}",
                                   int(round(size_pt * 4)))
        except OSError:                                   # pragma: no cover
            f = ImageFont.load_default()
        _FONTS[key] = f
    return _FONTS[key]


def text_lines(item, w_in, size_pt, *, bold=False):
    """Greedy word-wrap count for `item` in a `w_in` wide box at `size_pt`.

    `bold` applies to a plain string; a run list carries its own weights.
    Bold Arial is about 8 % wider, which is the difference between one line
    and two on a narrow card, so it must not be assumed away.
    """
    runs = [(item, bold)] if isinstance(item, str) else list(item)
    limit = w_in * 72.0 * 4.0                 # the font is rendered at 4x
    lines, cur = 1, 0.0
    for text, bold in runs:
        f = _font(size_pt, bold)
        for word in text.split(" "):
            if not word:
                continue
            adv = f.getlength(word + " ")
            if cur > 0 and cur + f.getlength(word) > limit:
                lines += 1
                cur = adv
            else:
                cur += adv
    return lines


def text_height(item, w_in, size_pt, line=1.18, *, bold=False):
    """Height in inches that `item` needs, including its line spacing."""
    return text_lines(item, w_in, size_pt, bold=bold) * size_pt * line / 72.0


def _block_height(items, tw, size, gap, line):
    h = 0.0
    for i, (_, _, body) in enumerate(items):
        h += max(BADGE_D, text_height(body, tw, size, line))
        if i:
            h += gap
    return h


def badges(slide, items, t=BODY_T, *, x=BADGE_X, tx=TEXT_X, tw=TEXT_W,
           size=SZ_BODY, colour=INK2, gap=0.26, line=1.18, bottom=None,
           floor=20.0, label=""):
    """A column of icon-badge bullets — his standard body block.

    `items` is a list of (icon, ring colour, body). The body is a string, or
    a run list for a bold lead-in: [("Trade off: ", True), ("…", False)].

    Each row is as tall as its own text, so a long sentence in a narrow column
    pushes the next badge down instead of printing over it. If the block still
    runs past `bottom`, the type steps down towards `floor` — and if even that
    is not enough the build says so on stderr, because the honest fix there is
    fewer words, not smaller ones.
    """
    if bottom is None:
        # every hslide() records where its source line starts, so a block laid
        # out on it stops above the citation rather than printing through it
        bottom = getattr(slide, "body_bottom", BODY_B)
    limit = bottom - t
    # Tighten the gaps before touching the type: his rows sit close together,
    # and 22 pt is the size the room reads. Only then step the size down.
    while gap > 0.15 and _block_height(items, tw, size, gap, line) > limit:
        gap -= 0.01
    while size > floor and _block_height(items, tw, size, gap, line) > limit:
        size -= 0.5
    over = _block_height(items, tw, size, gap, line) - limit
    if over > 0.03:
        print(f"  ! {label or 'badges'}: overruns by {over:.2f} in at "
              f"{size:g} pt — shorten the text", file=sys.stderr)
    y = t
    for ic, col, body in items:
        th = text_height(body, tw, size, line)
        rh = max(BADGE_D, th)
        badge(slide, x, y + rh / 2, ic, col)
        box, tf = B._frame(slide, tx, y + (rh - th) / 2 - 0.045, tw, th + 0.10)
        p = tf.paragraphs[0]
        p.line_spacing = line
        _runs(p, body, size, colour)
        y += rh + gap
    return y - gap


def mathline(slide, l, t, w, h, parts, size, *, align=PP_ALIGN.CENTER,
             color=INK, accent=RED):
    """One line of formula, as real editable text with real sub/superscripts.

    `parts` is a list of (text, kind); kind is "" for the base line, "sub" or
    "sup" for the shifted positions, and "hi" to colour a term in the accent.
    PowerPoint carries the shift in the run's `baseline` attribute, which
    python-pptx does not expose, so it is set on the XML directly.
    """
    box, tf = B._frame(slide, l, t, w, h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, kind in parts:
        run = B._para(p, text, size * (0.62 if kind in ("sub", "sup") else 1.0),
                      bold=True, color=accent if kind == "hi" else color,
                      font="Cambria Math")
        if kind == "sub":
            run.font._rPr.set("baseline", "-25000")
        elif kind == "sup":
            run.font._rPr.set("baseline", "30000")
    return box


def term_notes(slide, items, t, *, l=L, w=W, gap=0.24, size=13.0,
               head=13.5):
    """The legend under a formula: the symbol in red, its reading beneath.

    The note column is sized to the longest note, so one wordy term cannot
    push its own text out of its box while the others sit half empty.
    """
    n = len(items)
    cw = (w - gap * (n - 1)) / n
    nh = max(text_height(note, cw, size, 1.22) for _, note in items) + 0.06
    for i, (sym, note) in enumerate(items):
        x = l + i * (cw + gap)
        txt(slide, x, t, cw, 0.30, sym, head, bold=True, color=RED)
        txt(slide, x, t + 0.32, cw, nh, note, size, color=DIM, line=1.22)
    return t + 0.32 + nh


def sub_items(slide, lines, t, *, x=TEXT_X, w=TEXT_W, size=SZ_SUB,
              pitch=0.34):
    """The 16 pt detail lines he hangs under a bullet."""
    for i, line in enumerate(lines):
        txt(slide, x, t + i * pitch, w, pitch, line, size, color=DIM)
    return t + len(lines) * pitch


# ── slide scaffolding ───────────────────────────────────────────────────────
def hslide(prs, section, subject, source=None):
    """A content slide with his two-line title and a 10 pt source line.

    The source is bottom-anchored to `SRC_B` and sized to the lines it really
    needs, so a three-line citation grows upwards into the body instead of
    down into the master's footer rule. The slide carries the resulting
    `body_bottom`, which is what the body blocks lay out against.
    """
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
    tf = s.shapes.title.text_frame
    tf.text = ""
    p = tf.paragraphs[0]
    B._para(p, section, 22, bold=True, color=INK)
    if subject:
        br = p._p.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}br", {})
        p._p.append(br)
        B._para(p, subject, 22, bold=True, color=RED)
    bottom = SRC_B
    if source:
        sh = text_height(source, 11.2, SZ_SRC, 1.25)
        top = SRC_B - sh
        txt(s, L, top, 11.2, sh + 0.04, source, SZ_SRC, color=DIM, line=1.25)
        bottom = top
    s.body_bottom = bottom - 0.16
    return s


def build(prs, section, subject, source, items, *, draw=None, start=1,
          t=BODY_T, **kw):
    """Emit the same slide several times, one more bullet each time.

    This is the single most characteristic move in his decks: the audience
    never meets a wall of text, it meets one line at a time on a slide that
    otherwise does not move. `draw(slide)` paints whatever furniture the slide
    carries besides the bullets, and is called on every step so the geometry is
    identical — which is also what lets PowerPoint's Morph tween between them.
    """
    out = []
    kw.setdefault("label", f"{section} / {subject}")
    bottom = kw.pop("bottom", None)
    for n in range(start, len(items) + 1):
        s = hslide(prs, section, subject, source)
        if draw:
            draw(s)
        badges(s, items[:n], t, bottom=bottom or s.body_bottom, **kw)
        out.append(s)
    return out


def statement(prs, text, *, sub=None, size=40):
    """A title-less slide carrying one sentence — his section punctuation."""
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    hrule(s, L + 0.40, 2.90, 2.2, RED, 3.0)
    txt(s, L + 0.40, 3.15, 11.4, 2.0, text, size, bold=True, color=INK,
        line=1.16)
    if sub:
        txt(s, L + 0.40, 5.25, 11.4, 0.9, sub, SZ_LEAD, color=INK2, line=1.35)
    return s


def section_divider(prs, section, subject, sub=None):
    """A red full-bleed card announcing the next section."""
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    rect(s, 0, 0, SW, SH, RED)
    hrule(s, 1.05, 2.95, 2.0, RGBColor(0xE4, 0x9A, 0xA8), 3.0)
    txt(s, 1.05, 3.20, 11.0, 0.55, section, SZ_LEAD, bold=True,
        color=RGBColor(0xF2, 0xC8, 0xD1))
    txt(s, 1.05, 3.85, 11.0, 1.5, subject, 42, bold=True, color=WHITE,
        line=1.08)
    if sub:
        txt(s, 1.05, 5.35, 10.6, 0.9, sub, SZ_BODY,
            color=RGBColor(0xF7, 0xE0, 0xE5), line=1.35)
    return s


def figure_slide(prs, section, subject, source, path, *, items=None,
                 fig=(L, BODY_T, W, 3.55), t=None):
    """Figure across the top, badge bullets underneath."""
    s = hslide(prs, section, subject, source)
    pic(s, path, *fig)
    if items:
        badges(s, items, t if t is not None else fig[1] + fig[3] + 0.16,
               bottom=s.body_bottom, label=f"{section} / {subject}")
    return s


def split_slide(prs, section, subject, source, path, items, *,
                fig=(FIG_X, BODY_T, FIG_W, 5.10), tw=COL_W - 0.10,
                t=BODY_T + 0.10):
    """His two-column slide: text on the left, one large figure on the right."""
    s = hslide(prs, section, subject, source)
    pic(s, path, *fig)
    badges(s, items, t, x=0.95, tx=1.50, tw=tw, bottom=s.body_bottom,
           label=f"{section} / {subject}")
    return s


def contact(prs, *, name="Dr.-Ing. Lasse Bienzeisler",
            lines=("Technische Universität Braunschweig",
                   "Institut für Verkehr und Stadtbauwesen",
                   "+49 (0) 531 391 66806",
                   "l.bienzeisler@tu-braunschweig.de",
                   "github.com/TUBS-IVS/vroom-valhalla-lmd-hannover")):
    """The closing slide he ends every talk with — backup slides follow it."""
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    rect(s, L, 2.20, 0.09, 2.95, RED)
    txt(s, L + 0.42, 2.26, 10.5, 0.55, name, 30, bold=True, color=INK)
    txt(s, L + 0.42, 2.98, 11.0, len(lines) * SZ_BODY * 1.55 / 72.0 + 0.06,
        "\n".join(lines), SZ_BODY, color=INK2, line=1.55)
    return s
