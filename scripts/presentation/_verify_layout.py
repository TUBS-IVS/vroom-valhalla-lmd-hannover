"""Find the layout faults a generated deck is prone to, without opening it.

Generated slides fail in a small number of mechanical ways, and all of them are
measurable: a sentence needs more lines than its box is tall, two text boxes
claim the same rectangle, or a shape runs off the bottom of the slide into the
footer. This script measures each of those with the real font metrics and
reports the offending slides.

    python scripts/presentation/_verify_layout.py <deck.pptx> [--quiet]

It is a check, not a fix: it prints what is wrong and returns a non-zero exit
code so a build can be gated on it.
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):        # box-drawing and Greek in titles
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

from PIL import ImageFont
from pptx import Presentation
from pptx.util import Emu

SLIDE_H = 7.5
SLIDE_W = 13.333
FOOTER_T = 7.03      # the master's footer rule; nothing should reach it
_FONTS: dict = {}


def _font(size_pt: float, bold: bool):
    key = (round(size_pt, 1), bold)
    if key not in _FONTS:
        name = "arialbd.ttf" if bold else "arial.ttf"
        try:
            _FONTS[key] = ImageFont.truetype(rf"C:\Windows\Fonts\{name}",
                                             int(round(size_pt * 4)))
        except OSError:                                   # pragma: no cover
            _FONTS[key] = ImageFont.load_default()
    return _FONTS[key]


def _para_height(para, w_in: float, default_size: float) -> float:
    """Wrapped height of one paragraph, in inches."""
    runs = [(r.text, r.font.size.pt if r.font.size else default_size,
             bool(r.font.bold)) for r in para.runs if r.text]
    if not runs:
        return 0.0
    size = max(s for _, s, _ in runs)
    limit = w_in * 72.0 * 4.0
    lines, cur = 1, 0.0
    for text, sz, bold in runs:
        f = _font(sz, bold)
        for word in text.split(" "):
            if not word:
                continue
            if cur > 0 and cur + f.getlength(word) > limit:
                lines += 1
                cur = f.getlength(word + " ")
            else:
                cur += f.getlength(word + " ")
    spacing = para.line_spacing if isinstance(para.line_spacing, float) else 1.2
    after = (para.space_after.pt if para.space_after else 0.0) / 72.0
    return lines * size * spacing / 72.0 + after


def text_boxes(slide):
    """Every text-bearing shape with its rectangle and its needed height."""
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if sh.width is None or sh.height is None:
            continue
        l, t = Emu(sh.left).inches, Emu(sh.top).inches
        w, h = Emu(sh.width).inches, Emu(sh.height).inches
        tf = sh.text_frame
        pad = (Emu(tf.margin_left).inches if tf.margin_left else 0.0) + \
              (Emu(tf.margin_right).inches if tf.margin_right else 0.0)
        need = sum(_para_height(p, max(0.4, w - pad), 18.0)
                   for p in tf.paragraphs)
        out.append((sh, l, t, w, h, need))
    return out


def _overlap(a, b) -> float:
    """Area of the intersection of two (l, t, w, h) rectangles."""
    x = min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0])
    y = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
    return x * y if x > 0 and y > 0 else 0.0


def check(path: Path, *, tol: float = 0.08, quiet: bool = False) -> int:
    prs = Presentation(str(path))
    faults = 0
    for i, slide in enumerate(prs.slides, 1):
        boxes = text_boxes(slide)
        msgs = []
        for sh, l, t, w, h, need in boxes:
            if need > h + tol:
                msgs.append(f"overflow {need - h:+.2f}in  "
                            f"h={h:.2f} needs {need:.2f}  "
                            f"{sh.text_frame.text[:52]!r}")
            if t + max(h, need) > FOOTER_T:
                msgs.append(f"into footer, bottom {t + max(h, need):.2f}  "
                            f"{sh.text_frame.text[:52]!r}")
        # text-on-text collisions, using the height the text actually needs
        for j, (sa, la, ta, wa, ha, na) in enumerate(boxes):
            ra = (la, ta, wa, max(ha, na))
            for sb, lb, tb, wb, hb, nb in boxes[j + 1:]:
                rb = (lb, tb, wb, max(hb, nb))
                ov = _overlap(ra, rb)
                if ov > 0.12 * min(ra[2] * ra[3], rb[2] * rb[3]):
                    msgs.append(
                        f"collision {ov:.2f}in²  "
                        f"{sa.text_frame.text[:26]!r} × "
                        f"{sb.text_frame.text[:26]!r}")
        for sh in slide.shapes:
            if sh.height is None or sh.top is None or sh.left is None:
                continue
            bot = Emu(sh.top).inches + Emu(sh.height).inches
            right = Emu(sh.left).inches + Emu(sh.width).inches
            if bot > SLIDE_H + 0.02:
                msgs.append(f"off-slide bottom {bot:.2f} ({sh.shape_type})")
            if right > SLIDE_W + 0.02 or Emu(sh.left).inches < -0.02:
                msgs.append(f"off-slide horizontally "
                            f"{Emu(sh.left).inches:.2f}..{right:.2f} "
                            f"({sh.shape_type})")
        if msgs:
            faults += len(msgs)
            if not quiet:
                print(f"\nslide {i}")
                for m in msgs:
                    print(f"   {m}")
    print(f"\n{faults} fault(s) across {len(prs.slides)} slides")
    return faults


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("deck", type=Path)
    ap.add_argument("--quiet", action="store_true")
    ap.add_argument("--tol", type=float, default=0.08)
    a = ap.parse_args()
    return 1 if check(a.deck, tol=a.tol, quiet=a.quiet) else 0


if __name__ == "__main__":
    sys.exit(main())
